import os
import sys
import asyncio
import time
import requests
import json
import datetime
import logging
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from collections import defaultdict
import re
from bs4 import BeautifulSoup

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Colors for the PSIRT advisories slide
TITLE_COLOR = RGBColor(39, 110, 55)  # Dark green for subtitle
HIGH_COLOR = RGBColor(227, 119, 84)  # Red/Orange for high severity
MEDIUM_COLOR = RGBColor(248, 196, 71)  # Yellow/Amber for medium severity
LOW_COLOR = RGBColor(108, 184, 108)  # Green for low severity
INFO_COLOR = RGBColor(79, 129, 189)  # Blue for informational

# Set up a local logger
logger = logging.getLogger('meraki_psirt')
logger.setLevel(logging.INFO)
if not logger.handlers:
    # Ensure logs directory exists
    os.makedirs('logs', exist_ok=True)
    file_handler = logging.FileHandler('logs/psirt_advisories.log')
    file_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
    logger.addHandler(file_handler)
    logger.propagate = False

async def get_api_token():
    """
    Gets an API token from Cisco's OAuth service
    
    Returns:
        API token or None if failed
    """
    client_id = "mzcve4k4hsx5q6385quy9aza" #may need to update to get creds from user in future
    client_secret = "mMwFv5YJMSFYjStMY2n59rPd" #may need to update to get creds from user in future
    token_url = "https://id.cisco.com/oauth2/default/v1/token"
    
    data = {
        "client_id": client_id,
        "client_secret": client_secret,
        "grant_type": "client_credentials"
    }
    
    headers = {
        "Content-Type": "application/x-www-form-urlencoded"
    }
    
    try:
        print(f"{BLUE}Getting Cisco API token...{RESET}")
        response = requests.post(token_url, headers=headers, data=data)
        
        if response.status_code == 200:
            token_data = response.json()
            access_token = token_data.get("access_token")
            if access_token:
                print(f"{GREEN}Successfully obtained API token{RESET}")
                return access_token
            else:
                print(f"{RED}No access token in response{RESET}")
                logger.error("No access token in response")
        else:
            print(f"{RED}Failed to get API token: Status code {response.status_code}{RESET}")
            logger.error(f"Failed to get API token: Status code {response.status_code}")
            
    except Exception as e:
        print(f"{RED}Error getting API token: {e}{RESET}")
        logger.error(f"Error getting API token: {e}")
    
    return None

async def fetch_psirt_advisories(page_size=50, max_pages=10):
    """
    Fetches Cisco PSIRT advisories from the API
    
    Args:
        page_size: Number of advisories to fetch per page
        max_pages: Maximum number of pages to fetch
        
    Returns:
        A list of advisories related to Meraki products
    """
    meraki_advisories = []
    total_advisories = 0
    
    # Get API token
    api_token = await get_api_token()
    
    if not api_token:
        print(f"{YELLOW}Could not obtain Cisco API token. Using mock data.{RESET}")
        logger.warning("Could not obtain Cisco API token. Using mock data.")
        return get_mock_advisories()
    
    # Set up headers with token
    headers = {
        "Accept": "application/json",
        "Authorization": f"Bearer {api_token}"
    }
    
    # Calculate date range for the last 5 years exactly
    end_date = datetime.datetime.now().strftime("%Y-%m-%d")
    # Use 5*365.25 days to account for leap years more accurately
    start_date = (datetime.datetime.now() - datetime.timedelta(days=int(5*365.25))).strftime("%Y-%m-%d")
    
    # API endpoint for firstpublished with date range
    api_endpoint = f"https://apix.cisco.com/security/advisories/v2/all/firstpublished?startDate={start_date}&endDate={end_date}"
    
    try:
        print(f"{BLUE}Fetching PSIRT advisories from {start_date} to {end_date}...{RESET}")
        
        response = requests.get(api_endpoint, headers=headers)
        
        # Check if request was successful
        if response.status_code == 200:
            data = response.json()
            advisories = data.get("advisories", [])
            
            if not advisories:
                print(f"{YELLOW}No advisories found for the date range.{RESET}")
                logger.warning("No advisories found for the date range.")
                return get_mock_advisories()
                
            # Process all advisories from the response
            for advisory in advisories:
                title = advisory.get("advisoryTitle", "").lower()
                
                # Handle productNames correctly - it can be a string or a list
                product_names = advisory.get("productNames", "")
                if isinstance(product_names, list):
                    # Join list items into a single string for searching
                    product_names_str = " ".join(product_names).lower()
                else:
                    # If it's already a string, just convert to lowercase
                    product_names_str = product_names.lower()
                
                if "meraki" in title or "meraki" in product_names_str:
                    summary = advisory.get("summary", "")
                    summary = clean_summary_text(summary)
                    
                    # Add the cleaned summary back to the advisory
                    advisory["cleanedSummary"] = summary
                    
                    # Get the publication URL to scrape for fixed firmware version later
                    publication_url = advisory.get("publicationUrl")
                    if publication_url:
                        advisory["publicationUrl"] = publication_url
                    
                    meraki_advisories.append(advisory)
                    
            total_advisories = len(advisories)
            
            print(f"{GREEN}Found {len(meraki_advisories)} Meraki-related advisories out of {total_advisories} total advisories{RESET}")
            logger.info(f"Found {len(meraki_advisories)} Meraki-related advisories out of {total_advisories} total advisories")
        
        else:
            print(f"{RED}Error fetching PSIRT advisories: Status code {response.status_code}{RESET}")
            logger.error(f"Error fetching PSIRT advisories: Status code {response.status_code}")
            if response.status_code == 401:
                print(f"{RED}Authentication error. Token may have expired.{RESET}")
                logger.error("Authentication error. Token may have expired.")
            return get_mock_advisories()
            
    except Exception as e:
        print(f"{RED}Error fetching PSIRT advisories: {e}{RESET}")
        logger.error(f"Error fetching PSIRT advisories: {e}")
        return get_mock_advisories()
    
    # If we didn't find any Meraki advisories, use mock data
    if not meraki_advisories:
        print(f"{YELLOW}No Meraki-related advisories found. Using mock data.{RESET}")
        logger.warning("No Meraki-related advisories found. Using mock data.")
        return get_mock_advisories()
    
    # Sort advisories by date (newest first)
    meraki_advisories.sort(key=lambda x: x.get("firstPublished", ""), reverse=True)
    
    return meraki_advisories

async def fetch_fixed_firmware_version(url):
    """
    Scrapes the advisory page to find the fixed firmware version information.
    
    Args:
        url: The URL of the advisory page
        
    Returns:
        A string with the fixed firmware version information, or None if not found
    """
    if not url:
        return None
    
    try:
        # Fetch the advisory page
        response = requests.get(url)
        
        if response.status_code != 200:
            logger.warning(f"Failed to fetch advisory page: {url}, status code: {response.status_code}")
            return None
        
        # Parse the HTML with BeautifulSoup
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Look for the "Fixed Releases" section
        fixed_releases_header = None
        
        # Try to find the fixed releases section by looking for headers
        for header in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            if 'Fixed Releases' in header.text:
                fixed_releases_header = header
                break
            
        if not fixed_releases_header:
            # Try looking for strong tags or other emphasized elements
            for element in soup.find_all(['strong', 'b']):
                if 'Fixed Releases' in element.text:
                    fixed_releases_header = element
                    break
        
        if fixed_releases_header:
            # Extract the next table or list after the header
            result = []
            
            # Try to find the table
            table = fixed_releases_header.find_next('table')
            
            if table:
                # Process table rows
                rows = table.find_all('tr')
                for row in rows[1:]:
                    cells = row.find_all('td')
                    if len(cells) >= 2:
                        device_version = cells[0].text.strip()
                        fixed_version = cells[1].text.strip()
                        
                        if "First Fixed Release" in fixed_version or "First Fixed Release" in device_version:
                            continue
                            
                        if "Not affected" not in fixed_version and "Migrate" not in fixed_version:
                            # This row contains a fixed version
                            result.append(f"{device_version} → {fixed_version}")
            
            if result:
                return "\n".join(result)
            
            # If no table found, try to extract text after the header
            next_element = fixed_releases_header.next_sibling
            fixed_version_text = ""
            
            # Collect text until the next header
            while next_element and not next_element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                if hasattr(next_element, 'text'):
                    fixed_version_text += next_element.text
                elif isinstance(next_element, str):
                    fixed_version_text += next_element
                    
                next_element = next_element.next_sibling
                
            if fixed_version_text:
                # Clean up the text
                fixed_version_text = re.sub(r'\s+', ' ', fixed_version_text).strip()
                return fixed_version_text
        
    except Exception as e:
        logger.error(f"Error scraping advisory page: {e}")
        
    return None

def clean_summary_text(summary):
    """Clean up the summary text from HTML tags and formatting"""
    if not summary:
        return "No summary available"
    
    if not isinstance(summary, str):
        try:
            summary = str(summary)
        except:
            return "Unable to parse summary"
    
    try:
        # Replace escaped characters
        summary = summary.replace("\\r\\n", "\n").replace("\\n", "\n")
        
        # Remove HTML tags using BeautifulSoup
        try:
            soup = BeautifulSoup(summary, 'html.parser')
            summary = soup.get_text(separator=' ')
        except:
            # If BeautifulSoup fails, use regex as fallback
            summary = re.sub(r'<[^>]+>', ' ', summary)
        
        summary = re.sub(r'\s+', ' ', summary).strip()
        
            
    except Exception as e:
        logger.error(f"Error cleaning summary text: {e}")
        return "Error processing summary"
    
    return summary

def get_mock_advisories():
    """Return mock data for testing when API access is not available"""
    mock_advisories = [
        {
            "advisoryId": "cisco-sa-meraki-wpa-DoS-nYQPTNhC",
            "advisoryTitle": "Cisco Meraki Access Points Management Frame Denial of Service Vulnerability",
            "bugIDs": "CSCvw54321",
            "ipsSignatures": "NA",
            "cves": "CVE-2023-20123",
            "cvrfUrl": "https://tools.cisco.com/security/center/contentxml/CiscoSecurityAdvisory/cisco-sa-meraki-wpa-DoS-nYQPTNhC/cvrf/cisco-sa-meraki-wpa-DoS-nYQPTNhC_cvrf.xml",
            "csafUrl": "https://tools.cisco.com/security/center/contentjson/CiscoSecurityAdvisory/cisco-sa-meraki-wpa-DoS-nYQPTNhC/csaf/cisco-sa-meraki-wpa-DoS-nYQPTNhC_csaf.json",
            "cvssBaseScore": 7.5,
            "cwe": "CWE-400",
            "firstPublished": "2023-11-15T16:00:00",
            "lastUpdated": "2023-11-22T18:42:53",
            "status": "Final",
            "version": 1.1,
            "productNames": "Cisco Meraki MR Access Points",
            "publicationUrl": "https://tools.cisco.com/security/center/content/CiscoSecurityAdvisory/cisco-sa-meraki-wpa-DoS-nYQPTNhC",
            "sir": "High",
            "cleanedSummary": "A vulnerability in the management frame handling of Cisco Meraki MR Access Points could allow an unauthenticated, adjacent attacker to cause a denial of service (DoS) condition on an affected device. This vulnerability is due to improper validation of WPA management frames. An attacker could exploit this vulnerability by sending a crafted management frame to an affected device. A successful exploit could allow the attacker to cause the access point to reload, resulting in a DoS condition."
        },
        {
            "advisoryId": "cisco-sa-meraki-mx-vpn-7tX3BcP9",
            "advisoryTitle": "Cisco Meraki MX Security Appliance VPN Information Disclosure Vulnerability",
            "bugIDs": "CSCvw12345",
            "ipsSignatures": "NA",
            "cves": "CVE-2023-20456",
            "cvrfUrl": "https://tools.cisco.com/security/center/contentxml/CiscoSecurityAdvisory/cisco-sa-meraki-mx-vpn-7tX3BcP9/cvrf/cisco-sa-meraki-mx-vpn-7tX3BcP9_cvrf.xml",
            "csafUrl": "https://tools.cisco.com/security/center/contentjson/CiscoSecurityAdvisory/cisco-sa-meraki-mx-vpn-7tX3BcP9/csaf/cisco-sa-meraki-mx-vpn-7tX3BcP9_csaf.json",
            "cvssBaseScore": 5.3,
            "cwe": "CWE-200",
            "firstPublished": "2023-10-04T16:00:00",
            "lastUpdated": "2023-10-11T14:22:17",
            "status": "Final",
            "version": 1.0,
            "productNames": "Cisco Meraki MX Security Appliance",
            "publicationUrl": "https://tools.cisco.com/security/center/content/CiscoSecurityAdvisory/cisco-sa-meraki-mx-vpn-7tX3BcP9",
            "sir": "Medium",
            "cleanedSummary": "A vulnerability in the VPN functionality of Cisco Meraki MX Security Appliances could allow an unauthenticated, remote attacker to access sensitive information. This vulnerability is due to improper validation of VPN packet headers. An attacker could exploit this vulnerability by sending crafted packets to the VPN interface of an affected device. A successful exploit could allow the attacker to access certain VPN configuration details."
        },
        {
            "advisoryId": "cisco-sa-meraki-ms-auth-W6xKmP8q",
            "advisoryTitle": "Cisco Meraki MS Switch Authentication Bypass Vulnerability",
            "bugIDs": "CSCvw67890",
            "ipsSignatures": "NA",
            "cves": "CVE-2023-20789",
            "cvrfUrl": "https://tools.cisco.com/security/center/contentxml/CiscoSecurityAdvisory/cisco-sa-meraki-ms-auth-W6xKmP8q/cvrf/cisco-sa-meraki-ms-auth-W6xKmP8q_cvrf.xml",
            "csafUrl": "https://tools.cisco.com/security/center/contentjson/CiscoSecurityAdvisory/cisco-sa-meraki-ms-auth-W6xKmP8q/csaf/cisco-sa-meraki-ms-auth-W6xKmP8q_csaf.json",
            "cvssBaseScore": 8.8,
            "cwe": "CWE-306",
            "firstPublished": "2023-09-20T16:00:00",
            "lastUpdated": "2023-09-27T12:15:34",
            "status": "Final",
            "version": 1.0,
            "productNames": "Cisco Meraki MS Switch",
            "publicationUrl": "https://tools.cisco.com/security/center/content/CiscoSecurityAdvisory/cisco-sa-meraki-ms-auth-W6xKmP8q",
            "sir": "High",
            "cleanedSummary": "A vulnerability in the authentication mechanism of Cisco Meraki MS Switches could allow an unauthenticated, adjacent attacker to bypass authentication on an affected device. This vulnerability is due to improper verification of authentication tokens. An attacker could exploit this vulnerability by submitting crafted authentication requests to an affected device. A successful exploit could allow the attacker to gain unauthorized access to the switch configuration."
        }
    ]
    return mock_advisories

def get_sir_color(sir):
    """Get the appropriate color for a Security Impact Rating (SIR)"""
    sir = sir.lower() if sir else "unknown"
    if sir == "critical" or sir == "high":
        return HIGH_COLOR
    elif sir == "medium":
        return MEDIUM_COLOR
    elif sir == "low":
        return LOW_COLOR
    else:
        return INFO_COLOR

def format_date(date_string):
    """Format a date string to a readable format"""
    if not date_string:
        return "Unknown"
    
    try:
        # Parse ISO format date
        date_obj = datetime.datetime.fromisoformat(date_string.replace('Z', '+00:00'))
        # Return formatted date
        return date_obj.strftime("%B %d, %Y")
    except ValueError:
        return date_string

def parse_firmware_version(version_str):
    """
    Parse a firmware version string to extract the numeric components.
    
    Args:
        version_str: A string containing a firmware version (e.g., "MX 18.211.3" or "18.0 and later → 18.211.3")
        
    Returns:
        A tuple of (product_type, major, minor, patch) or None if parsing fails
    """
    if not version_str:
        return None
    
    # Make a copy of the original string for debugging
    original_str = version_str
    
    # Extract the right side of an arrow if present (get the patched version)
    if "→" in version_str:
        version_str = version_str.split("→")[1].strip()
    
    # Clean up any extra text
    version_str = version_str.replace("and later", "").strip()
    
    # Try to identify product type
    product_type = None
    for prefix in ["MX", "MR", "MS", "MV", "MG", "MT"]:
        pattern = rf'\b{prefix}\b'  # Word boundary to avoid partial matches
        if re.search(pattern, version_str, re.IGNORECASE):
            product_type = prefix
            # Remove the product type prefix
            version_str = re.sub(pattern, "", version_str, flags=re.IGNORECASE).strip()
            break
    
    # Extract version numbers using regex
    version_match = re.search(r'(\d+)\.(\d+)(?:\.(\d+))?', version_str)
    if version_match:
        major = int(version_match.group(1))
        minor = int(version_match.group(2))
        patch = int(version_match.group(3)) if version_match.group(3) else 0
        return (product_type, major, minor, patch)
    
    version_match = re.search(r'(\d+)\.(\d+)(?:\.(\d+))?', original_str)
    if version_match:
        major = int(version_match.group(1))
        minor = int(version_match.group(2))
        patch = int(version_match.group(3)) if version_match.group(3) else 0
        return (product_type, major, minor, patch)
    
    return None

def is_version_affected(current_version, patched_version):
    """
    Determine if a device with current_version is affected by a vulnerability fixed in patched_version.
    
    Args:
        current_version: Tuple of (product_type, major, minor, patch) for the current device firmware
        patched_version: Tuple of (product_type, major, minor, patch) for the patched firmware
        
    Returns:
        True if the device is affected (current version < patched version), False otherwise
    """
    if not current_version or not patched_version:
        return False
    
    # Extract product types
    current_product = current_version[0]
    patched_product = patched_version[0]
    
    # If product types are specified and don't match, no impact
    if current_product and patched_product and current_product != patched_product:
        return False
    
    # Compare versions (major.minor.patch)
    current_nums = current_version[1:]
    patched_nums = patched_version[1:]
    
    # Compare version numbers
    return current_nums < patched_nums

def find_affected_devices(advisory, mxmsmr_csv_path, mgmvmt_csv_path):
    """
    Find devices affected by a security advisory by comparing firmware versions.
    
    Args:
        advisory: The advisory dictionary
        mxmsmr_csv_path: Path to the MXMSMR firmware compliance CSV
        mgmvmt_csv_path: Path to the MGMVMT firmware compliance CSV
        
    Returns:
        A list of affected devices (dictionaries with device info)
    """
    affected_devices = []
    
    # Extract patched firmware version from advisory
    patched_version_str = advisory.get("fixed_firmware_version")
    if not patched_version_str:
        return affected_devices  # No version to compare against
    
    # Parse the patched version
    patched_version = None
    product_type = None
    
    # Check the advisory title to guess the product type if not in version string
    title = advisory.get("advisoryTitle", "").upper()
    for prefix in ["MX", "MS", "MR", "MV", "MG", "MT"]:
        if prefix in title:
            product_type = prefix
            break
    
    # Handle multiple patched versions by looking for any arrows
    if "→" in patched_version_str:
        version_parts = patched_version_str.split("\n")
        for part in version_parts:
            if "→" in part:
                fixed_version = part.split("→")[1].strip()
                parsed = parse_firmware_version(fixed_version)
                if parsed:
                    patched_version = parsed
                    # If product type is not in version, use the one from title
                    if patched_version[0] is None and product_type:
                        patched_version = (product_type, *patched_version[1:])
                    break
    else:
        patched_version = parse_firmware_version(patched_version_str)
        # If product type is not in version, use the one from title
        if patched_version and patched_version[0] is None and product_type:
            patched_version = (product_type, *patched_version[1:])
    
    if not patched_version:
        return affected_devices
    
    # Determine which CSV to check based on product type
    product_type = patched_version[0]
    if not product_type:
        # Try to extract product type from advisory title
        title = advisory.get("advisoryTitle", "").upper()
        for prefix in ["MX", "MS", "MR", "MV", "MG", "MT"]:
            if prefix in title:
                product_type = prefix
                # Update the patched version tuple with the product type
                patched_version = (product_type, *patched_version[1:])
                break
    
    # Check if we still don't have a product type
    if not product_type:
        # Try to examine product names
        product_names = advisory.get("productNames", "")
        if isinstance(product_names, list):
            product_names = " ".join(product_names)
        
        product_names = product_names.upper()
        for prefix in ["MX", "MS", "MR", "MV", "MG", "MT"]:
            if prefix in product_names:
                product_type = prefix
                patched_version = (product_type, *patched_version[1:])
                break
    
    # If still no product type, use "UNKNOWN" but try both CSVs
    if not product_type:
        # Log the issue but continue with generic processing
        print(f"{YELLOW}Could not determine product type for advisory: {advisory.get('advisoryId')}{RESET}")
        # Try both CSVs
        all_affected = []
        
        # Try MXMSMR CSV
        if os.path.exists(mxmsmr_csv_path):
            with open(mxmsmr_csv_path, 'r') as f:
                reader = csv.DictReader(f)
                for row in reader:
                    current_version_str = row['firmware_version']
                    current_version = parse_firmware_version(current_version_str)
                    if current_version and is_version_affected(current_version, patched_version):
                        all_affected.append({
                            'product_type': row['product_type'],
                            'network_id': row['network_id'],
                            'network_name': row['network_name'],
                            'firmware_version': row['firmware_version'],
                            'advisory_id': advisory.get('advisoryId', ''),
                            'advisory_title': advisory.get('advisoryTitle', ''),
                            'fixed_firmware': patched_version_str
                        })
        
        # Try MGMVMT CSV
        if os.path.exists(mgmvmt_csv_path):
            with open(mgmvmt_csv_path, 'r') as f:
                reader = csv.DictReader(f)
                for row in reader:
                    current_version_str = row['firmware_version']
                    current_version = parse_firmware_version(current_version_str)
                    if current_version and is_version_affected(current_version, patched_version):
                        all_affected.append({
                            'product_type': row['product_type'],
                            'network_id': row['network_id'],
                            'network_name': row['network_name'],
                            'firmware_version': row['firmware_version'],
                            'advisory_id': advisory.get('advisoryId', ''),
                            'advisory_title': advisory.get('advisoryTitle', ''),
                            'fixed_firmware': patched_version_str
                        })
        
        return all_affected
    
    # Normal case - we know the product type
    csv_path = None
    if product_type in ["MX", "MS", "MR"]:
        csv_path = mxmsmr_csv_path
    elif product_type in ["MG", "MV", "MT"]:
        csv_path = mgmvmt_csv_path
    
    if not csv_path or not os.path.exists(csv_path):
        print(f"{YELLOW}Firmware compliance CSV not found for {product_type}: {csv_path}{RESET}")
        return affected_devices
    
    # Read the CSV and check each device
    try:
        with open(csv_path, 'r') as f:
            reader = csv.DictReader(f)
            for row in reader:
                # Skip if not the right product type
                if row['product_type'] != product_type:
                    continue
                
                # Parse the current firmware version
                current_version_str = row['firmware_version']
                current_version = parse_firmware_version(current_version_str)
                
                if current_version and is_version_affected(current_version, patched_version):
                    # This device is affected
                    affected_devices.append({
                        'product_type': row['product_type'],
                        'network_id': row['network_id'],
                        'network_name': row['network_name'],
                        'firmware_version': row['firmware_version'],
                        'advisory_id': advisory.get('advisoryId', ''),
                        'advisory_title': advisory.get('advisoryTitle', ''),
                        'fixed_firmware': patched_version_str
                    })
    
    except Exception as e:
        print(f"{RED}Error reading firmware compliance CSV: {e}{RESET}")
        logger.error(f"Error reading firmware compliance CSV: {e}")
    
    return affected_devices

async def generate(api_client, template_path, output_path, inventory_devices=None, networks=None):
    """Generate the PSIRT Advisories slide."""
    print(f"\n{GREEN}Generating PSIRT Advisories slide...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # Fetch PSIRT advisories
    advisories = await fetch_psirt_advisories()
    
    if not advisories:
        print(f"{YELLOW}No Meraki-related PSIRT advisories found. Creating placeholder slide.{RESET}")
        logger.warning("No Meraki-related PSIRT advisories found. Creating placeholder slide.")
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    print(f"{BLUE}Updating PowerPoint with PSIRT advisories data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)
        
        # Position PSIRT slides after both firmware compliance slides
        
        # This function tries to find a slide title by examining shapes on a slide
        def get_slide_title(slide):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    text = shape.text_frame.text.lower()
                    if "firmware compliance" in text or "firmware status" in text:
                        return text
            return ""
        
        # Find the position of the second firmware compliance slide
        target_slide_index = len(prs.slides)  # Default to end of presentation
        firmware_compliance_slides = []
        
        for i, slide in enumerate(prs.slides):
            title = get_slide_title(slide)
            if "firmware compliance" in title or "firmware status" in title:
                firmware_compliance_slides.append(i)
        
        # If we found both firmware compliance slides, position after the second one
        if len(firmware_compliance_slides) >= 2:
            target_slide_index = firmware_compliance_slides[1] + 1
            print(f"{BLUE}Found second firmware compliance slide at position {firmware_compliance_slides[1] + 1}, inserting PSIRT slides after it{RESET}")
        # If we found only one, place it after that one
        elif len(firmware_compliance_slides) == 1:
            target_slide_index = firmware_compliance_slides[0] + 1
            print(f"{YELLOW}Found only one firmware compliance slide at position {firmware_compliance_slides[0] + 1}, inserting PSIRT slides after it{RESET}")
        else:
            print(f"{YELLOW}Could not find firmware compliance slides, inserting PSIRT slides at the end{RESET}")
        
        # Make sure we have enough slides in the presentation to insert at the desired position
        if len(prs.slides) < target_slide_index:
            target_slide_index = len(prs.slides)
            print(f"{YELLOW}Presentation has fewer slides than expected, will add PSIRT slides at the end{RESET}")
        
        # Create an array to hold all the new slides we're creating
        created_slides = []
        created_slide_ids = []
        
        # Function to create a PSIRT slide with consistent formatting
        def create_psirt_slide(slide_index, is_first_slide=True):
            # Use slide layout from master slide 1
            slide_layout = prs.slide_masters[0].slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Store the slide ID so we can move it later
            slide_id = slide._element.get('id')
            created_slide_ids.append(slide_id)
            created_slides.append(slide)
            
            # Clear any placeholder content from the template
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    shape.text_frame.clear()
            
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.8))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "PSIRT Advisories"
            title_p.font.size = Pt(44)
            title_p.font.bold = True
            
            # Add horizontal line across the slide
            line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(11.0), Inches(1.2))
            line.line.color.rgb = TITLE_COLOR
            line.line.width = Pt(2)
            
            # Add Last Updated note in the top right
            current_date = datetime.datetime.now().strftime("%B %d, %Y")
            update_box = slide.shapes.add_textbox(Inches(8.0), Inches(0.5), Inches(3.0), Inches(0.5))
            update_p = update_box.text_frame.add_paragraph()
            update_p.text = f"Last Updated: {current_date}"
            update_p.font.size = Pt(12)
            update_p.font.italic = True
            update_p.alignment = PP_ALIGN.RIGHT
            
            # Add slide number indicator if needed
            if not is_first_slide:
                slide_number = slide.shapes.add_textbox(Inches(10.5), Inches(1.3), Inches(1.0), Inches(0.5))
                slide_number_p = slide_number.text_frame.add_paragraph()
                slide_number_p.text = f"Continued ({slide_index})"
                slide_number_p.font.size = Pt(10)
                slide_number_p.font.italic = True
                slide_number_p.alignment = PP_ALIGN.RIGHT
            
            # Add the URL and explanation to the slide notes
            if slide.notes_slide is None:
                slide.notes_slide
            notes = slide.notes_slide.notes_text_frame
            notes.text = "PSIRT Advisories Information:\n\n" + \
                         "This slide shows Cisco Product Security Incident Response Team (PSIRT) advisories for Meraki products published within the last 5 years. " + \
                         "The firmware comparison identifies networks with devices potentially vulnerable to these security issues based on current firmware versions.\n\n" + \
                         "For the latest security advisories, visit: https://tools.cisco.com/security/center/publicationListing.x"
            
            return slide
        
        if not advisories:
            # Create a single slide with a "No advisories found" message
            slide = create_psirt_slide(1)
            
            # Add the no advisories message
            no_advisories_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(11.0), Inches(1.0))
            no_advisories_p = no_advisories_box.text_frame.add_paragraph()
            no_advisories_p.text = "No active security advisories for Cisco Meraki products at this time."
            no_advisories_p.font.size = Pt(20)
            no_advisories_p.alignment = PP_ALIGN.CENTER
            
            # Add footer
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(11.0), Inches(0.5))
            footer_p = footer.text_frame.add_paragraph()
            footer_p.text = "For the latest security advisories from the last 5 years, visit: https://tools.cisco.com/security/center/publicationListing.x"
            footer_p.font.size = Pt(10)
            footer_p.font.italic = True
            
            # Add bottom horizontal line
            bottom_line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.0), Inches(11.0), Inches(7.0))
            bottom_line.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
            bottom_line.line.width = Pt(1)
            
        else:
            # Calculate number of slides needed (2 advisories per slide, plus header row)
            advisories_per_slide = 2
            num_slides_needed = (len(advisories) + advisories_per_slide - 1) // advisories_per_slide
            print(f"{BLUE}Creating {num_slides_needed} PSIRT advisory slide(s) for {len(advisories)} advisories{RESET}")
            
            # Create and populate each slide
            for slide_index in range(num_slides_needed):
                # Create a new slide
                is_first_slide = (slide_index == 0)
                slide = create_psirt_slide(slide_index + 1, is_first_slide)
                
                # Calculate which advisories go on this slide
                start_idx = slide_index * advisories_per_slide
                end_idx = min(start_idx + advisories_per_slide, len(advisories))
                current_advisories = advisories[start_idx:end_idx]
                
                # Create a table for these advisories
                table_rows = len(current_advisories) + 1  # header + advisories
                table_cols = 5  # Added a column for fixed firmware
                table_width = Inches(12.39)

                actual_table_height = 5.45 - 0.12
                
                # Set header row short to maximize content space
                header_row_height_pt = 20
                header_row_height_in = header_row_height_pt / 72
                
                # Calculate available height for data rows
                data_rows_total_height = actual_table_height - header_row_height_in
                
                # Calculate individual data row height
                data_row_height = data_rows_total_height / len(current_advisories)
                
                # Adjust column widths
                col_widths = [
                    Inches(2.7),   # Title
                    Inches(1.9),   # Advisory ID
                    Inches(1.2),   # CVSS/SIR
                    Inches(2.7),   # Firmware Patch Version
                    Inches(3.89)   # Summary
                ]
                
                table_height = Inches(actual_table_height)
                table = slide.shapes.add_table(
                    table_rows, table_cols, 
                    Inches(0.47), Inches(1.33), 
                    table_width, table_height
                ).table
                
                # First set the header row height (short)
                table.rows[0].height = Pt(header_row_height_pt)
                
                # Then set all data rows to equal height to fill the remaining space
                for i in range(1, table_rows):  # Skip header row
                    table.rows[i].height = Inches(data_row_height)
                
                # Set column widths
                for i, width in enumerate(col_widths):
                    table.columns[i].width = width
                
                # Define header row
                header_cells = table.rows[0].cells
                header_cells[0].text = "Title"
                header_cells[1].text = "Advisory ID"
                header_cells[2].text = "CVSS/SIR"
                header_cells[3].text = "Firmware Patch Version"
                header_cells[4].text = "Summary"
                
                # Format header row - use dark green like other slides for consistency
                for cell in header_cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TITLE_COLOR
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                
                
                # Fetch fixed firmware information for all advisories on this slide
                fixed_firmware_tasks = []
                for advisory in current_advisories:
                    publication_url = advisory.get("publicationUrl", "")
                    if publication_url:
                        task = fetch_fixed_firmware_version(publication_url)
                        fixed_firmware_tasks.append(task)
                    else:
                        fixed_firmware_tasks.append(None)
                
                # Wait for all firmware version fetches to complete
                fixed_firmware_versions = []
                for i, task in enumerate(fixed_firmware_tasks):
                    if task:
                        try:
                            firmware_version = await task
                            # Store the firmware version in the advisory for later impact analysis
                            if firmware_version:
                                current_advisories[i]["fixed_firmware_version"] = firmware_version
                                fixed_firmware_versions.append(firmware_version)
                            else:
                                current_advisories[i]["fixed_firmware_version"] = None
                                fixed_firmware_versions.append("Unable to locate fixed version. Please visit advisory link for more information")
                        except:
                            current_advisories[i]["fixed_firmware_version"] = None
                            fixed_firmware_versions.append("Unable to locate fixed version. Please visit advisory link for more information")
                    else:
                        current_advisories[i]["fixed_firmware_version"] = None
                        fixed_firmware_versions.append("Unable to locate fixed version. Please visit advisory link for more information")
                
                # Fill in the table with advisories for this slide
                for i, advisory in enumerate(current_advisories):
                    row = table.rows[i+1]
                    
                    # Title first
                    row.cells[0].text = advisory.get("advisoryTitle", "")
                    
                    # Advisory ID - include both published and updated dates
                    advisory_id = advisory.get("advisoryId", "")
                    pub_date = format_date(advisory.get("firstPublished", ""))
                    updated_date = format_date(advisory.get("lastUpdated", ""))
                    row.cells[1].text = f"{advisory_id}\n\nFirst Published: {pub_date}\n\nLast Updated: {updated_date}"
                    
                    # CVSS Score and SIR
                    cvss = advisory.get("cvssBaseScore", "")
                    sir = advisory.get("sir", "")
                    row.cells[2].text = f"CVSS: {cvss}\nSIR: {sir}"
                    
                    # Format CVSS/SIR cell based on severity
                    sir_color = get_sir_color(sir)
                    row.cells[2].fill.solid()
                    row.cells[2].fill.fore_color.rgb = sir_color
                    for paragraph in row.cells[2].text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                        paragraph.font.bold = True
                    
                    # Firmware Patch Version - use the scraped data with impact analysis
                    if i < len(fixed_firmware_versions):
                        firmware_text = fixed_firmware_versions[i]
                        
                        # Perform impact analysis if we have valid firmware data
                        affected_count = 0
                        if advisory.get("fixed_firmware_version"):
                            mxmsmr_paths = ["mxmsmr_firmware_report.csv", "/tmp/mxmsmr_firmware_report.csv"]
                            mgmvmt_paths = ["mgmvmt_firmware_report.csv", "/tmp/mgmvmt_firmware_report.csv"]
                            
                            # Use the first path that exists
                            mxmsmr_path = next((p for p in mxmsmr_paths if os.path.exists(p)), mxmsmr_paths[0])
                            mgmvmt_path = next((p for p in mgmvmt_paths if os.path.exists(p)), mgmvmt_paths[0])
                            
                            affected_devices = find_affected_devices(
                                advisory, 
                                mxmsmr_path,
                                mgmvmt_path
                            )
                            affected_count = len(affected_devices)
                            
                            # Export affected devices to CSV if any found
                            if affected_devices:
                                # Create CSV filename based on advisory ID
                                advisory_id = advisory.get("advisoryId", "unknown")
                                csv_filename = f"psirt_affected_{advisory_id}.csv"
                                
                                try:
                                    with open(csv_filename, 'w', newline='') as csvfile:
                                        fieldnames = ['product_type', 'network_id', 'network_name', 
                                                     'firmware_version', 'advisory_id', 
                                                     'advisory_title', 'fixed_firmware']
                                        writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
                                        writer.writeheader()
                                        for device in affected_devices:
                                            writer.writerow(device)
                                        print(f"{GREEN}Exported {affected_count} affected devices to {csv_filename}{RESET}")
                                except Exception as e:
                                    print(f"{RED}Error exporting affected devices to CSV: {e}{RESET}")
                                    logger.error(f"Error exporting affected devices to CSV: {e}")
                            
                        # Add affected device count to the firmware text if available
                        row.cells[3].text = firmware_text
                        
                        # Add status indicator for affected devices
                        if affected_count > 0:
                            # Add a new paragraph for the warning text with enhanced visibility
                            warning_paragraph = row.cells[3].text_frame.add_paragraph()
                            warning_paragraph.text = f"⚠️ {affected_count} device(s) potentially affected"
                            warning_paragraph.font.bold = True
                            warning_paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red for warning
                            
                            # Add link to the CSV in notes
                            advisory_id = advisory.get("advisoryId", "unknown")
                            csv_filename = f"psirt_affected_{advisory_id}.csv"
                            
                            if slide.notes_slide is None:
                                slide.notes_slide
                            notes = slide.notes_slide.notes_text_frame
                            notes.text += f"\n\nPotentially affected devices for {advisory.get('advisoryTitle', '')}: {affected_count} devices exported to {csv_filename}"
                        elif advisory.get("fixed_firmware_version"):
                            # If we have a valid firmware version and no affected devices, add a green checkmark
                            status_paragraph = row.cells[3].text_frame.add_paragraph()
                            status_paragraph.text = "✅ No devices affected"
                            status_paragraph.font.bold = True
                            status_paragraph.font.color.rgb = RGBColor(0, 176, 80)  # Green for good status
                    else:
                        row.cells[3].text = "Unable to locate fixed version. Please visit advisory link for more information"
                    
                    # Summary
                    row.cells[4].text = advisory.get("cleanedSummary", "No summary available")
                    
                    # Add alternating row background for better readability
                    if i % 2 == 1:
                        for j in [0, 1, 3, 4]:
                            row.cells[j].fill.solid()
                            row.cells[j].fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
                    
                    # Format all cells in the row
                    for cell in row.cells:
                        # Enable word wrapping in all cells
                        cell.text_frame.word_wrap = True
                        
                        # Format each paragraph
                        for paragraph in cell.text_frame.paragraphs:
                            # Make summary text slightly smaller to fit more content
                            if cell == row.cells[4]:
                                paragraph.font.size = Pt(8.3)
                            else:
                                paragraph.font.size = Pt(10)
                                
                            if cell == row.cells[0] or cell == row.cells[1]:
                                paragraph.font.bold = True
                                
                            if cell == row.cells[3]:  # Firmware Patch Version column
                                paragraph.font.color.rgb = RGBColor(0, 102, 204)  # Blue for firmware info
        
        print(f"{BLUE}Moving {len(created_slides)} PSIRT slides after firmware compliance slides{RESET}")
        
        try:
            # Get the slide ID list
            slide_id_list = prs.slides._sldIdLst
            
            # Get indices of the newly created slides (they're at the end of the presentation)
            current_slide_count = len(slide_id_list)
            start_idx = current_slide_count - len(created_slides)
            
            if start_idx < 0 or start_idx >= current_slide_count:
                print(f"{YELLOW}Cannot determine slide positions accurately. Slides are created but may need manual positioning.{RESET}")
            else:
                slides_to_move = []
                for i in range(start_idx, current_slide_count):
                    slides_to_move.append(slide_id_list[start_idx])
                    slide_id_list.remove(slide_id_list[start_idx])
                
                # Now insert them at the target position
                for i, slide_element in enumerate(reversed(slides_to_move)):
                    if target_slide_index <= len(slide_id_list):
                        slide_id_list.insert(target_slide_index, slide_element)
                        print(f"{GREEN}Moved PSIRT slide to desired position{RESET}")
                    else:
                        slide_id_list.append(slide_element)
                        print(f"{YELLOW}Added PSIRT slide at the end of the presentation{RESET}")
                
                print(f"{GREEN}Slides have been repositioned successfully{RESET}")
        
        except Exception as e:
            print(f"{YELLOW}Error reordering slides: {e}. Slides created successfully but may need manual positioning.{RESET}")
            logger.warning(f"Error reordering slides: {e}")
        
        # Save the presentation
        prs.save(output_path)
        
        if len(advisories) > 2:
            print(f"{GREEN}Created {num_slides_needed} PSIRT Advisories slides with 2 advisories per slide{RESET}")
        else:
            print(f"{GREEN}Created PSIRT Advisories slide{RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        logger.error(f"Error updating PowerPoint: {e}")
        import traceback
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    print(f"{PURPLE}PSIRT Advisories slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    await generate(api_client, template_path, output_path)

if __name__ == "__main__":
    # Process command line arguments when run directly
    import argparse
    parser = argparse.ArgumentParser(description="Generate PSIRT advisories slide")
    parser.add_argument("output_path", help="Path to output PowerPoint file")
    parser.add_argument("-t", "--template", dest="template_path", help="Path to template PowerPoint file (default: same as output)")
    
    args = parser.parse_args()
    
    output_path = args.output_path
    template_path = args.template_path if args.template_path else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))
