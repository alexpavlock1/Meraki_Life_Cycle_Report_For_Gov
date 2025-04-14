import os
import sys
import asyncio
import time
import datetime
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Colors for the executive summary
GOOD_COLOR = RGBColor(0, 176, 80)       # Green
WARNING_COLOR = RGBColor(255, 192, 0)    # Amber/Yellow  
CRITICAL_COLOR = RGBColor(255, 0, 0)     # Red
ACCENT_COLOR = RGBColor(0, 112, 192)     # Blue
NEUTRAL_COLOR = RGBColor(89, 89, 89)     # Dark gray

def add_decorative_bubbles(slide):
    """
    Add three overlapping decorative bubbles to the top right corner of the slide
    for visual enhancement.
    """

    bubble_size = Inches(0.9)
    
    base_x = Inches(12.06)
    base_y = Inches(0.43)
    

    bubble_colors = [
        ACCENT_COLOR,         # Blue
        RGBColor(131, 206, 114),  # Light green (derived from GOOD_COLOR)
        RGBColor(255, 217, 102)   # Light yellow (derived from WARNING_COLOR)
    ]

    offsets = [
        (Inches(0.2), Inches(0.15)),
        (Inches(-0.5), Inches(0.15)),
        (Inches(-0.25), Inches(-0.35))
    ]
    
    bubbles = []
    for i, (offset_x, offset_y) in enumerate(offsets):
        bubble = slide.shapes.add_shape(
            9,  # OVAL shape
            base_x + offset_x, 
            base_y + offset_y, 
            bubble_size, 
            bubble_size
        )
        
        # Style the bubble
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = bubble_colors[i]
        
        # Make the bubble semi-transparent
        bubble.fill.transparency = 0.3
        
        # No outline
        bubble.line.fill.background()
        
        bubbles.append(bubble)
    
    return bubbles

def add_notes_to_slide(slide, notes_text):
    """
    Add notes to a PowerPoint slide that will only be visible to the presenter.
    """
    # Get or create the notes slide
    notes_slide = slide.notes_slide
    
    # Access the text frame of the notes slide
    notes_text_frame = notes_slide.notes_text_frame
    
    # Clear any existing notes (optional)
    for idx in range(len(notes_text_frame.paragraphs)):
        if idx > 0:  # Keep the first paragraph and modify it
            p = notes_text_frame.paragraphs[idx]
            p.text = ""
    
    if len(notes_text_frame.paragraphs) > 0:
        # Modify the first paragraph if it exists
        notes_text_frame.paragraphs[0].text = notes_text
    else:
        # Add a new paragraph if none exists
        p = notes_text_frame.add_paragraph()
        p.text = notes_text

# Status colors mapped to health levels
def get_status_color(status):
    if status in ['good', 'Good', 'Excellent']:
        return GOOD_COLOR
    elif status in ['warning', 'Warning', 'Fair', 'Satisfactory']:
        return WARNING_COLOR
    elif status in ['critical', 'Critical', 'Needs Attention', 'Critical Issues']:
        return CRITICAL_COLOR
    else:
        return NEUTRAL_COLOR

# Helper functions
def clean_slide(slide):
    """Remove everything from a slide except branding elements at the bottom."""
    # Identify shapes to remove (everything except bottom branding)
    shapes_to_remove = []
    for shape in slide.shapes:
        # Preserve shapes that appear to be logos/branding (bottom of slide)
        if shape.top >= Inches(6.5):
            continue
        # Keep bottom text
        if hasattr(shape, "text_frame") and shape.top >= Inches(6.5):
            continue
        # Remove everything else
        shapes_to_remove.append(shape)
    
    # Remove the identified shapes
    for shape in shapes_to_remove:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
        except Exception as e:
            print(f"{YELLOW}Could not remove shape: {e}{RESET}")

def move_slide_to_position(prs, old_index, new_index):
    """Move a slide from one position to another in the presentation."""
    if old_index == new_index:
        return
    
    slides = prs.slides._sldIdLst
    slide = slides[old_index]
    slides.remove(slide)
    slides.insert(new_index, slide)
def format_recommendation_text(text, max_chars_per_line=85):
    """
    Format recommendation text with appropriate line breaks to prevent text from running off slides.
    """
    # If text is already short enough, return as is
    if len(text) <= max_chars_per_line:
        return text
    
    # Find logical break points (after commas, periods, and certain conjunctions)
    words = text.split()
    current_line = ""
    formatted_text = ""
    
    # Go through each word
    for word in words:
        # Check if adding this word would exceed the line length
        if len(current_line) + len(word) + 1 > max_chars_per_line:
            # Add current line to formatted text
            formatted_text += current_line.strip() + "\n"
            current_line = word + " "
        else:
            # Add word to current line
            current_line += word + " "
    
    # Add any remaining text
    if current_line:
        formatted_text += current_line.strip()
    
    return formatted_text

def get_base_model(model):
    """Extract the base model prefix."""
    if not model:
        return None
    
    import re
    # Look for standard Meraki model patterns including Catalyst Wireless
    match = re.match(r'(MR\d+|MS\d+|MX\d+|MV\d+|MG\d+|MT\d+|Z\d+|CW\d+)', model)
    if match:
        return match.group(1)
    return model

def add_recommendations_to_slide(slide, recommendations, x, y, width, height):
    """
    Add prioritized recommendations to the slide, ensuring health-related recommendations
    appear first, followed by general and product recommendations.
    """
    # Recommendations title
    rec_title_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
    rec_title_p = rec_title_box.text_frame.add_paragraph()
    rec_title_p.text = "Recommendations"
    rec_title_p.font.size = Pt(12)
    rec_title_p.font.bold = True
    
    # Add Recommendations content
    rec_box = slide.shapes.add_textbox(x, y + Inches(0.3), width, height - Inches(0.3))
    rec_tf = rec_box.text_frame
    
    # Add recommendations (they're already prioritized from generate_insights)
    for i, recommendation in enumerate(recommendations[:5]):  # Limit to 5 recommendations
        rec_p = rec_tf.add_paragraph()
        
        # Identify if this is a NON-product adoption recommendation
        is_health_rec = any(term in recommendation.lower() for term in 
              ['firmware', 'upgrade', 'replacement', 'end-of-support', 'eol', 
               'critical', 'device health', 'lifecycle', 'health improvement',
               'implement', 'update', 'refresh', 'hardware', 'end-of-sale'])
               
        # Identify if this is a product adoption recommendation for advanced products
        is_advanced_product_adoption = any(product in recommendation for product in 
            ['Secure Connect', 'Umbrella', 'ThousandEyes', 'Spaces', 'XDR'])
            
        # Identify if this is a core product adoption recommendation (MX, MS, MR)
        is_core_product_recommendation = any(product in recommendation for product in 
            ['Deploy Meraki MX', 'Add Meraki MS', 'Implement Meraki wireless', 'Complete your Meraki deployment'])
        
        # Ensure text wrapping is enabled for this paragraph
        rec_tf.word_wrap = True
        rec_tf.auto_size = MSO_ANCHOR.TOP  # Anchor text at top and allow it to expand downward

        formatted_text = format_recommendation_text(recommendation)
        
        # Apply blue color to health recommendations and core product recommendations
        if (is_health_rec and not is_advanced_product_adoption) or is_core_product_recommendation:
            rec_p.text = f"• {formatted_text}"
            rec_p.font.bold = True  # Make important recommendations bold
            rec_p.font.color.rgb = ACCENT_COLOR  # Highlight important recommendations (blue)
        else:
            # Other recommendations in black (default)
            rec_p.text = f"• {formatted_text}"
        
        rec_p.font.size = Pt(10)
        rec_p.space_after = Pt(3)
    
    return rec_title_box, rec_box

def check_firmware_version_status(model_prefix, firmware_version, firmware_stats):
    """
    Check the status of a firmware version based on firmware stats.
    """
    if not firmware_stats or model_prefix not in firmware_stats:
        return 'unknown'
    
    # Get the latest firmware version
    latest_firmware = firmware_stats[model_prefix].get('latest', '')
    
    # If firmware version matches latest, it's good
    if firmware_version == latest_firmware:
        return 'good'
    
    # Check if firmware is in critical list
    critical_percentage = firmware_stats[model_prefix].get('Critical', 0) / max(firmware_stats[model_prefix].get('Total', 1), 1) * 100
    
    # If high percentage of critical firmware, classify as critical
    if critical_percentage >= 75:
        return 'critical'
    elif critical_percentage >= 25:
        return 'warning'
        
    # Otherwise, just a warning as it's not the latest
    return 'warning'

def categorize_device_health(device, firmware_stats, eol_data):
    """
    Categorize a device based on firmware compliance and EOL status.
    accounts for critical firmware status and weighs EOL more heavily
    """
    model = device.get('model', '')
    if not model:
        return 'unknown'
    
    # Initialize status flags
    is_eol = False
    approaching_eol = False
    end_of_sale_reached = False
    firmware_status = 'good'
    
    # Check firmware compliance
    if firmware_stats:
        model_prefix = model[:2].upper()
        if model_prefix in ['MX', 'MS', 'MR', 'MV', 'MG', 'MT'] and model_prefix in firmware_stats:
            device_firmware = device.get('firmware', '')
            if device_firmware:
                firmware_status = check_firmware_version_status(model_prefix, device_firmware, firmware_stats)
    
    # Check if device model appears in EOL data
    if eol_data:
        base_model = get_base_model(model)
        for eol_model in eol_data:
            if base_model and base_model.startswith(eol_model):
                # Check End of Support status
                eol_date = eol_data[eol_model].get('end_of_support')
                if eol_date:
                    try:
                        # Convert date string to datetime object
                        eol_datetime = datetime.strptime(eol_date, "%b %d, %Y")
                        current_date = datetime.now()
                        
                        # Calculate time difference in years
                        time_diff = (eol_datetime - current_date).days / 365.25
                        
                        if time_diff <= 0:
                            is_eol = True
                        elif time_diff <= 1:
                            approaching_eol = True
                    except:
                        pass
                
                # Check End of Sale status
                eos_date = eol_data[eol_model].get('end_of_sale')
                if eos_date:
                    try:
                        from datetime import datetime
                        eos_datetime = datetime.strptime(eos_date, "%b %d, %Y")
                        current_date = datetime.now()
                        
                        # Check if we've passed the end of sale date
                        if current_date > eos_datetime:
                            end_of_sale_reached = True
                    except:
                        pass
                break
    
    # Determine overall health
    if is_eol or firmware_status == 'critical':
        return 'critical'
    elif approaching_eol or end_of_sale_reached or firmware_status == 'warning':
        return 'warning'
    else:
        return 'good'
def add_score_deduction_explanation(slide, deduction_reasons):
    """
    Add explanations for score deductions to the slide using exact positioning.
    """
    if not deduction_reasons:
        return None  # No deductions to explain
    
    # Fixed positions and dimensions as specified
    container_x = Inches(0.45)          # Left position
    container_y = Inches(5.43)          # Top position
    container_width = Inches(3.8)       # Width
    container_height = Inches(1.24)     # Height
    
    # Create container for the deduction explanations
    container = slide.shapes.add_shape(1, container_x, container_y, container_width, container_height)
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Very light gray
    container.line.fill.solid()
    container.line.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray
    
    # Add title at exact position
    title_x = Inches(0.5)               # Left position: 0.5"
    title_y = Inches(5.18)              # Top position: 5.18"
    title_width = Inches(3.7)
    
    title_box = slide.shapes.add_textbox(title_x, title_y, title_width, Inches(0.25))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = "Score Deductions"
    title_p.font.size = Pt(10)
    title_p.font.bold = True
    
    # Add explanations at exact position
    explanation_x = Inches(0.5)         # Left position: 0.5"
    explanation_y = Inches(5.43)        # Top position: 5.43"
    explanation_width = Inches(3.6)
    explanation_height = Inches(1.2)
    
    explanation_box = slide.shapes.add_textbox(explanation_x, explanation_y, explanation_width, explanation_height)
    explanation_tf = explanation_box.text_frame
    explanation_tf.word_wrap = True
    
    for i, reason in enumerate(deduction_reasons[:5]):  # Limit to top 5 deductions to avoid overcrowding
        reason_p = explanation_tf.add_paragraph()
        reason_p.text = reason
        reason_p.font.size = Pt(8)
        reason_p.space_after = Pt(2)
        
    return container

def calculate_health_score(inventory_devices, firmware_stats, eol_data, dashboard_stats, products):
    """
    Calculate an overall network health score (0-100) using percentage-based thresholds.
    """
    score = 100  # Start with perfect score and deduct points
    deduction_reasons = []  # Track reasons for score deductions
    
    # Track device health categories
    device_health = {'good': 0, 'warning': 0, 'critical': 0, 'unknown': 0}
    
    # Track EOL categories specifically
    eol_status = {'end_of_support': 0, 'end_of_sale': 0, 'approaching_eol': 0, 'current': 0}
    
    # Analyze devices if available
    if inventory_devices:
        total_devices = len(inventory_devices)
        
        # Categorize each device
        for device in inventory_devices:
            health = categorize_device_health(device, firmware_stats, eol_data)
            device_health[health] += 1
            
            # Also track EOL status specifically
            model = device.get('model', '')
            if eol_data and model:
                base_model = get_base_model(model)
                if base_model:
                    for eol_model in eol_data:
                        if base_model.startswith(eol_model):
                            # Check End of Support
                            eol_date = eol_data[eol_model].get('end_of_support')
                            if eol_date:
                                try:
                                    eol_datetime = datetime.strptime(eol_date, "%b %d, %Y")
                                    current_date = datetime.now()
                                    
                                    if current_date > eol_datetime:
                                        eol_status['end_of_support'] += 1
                                        break
                                    
                                    # Check if approaching EOL
                                    time_diff = (eol_datetime - current_date).days / 365.25
                                    if time_diff <= 1:
                                        eol_status['approaching_eol'] += 1
                                        break
                                except:
                                    pass
                            
                            # Check End of Sale
                            eos_date = eol_data[eol_model].get('end_of_sale')
                            if eos_date:
                                try:
                                    from datetime import datetime
                                    eos_datetime = datetime.strptime(eos_date, "%b %d, %Y")
                                    current_date = datetime.now()
                                    
                                    if current_date > eos_datetime:
                                        eol_status['end_of_sale'] += 1
                                        break
                                except:
                                    pass
                            
                            eol_status['current'] += 1
                            break
        
        # Calculate percentage of devices in each category
        if total_devices > 0:
            # Calculate percentages
            critical_pct = (device_health['critical'] / total_devices) * 100
            warning_pct = (device_health['warning'] / total_devices) * 100
            eos_pct = (eol_status['end_of_support'] / total_devices) * 100
            eosale_pct = (eol_status['end_of_sale'] / total_devices) * 100
            
            # Use threshold-based penalties instead of linear scaling
            
            # Penalty for critical devices
            if critical_pct >= 25:
                deduction = 30  # Severe penalty for 25%+ critical devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_pct:.1f}% of devices have critical health status")
            elif critical_pct >= 15:
                deduction = 20  # Major penalty for 15-25% critical devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_pct:.1f}% of devices have critical health status")
            elif critical_pct >= 5:
                deduction = 10  # Moderate penalty for 5-15% critical devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_pct:.1f}% of devices have critical health status")
            elif critical_pct > 0:
                deduction = 5   # Minor penalty for >0-5% critical devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_pct:.1f}% of devices have critical health status")
            
            # Penalty for warning devices
            if warning_pct >= 40:
                deduction = 15  # Major penalty for 40%+ warning devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {warning_pct:.1f}% of devices have warning health status")
            elif warning_pct >= 25:
                deduction = 10  # Moderate penalty for 25-40% warning devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {warning_pct:.1f}% of devices have warning health status")
            elif warning_pct >= 10:
                deduction = 5   # Minor penalty for 10-25% warning devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {warning_pct:.1f}% of devices have warning health status")
            
            # Penalty for End of Support devices
            if eos_pct >= 20:
                deduction = 25  # Severe penalty for 20%+ EOS devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {eos_pct:.1f}% of devices have reached end-of-support")
            elif eos_pct >= 10:
                deduction = 15  # Major penalty for 10-20% EOS devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {eos_pct:.1f}% of devices have reached end-of-support")
            elif eos_pct >= 5:
                deduction = 10  # Moderate penalty for 5-10% EOS devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {eos_pct:.1f}% of devices have reached end-of-support")
            elif eos_pct > 0:
                deduction = 5   # Minor penalty for >0-5% EOS devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {eos_pct:.1f}% of devices have reached end-of-support")
            
            # Penalty for End of Sale devices
            if eosale_pct >= 30:
                deduction = 10  # Moderate penalty for 30%+ EOS devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {eosale_pct:.1f}% of devices have reached end-of-sale")
            elif eosale_pct >= 15:
                deduction = 5   # Minor penalty for 15-30% EOS devices
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {eosale_pct:.1f}% of devices have reached end-of-sale")
    
    # Check firmware criticality with threshold-based penalties
    if firmware_stats:
        total_fw_devices = 0
        critical_fw_devices = 0
        
        for device_type in ['MX', 'MS', 'MR', 'MV', 'MG', 'MT']:
            if device_type in firmware_stats:
                total_fw_devices += firmware_stats[device_type].get('Total', 0)
                critical_fw_devices += firmware_stats[device_type].get('Critical', 0)
        
        if total_fw_devices > 0:
            critical_fw_pct = (critical_fw_devices / total_fw_devices) * 100
            
            # Penalty for critical firmware
            if critical_fw_pct >= 25:
                deduction = 20  # Major penalty for 25%+ critical firmware
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_fw_pct:.1f}% of devices on critical firmware versions")
            elif critical_fw_pct >= 15:
                deduction = 15  # Moderate penalty for 15-25% critical firmware
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_fw_pct:.1f}% of devices on critical firmware versions")
            elif critical_fw_pct >= 5:
                deduction = 10  # Minor penalty for 5-15% critical firmware
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_fw_pct:.1f}% of devices on critical firmware versions")
            elif critical_fw_pct > 0:
                deduction = 5   # Minimal penalty for >0-5% critical firmware
                score -= deduction
                deduction_reasons.append(f"-{deduction} points: {critical_fw_pct:.1f}% of devices on critical firmware versions")
    
    # Check for missing key products
    if products:
        missing_core_products = [p for p in ['MX', 'MS', 'MR'] if not products.get(p, False)]
        missing_core_count = len(missing_core_products)
        
        missing_advanced_products = [p for p in ['Secure Connect', 'Umbrella Secure Internet Gateway'] 
                                      if not products.get(p, False)]
        missing_advanced_count = len(missing_advanced_products)
        
        # Deduct points for missing products
        if missing_core_count > 0:
            deduction = missing_core_count * 5
            score -= deduction
            deduction_reasons.append(f"-{deduction} points: Missing core products ({', '.join(missing_core_products)})")
        
        if missing_advanced_count > 0:
            deduction = missing_advanced_count * 2
            score -= deduction
            deduction_reasons.append(f"-{deduction} points: Missing advanced products ({', '.join(missing_advanced_products)})")
    
    # Client density check has been removed as requested
    
    # Ensure the score stays within 0-100 range
    final_score = max(0, min(100, round(score)))
    
    # Return both the score and the reasons for deductions
    return final_score, deduction_reasons

def get_health_assessment(score):
    """Convert a numerical health score to a text assessment."""
    if score >= 90:
        return "Excellent"
    elif score >= 80:
        return "Good"
    elif score >= 70:
        return "Satisfactory"
    elif score >= 60:
        return "Fair"
    elif score >= 40:
        return "Needs Attention"
    else:
        return "Critical Issues"

def is_end_of_sale_critical(device, eol_data):
    """
    Check if a device has reached end of sale.
    """
    if not eol_data:
        return False
        
    model = device.get('model', '')
    base_model = get_base_model(model)
    
    if not base_model:
        return False
        
    for eol_model in eol_data:
        if base_model.startswith(eol_model):
            eos_date = eol_data[eol_model].get('end_of_sale')
            if eos_date:
                try:
                    from datetime import datetime
                    eos_datetime = datetime.strptime(eos_date, "%b %d, %Y")
                    current_date = datetime.now()
                    
                    # Check if we've passed the end of sale date
                    if current_date > eos_datetime:
                        return True
                except:
                    pass
    
    return False

def is_end_of_support_critical(device, eol_data):
    """
    Check if a device has reached end of support or is within 1 year of end of support.
    """
    if not eol_data:
        return False
        
    model = device.get('model', '')
    base_model = get_base_model(model)
    
    if not base_model:
        return False
        
    for eol_model in eol_data:
        if base_model.startswith(eol_model):
            eol_date = eol_data[eol_model].get('end_of_support')
            if eol_date:
                try:
                    from datetime import datetime
                    eol_datetime = datetime.strptime(eol_date, "%b %d, %Y")
                    current_date = datetime.now()
                    
                    # Calculate time difference in years
                    time_diff = (eol_datetime - current_date).days / 365.25
                    
                    # Critical if already EOL or within 1 year
                    if time_diff <= 1:
                        return True
                except:
                    pass
    
    return False

def is_end_of_support_warning(device, eol_data):
    """
    Check if a device is within 1-2 years of end of support.
    """
    if not eol_data:
        return False
        
    model = device.get('model', '')
    base_model = get_base_model(model)
    
    if not base_model:
        return False
        
    for eol_model in eol_data:
        if base_model.startswith(eol_model):
            eol_date = eol_data[eol_model].get('end_of_support')
            if eol_date:
                try:
                    from datetime import datetime
                    eol_datetime = datetime.strptime(eol_date, "%b %d, %Y")
                    current_date = datetime.now()
                    
                    # Calculate time difference in years
                    time_diff = (eol_datetime - current_date).days / 365.25
                    
                    # Warning if within 1-2 years
                    if 1 < time_diff <= 2:
                        return True
                except:
                    pass
    
    return False

def generate_insights(inventory_devices, dashboard_stats, firmware_stats, eol_data, products, health_score=None):
    """Generate natural language insights based on the available data with prioritized recommendations."""
    insights = []
    health_recommendations = []  # Specifically for device health and lifecycle
    product_recommendations = []  # For product adoption
    general_recommendations = []  # For other recommendations

    # Comprehensive EOL and Firmware Status Summary
    if inventory_devices:
        total_devices = len(inventory_devices)
        
        # Count devices by health status for overall health assessment
        health_counts = {'good': 0, 'warning': 0, 'critical': 0, 'unknown': 0}
        
        # Count devices by EOL status specifically
        eol_status = {
            'past_eol': 0,      # Already past end-of-support
            'nearing_eol': 0,   # Within 12 months of end-of-support
            'approaching_eol': 0, # 12-24 months from end-of-support
            'current': 0        # More than 24 months from end-of-support
        }
        
        # Process each device for health and EOL status
        for device in inventory_devices:
            # Categorize overall health
            health = categorize_device_health(device, firmware_stats, eol_data)
            health_counts[health] += 1
            
            # Categorize EOL status specifically if EOL data is available
            if eol_data:
                model = device.get('model', '')
                base_model = get_base_model(model)
                
                if base_model:
                    eol_status_found = False
                    for eol_model in eol_data:
                        if base_model.startswith(eol_model):
                            eol_date = eol_data[eol_model].get('end_of_support')
                            if eol_date:
                                try:
                                    # Convert date string to datetime object
                                    eol_datetime = datetime.strptime(eol_date, "%b %d, %Y")
                                    current_date = datetime.now()
                                    
                                    # Calculate time difference in months
                                    time_diff = (eol_datetime - current_date).days / 30
                                    
                                    if time_diff <= 0:
                                        eol_status['past_eol'] += 1
                                    elif time_diff <= 12:
                                        eol_status['nearing_eol'] += 1
                                    elif time_diff <= 24:
                                        eol_status['approaching_eol'] += 1
                                    else:
                                        eol_status['current'] += 1
                                    
                                    eol_status_found = True
                                    break
                                except:
                                    pass
                    
                    # If no EOL status found, consider it current
                    if not eol_status_found:
                        eol_status['current'] += 1
                else:
                    # Unknown model defaults to current
                    eol_status['current'] += 1
        
        # Calculate percentages for reporting
        good_percentage = (health_counts['good'] / total_devices) * 100 if total_devices > 0 else 0
        warning_percentage = (health_counts['warning'] / total_devices) * 100 if total_devices > 0 else 0
        critical_percentage = (health_counts['critical'] / total_devices) * 100 if total_devices > 0 else 0
        
        # INSIGHT 1: Overall Device Health Summary
        if good_percentage >= 85:
            insights.append(f"Overall device health is excellent with {good_percentage:.1f}% of devices in good status.")
        elif good_percentage >= 70:
            insights.append(f"Overall device health is good with {good_percentage:.1f}% of devices in good status.")
        elif good_percentage >= 50:
            insights.append(f"Overall device health is fair with {good_percentage:.1f}% of devices in good status.")
        else:
            insights.append(f"Overall device health needs attention with only {good_percentage:.1f}% of devices in good status.")
        
        # INSIGHT 2: EOL Status Summary
        # First check end of sale status
        eol_sale_critical = sum(1 for device in inventory_devices if is_end_of_sale_critical(device, eol_data))
        eol_sale_critical_pct = (eol_sale_critical / total_devices) * 100 if total_devices > 0 else 0
        
        # Then check end of support status
        eol_support_critical = sum(1 for device in inventory_devices if is_end_of_support_critical(device, eol_data))
        eol_support_warning = sum(1 for device in inventory_devices if is_end_of_support_warning(device, eol_data))
        eol_support_critical_pct = (eol_support_critical / total_devices) * 100 if total_devices > 0 else 0
        eol_support_warning_pct = (eol_support_warning / total_devices) * 100 if total_devices > 0 else 0
        eol_support_good_pct = 100 - eol_support_critical_pct - eol_support_warning_pct
        
        # Calculate overall EOL devices
        eol_devices = eol_sale_critical  # Count devices past end of sale as EOL
        eol_devices_pct = (eol_devices / total_devices) * 100 if total_devices > 0 else 0
        
        # Generate appropriate EOL status insight based on actual data
        if eol_devices > 0:
            insights.append(f"EOL Status: {eol_devices} devices ({eol_devices_pct:.1f}%) have reached end-of-sale.")
            
            if eol_support_critical > 0:
                insights.append(f"Support Status: {eol_support_critical} devices ({eol_support_critical_pct:.1f}%) have reached or are within 1 year of end-of-support.")
            
            if eol_support_warning > 0:
                insights.append(f"Support Status: {eol_support_warning} devices ({eol_support_warning_pct:.1f}%) are within 2 years of end-of-support.")
        else:
            insights.append(f"EOL Status: Excellent lifecycle management with all devices on current hardware.")
        
        # Generate EOL recommendations based on EOL status
        if eol_sale_critical > 0:
            health_recommendations.append(f"Create a hardware refresh plan for {eol_sale_critical} devices ({eol_sale_critical_pct:.1f}%) that have reached end-of-sale.")
        
        if eol_support_critical > 0:
            health_recommendations.append(f"Prioritize replacement of {eol_support_critical} devices that have reached or will soon reach end-of-support.")
        
        if eol_support_warning > 0 and len(health_recommendations) < 2:
            health_recommendations.append(f"Plan for replacement of {eol_support_warning} devices approaching end-of-support within 2 years.")
    
    # 2. Detailed Firmware Compliance Analysis with prioritized recommendations
    if firmware_stats:
        # Aggregate firmware stats across all device types
        firmware_health = {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0}
        
        for device_type in ['MX', 'MS', 'MR']:
            if device_type in firmware_stats:
                stats = firmware_stats[device_type]
                firmware_health['Good'] += stats.get('Good', 0)
                firmware_health['Warning'] += stats.get('Warning', 0)
                firmware_health['Critical'] += stats.get('Critical', 0)
                firmware_health['Total'] += stats.get('Total', 0)
        
        # Generate comprehensive firmware compliance insight
        if firmware_health['Total'] > 0:
            good_fw_pct = (firmware_health['Good'] / firmware_health['Total']) * 100
            warning_fw_pct = (firmware_health['Warning'] / firmware_health['Total']) * 100
            critical_fw_pct = (firmware_health['Critical'] / firmware_health['Total']) * 100
            
            # INSIGHT 3: Overall Firmware Status
            if good_fw_pct >= 85:
                insights.append(f"Firmware Status: Excellent compliance with {good_fw_pct:.1f}% of devices running recommended firmware.")
            elif good_fw_pct >= 70:
                insights.append(f"Firmware Status: Good compliance with {good_fw_pct:.1f}% of devices running recommended firmware.")
            elif good_fw_pct >= 50:
                insights.append(f"Firmware Status: Fair compliance with {good_fw_pct:.1f}% of devices running recommended firmware.")
            else:
                insights.append(f"Firmware Status: Needs attention with only {good_fw_pct:.1f}% of devices running recommended firmware.")
            
            # Critical firmware issues by device type
            critical_firmware_issues = []
            for device_type in ['MX', 'MS', 'MR']:
                if device_type in firmware_stats:
                    stats = firmware_stats[device_type]
                    if 'Critical' in stats and 'Total' in stats and stats['Total'] > 0:
                        critical_percentage = (stats.get('Critical', 0) / stats['Total']) * 100
                        
                        if critical_percentage >= 20:
                            critical_firmware_issues.append(f"{device_type} ({critical_percentage:.1f}%)")
            
            if critical_firmware_issues:
                insights.append(f"Critical firmware compliance issues found in {', '.join(critical_firmware_issues)}.")
                health_recommendations.insert(0, f"Prioritize firmware upgrades for devices with critical compliance issues to address security vulnerabilities.")
        
        # Add specific recommendations based on firmware state (prioritize these)
        if firmware_health.get('Critical', 0) > 0:
            critical_count = firmware_health.get('Critical', 0)
            critical_pct = (critical_count / firmware_health.get('Total', 1)) * 100
            if not any("firmware upgrade" in r.lower() for r in health_recommendations):
                health_recommendations.insert(0, f"Upgrade {critical_count} devices with critical firmware versions to mitigate security risks.")
        
        if firmware_health.get('Warning', 0) > 0:
            warning_count = firmware_health.get('Warning', 0)
            warning_pct = (warning_count / firmware_health.get('Total', 1)) * 100
            if warning_pct > 25 and not any("firmware update" in r.lower() for r in health_recommendations):
                health_recommendations.append(f"Schedule firmware updates for {warning_count} devices running outdated versions.")
    
    # 3. Analyze inventory composition
    if inventory_devices:
        total_devices = len(inventory_devices)
        device_types = defaultdict(int)
        for device in inventory_devices:
            model = device.get('model', '')
            if model.startswith('MX'):
                device_types['MX'] += 1
            elif model.startswith('MS'):
                device_types['MS'] += 1
            elif model.startswith('MR') or model.startswith('CW'):
                device_types['MR'] += 1
            elif model.startswith('MV'):
                device_types['MV'] += 1
            elif model.startswith('MG'):
                device_types['MG'] += 1
            elif model.startswith('MT'):
                device_types['MT'] += 1
        
        # Generate insight about network composition
        network_composition = []
        for device_type, count in sorted(device_types.items(), key=lambda x: x[1], reverse=True):
            percentage = (count / total_devices) * 100
            if percentage >= 10:  # Only include significant device types
                network_composition.append(f"{device_type} ({count}, {percentage:.1f}%)")
        
        if network_composition:
            insights.append(f"Network consists primarily of {', '.join(network_composition)}.")
    
    # 4. Analyze client density and network capacity
    if dashboard_stats and 'avg_unique_clients_per_day' in dashboard_stats:
        avg_clients = dashboard_stats['avg_unique_clients_per_day']
        if avg_clients > 0:
            insights.append(f"Network serves an average of {avg_clients} unique clients per day.")
            
            if inventory_devices:
                mr_count = sum(1 for device in inventory_devices if device.get('model', '').startswith('MR'))
                if mr_count > 0:
                    clients_per_ap = avg_clients / mr_count
                    if clients_per_ap > 50:
                        insights.append(f"Very high client-to-AP ratio detected ({clients_per_ap:.1f} clients per AP).")
                        general_recommendations.append("Add more wireless access points to improve coverage and performance for your {avg_clients} daily clients.")
                    elif clients_per_ap > 30:
                        insights.append(f"High client-to-AP ratio ({clients_per_ap:.1f} clients per AP).")
                        general_recommendations.append("Consider expanding wireless capacity in high-density areas to maintain optimal performance.")
    
    # 5. Analyze product adoption opportunities (lower priority than health)
    if products:
        missing_products = [product for product, adopted in products.items() if not adopted]
        
        # Add recommendations for missing core products (highest priority)
        missing_core_products = [p for p in ['MX', 'MS', 'MR'] if not products.get(p, False)]
        if missing_core_products:
            # Create recommendations for each missing core product
            core_product_descriptions = {
                'MX': "Deploy Meraki MX security appliances to provide integrated security, SD-WAN capabilities, and simplified management through a single dashboard.",
                'MS': "Add Meraki MS switches to enable centralized network management with virtual stacking, automatic alerts, and deep visibility across your entire network.",
                'MR': "Implement Meraki wireless access points to deliver high-performance connectivity with integrated location analytics, automatic RF optimization, and built-in security capabilities."
            }
            
            # Add recommendations for each missing core product
            for product in missing_core_products:
                if product in core_product_descriptions:
                    product_recommendations.append(core_product_descriptions[product])
            
            # If multiple core products are missing, add a comprehensive recommendation
            if len(missing_core_products) > 1:
                product_recommendations.append(f"Complete your Meraki deployment by adding {', '.join(missing_core_products)} to gain the full benefits of a unified cloud-managed network architecture.")
        
        # Contextual adoption recommendations for advanced products
        if 'MX' in products and products['MX'] and 'Secure Connect' in missing_products:
            product_recommendations.append("Implement Secure Connect to enhance remote access security and simplify VPN management.")
        
        if 'MR' in products and products['MR'] and 'Umbrella Secure Internet Gateway' in missing_products:
            product_recommendations.append("Add Umbrella Secure Internet Gateway to strengthen threat protection at the DNS layer.")
        
        if 'MS' in products and products['MS'] and 'Thousand Eyes' in missing_products:
            product_recommendations.append("Deploy ThousandEyes to gain end-to-end visibility into network performance issues.")
        
        if 'MV' in products and products['MV'] and 'Spaces' in missing_products:
            product_recommendations.append("Leverage Spaces to derive business insights from your existing MV camera deployment.")

    # 6. Health score-based firmware recommendations
    if health_score is not None and firmware_stats:
        # CRITICAL FIX: Check device health percentage specifically before making recommendations
        device_health_pct = 0
        if inventory_devices:
            total_devices = len(inventory_devices)
            good_devices = sum(1 for device in inventory_devices if categorize_device_health(device, firmware_stats, eol_data) == 'good')
            device_health_pct = (good_devices / total_devices * 100) if total_devices > 0 else 0
        
        # If device health percentage is below 80%, always prioritize device health recommendations
        if device_health_pct < 80:
            # Check if we already have specific firmware recommendations
            has_firmware_rec = any("firmware" in r.lower() for r in health_recommendations)
            
            if not has_firmware_rec:
                # Add a targeted recommendation based on device health percentage
                if device_health_pct < 70:
                    health_recommendations.append(f"Improve device health status (currently {device_health_pct:.1f}%) by upgrading firmware and replacing devices with critical issues to enhance network reliability.")
                else:
                    health_recommendations.append(f"Address devices with warning status to improve overall device health from {device_health_pct:.1f}% to optimal levels.")
        
        # Only add general health score recommendations if no specific health recommendations exist
        elif len(health_recommendations) == 0:
            # For networks with excellent health (75%+)
            if health_score >= 75:
                has_warning_firmware = False
                for device_type in ['MX', 'MS', 'MR', 'MV', 'MG', 'MT']:
                    if device_type in firmware_stats and firmware_stats[device_type].get('Warning', 0) > 0:
                        has_warning_firmware = True
                        break
                
                if has_warning_firmware:
                    general_recommendations.append("Maintain your strong network health by implementing a regular firmware update schedule for devices on older but stable versions.")
                else:
                    general_recommendations.append("Preserve your excellent network health by continuing to keep firmware versions current across all device types.")
            
            # For networks with good health (60-74%)
            elif health_score >= 60:
                general_recommendations.append("Improve network health by implementing a structured firmware update program focusing on security and stability enhancements.")
            
            # For networks with fair health (40-59%)
            elif health_score >= 40:
                health_recommendations.append("Develop a firmware compliance strategy to progressively bring all devices to recommended versions, prioritizing security patches and critical bug fixes.")
    
    # Combine all recommendations with proper prioritization:
    # 1. Health and Lifecycle recommendations first
    # 2. General recommendations second
    # 3. Product recommendations last
    
    # Calculate device health percentage for final check
    device_health_pct = 0
    
    if inventory_devices:
        total_devices = len(inventory_devices)
        # Device health check
        good_devices = sum(1 for device in inventory_devices if categorize_device_health(device, firmware_stats, eol_data) == 'good')
        device_health_pct = (good_devices / total_devices * 100) if total_devices > 0 else 0
        
        # Calculate End of Sale percentage
        eol_sale_devices = sum(1 for device in inventory_devices if is_end_of_sale_critical(device, eol_data))
        eol_sale_pct = (eol_sale_devices / total_devices * 100) if total_devices > 0 else 0
        
        # Calculate End of Support percentages
        eol_support_critical = sum(1 for device in inventory_devices if is_end_of_support_critical(device, eol_data))
        eol_support_warning = sum(1 for device in inventory_devices if is_end_of_support_warning(device, eol_data))
        eol_support_pct = ((eol_support_critical + eol_support_warning) / total_devices * 100) if total_devices > 0 else 0
    
    # Final check to ensure we have appropriate health recommendations
    if device_health_pct < 80 and not any("device health" in r.lower() for r in health_recommendations):
        health_recommendations.append(f"Improve device health from {device_health_pct:.1f}% by upgrading firmware and replacing devices with critical issues.")
    
    # Check for EOL recommendations specifically if we have EOL issues
    if eol_sale_pct > 10 and not any("end-of-sale" in r.lower() for r in health_recommendations):
        health_recommendations.append(f"Address {eol_sale_pct:.1f}% of devices that have reached end-of-sale by planning hardware refresh.")
    
    if eol_support_pct > 10 and not any("end-of-support" in r.lower() for r in health_recommendations):
        health_recommendations.append(f"Plan replacement of devices approaching end-of-support to maintain network reliability.")
    
    # Limit health recommendations to 2 most important ones to avoid overwhelming
    if len(health_recommendations) > 2:
        health_recommendations = health_recommendations[:2]
    
    # Separate core product recommendations from other product recommendations
    core_product_recommendations = []
    other_product_recommendations = []
    
    for rec in product_recommendations:
        # Check if this is a core product recommendation
        if any(phrase in rec for phrase in ['Deploy Meraki MX', 'Add Meraki MS', 'Implement Meraki wireless', 'Complete your Meraki deployment']):
            core_product_recommendations.append(rec)
        else:
            other_product_recommendations.append(rec)
    
    # Prioritize recommendations in this order:
    # 1. Core product recommendations
    # 2. Health recommendations
    # 3. General recommendations
    # 4. Other product recommendations
    all_recommendations = core_product_recommendations + health_recommendations + general_recommendations + other_product_recommendations
    
    # Ensure we have at least one health recommendation if health is below optimal
    if device_health_pct < 80 and not any(any(term in r.lower() for term in ['device health', 'firmware', 'upgrade', 'health improvement']) for r in all_recommendations[:2]):
        # Add health recommendation but preserve any core product recommendations at the top
        if core_product_recommendations:
            all_recommendations.insert(len(core_product_recommendations), f"Address device health issues to improve from {device_health_pct:.1f}% by updating firmware and replacing problematic devices.")
        else:
            all_recommendations.insert(0, f"Address device health issues to improve from {device_health_pct:.1f}% by updating firmware and replacing problematic devices.")
    
    return insights, all_recommendations
def create_health_score_indicator(slide, x, y, width, health_score):
    """Create a modern, clean health score indicator."""
    # Health score section container
    container_height = Inches(1.0)
    container = slide.shapes.add_shape(1, x, y, width, container_height)
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Very light gray
    container.line.fill.solid()
    container.line.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray
    
    # Add title
    title_x = x + Inches(0.1)
    title_y = y + Inches(0.1)
    title_width = width - Inches(0.2)
    title_box = slide.shapes.add_textbox(title_x, title_y, title_width, Inches(0.25))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = "Network Health Score"
    title_p.font.size = Pt(11)
    title_p.font.bold = True
    
    # Add score number
    score_x = x + Inches(0.1)
    score_y = y + Inches(0.35)
    score_width = Inches(1.0)
    score_box = slide.shapes.add_textbox(score_x, score_y, score_width, Inches(0.5))
    score_p = score_box.text_frame.add_paragraph()
    score_p.text = str(health_score)
    score_p.font.size = Pt(32)
    score_p.font.bold = True
    
    # Determine color based on score
    if health_score >= 80:
        score_p.font.color.rgb = GOOD_COLOR
    elif health_score >= 60:
        score_p.font.color.rgb = WARNING_COLOR
    else:
        score_p.font.color.rgb = CRITICAL_COLOR
    
    # Add "out of 100" text
    out_of_box = slide.shapes.add_textbox(score_x + Inches(0.9), score_y + Inches(0.2), Inches(0.6), Inches(0.25))
    out_of_p = out_of_box.text_frame.add_paragraph()
    out_of_p.text = "/ 100"
    out_of_p.font.size = Pt(12)
    
    # Add assessment text
    assessment = get_health_assessment(health_score)
    assessment_x = x + Inches(1.6)
    assessment_y = y + Inches(0.47)
    assessment_width = width - Inches(1.7)
    assessment_box = slide.shapes.add_textbox(assessment_x, assessment_y, assessment_width, Inches(0.3))
    assessment_p = assessment_box.text_frame.add_paragraph()
    assessment_p.text = assessment
    assessment_p.font.size = Pt(16)
    assessment_p.font.bold = True
    
    # Determine color based on assessment
    assessment_p.font.color.rgb = get_status_color(assessment)
    
    return container

def add_insights_and_recommendations(slide, insights, recommendations, x, y, width, height):
    """Add insights and recommendations in a clean, professional format."""
    # Create a single content box
    content_box = slide.shapes.add_textbox(x, y, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    # Add insights section
    insights_title = content_tf.add_paragraph()
    insights_title.text = "Key Insights"
    insights_title.font.size = Pt(12)
    insights_title.font.bold = True
    insights_title.space_after = Pt(4)
    
    # Add insights
    for i, insight in enumerate(insights[:5]):  # Limit to 5 insights
        insight_p = content_tf.add_paragraph()
        insight_p.text = f"• {insight}"
        insight_p.font.size = Pt(10)
        insight_p.space_after = Pt(3)
    
    # Add some spacing between sections
    spacer = content_tf.add_paragraph()
    spacer.space_after = Pt(8)
    
    # Add recommendations section
    if recommendations:
        rec_title = content_tf.add_paragraph()
        rec_title.text = "Recommendations"
        rec_title.font.size = Pt(12)
        rec_title.font.bold = True
        rec_title.space_after = Pt(4)
        
        # Add recommendations
        for i, recommendation in enumerate(recommendations[:4]):  # Limit to 4 recommendations
            rec_p = content_tf.add_paragraph()
            rec_p.text = f"• {recommendation}"
            rec_p.font.size = Pt(10)
            rec_p.space_after = Pt(3)
    
    return content_box

def create_network_health_metrics(slide, device_health, firmware_stats, eol_data, inventory_devices, x, y, width, height):
    """Create simple, reliable network health metrics that won't cause PowerPoint rendering issues."""
    # Add section title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.25))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = "Network Health Metrics"
    title_p.font.size = Pt(12)
    title_p.font.bold = True
    
    # Calculate health metrics
    total_devices = sum(device_health.values())
    good_devices = device_health.get('good', 0)
    device_health_pct = (good_devices / total_devices * 100) if total_devices > 0 else 0
    
    # Calculate EOL statistics
    eol_stats = {'current': 0, 'approaching': 0, 'past': 0}
    
    if inventory_devices:
        for device in inventory_devices:
            model = device.get('model', '')
            base_model = get_base_model(model)
            
            if base_model and eol_data:
                eol_status_found = False
                for eol_model in eol_data:
                    if base_model.startswith(eol_model):
                        eol_date = eol_data[eol_model].get('end_of_support')
                        if eol_date:
                            try:
                                # Convert date string to datetime object
                                from datetime import datetime
                                eol_datetime = datetime.strptime(eol_date, "%b %d, %Y")
                                current_date = datetime.now()
                                
                                # Calculate time difference in months
                                time_diff = (eol_datetime - current_date).days / 30
                                
                                if time_diff <= 0:
                                    eol_stats['past'] += 1
                                elif time_diff <= 24:
                                    eol_stats['approaching'] += 1
                                else:
                                    eol_stats['current'] += 1
                                
                                eol_status_found = True
                                break
                            except:
                                # If date parsing fails, consider it current
                                eol_stats['current'] += 1
                                eol_status_found = True
                                break
                
                # If no EOL status found, consider it current
                if not eol_status_found:
                    eol_stats['current'] += 1
            else:
                # Unknown model defaults to current
                eol_stats['current'] += 1
    else:
        # No inventory devices, set current to total
        eol_stats['current'] = total_devices
    
    lifecycle_good = eol_stats['current']
    lifecycle_pct = (lifecycle_good / total_devices * 100) if total_devices > 0 else 0
    
    # Add Device Health metric - simple text with percentage
    health_label = slide.shapes.add_textbox(x, y + Inches(0.4), Inches(1.5), Inches(0.25))
    health_label_p = health_label.text_frame.add_paragraph()
    health_label_p.text = "Device Health"
    health_label_p.font.size = Pt(11)
    health_label_p.font.bold = True
    
    health_subtitle = slide.shapes.add_textbox(x, y + Inches(0.65), Inches(1.7), Inches(0.25))
    health_subtitle_p = health_subtitle.text_frame.add_paragraph()
    health_subtitle_p.text = "Overall operational status"
    health_subtitle_p.font.size = Pt(8)
    health_subtitle_p.font.italic = True
    
    health_value = slide.shapes.add_textbox(x, y + Inches(0.9), Inches(1.5), Inches(0.5))
    health_value_p = health_value.text_frame.add_paragraph()
    health_value_p.text = f"{device_health_pct:.1f}%"
    health_value_p.font.size = Pt(24)
    health_value_p.font.bold = True
    
    # Color based on percentage
    if device_health_pct >= 80:
        health_value_p.font.color.rgb = GOOD_COLOR
    elif device_health_pct >= 60:
        health_value_p.font.color.rgb = WARNING_COLOR
    else:
        health_value_p.font.color.rgb = CRITICAL_COLOR
    
    health_details = slide.shapes.add_textbox(x, y + Inches(1.4), Inches(2.0), Inches(0.25))
    health_details_p = health_details.text_frame.add_paragraph()
    health_details_p.text = f"{good_devices}/{total_devices} devices in good status"
    health_details_p.font.size = Pt(8)
    
    # Add Lifecycle metric - simple text with percentage
    lifecycle_label = slide.shapes.add_textbox(x + Inches(1.75), y + Inches(0.4), Inches(1.5), Inches(0.25))
    lifecycle_label_p = lifecycle_label.text_frame.add_paragraph()
    lifecycle_label_p.text = "Lifecycle"
    lifecycle_label_p.font.size = Pt(11)
    lifecycle_label_p.font.bold = True
    
    lifecycle_subtitle = slide.shapes.add_textbox(x + Inches(1.75), y + Inches(0.65), Inches(1.7), Inches(0.25))
    lifecycle_subtitle_p = lifecycle_subtitle.text_frame.add_paragraph()
    lifecycle_subtitle_p.text = "Hardware EOL status"
    lifecycle_subtitle_p.font.size = Pt(8)
    lifecycle_subtitle_p.font.italic = True
    
    lifecycle_value = slide.shapes.add_textbox(x + Inches(1.75), y + Inches(0.9), Inches(1.5), Inches(0.5))
    lifecycle_value_p = lifecycle_value.text_frame.add_paragraph()
    lifecycle_value_p.text = f"{lifecycle_pct:.1f}%"
    lifecycle_value_p.font.size = Pt(24)
    lifecycle_value_p.font.bold = True
    
    # Color based on percentage
    if lifecycle_pct >= 80:
        lifecycle_value_p.font.color.rgb = GOOD_COLOR
    elif lifecycle_pct >= 60:
        lifecycle_value_p.font.color.rgb = WARNING_COLOR
    else:
        lifecycle_value_p.font.color.rgb = CRITICAL_COLOR
    
    lifecycle_details = slide.shapes.add_textbox(x + Inches(1.75), y + Inches(1.4), Inches(2.0), Inches(0.25))
    lifecycle_details_p = lifecycle_details.text_frame.add_paragraph()
    lifecycle_details_p.text = f"{lifecycle_good}/{total_devices} devices on current hardware"
    lifecycle_details_p.font.size = Pt(8)
    
    return {
        'title_box': title_box,
        'device_health_pct': device_health_pct,
        'lifecycle_pct': lifecycle_pct,
        'eol_stats': eol_stats
    }

def create_simple_donut(slide, x, y, values, colors):
    """Create a simple donut chart without title or description."""

    chart_data = CategoryChartData()
    

    categories = []
    data_values = []
    
    for segment, value in values.items():
        if value > 0:
            categories.append(segment)
            data_values.append(value)
    
    chart_data.categories = categories
    chart_data.add_series('Status', data_values)
    
    # Set chart size
    chart_size = Inches(1.25)
    
    # Add donut chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, 
        x, 
        y, 
        chart_size, 
        chart_size,
        chart_data
    ).chart
    
    # Configure chart
    chart.has_title = False
    chart.has_legend = False
    
    # Set donut hole size
    plot = chart.plots[0]
    plot.has_data_labels = False
    plot.doughnut_hole_size = 50  # 50% of the radius
    
    # Color the segments
    for i, point in enumerate(plot.series[0].points):
        segment_name = categories[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[segment_name]
    
    return chart

def add_health_score_explanation(slide):
    """Add a small note explaining how health score is calculated."""

    note_box = slide.shapes.add_textbox(Inches(0.45), Inches(6.45), Inches(10), Inches(0.2))
    note_tf = note_box.text_frame
    note_p = note_tf.add_paragraph()
    note_p.text = "Network Health Score: Calculated based on firmware compliance, hardware lifecycle status, product adoption, and client density metrics."
    note_p.font.size = Pt(8)
    note_p.font.italic = True
    note_p.font.color.rgb = RGBColor(128, 128, 128)  # Gray
    
    return note_box

async def generate(api_client, template_path, output_path, 
                  inventory_devices=None, networks=None, 
                  dashboard_stats=None, firmware_stats=None, 
                  eol_data=None, products=None):
    """Generate the Executive Summary slide."""
    print(f"\n{GREEN}Generating Executive Summary slide...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # Calculate overall health score
    health_score, deduction_reasons = calculate_health_score(
        inventory_devices, firmware_stats, eol_data, dashboard_stats, products
    )
    
    # Generate insights and recommendations with the health score
    insights, recommendations = generate_insights(
        inventory_devices, dashboard_stats, firmware_stats, eol_data, products, health_score
    )
    
    # Create device health distribution data
    device_health = {'good': 0, 'warning': 0, 'critical': 0}
    if inventory_devices:
        for device in inventory_devices:
            health = categorize_device_health(device, firmware_stats, eol_data)
            if health in device_health:
                device_health[health] += 1
    
    # Create the slide
    try:
        prs = Presentation(output_path)
        
        # Find a suitable slide layout
        slide_layout = None
        for layout in prs.slide_layouts:
            if hasattr(layout, 'name') and layout.name and 'title' in layout.name.lower():
                slide_layout = layout
                break
        
        # If no suitable layout found, use the first one
        if not slide_layout and len(prs.slide_layouts) > 0:
            slide_layout = prs.slide_layouts[0]
        
        # If we have a layout, create the slide
        if slide_layout:
            # Add slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Clean the slide of any existing content
            clean_slide(slide)
            decorative_bubbles = add_decorative_bubbles(slide)
            # Move to position 1 (second slide, after title)
            new_slide_index = len(prs.slides) - 1
            move_slide_to_position(prs, new_slide_index, 1)
            
            # Add title with styling
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.0), Inches(11), Inches(0.4))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "Executive Summary"
            title_p.font.size = Pt(24)
            title_p.font.bold = True
            
            # Add horizontal line under the title
            line = slide.shapes.add_connector(1, Inches(0.5), Inches(0.75), Inches(11.0), Inches(0.75))
            line.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
            line.line.width = Pt(1.5)
            
            # COLUMN 1 (Left): Health score box
            left_col_x = Inches(0.5)
            left_col_width = Inches(3.5)
            
            # Create the health score container box
            container_height = Inches(1.0)
            container = slide.shapes.add_shape(1, left_col_x, Inches(1.21), left_col_width, container_height)
            container.fill.solid()
            container.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Very light gray
            container.line.fill.solid()
            container.line.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray

            # Add the Network Health Score title
            title_box = slide.shapes.add_textbox(left_col_x + Inches(0.1), Inches(1.1), left_col_width - Inches(0.2), Inches(0.25))
            title_p = title_box.text_frame.add_paragraph()
            title_p.text = "Network Health Score"
            title_p.font.size = Pt(11)
            title_p.font.bold = True

            # Add the score number
            score_box = slide.shapes.add_textbox(left_col_x + Inches(0.1), Inches(1.35), Inches(1.0), Inches(0.5))
            score_p = score_box.text_frame.add_paragraph()
            score_p.text = str(health_score)
            score_p.font.size = Pt(32)
            score_p.font.bold = True
                        
            if health_score >= 80:
                score_p.font.color.rgb = GOOD_COLOR
            elif health_score >= 60:
                score_p.font.color.rgb = WARNING_COLOR
            else:
                score_p.font.color.rgb = CRITICAL_COLOR

            # Add "out of 100" text
            out_of_box = slide.shapes.add_textbox(left_col_x + Inches(0.9), Inches(1.55), Inches(0.6), Inches(0.25))
            out_of_p = out_of_box.text_frame.add_paragraph()
            out_of_p.text = "/ 100"
            out_of_p.font.size = Pt(12)

            # Add assessment text
            assessment = get_health_assessment(health_score)
            assessment_box = slide.shapes.add_textbox(left_col_x + Inches(1.6), Inches(1.47), left_col_width - Inches(1.7), Inches(0.3))
            assessment_p = assessment_box.text_frame.add_paragraph()
            assessment_p.text = assessment
            assessment_p.font.size = Pt(16)
            assessment_p.font.bold = True
            assessment_p.font.color.rgb = get_status_color(assessment)
            
            # Add score deduction explanations with exact positioning
            if deduction_reasons:
                deduction_container = add_score_deduction_explanation(
                    slide, 
                    deduction_reasons
                )
            
            # Add Network Health Metrics below health score
            metrics_result = create_network_health_metrics(
                slide, device_health, firmware_stats, eol_data, inventory_devices,
                left_col_x, Inches(2.2), left_col_width, Inches(2.0)
            )
            
            # Extract values from the metrics function result
            eol_stats = metrics_result['eol_stats']
            device_health_pct = metrics_result['device_health_pct']
            lifecycle_pct = metrics_result['lifecycle_pct']
            
            # Calculate total devices for the donut charts
            total_devices = sum(device_health.values())
            
            # Create segment colors map
            segment_colors = {
                'good': GOOD_COLOR,
                'warning': WARNING_COLOR,
                'critical': CRITICAL_COLOR,
                'current': GOOD_COLOR,
                'approaching': WARNING_COLOR,
                'past': CRITICAL_COLOR
            }
            
            # Create Device Health donut chart
            health_donut = create_simple_donut(
                slide,
                Inches(0.6),
                Inches(4.19),
                device_health,
                segment_colors
            )
            
            # Create Lifecycle donut chart
            lifecycle_values = {
                'current': eol_stats['current'],
                'approaching': eol_stats['approaching'],
                'past': eol_stats['past']
            }

            lifecycle_donut = create_simple_donut(
                slide,
                Inches(2.5),
                Inches(4.19),
                lifecycle_values,
                segment_colors
            )
                        
            # COLUMN 2: Insights and Recommendations in the middle
            right_col_x = Inches(5.2)
            right_col_width = Inches(6.3)
            
            # Key Insights title
            insights_title_box = slide.shapes.add_textbox(right_col_x, Inches(1.5), right_col_width, Inches(0.3))
            insights_title_p = insights_title_box.text_frame.add_paragraph()
            insights_title_p.text = "Key Insights"
            insights_title_p.font.size = Pt(12)
            insights_title_p.font.bold = True
            
            # Add Key Insights content
            insights_box = slide.shapes.add_textbox(right_col_x, Inches(1.8), right_col_width, Inches(2.0))
            insights_tf = insights_box.text_frame
            
            # Add insights
            for i, insight in enumerate(insights[:5]):  # Limit to 5 insights
                insight_p = insights_tf.add_paragraph()
                insight_p.text = f"• {insight}"
                insight_p.font.size = Pt(10)
                insight_p.space_after = Pt(3)

            add_recommendations_to_slide(
                slide, 
                recommendations,
                right_col_x,
                Inches(3.8),
                right_col_width,
                Inches(2.3)
            )
            
            # Add health score explanation at the bottom
            add_health_score_explanation(slide)
            
            # Create comprehensive notes for the presenter
            health_score_notes = """
Network Health Score Logic:

The Network Health Score starts at 100 points and applies deductions based on percentage thresholds:

Device Health Deductions:
- Critical health devices: -5 to -30 points (higher % = larger deduction)
- Warning health devices: -5 to -15 points (higher % = larger deduction)

Lifecycle Deductions:
- End-of-support devices: -5 to -25 points (higher % = larger deduction)
- End-of-sale devices: -5 to -10 points (higher % = larger deduction)

Firmware Deductions:
- Critical firmware versions: -5 to -20 points (higher % = larger deduction)

Product Adoption Deductions:
- Missing core products (MX, MS, MR): -5 points each
- Missing advanced products (Secure Connect, Umbrella): -2 points each

Device Health Logic:
Devices are categorized as "Good," "Warning," or "Critical" based on:
1. Firmware Status:
   - Good: Running latest firmware 
   - Warning: Running outdated but stable firmware
   - Critical: Running versions with known security issues

2. Lifecycle Status:
   - Critical: Past end-of-support or within 1 year of end-of-support
   - Warning: Past end-of-sale or within 1-2 years of end-of-support
   - Good: Current hardware with >2 years until end-of-support

A device's overall health is determined by its most severe issue (firmware or lifecycle).
"""

            # Add the notes to the slide
            add_notes_to_slide(slide, health_score_notes)
            
            # Save the presentation
            prs.save(output_path)
            print(f"{GREEN}Added Executive Summary slide to the presentation (slide 2){RESET}")
            print(f"Saved presentation to {output_path}")  # Added confirmation message
        else:
            print(f"{RED}No suitable slide layout found in the presentation{RESET}")
        
    except Exception as e:
        print(f"{RED}Error creating Executive Summary slide: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    # Calculate execution time
    total_time = time.time() - start_time
    print(f"{PURPLE}Executive Summary slide created in {total_time:.2f} seconds{RESET}")
    
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some sample data for testing
    
    # Sample inventory devices
    inventory_devices = [
        {"model": "MX64", "firmware": "15.44.0", "networkId": "N1"},
        {"model": "MX84-HW", "firmware": "15.44.0", "networkId": "N1"},
        {"model": "MS220-8P", "firmware": "14.16.1", "networkId": "N2"},
        {"model": "MS220-24P", "firmware": "14.16.1", "networkId": "N2"},
        {"model": "MS320-48LP", "firmware": "14.16.1", "networkId": "N3"},
        {"model": "MR16", "firmware": "28.5", "networkId": "N4"},
        {"model": "MR34", "firmware": "29.5", "networkId": "N4"},
        {"model": "MV21", "firmware": "5.2", "networkId": "N5"},
        {"model": "MR84", "firmware": "29.5", "networkId": "N6"},
        {"model": "MS350-48LP", "firmware": "14.32.1", "networkId": "N7"}
    ]
    
    # Sample dashboard stats
    dashboard_stats = {
        "total_networks": 35,
        "total_inventory": 55,
        "total_active_nodes": 48,
        "total_unique_clients": 5423,
        "total_non_unique_clients": 32156,
        "avg_unique_clients_per_day": 387,
        "avg_non_unique_clients_per_day": 2296
    }
    
    # Sample firmware stats
    firmware_stats = {
        'MX': {
            'Good': 5, 'Warning': 2, 'Critical': 3, 'Total': 10,
            'Versions': {'16.14.1': 5, '15.44.0': 3, '14.53.1': 2},
            'latest': '16.14.1'
        },
        'MS': {
            'Good': 8, 'Warning': 10, 'Critical': 12, 'Total': 30,
            'Versions': {'14.32.1': 8, '14.16.1': 10, '13.28.0': 12},
            'latest': '14.32.1'
        },
        'MR': {
            'Good': 6, 'Warning': 4, 'Critical': 0, 'Total': 10,
            'Versions': {'29.5': 6, '28.5': 4},
            'latest': '29.5'
        }
    }
    
    # Sample EOL data
    eol_data = {
        "MX60": {"announcement": "Aug 30, 2021", "end_of_sale": "Aug 30, 2022", "end_of_support": "Aug 30, 2027"},
        "MX64": {"announcement": "Aug 30, 2021", "end_of_sale": "Aug 30, 2022", "end_of_support": "Aug 30, 2027"},
        "MR16": {"announcement": "Feb 28, 2017", "end_of_sale": "Feb 28, 2018", "end_of_support": "Feb 28, 2026"},
        "MS220": {"announcement": "Jun 8, 2020", "end_of_sale": "Jun 8, 2021", "end_of_support": "Jun 8, 2026"}
    }
    
    # Sample products adoption
    products = {
        'MX': True,
        'MS': True,
        'MR': True,
        'MG': False,
        'MV': True,
        'MT': False,
        'Secure Connect': False,
        'Umbrella Secure Internet Gateway': False,
        'Thousand Eyes': True,
        'Spaces': False,
        'XDR': False
    }
    
    # Generate the executive summary slide
    await generate(
        api_client, 
        template_path, 
        output_path, 
        inventory_devices,
        None,
        dashboard_stats,
        firmware_stats,
        eol_data,
        products
    )

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python executive_summary.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))