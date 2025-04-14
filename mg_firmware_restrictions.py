import os
import sys
import asyncio
import time
import re
import requests
from bs4 import BeautifulSoup
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# MG firmware version restrictions - Hardcoded fallback values
MG_FIRMWARE_RESTRICTIONS = {}

# Models that can run current firmware - Hardcoded fallback
MG_UNRESTRICTED_MODELS = [
    "MG21", "MG41", "MG51", "MG52"
]

# Model name normalizations for consistent counting
MG_MODEL_NORMALIZATIONS = {
    "MG21": "MG21",
    "MG41": "MG41",
    "MG51": "MG51",
    "MG52": "MG52"
}

# Last updated date - fallback value
RESTRICTIONS_LAST_UPDATED = "Mar 11, 2025"

def get_firmware_restrictions_from_doc():
    """
    Attempt to fetch MG firmware restrictions from documentation.
    
    Returns:
        tuple: (firmware_restrictions dict, unrestricted_models list, last_updated string, is_from_doc bool)
    """
    try:
        # Attempt to fetch documentation
        #print(f"{BLUE}Attempting to fetch MG firmware information from documentation{RESET}")
        
        # Default fallback data
        fallback_restrictions = MG_FIRMWARE_RESTRICTIONS
        fallback_unrestricted = MG_UNRESTRICTED_MODELS
        
        # Use the correct URL for firmware information
        doc_url = "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions"
        
        # Add User-Agent header to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Make the request with a timeout and headers
        response = requests.get(doc_url, timeout=15, headers=headers)
        response.raise_for_status()
        
        # Get the raw HTML content
        html_content = response.text
        
        # TARGETED APPROACH: Extract date from meta tags and schema.org data
        last_updated = None
        
        # Look for meta tag with article:modified_time
        meta_match = re.search(r'<meta\s+property="article:modified_time"\s+content="([^"]+)"', html_content)
        if meta_match:
            iso_date = meta_match.group(1)
            # Convert ISO date to readable format
            try:
                import datetime
                dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                #print(f"{GREEN}Found last updated date in meta tag: '{last_updated}'{RESET}")
            except Exception as e:
                # If datetime conversion fails, use the raw date
                print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                last_updated = iso_date
        
        # If not found in meta tags, look for dateModified in JSON-LD
        if not last_updated:
            schema_match = re.search(r'"dateModified":"([^"]+)"', html_content)
            if schema_match:
                iso_date = schema_match.group(1)
                # Convert ISO date to readable format
                try:
                    import datetime
                    dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                    last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                    #print(f"{GREEN}Found last updated date in schema.org data: '{last_updated}'{RESET}")
                except Exception as e:
                    # If datetime conversion fails, use the raw date
                    print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                    last_updated = iso_date
        
        # Now parse the HTML with BeautifulSoup for firmware restrictions
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Initialize collections for firmware data
        firmware_restrictions = {}  # model -> max firmware version
        unrestricted_models = []    # models that can run "Current" firmware
        
        # APPROACH #1: Look for tables with firmware information
        #print(f"{BLUE}Scanning tables for MG firmware information...{RESET}")
        
        tables = soup.find_all('table')
        
        for table in tables:
            # Check if this table might contain MG firmware information
            table_text = table.get_text().lower()
            if ('mg' in table_text and 'firmware' in table_text) or ('cellular gateway' in table_text and 'firmware' in table_text):
                #print(f"{BLUE}Found table with MG and firmware mentions{RESET}")
                
                # Check table headers to understand structure
                headers = []
                rows = table.find_all('tr')
                
                if rows:
                    header_cells = rows[0].find_all(['th', 'td'])
                    headers = [cell.get_text().strip().lower() for cell in header_cells]
                    #print(f"{BLUE}Table headers: {headers}{RESET}")
                
                # Find the relevant columns
                product_col = None
                max_firmware_col = None
                
                for i, header in enumerate(headers):
                    if any(term in header for term in ['product', 'model', 'gateway', 'device']):
                        product_col = i
                    if any(term in header for term in ['maximum', 'max', 'firmware restriction']):
                        max_firmware_col = i
                    elif 'firmware' in header and any(term in header for term in ['maximum', 'max', 'restriction']):
                        max_firmware_col = i
                
                # If we couldn't identify columns but "maximum runnable firmware" is in headers
                if product_col is None and max_firmware_col is None:
                    if 'maximum runnable firmware' in headers:
                        max_firmware_col = headers.index('maximum runnable firmware')
                        # Look for product/model column - often the first column
                        if 'product' in headers:
                            product_col = headers.index('product')
                        else:
                            product_col = 0
                
                # If we identified the needed columns, extract the data
                if product_col is not None and max_firmware_col is not None:
                    #print(f"{GREEN}Found table with product (col {product_col}) and max firmware (col {max_firmware_col}) columns{RESET}")
                    
                    for row in rows[1:]:
                        cells = row.find_all(['td', 'th'])
                        
                        if len(cells) > max(product_col, max_firmware_col):
                            product_text = cells[product_col].get_text().strip()
                            max_firmware_text = cells[max_firmware_col].get_text().strip().lower()
                            
                            # Extract the base model (e.g., MG21 from MG21-HW)
                            mg_models = re.findall(r'(MG\d+\w*)', product_text, re.IGNORECASE)
                            
                            for model in mg_models:
                                # Check if this model has a firmware restriction or can run "Current"
                                if any(term in max_firmware_text for term in ['current', 'latest', 'newest', 'unrestricted']):
                                    if model not in unrestricted_models:
                                        unrestricted_models.append(model)
                                        #print(f"{GREEN}Found unrestricted model: {model} (can run Current firmware){RESET}")
                                else:
                                    # Extract version number
                                    version_match = re.search(r'(\d+(?:\.\d+)*)', max_firmware_text)
                                    if version_match:
                                        version = version_match.group(1)
                                        if version not in firmware_restrictions:
                                            firmware_restrictions[version] = []
                                        
                                        if model not in firmware_restrictions[version]:
                                            firmware_restrictions[version].append(model)
                                            #print(f"{GREEN}Found restriction: {model} -> MG {version}{RESET}")
        
        # APPROACH #2: Look for MG models and firmware mentions outside tables
        if not firmware_restrictions and not unrestricted_models:
            #print(f"{BLUE}Looking for MG firmware information in page text...{RESET}")
            
            # Get page text for searching
            page_text = soup.get_text()
            
            # Search for MG models followed by firmware info
            model_firmware_pattern = re.compile(r'(MG\d+\w*).*?(?:maximum|restricted to|cannot run beyond).*?(?:firmware|version).*?(?:(current|latest)|(?:MG)?\s*(\d+(?:\.\d+)?))', re.IGNORECASE)
            
            for match in model_firmware_pattern.finditer(page_text):
                model = match.group(1)  # The MG model
                is_current = match.group(2)  # "current" or "latest" if matched
                version = match.group(3)  # Version number if matched
                
                if is_current:
                    # This model can run current firmware
                    if model not in unrestricted_models:
                        unrestricted_models.append(model)
                        #print(f"{GREEN}Found unrestricted model (text): {model} (can run Current firmware){RESET}")
                elif version:
                    # This model has a firmware restriction
                    if version not in firmware_restrictions:
                        firmware_restrictions[version] = []
                    
                    if model not in firmware_restrictions[version]:
                        firmware_restrictions[version].append(model)
                        #print(f"{GREEN}Found restriction (text): {model} -> MG {version}{RESET}")
        
        # If no MG models found in the restrictions or unrestricted list, add our fallback models as unrestricted
        if not unrestricted_models and not firmware_restrictions:
            print(f"{YELLOW}No MG model information found in documentation, using fallback{RESET}")
            unrestricted_models = MG_UNRESTRICTED_MODELS
        
        # If we found useful data, return it
        if firmware_restrictions or unrestricted_models:
            # Print summary of findings
            # print(f"{GREEN}Successfully parsed MG firmware information from documentation{RESET}")
            
            if firmware_restrictions:
                # print(f"Found {len(firmware_restrictions)} firmware restrictions:")
                # for version, models in sorted(firmware_restrictions.items(), key=lambda x: float(x[0].split('.')[0]), reverse=True):
                #     print(f"  - MG {version}: {len(models)} models - {', '.join(sorted(models))}")
                pass
            
            if unrestricted_models:
                # print(f"Found {len(unrestricted_models)} unrestricted models that can run Current firmware:")
                # print(f"  - {', '.join(sorted(unrestricted_models))}")
                pass
            
            return firmware_restrictions, unrestricted_models, last_updated, True
        else:
            # print(f"{YELLOW}Could not parse firmware information from documentation, using fallback{RESET}")
            pass
            return fallback_restrictions, fallback_unrestricted, None, False
            
    except Exception as e:
        # print(f"{RED}Error fetching/parsing documentation: {e}{RESET}")
        # import traceback
        # traceback.print_exc()
        
        # Use fallback values but no fallback date
        # print(f"{YELLOW}Using fallback firmware information{RESET}")
        return MG_FIRMWARE_RESTRICTIONS, MG_UNRESTRICTED_MODELS, None, False


# Helper function to safely check if a line has a specific RGB color
def has_rgb_color(shape, target_rgb):
    """Check if shape has a line with the target RGB color, safely handling None cases."""
    if not hasattr(shape, 'line'):
        return False
    
    # Check if shape has a line
    line = shape.line
    if line is None:
        return False
    
    # Check if line has a color
    if not hasattr(line, 'color') or line.color is None:
        return False
        
    # Check if color has rgb attribute
    if not hasattr(line.color, 'rgb') or line.color.rgb is None:
        return False
    
    # Finally check the RGB value
    return line.color.rgb == target_rgb

# Helper function to extract base model
def get_base_model(model):
    """Extract the base model (e.g., MG21 from MG21-HW)."""
    if not model:
        return None
        
    # Handle different cases
    model = model.strip().upper()
    
    # Extract the base model
    base_match = re.search(r'(MG\d+)', model, re.IGNORECASE)
    return base_match.group(1) if base_match else model

# Helper function to normalize model names
def normalize_model_name(model):
    """Normalize the model name for consistent counting."""
    if not model:
        return None
        
    # Handle different cases
    model = model.strip().upper()
    
    # Check if model is in our normalization mapping
    for pattern, normalized in MG_MODEL_NORMALIZATIONS.items():
        if pattern.upper() == model:
            return normalized
        
    # If not found in mapping, use the base model
    base_model = get_base_model(model)
    return base_model if base_model else model

# Helper function to check if model has firmware restriction
def get_model_firmware_version(model, firmware_restrictions, unrestricted_models):
    """
    Determine if a model has a firmware restriction, and if so, which version.
    
    Args:
        model: The full model string (e.g., MG21-HW)
        firmware_restrictions: Dict of firmware versions and their restricted models
        unrestricted_models: List of models that can run Current firmware
        
    Returns:
        str or None: The firmware version restriction or None if unrestricted
    """
    # Extract the base model
    base_model = get_base_model(model)
    
    if not base_model:
        return None  # Not a recognizable model
    
    # Check if the base model is in the unrestricted list
    for um in unrestricted_models:
        um_upper = um.upper()
        if base_model.upper() == um_upper or base_model.upper().startswith(um_upper):
            return None  # This model is unrestricted
    
    # Check if the base model is in any firmware restriction list
    for version, models in firmware_restrictions.items():
        for rm in models:
            rm_upper = rm.upper()
            if base_model.upper() == rm_upper or base_model.upper().startswith(rm_upper):
                return version  # Return the firmware restriction version
    
    # If not explicitly listed in either restricted or unrestricted, treat as unrestricted
    return None

async def generate(api_client, template_path, output_path, inventory_devices=None):
    """Generate the MG Firmware Restrictions slide."""
    print(f"\n{GREEN}Generating MG Firmware Restrictions slide (Slide 7)...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # If inventory_devices is provided, use it
    if not inventory_devices:
        print(f"{RED}No inventory data provided{RESET}")
        return
    
    #print(f"{BLUE}Using inventory data provided from slide 1{RESET}")
    
    # Get firmware restrictions from documentation (or use hardcoded fallback)
    firmware_restrictions, unrestricted_models, last_updated_date, is_from_doc = get_firmware_restrictions_from_doc()
    
    # Log the source of firmware restrictions
    if is_from_doc:
        #print(f"{GREEN}Using MG firmware information from documentation (last updated: {last_updated_date}){RESET}")
        pass
    else:
        print(f"{YELLOW}Using fallback MG firmware information - documentation unavailable{RESET}")
    
    # Process MG device data
    process_start_time = time.time()
    #print(f"{PURPLE}[{time.strftime('%H:%M:%S')}] Processing MG device data...{RESET}")
    
    # Filter only MG devices
    mg_devices = [device for device in inventory_devices 
                 if device.get('model', '').upper().startswith('MG')]
    
    # Display the firmware restrictions data for verification
    #print(f"{BLUE}Firmware restrictions data:{RESET}")
    for version, models in firmware_restrictions.items():
        #print(f"  - MG {version}: {', '.join(sorted(models))}")
        pass
    
    if unrestricted_models:
        #print(f"{BLUE}Unrestricted models:{RESET}")
        #print(f"  - {', '.join(sorted(unrestricted_models))}")
        pass
    
    # Count devices by firmware version and model
    restricted_devices = {}
    unrestricted_devices = {}
    total_mg_devices = len(mg_devices)
    
    # Group devices by their firmware restriction and model
    for device in mg_devices:
        model = device.get('model', 'unknown')
        
        # Skip if not an MG device (shouldn't happen due to filtering above, but just in case)
        if not model or not model.upper().startswith('MG'):
            continue
        
        # Normalize the model name for consistent counting
        normalized_model = normalize_model_name(model)
        
        # Get the firmware restriction for this model
        restricted_version = get_model_firmware_version(model, firmware_restrictions, unrestricted_models)
        
        if restricted_version:
            # This model has a firmware restriction
            if restricted_version not in restricted_devices:
                restricted_devices[restricted_version] = {}
            
            if normalized_model not in restricted_devices[restricted_version]:
                restricted_devices[restricted_version][normalized_model] = 0
            
            restricted_devices[restricted_version][normalized_model] += 1
        else:
            # This model doesn't have a specific restriction (is "Current")
            if normalized_model not in unrestricted_devices:
                unrestricted_devices[normalized_model] = 0
            
            unrestricted_devices[normalized_model] += 1
    
    #Print statistics for verification
    #print(f"{BLUE}MG Device Statistics:{RESET}")
    #print(f"Total MG devices found: {total_mg_devices}")
    
    for version in restricted_devices:
        device_count = sum(restricted_devices[version].values())
        #print(f"MG {version}: {device_count} devices")
        # Print device models
        for model, count in sorted(restricted_devices[version].items()):
            #print(f"  - {model}: {count}")
            pass
    
    unrestricted_count = sum(unrestricted_devices.values())
    #print(f"Not Firmware Restricted: {unrestricted_count} devices")
    # Print unrestricted device models
    for model, count in sorted(unrestricted_devices.items()):
        #print(f"  - {model}: {count}")
        pass
    
    process_time = time.time() - process_start_time
    print(f"{BLUE}MG data processing completed in {process_time:.2f} seconds{RESET}")
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    print(f"{BLUE}Updating PowerPoint with MG data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)
        
        # Use slide 7 (index 6)
        # If the slide doesn't exist, add it
        if len(prs.slides) < 7:
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides[6]
        
        # Clear existing shapes except for title
        title_shape = None
        teal_line = None
        black_line = None
        
        # Look for existing title and lines
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and "MG Firmware Restrictions" in shape.text_frame.text:
                title_shape = shape
                continue
                
            if has_rgb_color(shape, RGBColor(80, 200, 192)):
                teal_line = shape
                continue
                
            if has_rgb_color(shape, RGBColor(0, 0, 0)):
                black_line = shape
                continue
        
        # Create title if it doesn't exist
        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "MG Firmware Restrictions"
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            # print(f"{YELLOW}Added new title: 'MG Firmware Restrictions'{RESET}")
        else:
            # print(f"{BLUE}Found existing textbox title: 'MG Firmware Restrictions'{RESET}")
            pass
        
        # Remove all shapes except title and lines
        shapes_to_preserve = [title_shape, teal_line, black_line]
        shapes_to_remove = []
        
        for shape in slide.shapes:
            if shape not in shapes_to_preserve and shape is not None:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                if hasattr(shape, '_sp'):
                    sp = shape._sp
                    sp.getparent().remove(sp)
            except Exception as e:
                # print(f"{RED}Error removing shape: {e}{RESET}")
                pass
        
        # print(f"{BLUE}Removing {len(shapes_to_remove)} shapes while preserving title and line{RESET}")
        
        # Add last updated date with data source indicator
        update_text = f"Firmware restriction last updated {last_updated_date}"
        if not is_from_doc:
            update_text += " (using fallback data)"
            
        update_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
        update_tf = update_box.text_frame
        update_p = update_tf.add_paragraph()
        update_p.text = update_text
        update_p.font.size = Pt(10)
        update_p.font.italic = True
        
        # Add an explanatory note to define to user what "firmware restrictions" means
        explanation_text = "Note: These values represent the maximum firmware versions these devices can run."
        explanation_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(8), Inches(0.3))
        explanation_tf = explanation_box.text_frame
        explanation_p = explanation_tf.add_paragraph()
        explanation_p.text = explanation_text
        explanation_p.font.size = Pt(10)
        explanation_p.font.italic = True
        
        # Setup style settings
        header_size = Pt(16)
        item_size = Pt(12)
        
        # Current Y position for content
        current_y = Inches(1.9)
        
        left_col_x = Inches(0.5)
        right_col_x = Inches(4.75)
        
        # Left Column - Not Firmware Restricted
        header = slide.shapes.add_textbox(left_col_x, current_y, Inches(4), Inches(0.3))
        tf = header.text_frame
        p = tf.add_paragraph()
        p.text = "Not Firmware Restricted"
        p.font.size = header_size
        p.font.bold = True
        
        # Y position for content
        left_content_y = current_y + Inches(0.5)
        
        # Add "Cellular Gateways" header
        gateways_header = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
        tf = gateways_header.text_frame
        p = tf.add_paragraph()
        p.text = "Cellular Gateways:"
        p.font.size = item_size
        p.font.bold = True
        
        left_content_y += Inches(0.3)
        
        # Create the unrestricted models display
        if unrestricted_devices:
            # Create model lines with good formatting
            model_lines = []
            current_line = ""
            
            for model, count in sorted(unrestricted_devices.items()):
                model_text = f"{model} ({count})"
                
                # Check if adding this would make the line too long
                if current_line and len(current_line) + len(model_text) + 2 > 40:
                    model_lines.append(current_line)
                    current_line = model_text
                else:
                    if current_line:
                        current_line += ", " + model_text
                    else:
                        current_line = model_text
            
            if current_line:
                model_lines.append(current_line)
            
            # Add each line to the slide
            for line in model_lines:
                item = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                tf = item.text_frame
                p = tf.add_paragraph()
                p.text = line
                p.font.size = item_size
                p.alignment = PP_ALIGN.LEFT
                
                left_content_y += Inches(0.25)

        # If no MG devices were found with inventory, show the fallback models  
        elif total_mg_devices == 0:
            fallback_line = ", ".join(MG_UNRESTRICTED_MODELS)
            item = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
            tf = item.text_frame
            p = tf.add_paragraph()
            p.text = fallback_line
            p.font.size = item_size
            p.alignment = PP_ALIGN.LEFT
            
            left_content_y += Inches(0.4)
            
            # Add note about no devices found
            note = slide.shapes.add_textbox(Inches(2), left_content_y, Inches(6), Inches(0.25))
            tf = note.text_frame
            p = tf.add_paragraph()
            p.text = "Note: No MG devices found in current inventory"
            p.font.size = item_size
            p.font.italic = True
            p.alignment = PP_ALIGN.LEFT
        
        # Check if we have any device with firmware restrictions (should be none per fallback)
        if restricted_devices:
            # Sort firmware versions in reverse order (newer versions first)
            sorted_versions = sorted(restricted_devices.keys(), 
                                    key=lambda x: [float(n) for n in x.split('.')], 
                                    reverse=True)
            
            if sorted_versions:
                right_content_y = current_y
                
                # Process each version in the right column
                for version_index, version in enumerate(sorted_versions):
                    header = slide.shapes.add_textbox(right_col_x, right_content_y, Inches(4), Inches(0.3))
                    tf = header.text_frame
                    p = tf.add_paragraph()
                    p.text = f"MG {version}"
                    p.font.size = header_size
                    p.font.bold = True
                    
                    right_content_y += Inches(0.4)
                    
                    subtitle = slide.shapes.add_textbox(right_col_x + Inches(0.15), right_content_y, Inches(4), Inches(0.25))
                    tf = subtitle.text_frame
                    p = tf.add_paragraph()
                    p.text = "Cellular Gateways:"
                    p.font.size = item_size
                    p.font.bold = True
                    
                    right_content_y += Inches(0.3)
                    
                    # Create formatted model lines
                    model_lines = []
                    current_line = ""
                    
                    for model, count in sorted(restricted_devices[version].items()):
                        model_text = f"{model} ({count})"
                        
                        # Check if adding this would make the line too long
                        if current_line and len(current_line) + len(model_text) + 2 > 40:
                            model_lines.append(current_line)
                            current_line = model_text
                        else:
                            if current_line:
                                current_line += ", " + model_text
                            else:
                                current_line = model_text
                    
                    if current_line:
                        model_lines.append(current_line)
                    
                    # Add each line to the slide
                    for line in model_lines:
                        item = slide.shapes.add_textbox(right_col_x + Inches(0.15), right_content_y, Inches(4), Inches(0.25))
                        tf = item.text_frame
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = item_size
                        p.alignment = PP_ALIGN.LEFT
                        right_content_y += Inches(0.25)
                    
                    # Add spacing between versions
                    right_content_y += Inches(0.3)
        
        # Add total count at the bottom right
        total_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(3), Inches(0.4))
        tf = total_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Total MG Devices: {total_mg_devices}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Updated MG slide (Slide 7) with proper firmware categorization{RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    print(f"{PURPLE}MG Firmware Restrictions slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some sample inventory devices for testing
    inventory_devices = [
        {"model": "MG21", "firmware": "1.0", "networkId": "N1"},
        {"model": "MG41", "firmware": "1.0", "networkId": "N2"},
        {"model": "MG51", "firmware": "1.0", "networkId": "N3"},
        {"model": "MG52", "firmware": "1.0", "networkId": "N4"}
    ]
    
    await generate(api_client, template_path, output_path, inventory_devices=inventory_devices)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python slide8.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))