import os
import sys
import argparse
import asyncio
import datetime
import time
import json
import subprocess
import logging
from pptx import Presentation
import shutil

# Configure root logger to prevent debug messages from appearing in console
logging.basicConfig(level=logging.WARNING)

# Import slide modules
try:
    import clients
    CLIENTS_AVAILABLE = True
except ImportError:
    print("Warning: clients.py module not found. Dashboard summary slide will not be updated.")
    CLIENTS_AVAILABLE = False

try:
    import mx_firmware_restrictions
    MX_FIRMWARE_AVAILABLE = True
except ImportError:
    print("Warning: mx_firmware_restrictions.py module not found. MX firmware slide will not be updated.")
    MX_FIRMWARE_AVAILABLE = False

try:
    import ms_firmware_restrictions
    MS_FIRMWARE_AVAILABLE = True
except ImportError:
    print("Warning: ms_firmware_restrictions.py module not found. MS firmware slide will not be updated.")
    MS_FIRMWARE_AVAILABLE = False

try:
    import mr_firmware_restrictions
    MR_FIRMWARE_AVAILABLE = True
except ImportError:
    print("Warning: mr_firmware_restrictions.py module not found. MR firmware slide will not be updated.")
    MR_FIRMWARE_AVAILABLE = False

try:
    import mv_firmware_restrictions
    MV_FIRMWARE_AVAILABLE = True
except ImportError:
    print("Warning: mv_firmware_restrictions.py module not found. MV firmware slide will not be updated.")
    MV_FIRMWARE_AVAILABLE = False

try:
    import mg_firmware_restrictions
    MG_FIRMWARE_AVAILABLE = True
except ImportError:
    print("Warning: mg_firmware_restrictions.py module not found. MG firmware slide will not be updated.")
    MG_FIRMWARE_AVAILABLE = False

try:
    import firmware_compliance_mxmsmr
    FIRMWARE_COMPLIANCE_MXMSMR_AVAILABLE = True
except ImportError:
    print("Warning: firmware_compliance_mxmsmr.py module not found. MX/MS/MR firmware compliance slide will not be updated.")
    FIRMWARE_COMPLIANCE_MXMSMR_AVAILABLE = False

try:
    import firmware_compliance_mgmvmt
    FIRMWARE_COMPLIANCE_MGMVMT_AVAILABLE = True
except ImportError:
    print("Warning: firmware_compliance_mgmvmt.py module not found. MG/MV/MT firmware compliance slide will not be updated.")
    FIRMWARE_COMPLIANCE_MGMVMT_AVAILABLE = False

try:
    import end_of_life
    END_OF_LIFE_AVAILABLE = True
except ImportError:
    print("Warning: end_of_life.py module not found. End of Life slides will not be updated.")
    END_OF_LIFE_AVAILABLE = False

# Import the adoption module
try:
    import adoption
    ADOPTION_AVAILABLE = True
except ImportError:
    print("Warning: adoption.py module not found. Meraki Product Adoption slide will not be added.")
    ADOPTION_AVAILABLE = False

# Import the new executive summary module
try:
    import executive_summary
    EXECUTIVE_SUMMARY_AVAILABLE = True
except ImportError:
    print("Warning: executive_summary.py module not found. Executive Summary slide will not be added.")
    EXECUTIVE_SUMMARY_AVAILABLE = False
try:
    import predictive_lifecycle
    PREDICTIVE_LIFECYCLE_AVAILABLE = True
except ImportError:
    print("Warning: predictive_lifecycle.py module not found.")
    PREDICTIVE_LIFECYCLE_AVAILABLE = False

# Import the PSIRT advisories module
try:
    import psirt_advisories
    PSIRT_ADVISORIES_AVAILABLE = True
except ImportError:
    print("Warning: psirt_advisories.py module not found. PSIRT Advisories slide will not be added.")
    PSIRT_ADVISORIES_AVAILABLE = False
    
# Constants
TEMPLATE_PATH = "template.pptx"  # Path to template PPTX file, if available
OUTPUT_PATH = "meraki_report.pptx"  # Default output path

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

def print_progress_bar(progress, total, prefix='Progress:', suffix='Complete', length=50, fill='█'):
    """
    Print a progress bar to the terminal.
    
    Args:
        progress: Current progress value
        total: Total value for 100% progress
        prefix: Text before the progress bar
        suffix: Text after the progress bar
        length: Character length of the progress bar
        fill: Character to use for the filled portion
    """
    percent = f"{100 * (progress / float(total)):.1f}%"
    filled_length = int(length * progress // total)
    bar = fill * filled_length + '-' * (length - filled_length)
    # Clear the current line and print the progress bar
    sys.stdout.write(f"\r{BLUE}{prefix} |{GREEN}{bar}{BLUE}| {percent} {suffix}{RESET}")
    sys.stdout.flush()
    # Print a new line if progress is complete
    if progress == total:
        print()

def delete_template_slide_3(output_path):
    """
    Delete slide 3 which is just inserted with the template.
    
    Args:
        output_path: Path to the PowerPoint file
    """
    try:
        # Load presentation
        #print(f"{BLUE}Opening PowerPoint to remove slide 3...{RESET}")
        prs = Presentation(output_path)
        
        # Make sure we have enough slides
        if len(prs.slides) >= 3:
            # Get the slide 3 (index 2)
            slide = prs.slides[2]
            slide_id = slide.slide_id
            
            # Get the parent element (slide list)
            slides = prs.slides._sldIdLst
            
            # Find and remove the slide
            for i, slide_element in enumerate(slides):
                if slide_element.id == slide_id:
                    slides.remove(slide_element)
                    #print(f"{GREEN}Removed template slide 3{RESET}")
                    break
            
            # Save the updated presentation
            prs.save(output_path)
            #print(f"{GREEN}PowerPoint saved with slide 3 removed{RESET}")
        else:
            print(f"{YELLOW}Slide 3 not found in the presentation{RESET}")
    
    except Exception as e:
        print(f"{RED}Error removing slide 3: {e}{RESET}")
        import traceback
        traceback.print_exc()

def delete_slides_for_missing_devices(output_path, device_types):
    """
    Delete slides for device types that aren't present in inventory.
    
    Args:
        output_path: Path to the PowerPoint file
        device_types: Dictionary of device types with boolean values indicating presence
    """
    try:
        # Load presentation
        #print(f"{BLUE}Opening PowerPoint to remove unnecessary slides...{RESET}")
        prs = Presentation(output_path)
        
        # Map of slide indices to check and delete
        # Note: These are 0-based indices (after slide 3 is removed)
        slides_to_check = {
            2: ("MX", device_types.get('has_mx_devices', False)),
            3: ("MS", device_types.get('has_ms_devices', False)),
            4: ("MR", device_types.get('has_mr_devices', False)),
            5: ("MV", device_types.get('has_mv_devices', False)),
            6: ("MG", device_types.get('has_mg_devices', False))
        }
        
        # We need to delete from higher indices to lower to avoid index shifting
        slides_to_delete = [(idx, device_type) for idx, (device_type, exists) in slides_to_check.items() 
                            if not exists]
        
        # Sort in reverse order so we delete from the end first
        slides_to_delete.sort(reverse=True)
        
        if slides_to_delete:
            #print(f"{BLUE}Will remove {len(slides_to_delete)} slides for missing device types:{RESET}")
            for idx, device_type in slides_to_delete:
                #print(f"  - Slide {idx + 1} ({device_type} Firmware Restrictions)")
                pass
            
            # Delete slides
            for idx, device_type in slides_to_delete:
                # Check if the index is still valid
                if idx < len(prs.slides):
                    # Get the XML element for the slide
                    slide = prs.slides[idx]
                    slide_id = slide.slide_id
                    
                    # Get the parent element (slide list)
                    slides = prs.slides._sldIdLst
                    
                    # Find and remove the slide
                    for i, slide_element in enumerate(slides):
                        if slide_element.id == slide_id:
                            slides.remove(slide_element)
                            #print(f"{GREEN}Removed slide {idx + 1} ({device_type} Firmware Restrictions){RESET}")
                            break
                else:
                    print(f"{YELLOW}Slide index {idx + 1} is out of range, skipping{RESET}")
            
            # Save the updated presentation
            prs.save(output_path)
            #print(f"{GREEN}PowerPoint saved with unnecessary slides removed{RESET}")
        else:
            print(f"{BLUE}All device types present, no slides need to be removed{RESET}")
    
    except Exception as e:
        print(f"{RED}Error removing slides: {e}{RESET}")
        import traceback
        traceback.print_exc()

async def main():
    """Main orchestration function."""
    # Start timer
    start_time = time.time()
    
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="Generate Meraki Dashboard Report in PowerPoint")
    parser.add_argument("-o", required=True, nargs='+', 
                        help="Space-separated list of Meraki organization IDs")
    parser.add_argument("-d", "--days", type=int, default=14,
                        help="Number of days to look back for client data (1-31, default: 14)")
    parser.add_argument("--output", default=OUTPUT_PATH,
                        help=f"Output path for PowerPoint file (default: {OUTPUT_PATH})")
    parser.add_argument("--template", default=TEMPLATE_PATH,
                        help=f"Path to PowerPoint template (default: {TEMPLATE_PATH})")
    parser.add_argument("--slides", default="all",
                        help="Comma-separated list of slide types to generate (default: all). Valid values: dashboard, mx, ms, mr, mv, mg, compliance-mxmsmr, compliance-mgmvmt, eol-summary, eol-detail, product-adoption, executive-summary, predictive-lifecycle, psirt-advisories")
    parser.add_argument("--debug", action="store_true",
                        help="Enable verbose debug output")
    parser.add_argument("--keep-all-slides", action="store_true",
                        help="Don't remove slides for missing device types")
    parser.add_argument("--secure-connect", action="store_true",
                        help="Indicate that the organization has Secure Connect deployed")
    parser.add_argument("--umbrella", action="store_true",
                        help="Indicate that the organization has Umbrella deployed")
    parser.add_argument("--thousand-eyes", action="store_true",
                        help="Indicate that the organization has Thousand Eyes deployed")
    parser.add_argument("--spaces", action="store_true",
                        help="Indicate that the organization has Spaces deployed")
    parser.add_argument("--xdr", action="store_true",
                        help="Indicate that the organization has XDR deployed")
    parser.add_argument("--no-progress-bar", action="store_true",
                        help="Disable progress bar display")
    parser.add_argument("--no-csv-export", action="store_true",
                        help="Disable automatic export of firmware compliance data to CSV files")
    
    args = parser.parse_args()
    
    # Enable debug mode if requested
    debug_mode = args.debug
    
    # Setup progress tracking
    use_progress_bar = not args.no_progress_bar
    progress_current = 0
    progress_total = 0  # Will be calculated based on slides to generate
    
    # Validate days input
    if args.days < 1 or args.days > 31:
        print("Error: Days must be between 1 and 31. Using default (14).")
        args.days = 14
    
    # Update paths based on arguments
    output_path = args.output
    template_path = args.template
    
    # Define the mapping between slide names and their corresponding indices/identifiers
    slide_mapping = {
        'dashboard': 2,
        'mx': 3,
        'ms': 4, 
        'mr': 5,
        'mv': 6,
        'mg': 7,
        'compliance-mxmsmr': 8,
        'compliance-mgmvmt': 9,
        'eol-summary': 10,
        'eol-detail': 11,
        'psirt-advisories': 12,  # PSIRT comes after both firmware compliance slides and EOL slides
        'product-adoption': 'product_adoption',
        'executive-summary': 'executive_summary',
        'predictive-lifecycle': 'predictive_lifecycle'
    }
    
    # Reverse mapping for validation
    reverse_mapping = {
        2: 'dashboard',
        3: 'mx',
        4: 'ms',
        5: 'mr',
        6: 'mv',
        7: 'mg',
        8: 'compliance-mxmsmr',
        9: 'compliance-mgmvmt',
        10: 'eol-summary',
        11: 'eol-detail',
        12: 'psirt-advisories',
        'product_adoption': 'product-adoption',
        'executive_summary': 'executive-summary',
        'predictive_lifecycle': 'predictive-lifecycle'
    }
    
    # Determine which slides are available
    available_slides = []
    if CLIENTS_AVAILABLE:
        available_slides.append(slide_mapping['dashboard'])
    if MX_FIRMWARE_AVAILABLE:
        available_slides.append(slide_mapping['mx'])
    if MS_FIRMWARE_AVAILABLE:
        available_slides.append(slide_mapping['ms'])
    if MR_FIRMWARE_AVAILABLE:
        available_slides.append(slide_mapping['mr'])
    if MV_FIRMWARE_AVAILABLE:
        available_slides.append(slide_mapping['mv'])
    if MG_FIRMWARE_AVAILABLE:
        available_slides.append(slide_mapping['mg'])
    if FIRMWARE_COMPLIANCE_MXMSMR_AVAILABLE:
        available_slides.append(slide_mapping['compliance-mxmsmr'])
    if FIRMWARE_COMPLIANCE_MGMVMT_AVAILABLE:
        available_slides.append(slide_mapping['compliance-mgmvmt'])
    if END_OF_LIFE_AVAILABLE:
        available_slides.append(slide_mapping['eol-summary'])
        available_slides.append(slide_mapping['eol-detail'])
    if PSIRT_ADVISORIES_AVAILABLE:
        available_slides.append(slide_mapping['psirt-advisories'])
        
    # Create a list for additional special slides
    special_slides = []
    if ADOPTION_AVAILABLE:
        special_slides.append(slide_mapping['product-adoption'])
    if EXECUTIVE_SUMMARY_AVAILABLE:
        special_slides.append(slide_mapping['executive-summary'])
    if PREDICTIVE_LIFECYCLE_AVAILABLE:
        special_slides.append(slide_mapping['predictive-lifecycle'])
        
    if args.slides.lower() == 'all':
        slides_to_generate = available_slides + special_slides
    else:
        try:
            requested_slides = []
            for slide in args.slides.split(','):
                slide = slide.strip().lower()
                if slide in slide_mapping:
                    requested_slides.append(slide_mapping[slide])
                else:
                    # Try to help with partial matches
                    possible_matches = [name for name in slide_mapping.keys() if slide in name]
                    if possible_matches:
                        print(f"{YELLOW}Slide type '{slide}' not found. Did you mean one of these: {', '.join(possible_matches)}?{RESET}")
                    else:
                        print(f"{YELLOW}Slide type '{slide}' not recognized.{RESET}")
            
            # If we have valid requested slides, use them; otherwise fall back to all available slides
            if requested_slides:
                slides_to_generate = requested_slides
            else:
                print(f"{YELLOW}No valid slide types specified. Using all available slides.{RESET}")
                slides_to_generate = available_slides + special_slides
        except Exception as e:
            print(f"{RED}Error parsing slide types: {e}. Using all available slides.{RESET}")
            slides_to_generate = available_slides + special_slides
    
    if not slides_to_generate:
        #print(f"{RED}No valid slides specified or no slide modules available. Exiting.{RESET}")
        return

    print(f"\n{BLUE}Starting Meraki Dashboard Report Generation{RESET}")
    
    # Set up progress tracking
    if use_progress_bar:
        # Set up a fixed total of 100 steps for easier percentage tracking
        progress_total = 100
        
        # Initialize progress
        progress_current = 0
        print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Variables to store all collected data
    dashboard_stats = None
    all_inventory_devices = []
    org_names = {}  # Store organization names
    
    # Variables to track device types
    device_types = {
        'has_mx_devices': False,
        'has_ms_devices': False,
        'has_mr_devices': False,
        'has_mv_devices': False,
        'has_mg_devices': False
    }
    
    # PHASE 1: Data Collection
    # First collect data for all slides without updating PowerPoint
    if 2 in slides_to_generate and CLIENTS_AVAILABLE:
        data_start_time = time.time()
        print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Collecting dashboard data...{RESET}")
        
        try:
            # Import functions from clients.py
            from clients import get_inventory_devices, get_api_key, AdaptiveRateLimiter
            from clients import get_networks, get_dashboard_stats, filter_incompatible_networks
            from clients import get_client_stats, get_organization_names
            import meraki.aio
            
            # Get API key and create rate limiter
            api_key = get_api_key()
            rate_limiter = AdaptiveRateLimiter()
            
            # Set up Meraki client - IMPORTANT: Keep Government API base URL
            async with meraki.aio.AsyncDashboardAPI(
                api_key=api_key,
                suppress_logging=True,
                maximum_retries=3,
                base_url="https://api.gov-meraki.com/api/v1"
            ) as aiomeraki:
                # Get organization names
                print(f"{BLUE}Getting organization names...{RESET}")
                org_names = await get_organization_names(aiomeraki, args.o, rate_limiter)
                
                # Get networks
                print(f"{BLUE}Getting networks...{RESET}")
                all_networks = []
                for org_id in args.o:
                    try:
                        networks = await get_networks(aiomeraki, org_id, rate_limiter)
                        all_networks.extend(networks)
                    except Exception as e:
                        print(f"{RED}Error retrieving networks for org {org_id}: {e}{RESET}")
                
                # Get network IDs
                network_ids = [network['id'] for network in all_networks]
                
                # Filter incompatible networks
                valid_network_ids = await filter_incompatible_networks(network_ids, all_networks)
                
                # Get dashboard statistics
                print(f"{BLUE}Getting dashboard statistics...{RESET}")
                dash_stats = await get_dashboard_stats(aiomeraki, args.o, rate_limiter)
                
                # Get client statistics
                print(f"{BLUE}Getting client statistics...This may take some time. Please be patient.{RESET}")
                
                # Update progress bar before getting client statistics
                if use_progress_bar:
                    # Set progress to 4%
                    progress_current = 4
                    print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
                client_stats = await get_client_stats(aiomeraki, valid_network_ids, rate_limiter, args.days)
                
                # Combine all stats
                dashboard_stats = {**dash_stats, **client_stats}
                
                # Get inventory devices for firmware restriction slides
                print(f"{BLUE}Getting inventory devices...{RESET}")
                
                # Update progress bar before getting inventory devices
                if use_progress_bar:
                    # Set progress to 8%
                    progress_current = 8
                    print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
                for org_id in args.o:
                    try:
                        devices = await get_inventory_devices(aiomeraki, org_id, rate_limiter)
                        all_inventory_devices.extend(devices)
                    except Exception as e:
                        print(f"{RED}Error retrieving inventory for org {org_id}: {e}{RESET}")
                
                # Detect which device types are present in inventory
                device_types['has_mx_devices'] = any(device.get('model', '').upper().startswith('MX') for device in all_inventory_devices)
                device_types['has_ms_devices'] = any(device.get('model', '').upper().startswith('MS') for device in all_inventory_devices)
                device_types['has_mr_devices'] = any(device.get('model', '').upper().startswith('MR') or device.get('model', '').upper().startswith('CW') for device in all_inventory_devices)
                device_types['has_mv_devices'] = any(device.get('model', '').upper().startswith('MV') for device in all_inventory_devices)
                device_types['has_mg_devices'] = any(device.get('model', '').upper().startswith('MG') for device in all_inventory_devices)
        
        except Exception as e:
            print(f"{RED}Error collecting data: {e}{RESET}")
            import traceback
            traceback.print_exc()
        
        data_time = time.time() - data_start_time
        print(f"{PURPLE}Data collection completed in {data_time:.2f} seconds{RESET}")
        
        # Update progress bar
        if use_progress_bar:
            progress_current = 12.5
            print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # PHASE 2: PowerPoint Updates
    # Now update PowerPoint with the collected data
    
    # First, copy template to output if needed
    if template_path != output_path and not os.path.exists(output_path):
        try:
            import shutil
            shutil.copy2(template_path, output_path)
        except Exception as e:
            print(f"{RED}Error copying template to output: {e}{RESET}")
    
    # Delete template slide 3 before adding content
    delete_template_slide_3(output_path)
    
    # Update progress bar
    if use_progress_bar:
        progress_current = 15
        print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Update slide 2 (Dashboard Summary) using the update_clients.py script
    if dashboard_stats and 2 in slides_to_generate and CLIENTS_AVAILABLE:
        slide_start_time = time.time()
        print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 2...{RESET}")
        
        try:
            # Convert dashboard_stats to JSON
            stats_json = json.dumps(dashboard_stats)
            
            # Also convert org_names to JSON for the update script
            org_names_json = json.dumps(org_names)
            
            # Call the update_slides.py script with the stats, days, and org_names
            result = subprocess.run(
                ["python3", "update_clients.py", template_path, output_path, stats_json, str(args.days), org_names_json],
                capture_output=True,
                text=True
            )
            
            # Print output from the script
            print(result.stdout)

            if result.stderr:
                print(f"{RED}Errors from update script:{RESET}")
                print(result.stderr)

            if result.returncode == 0:
                pass
            else:
                print(f"{RED}Error updating dashboard summary{RESET}")
                
        except Exception as e:
            print(f"{RED}Error updating slide 2: {e}{RESET}")
            import traceback
            traceback.print_exc()
        
        slide_time = time.time() - slide_start_time
        print(f"{PURPLE}Slide 2 update completed in {slide_time:.2f} seconds{RESET}")
        
        # Update progress bar
        if use_progress_bar:
            progress_current = 18.8
            print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Update slide 3 (MX Firmware Restrictions) using mx_firmware_restrictions.py
    if all_inventory_devices and 3 in slides_to_generate and MX_FIRMWARE_AVAILABLE:
        if device_types['has_mx_devices']:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 3...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call mx_firmware_restrictions's generate function
                if hasattr(mx_firmware_restrictions, 'generate'):
                    await mx_firmware_restrictions.generate(
                        api_client,
                        output_path,  # Use potentially already updated file as template
                        output_path,
                        inventory_devices=all_inventory_devices
                    )
                    print(f"{GREEN}Updated MX firmware restrictions in PowerPoint{RESET}")
                else:
                    print(f"{RED}mx_firmware_restrictions.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 3: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 3 update completed in {slide_time:.2f} seconds{RESET}")
            
            # Update progress bar
            if use_progress_bar:
                # After slide 3, progress to 25%
                progress_current = 25
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
        else:
            print(f"\n{YELLOW}Skipping slide 3 - No MX devices found in inventory{RESET}")
            
            # Update progress bar even if we skip
            if use_progress_bar and 3 in slides_to_generate:
                # After slide 3 (even if skipped), progress to 25%
                progress_current = 25
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Update slide 4 (MS Firmware Restrictions) using ms_firmware_restrictions.py
    if all_inventory_devices and 4 in slides_to_generate and MS_FIRMWARE_AVAILABLE:
        if device_types['has_ms_devices']:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 4...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call ms_firmware_restrictions's generate function
                if hasattr(ms_firmware_restrictions, 'generate'):
                    await ms_firmware_restrictions.generate(
                        api_client,
                        output_path,  # Use potentially already updated file as template
                        output_path,
                        inventory_devices=all_inventory_devices
                    )
                    print(f"{GREEN}Updated MS firmware restrictions in PowerPoint{RESET}")
                else:
                    print(f"{RED}ms_firmware_restrictions.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 4: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 4 update completed in {slide_time:.2f} seconds{RESET}")
            
            if use_progress_bar:
                progress_current = 31.2
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
        else:
            print(f"\n{YELLOW}Skipping slide 4 - No MS devices found in inventory{RESET}")
            
            if use_progress_bar and 4 in slides_to_generate:
                progress_current = 31.2
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Update slide 5 (MR Firmware Restrictions) using mr_firmware_restrictions.py
    if all_inventory_devices and 5 in slides_to_generate and MR_FIRMWARE_AVAILABLE:
        if device_types['has_mr_devices']:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 5...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call mr_firmware_restrictions's generate function
                if hasattr(mr_firmware_restrictions, 'generate'):
                    await mr_firmware_restrictions.generate(
                        api_client,
                        output_path,  # Use potentially already updated file as template
                        output_path,
                        inventory_devices=all_inventory_devices
                    )
                    print(f"{GREEN}Updated MR firmware restrictions in PowerPoint{RESET}")
                else:
                    print(f"{RED}mr_firmware_restrictions.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 5: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 5 update completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 5 - No MR devices found in inventory{RESET}")
    
    # Update slide 6 (MV Firmware Restrictions) using mv_firmware_restrictions.py
    if all_inventory_devices and 6 in slides_to_generate and MV_FIRMWARE_AVAILABLE:
        if device_types['has_mv_devices']:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 6...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call mv_firmware_restrictions's generate function
                if hasattr(mv_firmware_restrictions, 'generate'):
                    await mv_firmware_restrictions.generate(
                        api_client,
                        output_path,
                        output_path,
                        inventory_devices=all_inventory_devices
                    )
                    print(f"{GREEN}Updated MV firmware restrictions in PowerPoint{RESET}")
                else:
                    print(f"{RED}mv_firmware_restrictions.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 6: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 6 update completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 6 - No MV devices found in inventory{RESET}")
    
    # Update slide 7 (MG Firmware Restrictions) using mg_firmware_restrictions.py
    if all_inventory_devices and 7 in slides_to_generate and MG_FIRMWARE_AVAILABLE:
        if device_types['has_mg_devices']:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 7...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call mg_firmware_restrictions's generate function
                if hasattr(mg_firmware_restrictions, 'generate'):
                    await mg_firmware_restrictions.generate(
                        api_client,
                        output_path,  # Use potentially already updated file as template
                        output_path,
                        inventory_devices=all_inventory_devices
                    )
                    print(f"{GREEN}Updated MG firmware restrictions in PowerPoint{RESET}")
                else:
                    print(f"{RED}mg_firmware_restrictions.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 7: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 7 update completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 7 - No MG devices found in inventory{RESET}")
    
    # Update slide 8 (Firmware Compliance MX/MS/MR) using firmware_compliance_mxmsmr.py
    if 8 in slides_to_generate and FIRMWARE_COMPLIANCE_MXMSMR_AVAILABLE:
        # This slide requires networks data, which should be available from clients.py
        if 'all_networks' in locals() and all_networks:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 8...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                if hasattr(firmware_compliance_mxmsmr, 'generate'):
                    await firmware_compliance_mxmsmr.generate(
                        api_client,
                        output_path,
                        output_path,
                        networks=all_networks,
                        export_csv=not args.no_csv_export
                    )
                    print(f"{GREEN}Updated Firmware Compliance MX/MS/MR slide in PowerPoint{RESET}")
                else:
                    print(f"{RED}firmware_compliance_mxmsmr.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 8: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 8 update completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 8 - No networks data available{RESET}")
    
    # Update slide 9 (Firmware Compliance MG/MV/MT) using firmware_compliance_mgmvmt.py
    if 9 in slides_to_generate and FIRMWARE_COMPLIANCE_MGMVMT_AVAILABLE:
        # This slide requires networks data, which should be available from clients.py
        if 'all_networks' in locals() and all_networks:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 9...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                if hasattr(firmware_compliance_mgmvmt, 'generate'):
                    await firmware_compliance_mgmvmt.generate(
                        api_client,
                        output_path,
                        output_path,
                        networks=all_networks,
                        export_csv=not args.no_csv_export
                    )
                    print(f"{GREEN}Updated Firmware Compliance MG/MV/MT slide in PowerPoint{RESET}")
                else:
                    print(f"{RED}firmware_compliance_mgmvmt.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 9: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 9 update completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 9 - No networks data available{RESET}")
    
    # Update slide 10 (End of Life Products) using end_of_life.py
    if 10 in slides_to_generate and END_OF_LIFE_AVAILABLE:
        # This slide requires inventory devices data, which should be available
        if all_inventory_devices:
            slide_start_time = time.time()
                # Silenced to reduce terminal output
            # print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating slide 10...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call end_of_life's generate function
                if hasattr(end_of_life, 'generate'):
                    await end_of_life.generate(
                        api_client,
                        output_path,
                        output_path,
                        inventory_devices=all_inventory_devices,
                        networks=all_networks if 'all_networks' in locals() else None
                    )
                    print(f"{GREEN}Updated End of Life Products slide in PowerPoint{RESET}")
                else:
                    print(f"{RED}end_of_life.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error updating slide 10: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 10 update completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 10 - No inventory device data available{RESET}")

    # Update slide 11 (Device Models and EOL Dates) using end_of_life.py
    if 11 in slides_to_generate and END_OF_LIFE_AVAILABLE:
        # This slide requires inventory devices data, which should be available
        if all_inventory_devices:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Creating slide 11...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call end_of_life's generate_detail_slide function
                if hasattr(end_of_life, 'generate_detail_slide'):
                    await end_of_life.generate_detail_slide(
                        api_client,
                        output_path,  # Use potentially already updated file as template
                        output_path,
                        inventory_devices=all_inventory_devices,
                        networks=all_networks if 'all_networks' in locals() else None
                    )
                    print(f"{GREEN}Created Device Models and EOL Dates slide in PowerPoint{RESET}")
                else:
                    print(f"{RED}end_of_life.py doesn't have generate_detail_slide function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error creating slide 11: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Slide 11 creation completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping slide 11 - No inventory device data available{RESET}")
    
    # Add the PSIRT Advisories slide - This has been moved to execute AFTER both firmware compliance slides
    if 12 in slides_to_generate and PSIRT_ADVISORIES_AVAILABLE:
        psirt_start_time = time.time()
        print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Creating PSIRT Advisories slide...{RESET}")
        
        try:
            # Create simple API client
            class SimpleApiClient:
                def __init__(self, org_ids):
                    self.org_ids = org_ids
                    self.dashboard = None
            
            api_client = SimpleApiClient(args.o)
            
            # Call psirt_advisories's generate function
            if hasattr(psirt_advisories, 'generate'):
                await psirt_advisories.generate(
                    api_client,
                    output_path,
                    output_path,
                    inventory_devices=all_inventory_devices,
                    networks=all_networks if 'all_networks' in locals() else None
                )
                print(f"{GREEN}Created PSIRT Advisories slide in PowerPoint{RESET}")
            else:
                print(f"{RED}psirt_advisories.py doesn't have generate function{RESET}")
        
        except Exception as e:
            print(f"{RED}Error creating PSIRT Advisories slide: {e}{RESET}")
            import traceback
            traceback.print_exc()
        
        psirt_time = time.time() - psirt_start_time
        print(f"{PURPLE}PSIRT Advisories slide creation completed in {psirt_time:.2f} seconds{RESET}")
        
        if use_progress_bar:
            progress_current = 35
            print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Create Meraki Product Adoption slide using adoption.py
    products_adoption_data = None
    if 'product_adoption' in slides_to_generate and ADOPTION_AVAILABLE:
        # This slide requires inventory devices data, which should be available
        if all_inventory_devices:
            slide_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Creating Meraki Product Adoption slide...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Create manual configuration dictionary - this can be loaded from a config file or args
                manual_config = {
                    'Secure Connect': False,
                    'Umbrella Secure Internet Gateway': False,
                    'Thousand Eyes': False, 
                    'Spaces': False,
                    'XDR': False
                }
                
                # Check for manual override arguments
                if hasattr(args, 'secure_connect'):
                    manual_config['Secure Connect'] = args.secure_connect
                if hasattr(args, 'umbrella'):
                    manual_config['Umbrella Secure Internet Gateway'] = args.umbrella
                if hasattr(args, 'thousand_eyes'):
                    manual_config['Thousand Eyes'] = args.thousand_eyes
                if hasattr(args, 'spaces'):
                    manual_config['Spaces'] = args.spaces
                if hasattr(args, 'xdr'):
                    manual_config['XDR'] = args.xdr
                
                # Store products adoption data for potential use in executive summary
                products_adoption_data = {
                    'MX': device_types['has_mx_devices'],
                    'MS': device_types['has_ms_devices'],
                    'MR': device_types['has_mr_devices'],
                    'MG': device_types['has_mg_devices'],
                    'MV': device_types['has_mv_devices'],
                    'MT': any(device.get('model', '').upper().startswith('MT') for device in all_inventory_devices),
                    'Secure Connect': manual_config['Secure Connect'],
                    'Umbrella Secure Internet Gateway': manual_config['Umbrella Secure Internet Gateway'],
                    'Thousand Eyes': manual_config['Thousand Eyes'],
                    'Spaces': manual_config['Spaces'],
                    'XDR': manual_config['XDR']
                }
                
                # Call adoption's generate function
                if hasattr(adoption, 'generate'):
                    await adoption.generate(
                        api_client,
                        output_path,
                        output_path,
                        inventory_devices=all_inventory_devices,
                        networks=all_networks if 'all_networks' in locals() else None,
                        manual_config=manual_config
                    )
                    print(f"{GREEN}Created Meraki Product Adoption slide in PowerPoint{RESET}")
                else:
                    print(f"{RED}adoption.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error creating Meraki Product Adoption slide: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            slide_time = time.time() - slide_start_time
            print(f"{PURPLE}Meraki Product Adoption slide creation completed in {slide_time:.2f} seconds{RESET}")
        else:
            print(f"\n{YELLOW}Skipping Meraki Product Adoption slide - No inventory device data available{RESET}")
    
    # After all slides have been updated, delete unnecessary slides
    if all_inventory_devices and not args.keep_all_slides:
        delete_start_time = time.time()
        print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Removing slides for missing device types...{RESET}")
        
        # Delete slides for missing device types
        delete_slides_for_missing_devices(output_path, device_types)
        
        delete_time = time.time() - delete_start_time
        print(f"{PURPLE}Slide deletion completed in {delete_time:.2f} seconds{RESET}")
        
        if use_progress_bar:
            progress_current = 37.5
            print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Update title slide with organization names if not done through update_slides.py
    if org_names and not (2 in slides_to_generate and CLIENTS_AVAILABLE):
        title_start_time = time.time()
        print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Updating title slide...{RESET}")
        
        title_time = time.time() - title_start_time
        print(f"{PURPLE}Title slide update completed in {title_time:.2f} seconds{RESET}")
            
    # Finally, after ALL other slides are generated, create the Executive Summary slide
    firmware_compliance_data = None
    eol_data = None
    
    # Try to extract firmware compliance data from slides 8-9
    firmware_compliance_data = {}
    
    # First try to load from the JSON files created by the firmware compliance scripts
    try:
        # Try to load MXMSMR data
        if os.path.exists('mxmsmr_firmware_stats.json'):
            with open('mxmsmr_firmware_stats.json', 'r') as f:
                mxmsmr_data = json.load(f)
                mxmsmr_stats = mxmsmr_data.get('firmware_stats', {})
                mxmsmr_latest = mxmsmr_data.get('latest_versions', {})
                
                # Add to combined data
                for device_type in ['MX', 'MS', 'MR']:
                    if device_type in mxmsmr_stats:
                        firmware_compliance_data[device_type] = mxmsmr_stats[device_type]
                        # Ensure latest firmware version is included
                        if device_type in mxmsmr_latest:
                            firmware_compliance_data[device_type]['latest'] = mxmsmr_latest[device_type]
            
        # Try to load MGMVMT data
        if os.path.exists('mgmvmt_firmware_stats.json'):
            with open('mgmvmt_firmware_stats.json', 'r') as f:
                mgmvmt_data = json.load(f)
                mgmvmt_stats = mgmvmt_data.get('firmware_stats', {})
                mgmvmt_latest = mgmvmt_data.get('latest_versions', {})
                
                # Add to combined data
                for device_type in ['MG', 'MV', 'MT']:
                    if device_type in mgmvmt_stats:
                        firmware_compliance_data[device_type] = mgmvmt_stats[device_type]
                        # Ensure latest firmware version is included
                        if device_type in mgmvmt_latest:
                            firmware_compliance_data[device_type]['latest'] = mgmvmt_latest[device_type]
                        
    except Exception as e:
        print(f"{YELLOW}Error loading firmware data from JSON files: {e}{RESET}")
    
    # Fallback: Try to get data from global variables in the modules
    if not firmware_compliance_data and FIRMWARE_COMPLIANCE_MXMSMR_AVAILABLE:
        try:
            # Get data from global variable in firmware_compliance_mxmsmr module
            from firmware_compliance_mxmsmr import firmware_stats_mxmsmr
            for device_type in ['MX', 'MS', 'MR']:
                if device_type in firmware_stats_mxmsmr:
                    if device_type not in firmware_compliance_data:
                        firmware_compliance_data[device_type] = firmware_stats_mxmsmr[device_type]
        except Exception as e:
            print(f"{YELLOW}Could not extract firmware compliance data from MXMSMR module: {e}{RESET}")
            
    if not firmware_compliance_data and FIRMWARE_COMPLIANCE_MGMVMT_AVAILABLE:
        try:
            # Get data from global variable in firmware_compliance_mgmvmt module
            from firmware_compliance_mgmvmt import firmware_stats_mgmvmt
            for device_type in ['MG', 'MV', 'MT']:
                if device_type in firmware_stats_mgmvmt:
                    if device_type not in firmware_compliance_data:
                        firmware_compliance_data[device_type] = firmware_stats_mgmvmt[device_type]
        except Exception as e:
            print(f"{YELLOW}Could not extract firmware compliance data from MGMVMT module: {e}{RESET}")
            
    # Last resort fallback if all else fails
    if not firmware_compliance_data:
        print(f"{YELLOW}Could not extract firmware compliance data, creating empty template{RESET}")
        firmware_compliance_data = {
            'MX': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'latest': None},
            'MS': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'latest': None},
            'MR': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'latest': None},
            'MG': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'latest': None},
            'MV': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'latest': None},
            'MT': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'latest': None}
        }
    
    # Try to extract EOL data from slides 10-11
    if END_OF_LIFE_AVAILABLE:
        try:
            # Get the actual EOL data from the documentation
            from end_of_life import get_eol_info_from_doc
            eol_data, last_updated, is_from_doc = get_eol_info_from_doc()

        except Exception as e:
            print(f"{YELLOW}Could not fetch EOL data: {e}, using fallback{RESET}")
            # Only use fallback data if fetch fails
            try:
                from end_of_life import EOL_FALLBACK_DATA
                eol_data = EOL_FALLBACK_DATA
            except Exception as e2:
                print(f"{YELLOW}Could not import EOL data: {e2}, using None{RESET}")
    
    if 'executive_summary' in slides_to_generate and EXECUTIVE_SUMMARY_AVAILABLE:
        if all_inventory_devices:
            exec_summary_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Creating Executive Summary slide...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Call the executive summary's generate function
                if hasattr(executive_summary, 'generate'):
                    await executive_summary.generate(
                        api_client,
                        output_path,
                        output_path,
                        inventory_devices=all_inventory_devices,
                        networks=all_networks if 'all_networks' in locals() else None,
                        dashboard_stats=dashboard_stats,
                        firmware_stats=firmware_compliance_data,
                        eol_data=eol_data,
                        products=products_adoption_data
                    )
                    print(f"{GREEN}Created Executive Summary slide in PowerPoint (inserted at position 2){RESET}")
                else:
                    print(f"{RED}executive_summary.py doesn't have generate function{RESET}")
            
            except Exception as e:
                print(f"{RED}Error creating Executive Summary slide: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            exec_summary_time = time.time() - exec_summary_start_time
            print(f"{PURPLE}Executive Summary slide creation completed in {exec_summary_time:.2f} seconds{RESET}")
            
            # Update progress bar - Executive Summary complete
            if use_progress_bar:
                # After Executive Summary, progress to 43.8%
                progress_current = 43.8
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
        else:
            print(f"\n{YELLOW}Skipping Executive Summary slide - No inventory device data available{RESET}")
            
            # Update progress bar even if we skip
            if use_progress_bar and 'executive_summary' in slides_to_generate:
                # After Executive Summary (even if skipped), progress to 43.8%
                progress_current = 43.8
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    if 'predictive_lifecycle' in slides_to_generate and PREDICTIVE_LIFECYCLE_AVAILABLE:
        # This slide uses inventory devices, EOL data, and networks
        if all_inventory_devices:
            lifecycle_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Creating Predictive Lifecycle Management slides...{RESET}")
            
            try:
                # Create simple API client
                class SimpleApiClient:
                    def __init__(self, org_ids):
                        self.org_ids = org_ids
                        self.dashboard = None
                
                api_client = SimpleApiClient(args.o)
                
                # Extract EOL data if available
                eol_data = None
                if END_OF_LIFE_AVAILABLE:
                    try:
                        # Get the actual EOL data from the documentation
                        from end_of_life import get_eol_info_from_doc
                        eol_data, last_updated, is_from_doc = get_eol_info_from_doc()
                        
                        pass
                    except Exception as e:
                        print(f"{YELLOW}Could not fetch EOL data: {e}, using fallback{RESET}")
                        # Only use fallback data if fetch fails
                        try:
                            from end_of_life import EOL_FALLBACK_DATA
                            eol_data = EOL_FALLBACK_DATA
                        except Exception as e2:
                            print(f"{YELLOW}Could not import EOL data: {e2}, using None{RESET}")
                                
                if use_progress_bar:
                    progress_current = 66
                    print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
                
                # Generate the predictive lifecycle slides
                # Note: not specifying a position, so they'll be added at the end
                await predictive_lifecycle.generate(
                    api_client,
                    output_path,
                    output_path,
                    inventory_devices=all_inventory_devices,
                    networks=all_networks if 'all_networks' in locals() else None,
                    eol_data=eol_data
                )
                
                if use_progress_bar:
                    progress_current = 75
                    print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
                
                print(f"{GREEN}Added Predictive Lifecycle Management slides to the end of the presentation{RESET}")
                
            except Exception as e:
                print(f"{RED}Error creating Predictive Lifecycle Management slides: {e}{RESET}")
                import traceback
                traceback.print_exc()
            
            lifecycle_time = time.time() - lifecycle_start_time
            print(f"{PURPLE}Predictive Lifecycle slides creation completed in {lifecycle_time:.2f} seconds{RESET}")
            

            if use_progress_bar:
                progress_current = 95
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
        else:
            print(f"\n{YELLOW}Skipping Predictive Lifecycle slides - No inventory device data available{RESET}")
            

            if use_progress_bar and 'predictive_lifecycle' in slides_to_generate:
                progress_current = 95
                print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    # Calculate total script execution time
    total_time = time.time() - start_time
    
    if use_progress_bar and progress_current < 100:
        progress_current = 100
        print_progress_bar(progress_current, progress_total, prefix='Overall Progress:', suffix='Complete')
    
    print(f"\n{PURPLE}Total script execution time: {total_time:.2f} seconds{RESET}")
    print(f"\n{BLUE}Dashboard Report created successfully at {output_path}{RESET}")

def run_individual_slide(slide_type):
    """Helper function to run a single slide generator for debugging."""
    # Create a mapping between slide types and their respective modules
    debug_mapping = {
        'dashboard': (clients, 'main_async', [["123456"], 7, TEMPLATE_PATH, OUTPUT_PATH]),
        'mx': (mx_firmware_restrictions, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'ms': (ms_firmware_restrictions, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'mr': (mr_firmware_restrictions, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'mv': (mv_firmware_restrictions, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'mg': (mg_firmware_restrictions, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'compliance-mxmsmr': (firmware_compliance_mxmsmr, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'compliance-mgmvmt': (firmware_compliance_mgmvmt, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'eol-summary': (end_of_life, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'eol-detail': (end_of_life, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH, True]),
        'product-adoption': (adoption, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH, {
            'Secure Connect': True,
            'Umbrella Secure Internet Gateway': False,
            'Thousand Eyes': True,
            'Spaces': False,
            'XDR': False
        }]),
        'executive-summary': (executive_summary, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'predictive-lifecycle': (predictive_lifecycle, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH]),
        'psirt-advisories': (psirt_advisories, 'main_async', [["123456"], TEMPLATE_PATH, OUTPUT_PATH])
    }
    
    # Check if the slide type is in our mapping
    if slide_type in debug_mapping:
        module, function_name, args = debug_mapping[slide_type]
        
        # Check if the module is available
        module_name = module.__name__
        module_var_name = f"{module_name.upper()}_AVAILABLE"
        if module_name in globals() and globals().get(module_var_name, False):
            print(f"{YELLOW}Running {module_name}.py directly for debugging{RESET}")
            # Call the appropriate function with the args
            function = getattr(module, function_name)
            asyncio.run(function(*args))
        else:
            print(f"{RED}Module {module_name} is not available{RESET}")
    else:
        print(f"{RED}Invalid slide type: {slide_type}. Valid types are: {', '.join(debug_mapping.keys())}{RESET}")

if __name__ == "__main__":
    # Check for special debug flags
    if len(sys.argv) > 1 and sys.argv[1] == "--debug-clients":
        run_individual_slide('dashboard')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-mx":
        run_individual_slide('mx')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-ms":
        run_individual_slide('ms')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-mr":
        run_individual_slide('mr')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-mv":
        run_individual_slide('mv')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-mg":
        run_individual_slide('mg')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-compliance-mxmsmr":
        run_individual_slide('compliance-mxmsmr')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-compliance-mgmvmt":
        run_individual_slide('compliance-mgmvmt')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-eol-summary":
        run_individual_slide('eol-summary')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-eol-detail":
        run_individual_slide('eol-detail')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-adoption":
        run_individual_slide('product-adoption')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-executive-summary":
        run_individual_slide('executive-summary')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-predictive-lifecycle":
        run_individual_slide('predictive-lifecycle')
    elif len(sys.argv) > 1 and sys.argv[1] == "--debug-psirt-advisories":
        run_individual_slide('psirt-advisories')
    # Allow direct slide type debugging
    elif len(sys.argv) > 2 and sys.argv[1] == "--debug-slide":
        run_individual_slide(sys.argv[2])
    else:
        asyncio.run(main())
