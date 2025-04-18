import os
import sys
import asyncio
import time
import re
import requests
from bs4 import BeautifulSoup
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# MX firmware version restrictions - Hardcoded fallback values
MX_FIRMWARE_RESTRICTIONS = {
    "18.107.10": ["MX64", "MX65", "MX84", "MX100"],
    "16.16.9": ["MX400", "MX600"],
    "14.56": ["MX60", "MX80", "MX90", "Z1"]
}

# Models that can run current firmware - Hardcoded fallback
MX_UNRESTRICTED_MODELS = [
    "MX67", "MX68", "MX75", "MX85", "MX95", "MX105", 
    "MX250", "MX450", "Z4", "vMX"
]

# Model name normalizations for consistent counting
MX_MODEL_NORMALIZATIONS = {
    # MX64 variants
    "MX64": "MX64", "MX64W": "MX64",
    # MX65 variants
    "MX65": "MX65", "MX65W": "MX65",
    # MX67 variants
    "MX67": "MX67", "MX67W": "MX67", "MX67C": "MX67",
    # MX68 variants
    "MX68": "MX68", "MX68W": "MX68", "MX68CW": "MX68",
    # vMX variants
    "vMX": "vMX", "VMX": "vMX", 
    "vMX100": "vMX", "VMX100": "vMX",
    "vMX-S": "vMX", "VMX-S": "vMX",
    "vMX-M": "vMX", "VMX-M": "vMX",
    "vMX-L": "vMX", "VMX-L": "vMX",
    "vMX-XL": "vMX", "VMX-XL": "vMX"
}

def extract_last_updated_date(soup):
    """
    A more aggressive function to extract the last updated date from Meraki documentation.
    
    Args:
        soup: BeautifulSoup object of the parsed HTML
        
    Returns:
        str or None: The extracted date or None if not found
    """
    last_updated = None
    
    # APPROACH 1: Look for the exact "Last updated" text with asterisks (Markdown style)
    raw_html = str(soup)
    markdown_patterns = [
        r'\*\*Last updated\*\*\s*(Mar\s+\d+,?\s+2025)',
        r'\*\*Last updated\*\*\s*(March\s+\d+,?\s+2025)'
    ]
    
    for pattern in markdown_patterns:
        match = re.search(pattern, raw_html)
        if match:
            last_updated = match.group(1)
            # print(f"{GREEN}Found last updated date in Markdown: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 2: Look in the metadata section at the top of the page
    # Find any element that might contain metadata information
    meta_sections = soup.select('.page-metadata, .doc-metadata, .metadata, header p, .last-updated')
    for section in meta_sections:
        section_text = section.get_text()
        # Look for variations of "Last updated" followed by a date
        date_match = re.search(r'(?:Last\s+updated|Updated)(?:\s+on)?:?\s*(\w+\s+\d+,?\s+\d{4})', section_text, re.IGNORECASE)
        if date_match:
            last_updated = date_match.group(1)
            # print(f"{GREEN}Found last updated date in metadata: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 3: Look in the first few elements after the title
    # This targets the common pattern where the date appears right after the title
    title = soup.find('h1')
    if title:
        # Check the next few siblings of the title
        next_elements = []
        current = title.next_sibling
        for _ in range(5):
            if current:
                if hasattr(current, 'get_text'):
                    next_elements.append(current)
                current = current.next_sibling
        
        # Also get the parent and check its children
        parent = title.parent
        if parent:
            next_elements.extend(parent.find_all(recursive=False))
        
        # Check these elements for the date
        for elem in next_elements:
            if hasattr(elem, 'get_text'):
                elem_text = elem.get_text()
                date_match = re.search(r'(?:Last\s+updated|Updated)(?:\s+on)?:?\s*(\w+\s+\d+,?\s+\d{4})', elem_text, re.IGNORECASE)
                if date_match:
                    last_updated = date_match.group(1)
                    # print(f"{GREEN}Found last updated date near title: '{last_updated}'{RESET}")
                    return last_updated
    
    # APPROACH 4: Direct targeted extraction - look for specific March 2025 date
    specific_date_patterns = [
        r'(Mar\s+11,?\s+2025)',
        r'(March\s+11,?\s+2025)'
    ]
    
    for pattern in specific_date_patterns:
        match = re.search(pattern, raw_html)
        if match:
            last_updated = match.group(1)
            # print(f"{GREEN}Found specific date: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 5: Look for any date in the first part of the page
    # Find all text nodes in the first part of the document
    body = soup.find('body')
    if body:
        # Get the first ~20% of the HTML content
        first_part = str(body)[:int(len(str(body))*0.2)]
        # Look for any date with Mar/March 2025
        date_matches = re.findall(r'((?:Mar|March)\s+\d+,?\s+2025)', first_part)
        if date_matches:
            last_updated = date_matches[0]
            # print(f"{GREEN}Found date in first part of page: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 6: Look for "Last updated" line in any text node
    for tag in soup.find_all(['p', 'div', 'span']):
        if tag.string:
            text = tag.string.strip()
            if "Last updated" in text and "2025" in text:
                date_match = re.search(r'Last updated:?\s*(\w+\s+\d+,?\s+\d{4})', text, re.IGNORECASE)
                if date_match:
                    last_updated = date_match.group(1)
                    # print(f"{GREEN}Found last updated in clean text node: '{last_updated}'{RESET}")
                    return last_updated
    
    # print(f"{YELLOW}Could not find the last updated date using any extraction method{RESET}")
    return None

def get_firmware_restrictions_from_doc():
    """
    Attempt to fetch MX firmware restrictions from documentation.
    
    Returns:
        tuple: (firmware_restrictions dict, unrestricted_models list, last_updated string, is_from_doc bool)
    """
    try:
        # Attempt to fetch documentation
        # print(f"{BLUE}Attempting to fetch MX firmware information from documentation{RESET}")
        
        # Default fallback data for restrictions and models only (not date)
        fallback_restrictions = MX_FIRMWARE_RESTRICTIONS
        fallback_unrestricted = MX_UNRESTRICTED_MODELS
        
        # Use the correct URL for firmware information
        doc_url = "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions"
        
        # Add User-Agent header to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Make the request with a timeout and headers
        response = requests.get(doc_url, timeout=15, headers=headers)
        response.raise_for_status()
        
        # Get the raw HTML content
        html_content = response.text
        
        # TARGETED APPROACH: Extract date from meta tags and schema.org data
        last_updated = None
        
        # Look for meta tag with article:modified_time
        meta_match = re.search(r'<meta\s+property="article:modified_time"\s+content="([^"]+)"', html_content)
        if meta_match:
            iso_date = meta_match.group(1)
            # Convert ISO date to readable format
            try:
                import datetime
                dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                # print(f"{GREEN}Found last updated date in meta tag: '{last_updated}'{RESET}")
            except Exception as e:
                # If datetime conversion fails, use the raw date
                print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                last_updated = iso_date
        
        # If not found in meta tags, look for dateModified in JSON-LD
        if not last_updated:
            schema_match = re.search(r'"dateModified":"([^"]+)"', html_content)
            if schema_match:
                iso_date = schema_match.group(1)
                # Convert ISO date to readable format
                try:
                    import datetime
                    dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                    last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                    # print(f"{GREEN}Found last updated date in schema.org data: '{last_updated}'{RESET}")
                except Exception as e:
                    # If datetime conversion fails, use the raw date
                    print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                    last_updated = iso_date
                    
        # Now parse the HTML with BeautifulSoup for firmware restrictions
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Initialize collections for firmware data
        firmware_restrictions = {}  # model -> max firmware version
        unrestricted_models = []    # models that can run "Current" firmware
        
        # APPROACH #1: Look for tables with firmware information
        # print(f"{BLUE}Scanning tables for MX firmware information...{RESET}")
        
        tables = soup.find_all('table')
        
        for table in tables:
            # Check if this table might contain MX firmware information
            table_text = table.get_text().lower()
            if ('mx' in table_text and 'firmware' in table_text) or ('security appliance' in table_text and 'firmware' in table_text):
                # print(f"{BLUE}Found table with MX and firmware mentions{RESET}")
                pass
                
                # Check table headers to understand structure
                headers = []
                rows = table.find_all('tr')
                
                if rows:
                    header_cells = rows[0].find_all(['th', 'td'])
                    headers = [cell.get_text().strip().lower() for cell in header_cells]
                    # print(f"{BLUE}Table headers: {headers}{RESET}")
                
                # Find the relevant columns
                product_col = None
                max_firmware_col = None
                
                for i, header in enumerate(headers):
                    if any(term in header for term in ['product', 'model', 'appliance', 'device']):
                        product_col = i
                    if any(term in header for term in ['maximum', 'max', 'firmware restriction']):
                        max_firmware_col = i
                    elif 'firmware' in header and any(term in header for term in ['maximum', 'max', 'restriction']):
                        max_firmware_col = i
                
                # If we couldn't identify columns but "maximum runnable firmware" is in headers
                if product_col is None and max_firmware_col is None:
                    if 'maximum runnable firmware' in headers:
                        max_firmware_col = headers.index('maximum runnable firmware')
                        if 'product' in headers:
                            product_col = headers.index('product')
                        else:
                            product_col = 0
                
                # If we identified the needed columns, extract the data
                if product_col is not None and max_firmware_col is not None:
                    # print(f"{GREEN}Found table with product (col {product_col}) and max firmware (col {max_firmware_col}) columns{RESET}")
                    pass
                    
                    for row in rows[1:]:
                        cells = row.find_all(['td', 'th'])
                        
                        if len(cells) > max(product_col, max_firmware_col):
                            product_text = cells[product_col].get_text().strip()
                            max_firmware_text = cells[max_firmware_col].get_text().strip().lower()
                            
                            # Extract the base model (e.g., MX64 from MX64W)
                            mx_models = re.findall(r'(MX\d+\w*|Z\d+\w*|vMX\w*)', product_text, re.IGNORECASE)
                            
                            for model in mx_models:
                                # Check if this model has a firmware restriction or can run "Current"
                                if any(term in max_firmware_text for term in ['current', 'latest', 'newest', 'unrestricted']):
                                    if model not in unrestricted_models:
                                        unrestricted_models.append(model)
                                        # print(f"{GREEN}Found unrestricted model: {model} (can run Current firmware){RESET}")
                                else:
                                    # Extract version number
                                    version_match = re.search(r'(\d+(?:\.\d+)*)', max_firmware_text)
                                    if version_match:
                                        version = version_match.group(1)
                                        if version not in firmware_restrictions:
                                            firmware_restrictions[version] = []
                                        
                                        if model not in firmware_restrictions[version]:
                                            firmware_restrictions[version].append(model)
                                            # print(f"{GREEN}Found restriction: {model} -> MX {version}{RESET}")
        
        # APPROACH #2: Look for MX models and firmware mentions outside tables
        if not firmware_restrictions and not unrestricted_models:
            # print(f"{BLUE}Looking for MX firmware information in page text...{RESET}")
            pass
            
            # Look for specific patterns like "MX models that can support the latest firmware"
            page_text = soup.get_text()
            
            # Search for MX models followed by firmware info
            model_firmware_pattern = re.compile(r'(MX\d+\w*|Z\d+|vMX\w*).*?(?:maximum|restricted to|cannot run beyond).*?(?:firmware|version).*?(?:(current|latest)|(?:MX)?\s*(\d+(?:\.\d+)?))', re.IGNORECASE)
            
            for match in model_firmware_pattern.finditer(page_text):
                model = match.group(1)  # The MX model
                is_current = match.group(2)  # "current" or "latest" if matched
                version = match.group(3)  # Version number if matched
                
                if is_current:
                    # This model can run current firmware
                    if model not in unrestricted_models:
                        unrestricted_models.append(model)
                        # print(f"{GREEN}Found unrestricted model (text): {model} (can run Current firmware){RESET}")
                elif version:
                    # This model has a firmware restriction
                    if version not in firmware_restrictions:
                        firmware_restrictions[version] = []
                    
                    if model not in firmware_restrictions[version]:
                        firmware_restrictions[version].append(model)
                        # print(f"{GREEN}Found restriction (text): {model} -> MX {version}{RESET}")
        
        # If we found useful data, return it
        if firmware_restrictions or unrestricted_models:
            # Print summary of findings
            # print(f"{GREEN}Successfully parsed MX firmware information from documentation{RESET}")
            
            if firmware_restrictions:
                # print(f"Found {len(firmware_restrictions)} firmware restrictions:")
                # for version, models in sorted(firmware_restrictions.items(), key=lambda x: float(x[0].split('.')[0]), reverse=True):
                #     print(f"  - MX {version}: {len(models)} models - {', '.join(sorted(models))}")
                pass
            
            if unrestricted_models:
                # print(f"Found {len(unrestricted_models)} unrestricted models that can run Current firmware:")
                # print(f"  - {', '.join(sorted(unrestricted_models))}")
                pass
            
            return firmware_restrictions, unrestricted_models, last_updated, True
        else:
            # print(f"{YELLOW}Could not parse firmware information from documentation, using fallback{RESET}")
            return fallback_restrictions, fallback_unrestricted, None, False
            
    except Exception as e:
        print(f"{RED}Error fetching/parsing documentation: {e}{RESET}")
        import traceback
        traceback.print_exc()
        
        # Use fallback values but no fallback date
        print(f"{YELLOW}Using fallback firmware information{RESET}")
        return MX_FIRMWARE_RESTRICTIONS, MX_UNRESTRICTED_MODELS, None, False

def has_rgb_color(shape, target_rgb):
    """Check if shape has a line with the target RGB color, safely handling None cases."""
    if not hasattr(shape, 'line'):
        return False

    line = shape.line
    if line is None:
        return False
    
    if not hasattr(line, 'color') or line.color is None:
        return False

    if not hasattr(line.color, 'rgb') or line.color.rgb is None:
        return False
    
    return line.color.rgb == target_rgb

# Helper function to extract base model
def get_base_model(model):
    """Extract the base model (e.g., MX64 from MX64W)."""
    if not model:
        return None
        
    # Handle different cases
    model = model.strip().upper()
    
    # Extract the base model using regex
    base_match = re.search(r'(MX\d+|Z\d+|VMX[-\w]*)', model, re.IGNORECASE)
    return base_match.group(1) if base_match else model

# Helper function to normalize model names
def normalize_model_name(model):
    """Normalize the model name for consistent counting."""
    if not model:
        return None
        
    # Handle different cases
    model = model.strip().upper()
    
    # Check if model is in our normalization mapping
    for pattern, normalized in MX_MODEL_NORMALIZATIONS.items():
        if pattern.upper() == model:
            return normalized
        
    # If not found in mapping, use the base model
    base_model = get_base_model(model)
    return base_model if base_model else model

def get_model_firmware_version(model, firmware_restrictions, unrestricted_models):
    """
    Determine if a model has a firmware restriction, and if so, which version.
    
    Args:
        model: The full model string (e.g., MX64W)
        firmware_restrictions: Dict of firmware versions and their restricted models
        unrestricted_models: List of models that can run Current firmware
        
    Returns:
        str or None: The firmware version restriction or None if unrestricted
    """
    # Extract the base model
    base_model = get_base_model(model)
    
    if not base_model:
        return None  # Not a recognizable model
    
    # Check if the base model is in the unrestricted list
    for um in unrestricted_models:
        um_upper = um.upper()
        if base_model.upper() == um_upper or base_model.upper().startswith(um_upper):
            return None  # This model is unrestricted
    
    # Check if the base model is in any firmware restriction list
    for version, models in firmware_restrictions.items():
        for rm in models:
            rm_upper = rm.upper()
            if base_model.upper() == rm_upper or base_model.upper().startswith(rm_upper):
                return version  # Return the firmware restriction version
    
    # If not explicitly listed in either restricted or unrestricted, treat as unrestricted
    return None

async def generate(api_client, template_path, output_path, inventory_devices=None):
    """Generate the MX Firmware Restrictions slide."""
    print(f"\n{GREEN}Generating MX Firmware Restrictions slide (Slide 3)...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # If inventory_devices is provided, use it
    # Otherwise, would need to fetch it (not implemented here)
    if not inventory_devices:
        print(f"{RED}No inventory data provided{RESET}")
        return
    
    # print(f"{BLUE}Using inventory data provided from slide 1{RESET}")
    
    # Get firmware restrictions from documentation (or use hardcoded fallback)
    firmware_restrictions, unrestricted_models, last_updated_date, is_from_doc = get_firmware_restrictions_from_doc()
    
    # Log the source of firmware restrictions
    if is_from_doc:
        if last_updated_date:
            # print(f"{GREEN}Using MX firmware information from documentation (last updated: {last_updated_date}){RESET}")
            pass
        else:
            # print(f"{GREEN}Using MX firmware information from documentation (no date found){RESET}")
            pass
    else:
        # print(f"{YELLOW}Using fallback MX firmware information - documentation unavailable{RESET}")
        pass
    
    # Process MX device data
    process_start_time = time.time()
    #print(f"{PURPLE}[{time.strftime('%H:%M:%S')}] Processing MX device data...{RESET}")
    
    # Filter only MX devices, Z-Series, and vMX
    mx_devices = [device for device in inventory_devices 
                 if device.get('model', '').upper().startswith('MX') or
                    device.get('model', '').upper().startswith('VMX') or
                    device.get('model', '').upper().startswith('Z')]
    
    # Display the firmware restrictions data for verification
    # print(f"{BLUE}Firmware restrictions data:{RESET}")
    # for version, models in firmware_restrictions.items():
    #     print(f"  - MX {version}: {', '.join(sorted(models))}")
    # 
    # if unrestricted_models:
    #     print(f"{BLUE}Unrestricted models:{RESET}")
    #     print(f"  - {', '.join(sorted(unrestricted_models))}")
    
    # Count devices by firmware version and model
    restricted_devices = {}  # Devices with firmware restrictions
    unrestricted_devices = {}  # Devices without firmware restrictions
    total_mx_devices = len(mx_devices)
    
    # Group devices by their firmware restriction and model
    for device in mx_devices:
        model = device.get('model', 'unknown')
        
        # Skip if not an MX device (shouldn't happen due to filtering above, but just in case)
        if not model or not (model.upper().startswith('MX') or model.upper().startswith('VMX') or model.upper().startswith('Z')):
            continue
        
        # Normalize the model name for consistent counting
        normalized_model = normalize_model_name(model)
        
        # Get the firmware restriction for this model
        restricted_version = get_model_firmware_version(model, firmware_restrictions, unrestricted_models)
        
        if restricted_version:
            # This model has a firmware restriction
            if restricted_version not in restricted_devices:
                restricted_devices[restricted_version] = {}
            
            if normalized_model not in restricted_devices[restricted_version]:
                restricted_devices[restricted_version][normalized_model] = 0
            
            restricted_devices[restricted_version][normalized_model] += 1
        else:
            # This model doesn't have a specific restriction (is "Current")
            if normalized_model not in unrestricted_devices:
                unrestricted_devices[normalized_model] = 0
            
            unrestricted_devices[normalized_model] += 1
    
    # Print statistics for verification
    # print(f"{BLUE}MX Device Statistics:{RESET}")
    # print(f"Total MX devices found: {total_mx_devices}")
    # 
    # for version in restricted_devices:
    #     device_count = sum(restricted_devices[version].values())
    #     print(f"MX {version}: {device_count} devices")
    #     # Print device models
    #     for model, count in sorted(restricted_devices[version].items()):
    #         print(f"  - {model}: {count}")
    # 
    # unrestricted_count = sum(unrestricted_devices.values())
    # print(f"Not Firmware Restricted: {unrestricted_count} devices")
    # # Print unrestricted device models
    # for model, count in sorted(unrestricted_devices.items()):
    #     print(f"  - {model}: {count}")
    
    process_time = time.time() - process_start_time
    #print(f"{BLUE}MX data processing completed in {process_time:.2f} seconds{RESET}")
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    #print(f"{BLUE}Updating PowerPoint with MX data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)
        
        # Use slide 3 (index 2)
        # If the slide doesn't exist, add it
        if len(prs.slides) < 3:
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides[2]
        
        # Clear existing shapes except for title
        title_shape = None
        teal_line = None
        black_line = None
        
        # Look for existing title and lines
        for shape in slide.shapes:
            # Find title
            if hasattr(shape, "text_frame") and "MX Firmware Restrictions" in shape.text_frame.text:
                title_shape = shape
                continue
                
            # Find teal horizontal line
            if has_rgb_color(shape, RGBColor(80, 200, 192)):
                teal_line = shape
                continue
                
            # Find black horizontal line
            if has_rgb_color(shape, RGBColor(0, 0, 0)):
                black_line = shape
                continue
        
        # Create title if it doesn't exist
        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "MX Firmware Restrictions"
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            # print(f"{YELLOW}Added new title: 'MX Firmware Restrictions'{RESET}")
        else:
            # print(f"{BLUE}Found existing textbox title: 'MX Firmware Restrictions'{RESET}")
            pass
        
        # Remove all shapes except title and lines
        shapes_to_preserve = [title_shape, teal_line, black_line]
        shapes_to_remove = []
        
        for shape in slide.shapes:
            if shape not in shapes_to_preserve and shape is not None:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                if hasattr(shape, '_sp'):
                    sp = shape._sp
                    sp.getparent().remove(sp)
            except Exception as e:
                print(f"{RED}Error removing shape: {e}{RESET}")
        
        # print(f"{BLUE}Removing {len(shapes_to_remove)} shapes while preserving title and line{RESET}")
        
        # Add last updated date only if we have a date
        if last_updated_date:
            update_text = f"Firmware restriction last updated {last_updated_date}"
            update_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
            update_tf = update_box.text_frame
            update_p = update_tf.add_paragraph()
            update_p.text = update_text
            update_p.font.size = Pt(10)
            update_p.font.italic = True
        
        # Add an explanatory note to clarify what "firmware restrictions" means
        explanation_text = "Note: These values represent the maximum firmware versions these devices can run."
        explanation_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(8), Inches(0.3))
        explanation_tf = explanation_box.text_frame
        explanation_p = explanation_tf.add_paragraph()
        explanation_p.text = explanation_text
        explanation_p.font.size = Pt(10)
        explanation_p.font.italic = True
        
        # Setup style settings
        header_size = Pt(16)
        item_size = Pt(12)
        
        # Current Y position for content
        current_y = Inches(1.9)
        
        # Define column positions
        left_col_x = Inches(0.5)
        right_col_x = Inches(4.75)
        
        # Left Column - Not Firmware Restricted
        if unrestricted_devices:
            # Add header for unrestricted models
            header = slide.shapes.add_textbox(left_col_x, current_y, Inches(4), Inches(0.3))
            tf = header.text_frame
            p = tf.add_paragraph()
            p.text = "Not Firmware Restricted"
            p.font.size = header_size
            p.font.bold = True
            
            # Y position for content
            left_content_y = current_y + Inches(0.5)
            
            # Group models by type (MX, vMX, Z-Series)
            model_groups = {
                "Security Appliances": [],
                "Virtual Appliances": [],
                "Z-Series": []
            }
            
            for model, count in sorted(unrestricted_devices.items()):
                model_upper = model.upper()
                if model_upper.startswith('VMX') or 'VMX' in model_upper:
                    model_groups["Virtual Appliances"].append((model, count))
                elif model_upper.startswith('Z'):
                    model_groups["Z-Series"].append((model, count))
                else:
                    model_groups["Security Appliances"].append((model, count))
            
            # Add each group of models
            for group_title, models in model_groups.items():
                if not models:
                    continue
                    
                # Add group header
                group_header = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                tf = group_header.text_frame
                p = tf.add_paragraph()
                p.text = group_title + ":"
                p.font.size = item_size
                p.font.bold = True
                
                left_content_y += Inches(0.3)
                
                # Create model lines with good formatting
                model_lines = []
                current_line = ""
                
                for model, count in sorted(models):
                    model_text = f"{model} ({count})"
                    
                    # Check if adding this would make the line too long
                    if current_line and len(current_line) + len(model_text) + 2 > 40:
                        model_lines.append(current_line)
                        current_line = model_text
                    else:
                        if current_line:
                            current_line += ", " + model_text
                        else:
                            current_line = model_text
                
                if current_line:
                    model_lines.append(current_line)
                
                # Add each line to the slide
                for line in model_lines:
                    item = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                    tf = item.text_frame
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = item_size
                    
                    left_content_y += Inches(0.25)
                
                # Add space between groups
                left_content_y += Inches(0.1)

        sorted_versions = sorted(restricted_devices.keys(), 
                                key=lambda x: [int(n) for n in x.split('.')], 
                                reverse=True)
        
        if sorted_versions:
            right_content_y = current_y
            
            # Process each version in the right column
            for version_index, version in enumerate(sorted_versions):
                # Add firmware version header
                header = slide.shapes.add_textbox(right_col_x, right_content_y, Inches(4), Inches(0.3))
                tf = header.text_frame
                p = tf.add_paragraph()
                p.text = f"MX {version}"
                p.font.size = header_size
                p.font.bold = True
                
                right_content_y += Inches(0.4)
                
                # Add subtitle
                subtitle = slide.shapes.add_textbox(right_col_x + Inches(0.15), right_content_y, Inches(4), Inches(0.25))
                tf = subtitle.text_frame
                p = tf.add_paragraph()
                p.text = "Security Appliances:"
                p.font.size = item_size
                p.font.bold = True
                
                right_content_y += Inches(0.3)
                
                # Group models by type
                model_groups = {
                    "MX": [],
                    "Z-Series": []
                }
                
                for model, count in sorted(restricted_devices[version].items()):
                    if model.upper().startswith('Z'):
                        model_groups["Z-Series"].append((model, count))
                    else:
                        model_groups["MX"].append((model, count))
                
                # Add each model group
                for group_name, models in model_groups.items():
                    if not models:
                        continue
                    
                    # Create formatted model lines
                    model_lines = []
                    current_line = ""
                    
                    for model, count in sorted(models):
                        model_text = f"{model} ({count})"
                        
                        # Check if adding this would make the line too long
                        if current_line and len(current_line) + len(model_text) + 2 > 40:
                            model_lines.append(current_line)
                            current_line = model_text
                        else:
                            if current_line:
                                current_line += ", " + model_text
                            else:
                                current_line = model_text
                    
                    if current_line:
                        model_lines.append(current_line)
                    
                    for line in model_lines:
                        item = slide.shapes.add_textbox(right_col_x + Inches(0.15), right_content_y, Inches(4), Inches(0.25))
                        tf = item.text_frame
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = item_size
                        
                        right_content_y += Inches(0.25)
                
                right_content_y += Inches(0.3)
        
        # Add total count at the bottom right
        total_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(3), Inches(0.4))
        tf = total_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Total MX Devices: {total_mx_devices}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT
        
        # Add URL to slide notes (visible only to the presenter)
        documentation_url = "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions#MX"
        
        if hasattr(slide, 'notes_slide'):
            notes = slide.notes_slide
        else:
            notes = slide.notes_slide = prs.notes_master.clone_master_slide()
        
        # Clear any existing notes
        for shape in notes.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()
        
        # Add the URL to the slide notes
        notes_text_frame = notes.notes_text_frame
        note_p = notes_text_frame.add_paragraph()
        note_p.text = f"Source: {documentation_url}"
        note_p.font.size = Pt(12)
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Updated MX slide (Slide 3) with proper firmware categorization{RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    print(f"{PURPLE}MX Firmware Restrictions slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some sample inventory devices for testing
    inventory_devices = [
        {"model": "MX64", "firmware": "15.44.0", "networkId": "N1"},
        {"model": "MX64W", "firmware": "15.44.0", "networkId": "N2"},
        {"model": "MX65", "firmware": "15.44.0", "networkId": "N3"},
        {"model": "MX84", "firmware": "15.44.0", "networkId": "N4"},
        {"model": "MX100", "firmware": "15.44.0", "networkId": "N5"},
        {"model": "MX250", "firmware": "17.6.0", "networkId": "N6"},
        {"model": "MX450", "firmware": "17.6.0", "networkId": "N7"},
        {"model": "MX67", "firmware": "17.6.0", "networkId": "N8"},
        {"model": "MX68", "firmware": "17.6.0", "networkId": "N9"},
        {"model": "vMX100", "firmware": "17.6.0", "networkId": "N10"},
        {"model": "Z3", "firmware": "17.6.0", "networkId": "N11"}
    ]
    
    await generate(api_client, template_path, output_path, inventory_devices=inventory_devices)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python mx_firmware_restrictions.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))
