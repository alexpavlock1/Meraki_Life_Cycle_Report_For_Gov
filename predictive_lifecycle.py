import os
import sys
import asyncio
import time
import datetime
import numpy as np
import pandas as pd
from dateutil.relativedelta import relativedelta
from collections import defaultdict
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
import requests
from bs4 import BeautifulSoup
import re
import json
import logging
from concurrent.futures import ThreadPoolExecutor
from urllib.parse import urljoin
import random
# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color
try:
    from end_of_life import get_eol_info_from_doc, get_base_model, is_model_eol
    #print(f"{GREEN}Successfully imported EOL functions from end_of_life.py{RESET}")
except ImportError:
    print(f"{YELLOW}Could not import from end_of_life.py, using internal EOL functions{RESET}")
    
    # Define these functions directly if end_of_life.py isn't available
    def get_eol_info_from_doc():
        """
        Attempt to fetch EOL information from documentation.
        
        Returns:
            tuple: (eol_data dict, last_updated string, is_from_doc bool)
        """
        try:
            from end_of_life import EOL_FALLBACK_DATA, EOL_LAST_UPDATED
        except ImportError:
            # Define fallback data if even the constants aren't available
            EOL_FALLBACK_DATA = {
                # MX Devices
                "MX60": {"announcement": "Jul 10, 2015", "end_of_sale": "Oct 24, 2015", "end_of_support": "Oct 24, 2022"},
                "MX64": {"announcement": "Jan 26, 2022", "end_of_sale": "July 26, 2022", "end_of_support": "July 26, 2027"},
                "MX65": {"announcement": "Nov 20, 2018", "end_of_sale": "May 28, 2019", "end_of_support": "May 28, 2026"},
                "MX80": {"announcement": "Jan 26, 2016", "end_of_sale": "Aug 30, 2016", "end_of_support": "Aug 30, 2023"},
                "MX84": {"announcement": "Aug 10, 2021", "end_of_sale": "Oct 31, 2021", "end_of_support": "Oct 31, 2026"},
                "MX90": {"announcement": "Nov 5, 2013", "end_of_sale": "Apr 26, 2014", "end_of_support": "Apr 26, 2021"},
                "MX100": {"announcement": "Aug 10, 2021", "end_of_sale": "Feb 1, 2022", "end_of_support": "Feb 1, 2027"},
                "MX400": {"announcement": "Feb 28, 2018", "end_of_sale": "May 20, 2018", "end_of_support": "May 20, 2025"},
                "MX600": {"announcement": "Feb 28, 2018", "end_of_sale": "May 20, 2018", "end_of_support": "May 20, 2025"},
                "Z1": {"announcement": "April 27, 2018", "end_of_sale": "July 27, 2018", "end_of_support": "July 27, 2025"},
                "Z3": {"announcement": "Mar 4, 2024", "end_of_sale": "Sep 4, 2024", "end_of_support": "Sep 4, 2029"},
                "Z3C-HW-NA": {"announcement": "Mar 4, 2024", "end_of_sale": "Sep 4, 2024", "end_of_support": "Sep 4, 2029"},
                "Z3C-HW-WW": {"announcement": "Aug 11, 2023", "end_of_sale": "Feb 11, 2024", "end_of_support": "Feb 11, 2029"},
                "MX60W": {"announcement": "Jul 10, 2015", "end_of_sale": "Oct 24, 2015", "end_of_support": "Oct 24, 2022"},
                "MX64W": {"announcement": "Jan 26, 2022", "end_of_sale": "July 26, 2022", "end_of_support": "July 26, 2027"},
                "MX70": {"announcement": "Jan 18, 2012", "end_of_sale": "Mar 31, 2012", "end_of_support": "Mar 31, 2017"},
                }
            EOL_LAST_UPDATED = "April 4th, 2025"
            
        # Fallback data as defined in end_of_life.py
        return EOL_FALLBACK_DATA, EOL_LAST_UPDATED, False
    
    def get_base_model(model):
        """Extract the base model (e.g., MX64 from MX64-HW or MS220-8P from MS220-8P)."""
        base_match = re.match(r'(MR\d+|MS\d+|MX\d+|MV\d+|MG\d+|MT\d+|Z\d+|CW\d+)', model)
        if base_match:
            return base_match.group(1)
        
        # Try more specific regex patterns for each model line
        specific_patterns = [
            r'(MS\d+)-\d+',  # MS220-8P -> MS220
            r'(MX\d+)-\w+',  # MX64-HW -> MX64
            r'(MR\d+)-\w+',  # MR42-HW -> MR42
            r'(MV\d+)-\w+',  # MV12WE -> MV12
            r'(MG\d+)-\w+',  # MG21-HW -> MG21
            r'(MT\d+)-\w+',  # MT10-HW -> MT10
            r'(CW\d+)\w*'    # CW9166I -> CW9166
        ]
        
        for pattern in specific_patterns:
            match = re.match(pattern, model)
            if match:
                return match.group(1)
        
        return None
    
    def is_model_eol(model, eol_data):
        """
        Check if a model is in the EOL list and return its EOL info.
        """
        # Extract base model from full model string
        base_model = get_base_model(model)
        
        if not base_model:
            return None
        
        # Check exact model match
        if base_model in eol_data:
            return eol_data[base_model]
        
        # Try to match by model family (e.g., MS220-8P should match MS220)
        for eol_model in eol_data:
            if base_model.startswith(eol_model):
                return eol_data[eol_model]
        
        # Not found in EOL list
        return None
# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger("meraki_lifecycle")

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Colors for visualizations
HIGH_RISK_COLOR = RGBColor(229, 84, 81)     # Red
MEDIUM_RISK_COLOR = RGBColor(255, 192, 0)   # Amber
LOW_RISK_COLOR = RGBColor(112, 173, 71)     # Green
TEXT_COLOR = RGBColor(68, 68, 68)           # Dark Gray
ACCENT_COLOR = RGBColor(0, 112, 192)        # Blue

class RhinoPriceCatalog:
    """Class to fetch Meraki device pricing from Rhino Networks website."""
    
    def __init__(self, cache_file="meraki_rhino_prices_cache.json", cache_expiry_days=7):
        """Initialize the price catalog with Rhino Networks web scraping support."""
        self.base_url = "https://www.rhinonetworks.com"
        self.category_urls = {
            'MX': f"{self.base_url}/shop/category/security-appliances",
            'MS': f"{self.base_url}/shop/category/switching",
            'MR': f"{self.base_url}/shop/category/wireless-access-points",
            'CW': f"{self.base_url}/shop/category/wireless-access-points",
            'MV': f"{self.base_url}/shop/category/security-cameras",
            'Z': f"{self.base_url}/shop/category/security-appliances",
            'MT': f"{self.base_url}/series/device/environmental-sensors",
            'MG': f"{self.base_url}/series/device/wireless-wan",
            'License': f"{self.base_url}/shop/category/licenses"
        }
        self.cache_file = cache_file
        self.cache_expiry_days = cache_expiry_days
        
        # Logging setup
        self.logger = logging.getLogger("meraki_rhino_price_catalog")
        if not self.logger.handlers:
            handler = logging.StreamHandler()
            formatter = logging.Formatter("%(asctime)s [%(levelname)s] %(message)s")
            handler.setFormatter(formatter)
            self.logger.addHandler(handler)
            self.logger.setLevel(logging.INFO)
        
        # ANSI color codes for terminal output
        self.BLUE = '\033[94m'
        self.GREEN = '\033[92m'
        self.YELLOW = '\033[93m'
        self.RED = '\033[91m'
        self.RESET = '\033[0m'
        
        # Flag to track whether we're using fallback pricing
        self.using_fallback_pricing = False
        
        # Try to load prices from cache
        self.prices = self.load_cached_prices()
        
        # If no cached prices, try scraping
        if not self.prices:
            try:
                #self.logger.info(f"{self.BLUE}No valid cache found, fetching prices from Rhino Networks...{self.RESET}")
                
                # First try the direct URL approach (more reliable)
                if not self.scrape_direct_product_urls():
                    # If direct approach fails, try category pages
                    #self.logger.info(f"{self.BLUE}Direct URL approach failed, trying category pages...{self.RESET}")
                    self.fetch_all_prices()
                    
                    # Check if we got any products
                    total_prices = sum(len(family_prices) for family_prices in self.prices.values())
                    
                    # If no products found, use fallback prices
                    if total_prices == 0:
                        self.logger.warning(f"{self.YELLOW}No products found during scraping. Using fallback prices.{self.RESET}")
                        self.prices = self.get_fallback_prices()
                        self.using_fallback_pricing = True
                
            except Exception as e:
                self.logger.error(f"{self.RED}Failed to fetch prices from Rhino Networks: {e}{self.RESET}")
                self.prices = self.get_fallback_prices()
                self.using_fallback_pricing = True
                self.logger.warning(f"{self.YELLOW}Using fallback price estimates. Pricing information may not be accurate.{self.RESET}")
        
        # Track price lookup misses for reporting
        self.price_misses = set()
    def scrape_device_license_costs(self):
        """
        Scrape 1-year license costs for Meraki devices from the Rhino Networks license page.
        Maps license costs to corresponding hardware models.
        Focuses on base Enterprise licensing only.
        """
        #self.logger.info(f"{self.BLUE}Scraping 1-year license costs for Meraki devices...{self.RESET}")
        
        # Initialize license cost dictionary
        license_costs = {}
        
        # List of license pages to scrape (base URL + pagination)
        base_url = "https://www.rhinonetworks.com/shop/category/licenses"
        page_urls = [f"{base_url}?page={i}" for i in range(7)]  # Pages 0-6
        
        # Define product types to skip
        skip_product_keywords = [
            "display", 
            "insight", 
            "cloud archive", 
            "sd-wan plus",
            "secure sd-wan plus",
            "advanced security", 
            "per device sd-wan"
        ]
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
        }
        
        session = requests.Session()
        
        for page_url in page_urls:
            try:
                #self.logger.info(f"{self.BLUE}Fetching license page: {page_url}{self.RESET}")
                
                # Add delay to avoid being blocked
                time.sleep(random.uniform(1, 3))
                
                response = session.get(page_url, headers=headers)
                if response.status_code != 200:
                    #self.logger.warning(f"{self.YELLOW}Failed to fetch license page, status: {response.status_code}{self.RESET}")
                    continue
                    
                # Parse the HTML
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Find all product cards (updated selector to match current HTML)
                product_cards = soup.select('article.product__card')
                #self.logger.info(f"{self.GREEN}Found {len(product_cards)} license products on page{self.RESET}")
                
                # Save the first page for debugging
                #if page_url == base_url:
                    #with open('license_page_debug.html', 'w', encoding='utf-8') as f:
                        #f.write(response.text)
                        #self.logger.info(f"{self.GREEN}Saved HTML to license_page_debug.html for debugging{self.RESET}")
                
                # Process each product card
                for card in product_cards:
                    try:
                        # Get the title element to extract the license info
                        title_elem = card.select_one('.product__card__title h3')
                        if not title_elem:
                            continue
                            
                        # Get the title text
                        title_text = title_elem.text.strip()
                        
                        # Skip products we don't want to process
                        should_skip = any(keyword.lower() in title_text.lower() for keyword in skip_product_keywords)
                        if should_skip:
                            #self.logger.info(f"{self.YELLOW}Skipping product: {title_text} (excluded type){self.RESET}")
                            continue
                        
                        #self.logger.info(f"{self.BLUE}Processing product: {title_text}{self.RESET}")
                        
                        # Updated regex patterns to match current website format
                        # Pattern 1: "Meraki MS120-24 License" -> MS120-24
                        # Pattern 2: "Meraki C9300-24 Enterprise License" -> C9300-24, ENT
                        # Pattern 3: "Meraki MG21 License" -> MG21
                        # Pattern 4: "Meraki MX85 License" -> MX85
                        license_patterns = [
                            # Try to extract model and license type like "Meraki MR Advanced License"
                            r'Meraki ([A-Z]+) (Advanced|Enterprise) License',
                            # Try to extract model and license type like "Meraki C9300-24 Enterprise License" 
                            r'Meraki ([A-Z0-9]+-[0-9]+(?:[A-Z]+)?) (Advanced|Enterprise) License',
                            # Try to extract model like "Meraki MS120-24 License"
                            r'Meraki ([A-Z][A-Z0-9]+-?[0-9]*(?:[A-Z]+)?) License',
                            # Specifically for MX models
                            r'Meraki (MX[0-9]+[A-Z]*)',
                            # Legacy patterns for backward compatibility
                            r'LIC-([A-Z0-9]+)-([A-Z]+)-1YR',
                            r'LIC-([A-Z0-9]+)-1YR',
                            r'LICENSE-1YR-([A-Z0-9]+)',
                            r'1 ?YR ?(?:LICENSE)? ?(?:FOR)? ?([A-Z0-9-]+)',
                        ]
                        
                        hardware_model = None
                        license_type = "ENT"  # Default to Enterprise license type
                        
                        for pattern in license_patterns:
                            match = re.search(pattern, title_text, re.IGNORECASE)
                            if match:
                                #self.logger.info(f"{self.GREEN}Pattern match: {pattern} for {title_text}{self.RESET}")
                                if len(match.groups()) == 2:
                                    hardware_model = match.group(1).upper()
                                    # Convert "Enterprise" to "ENT" and "Advanced" to "ADV"
                                    license_type_text = match.group(2).upper()
                                    if license_type_text == "ENTERPRISE":
                                        license_type = "ENT"
                                    elif license_type_text == "ADVANCED":
                                        license_type = "ADV"
                                    else:
                                        license_type = license_type_text
                                else:
                                    hardware_model = match.group(1).upper()
                                break
                        
                        # If we couldn't extract with regex, try manual parsing
                        if not hardware_model:
                            #self.logger.info(f"{self.YELLOW}No regex match for {title_text}, trying manual parsing{self.RESET}")
                            # Look for common model prefixes in the title
                            for prefix in ["MX", "MS", "MR", "MV", "MT", "MG", "Z", "C9"]:
                                if prefix in title_text:
                                    model_match = re.search(f"{prefix}[0-9]+" + r"[A-Z0-9-]*", title_text)
                                    if model_match:
                                        hardware_model = model_match.group(0)
                                        #self.logger.info(f"{self.GREEN}Manual extraction found model: {hardware_model}{self.RESET}")
                                        break
                        
                        # Skip if we couldn't identify the hardware model
                        if not hardware_model:
                            #self.logger.warning(f"{self.YELLOW}Could not extract model from: {title_text}{self.RESET}")
                            continue
                        
                        # Get the price element
                        price_elem = card.select_one('.product__card__price')
                        if not price_elem:
                            #self.logger.warning(f"{self.YELLOW}No price element found for {hardware_model}{self.RESET}")
                            continue
                        
                        # Extract the discounted price (in the <strong> tag)
                        strong_price = price_elem.select_one('strong')
                        if not strong_price:
                            #self.logger.warning(f"{self.YELLOW}No strong price element found for {hardware_model}{self.RESET}")
                            continue
                        
                        # Extract the price value - allow for commas in prices like $1,199.99
                        price_text = strong_price.text.strip()
                        #self.logger.info(f"{self.BLUE}Found price text: {price_text} for {hardware_model}{self.RESET}")
                        
                        price_match = re.search(r'\$([0-9,.]+)', price_text)
                        if not price_match:
                            #self.logger.warning(f"{self.YELLOW}Could not extract price from: {price_text} for {hardware_model}{self.RESET}")
                            continue
                            
                        price = float(price_match.group(1).replace(',', ''))
                        
                        # Skip Advanced licenses if we're only interested in Enterprise
                        if license_type == "ADV":
                            #self.logger.info(f"{self.YELLOW}Skipping Advanced license for {hardware_model} (focusing on Enterprise only){self.RESET}")
                            continue
                        
                        # Store in license_costs dictionary
                        # Key format: (model, license_type)
                        license_costs[(hardware_model, license_type)] = price
                        #self.logger.info(f"{self.GREEN}Found 1YR {license_type} license for {hardware_model}: ${price}{self.RESET}")
                        
                    except Exception as e:
                        self.logger.error(f"{self.RED}Error extracting license info: {e}{self.RESET}")

                        self.logger.error(f"{self.RED}{traceback.format_exc()}{self.RESET}")
                        continue
            
            except Exception as e:
                self.logger.error(f"{self.RED}Error fetching license page {page_url}: {e}{self.RESET}")
                continue
        
        # Add fallback license prices for common models that might be missing
        if len(license_costs) < 20:  # If we didn't find many licenses, add fallbacks
            self.logger.warning(f"{self.YELLOW}Found only {len(license_costs)} licenses, adding fallbacks{self.RESET}")
            self._add_fallback_license_costs(license_costs)
        
        #self.logger.info(f"{self.GREEN}Completed license cost scraping. Found {len(license_costs)} license costs.{self.RESET}")
        return license_costs

    def _add_fallback_license_costs(self, license_costs):
        """
        Add fallback license costs for common models that might be missing.
        """
        fallback_costs = [
            # MX models
            {"model": "MX67", "license_type": "ENT", "price": 195.00},
            {"model": "MX68", "license_type": "ENT", "price": 250.00},
            {"model": "MX75", "license_type": "ENT", "price": 300.00},
            {"model": "MX84", "license_type": "ENT", "price": 495.00},
            {"model": "MX85", "license_type": "ENT", "price": 550.00},
            {"model": "MX95", "license_type": "ENT", "price": 750.00},
            {"model": "MX105", "license_type": "ENT", "price": 995.00},
            {"model": "MX250", "license_type": "ENT", "price": 2995.00},
            {"model": "MX450", "license_type": "ENT", "price": 5995.00},
            
            # MS models (switches by port count)
            {"model": "MS120-8", "license_type": "ENT", "price": 50.00},
            {"model": "MS120-24", "license_type": "ENT", "price": 95.00},
            {"model": "MS120-48", "license_type": "ENT", "price": 195.00},
            {"model": "MS210-24", "license_type": "ENT", "price": 125.00},
            {"model": "MS210-48", "license_type": "ENT", "price": 225.00},
            {"model": "MS225-24", "license_type": "ENT", "price": 150.00},
            {"model": "MS225-48", "license_type": "ENT", "price": 250.00},
            {"model": "MS250-24", "license_type": "ENT", "price": 175.00},
            {"model": "MS250-48", "license_type": "ENT", "price": 295.00},
            {"model": "MS350-24", "license_type": "ENT", "price": 225.00},
            {"model": "MS350-48", "license_type": "ENT", "price": 345.00},
            {"model": "MS390-24", "license_type": "ENT", "price": 295.00},
            {"model": "MS390-48", "license_type": "ENT", "price": 495.00},
            {"model": "MS410-16", "license_type": "ENT", "price": 350.00},
            {"model": "MS410-32", "license_type": "ENT", "price": 650.00},
            {"model": "MS425-16", "license_type": "ENT", "price": 450.00},
            {"model": "MS425-32", "license_type": "ENT", "price": 850.00},
            
            # MR models (wireless)
            {"model": "MR36", "license_type": "ENT", "price": 150.00},
            {"model": "MR44", "license_type": "ENT", "price": 175.00},
            {"model": "MR46", "license_type": "ENT", "price": 195.00},
            {"model": "MR56", "license_type": "ENT", "price": 225.00},
            {"model": "MR57", "license_type": "ENT", "price": 250.00},
            
            # MV models (cameras) - using average prices
            {"model": "MV2", "license_type": "ENT", "price": 195.00},
            {"model": "MV32", "license_type": "ENT", "price": 225.00},
            {"model": "MV72", "license_type": "ENT", "price": 250.00},
            
            # Z models
            {"model": "Z3", "license_type": "ENT", "price": 95.00},
            {"model": "Z4", "license_type": "ENT", "price": 125.00},
        ]
        
        # Add fallbacks only for models not already in the dictionary
        for item in fallback_costs:
            model = item["model"]
            license_type = item["license_type"]
            price = item["price"]
            
            # Check if we already have this specific model+license
            if (model, license_type) not in license_costs:
                license_costs[(model, license_type)] = price
                self.logger.info(f"{self.BLUE}Added fallback 1YR {license_type} license for {model}: ${price}{self.RESET}")

    def discover_new_models(self):
        """
        Enhanced method to scan Rhino Networks category pages to discover new Meraki models.
        
        Returns:
            dict: A dictionary mapping model names to their product URLs
        """
        #self.logger.info(f"{self.BLUE}Starting dynamic model discovery...{self.RESET}")
        
        discovered_models = {}
        
        # Headers for requests - Use a more realistic browser User-Agent
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
            'Accept-Language': 'en-US,en;q=0.9',
            'Referer': 'https://www.rhinonetworks.com/',
            'Cache-Control': 'no-cache',
        }
        
        session = requests.Session()
        
        # Process each category page
        for family, category_url in self.category_urls.items():
            #self.logger.info(f"{self.BLUE}Scanning category: {family} at {category_url}{self.RESET}")
            
            current_url = category_url
            page_num = 1
            
            # Process all pages in this category
            while current_url:
                try:
                    #self.logger.info(f"{self.BLUE}Scanning page {page_num} for {family}...{self.RESET}")
                    
                    # Add random delay to avoid being blocked
                    time.sleep(random.uniform(2, 4))
                    
                    # Fetch the page
                    response = session.get(current_url, headers=headers)
                    if response.status_code != 200:
                        self.logger.warning(f"{self.YELLOW}Failed to fetch page {current_url}, status: {response.status_code}{self.RESET}")
                        break
                    
                    # Save a sample page for debugging
                    #if page_num == 1:
                        #with open(f'rhino_{family}_page1.html', 'w', encoding='utf-8') as f:
                            #f.write(response.text)
                            #self.logger.info(f"{self.GREEN}Saved page HTML to rhino_{family}_page1.html for debugging{self.RESET}")
                    
                    # Parse the page
                    soup = BeautifulSoup(response.text, 'html.parser')
                    
                    # Find all product links - try multiple possible selectors
                    product_links = []
                    
                    # Try various selectors based on Rhino's site structure
                    # Product cards
                    cards = soup.select('.product__card') or soup.select('article.card') or soup.select('.card')
                    for card in cards:
                        parent_link = card.find_parent('a')
                        if parent_link and 'href' in parent_link.attrs:
                            product_links.append(parent_link['href'])
                    
                    # Direct links to product pages (expanded selectors)
                    link_selectors = [
                        'a[href*="/product/device/meraki-"]',
                        'a[href*="/product/license/meraki-"]',
                        'a.card__link',
                        '.product-card a',
                        '.product a',
                        'a[href*="meraki"]'
                    ]
                    
                    for selector in link_selectors:
                        links = soup.select(selector)
                        for link in links:
                            if 'href' in link.attrs:
                                product_links.append(link['href'])
                    
                    # Extract product names from headings
                    product_headings = soup.select('.product-title') or soup.select('h2.card__title') or soup.select('h3')
                    for heading in product_headings:
                        # Look for Meraki model patterns
                        text = heading.text.strip()
                        model_patterns = [
                            r'(MR[0-9]+[A-Z]?)',
                            r'(MS[0-9]+-[0-9]+[A-Z]*)',
                            r'(MX[0-9]+[A-Z]?)',
                            r'(MV[0-9]+[A-Z]?)',
                            r'(MG[0-9]+[A-Z]?)',
                            r'(MT[0-9]+[A-Z]?)',
                            r'(Z[0-9]+[A-Z]?)',
                            r'(CW[0-9]+[A-Z]?)',
                            r'(C[0-9]+[A-Z]*-[0-9]+[A-Z]*)'
                        ]
                        
                        for pattern in model_patterns:
                            match = re.search(pattern, text)
                            if match:
                                model = match.group(1)
                                # Find parent link if available
                                parent = heading.find_parent('a')
                                if parent and 'href' in parent.attrs:
                                    product_links.append(parent['href'])
                                    #self.logger.info(f"{self.GREEN}Found model in heading: {model} - {parent['href']}{self.RESET}")
                    
                    # Make links unique
                    product_links = list(set(product_links))
                    
                    # Process each product link
                    for link in product_links:
                        # Make sure the link is absolute
                        full_link = urljoin(self.base_url, link)
                        
                        # Extract the model from the URL
                        model_match = re.search(r'/meraki-([a-zA-Z0-9-]+)(?:/|\?|$)', full_link)
                        if model_match:
                            model_name = model_match.group(1).upper()
                            
                            # Add prefix to model name if needed to match format (MX, MS, MR, etc.)
                            if not re.match(r'^(MR|MS|MX|MV|MG|MT|Z|LIC|CW|C[0-9])', model_name):
                                # Try to determine family from URL structure
                                family_prefix = None
                                if 'security-appliances' in full_link:
                                    if model_name.startswith('Z'):
                                        family_prefix = 'Z'
                                    else:
                                        family_prefix = 'MX'
                                elif 'switching' in full_link:
                                    family_prefix = 'MS'
                                elif 'wireless' in full_link or 'access-points' in full_link:
                                    family_prefix = 'MR'
                                elif 'cameras' in full_link:
                                    family_prefix = 'MV'
                                elif 'sensors' in full_link:
                                    family_prefix = 'MT'
                                elif 'cellular' in full_link or 'wan' in full_link:
                                    family_prefix = 'MG'
                                elif 'catalyst' in full_link:
                                    family_prefix = 'C'
                                
                                if family_prefix:
                                    model_name = f"{family_prefix}{model_name}"
                            
                            # Make model name uppercase for consistency
                            model_name = model_name.upper()
                            
                            # Store in our results
                            discovered_models[model_name] = full_link
                            #self.logger.info(f"{self.GREEN}Discovered model: {model_name} at {full_link}{self.RESET}")
                    
                    # Check for next page - extended selectors to find pagination
                    next_page = None
                    pagination_selectors = [
                        '.pagination', 
                        '.pager',
                        'nav[aria-label="pagination"]',
                        '.page-numbers'
                    ]
                    
                    for selector in pagination_selectors:
                        pagination = soup.select_one(selector)
                        if pagination:
                            # Find the "next" link
                            next_links = pagination.select('a')
                            for link in next_links:
                                if ('next' in link.get('class', []) or 
                                    '›' in link.text or 
                                    'Next' in link.text or 
                                    'next' in link.text.lower() or
                                    'arrow-right' in str(link)):
                                    href = link.get('href')
                                    if href:
                                        next_page = urljoin(self.base_url, href)
                                        break
                    
                    # Move to next page if available
                    if next_page:
                        current_url = next_page
                        page_num += 1
                        # Be kind to the server
                        time.sleep(random.uniform(2, 3))
                    else:
                        break
                        
                except Exception as e:
                    self.logger.error(f"{self.RED}Error scanning page {current_url}: {e}{self.RESET}")
                    break
        
        #self.logger.info(f"{self.GREEN}Model discovery complete. Found {len(discovered_models)} Meraki models.{self.RESET}")
        return discovered_models

    def get_license_product_urls(self):
        """
        Get the correct URLs for license products.
        This handles VMX, Insight, and other Meraki license types.
        """
        # URLs for Insight licenses
        insight_licenses = [
            {"url": "https://www.rhinonetworks.com/product/license/meraki-insight-license-xsmall", "model": "INSIGHT-LICENSE-XSMALL"},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-insight-license-small", "model": "INSIGHT-LICENSE-SMALL"},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-insight-license-medium", "model": "INSIGHT-LICENSE-MEDIUM"},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-insight-license-large", "model": "INSIGHT-LICENSE-LARGE"},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-insight-license-xlarge", "model": "INSIGHT-LICENSE-XLARGE"}
        ]
        
        # URLs for VMX licenses with -license suffix
        vmx_licenses = [
            {"url": "https://www.rhinonetworks.com/product/license/meraki-vmx-small-license", "model": "VMX-SMALL-LICENSE", "fallback_price": 395.00},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-vmx-medium-license", "model": "VMX-MEDIUM-LICENSE", "fallback_price": 995.00},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-vmx-large-license", "model": "VMX-LARGE-LICENSE", "fallback_price": 3995.00},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-vmx-xlarge-license", "model": "VMX-XLARGE-LICENSE", "fallback_price": 8995.00}
        ]
        
        # Other common license types
        other_licenses = [
            #{"url": "https://www.rhinonetworks.com/product/license/meraki-enterprise-license", "model": "LIC-ENT"},
            #{"url": "https://www.rhinonetworks.com/product/license/meraki-advanced-license", "model": "LIC-ADV"},
            {"url": "https://www.rhinonetworks.com/product/license/meraki-systems-manager-enterprise-licenses", "model": "LIC-SME"}
        ]
        
        # Combine all license types
        return insight_licenses + vmx_licenses + other_licenses
    def add_standard_license_prices(self):

        standard_licenses = [
            # Enterprise License options (various durations)
            {"model": "LIC-ENT-1YR", "price": 150.00},
            {"model": "LIC-ENT-3YR", "price": 380.00},
            {"model": "LIC-ENT-5YR", "price": 600.00},
            {"model": "LIC-ENT-7YR", "price": 840.00},
            {"model": "LIC-ENT-10YR", "price": 1200.00},
            
            # Advanced License options
            {"model": "LIC-ADV-1YR", "price": 200.00},
            {"model": "LIC-ADV-3YR", "price": 500.00},
            {"model": "LIC-ADV-5YR", "price": 800.00},
            {"model": "LIC-ADV-7YR", "price": 1120.00},
            {"model": "LIC-ADV-10YR", "price": 1600.00}
        ]
        
        count = 0
        for item in standard_licenses:
            model = item["model"]
            price = item["price"]
            family = "License"
            
            if family in self.prices and model not in self.prices[family]:
                self.prices[family][model] = price
                count += 1
                #self.logger.info(f"{self.GREEN}Added standard license price for {model}: ${price}{self.RESET}")
        
        return count
    def add_vmx_license_prices(self):
        """
        Add prices for VMX licenses if they couldn't be scraped from the website.
        These products might not be listed on the website anymore or might have different URLs.
        """
        vmx_prices = [
            # Primary license model names with -LICENSE
            {"model": "VMX-SMALL-LICENSE", "price": 395.00},
            {"model": "VMX-MEDIUM-LICENSE", "price": 995.00},
            {"model": "VMX-LARGE-LICENSE", "price": 3995.00},
            {"model": "VMX-XLARGE-LICENSE", "price": 8995.00},
            
            # Also add aliases without -LICENSE for compatibility
            {"model": "VMX-SMALL", "price": 395.00},
            {"model": "VMX-MEDIUM", "price": 995.00},
            {"model": "VMX-LARGE", "price": 3995.00},
            {"model": "VMX-XLARGE", "price": 8995.00}
        ]
        
        count = 0
        for item in vmx_prices:
            model = item["model"]
            price = item["price"]
            family = "License"  # VMX licenses go in the License family
            
            if family in self.prices and model not in self.prices[family]:
                self.prices[family][model] = price
                count += 1
                #self.logger.info(f"{self.GREEN}Added VMX license price for {model}: ${price}{self.RESET}")
        
        return count


    def scrape_license_products(self):

        #self.logger.info(f"{self.BLUE}Scraping license product prices...{self.RESET}")
        
        # Get license product URLs
        license_products = self.get_license_product_urls()
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
        }
        
        session = requests.Session()
        total_products = 0
        
        # Process each license product
        for product in license_products:
            try:
                model = product["model"]
                url = product["url"]
                
                #self.logger.info(f"{self.BLUE}Fetching license: {model} from {url}{self.RESET}")
                
                # Add delay to avoid being blocked
                time.sleep(random.uniform(1, 2))
                
                response = session.get(url, headers=headers)
                
                # If not a successful response, try once more
                if response.status_code != 200:
                    self.logger.warning(f"{self.YELLOW}Retry for {model}, status: {response.status_code}{self.RESET}")
                    time.sleep(2)  # Longer delay for retry
                    response = session.get(url, headers=headers)
                
                if response.status_code != 200:
                    self.logger.warning(f"{self.YELLOW}Failed to fetch {model}, status: {response.status_code}{self.RESET}")
                    
                    # If this is a VMX license with a fallback price, use it
                    if "fallback_price" in product:
                        self.prices["License"][model] = product["fallback_price"]
                        total_products += 1
                        self.logger.info(f"{self.GREEN}Using fallback price for {model}: ${product['fallback_price']}{self.RESET}")
                    
                    continue
                    
                # Save first product for debugging
                #if product == license_products[0]:
                    #with open(f'rhino_license_{model}.html', 'w', encoding='utf-8') as f:
                        #f.write(response.text)
                        #self.logger.info(f"{self.GREEN}Saved license page to rhino_license_{model}.html{self.RESET}")
                
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # License pages may have different price elements
                price_selectors = [
                    '.product__price',
                    '.price', 
                    '.product-price',
                    'span.amount',
                    '.woocommerce-Price-amount',
                    '[itemprop="price"]'
                ]
                
                price_found = False
                for selector in price_selectors:
                    price_elems = soup.select(selector)
                    for price_elem in price_elems:
                        price_text = price_elem.text.strip()
                        price_match = re.search(r'\$([0-9,.]+)', price_text)
                        if price_match:
                            price = float(price_match.group(1).replace(',', ''))
                            
                            # Skip extremely high prices (likely not the base price)
                            if price > 30000:
                                continue
                                
                            # Add to our prices dictionary
                            self.prices["License"][model] = price
                            total_products += 1
                            price_found = True
                            #self.logger.info(f"{self.GREEN}Added license price for {model}: ${price}{self.RESET}")
                            break
                    
                    if price_found:
                        break
                
                # If we couldn't find a price but there's a fallback, use it
                if not price_found and "fallback_price" in product:
                    self.prices["License"][model] = product["fallback_price"]
                    total_products += 1
                    self.logger.info(f"{self.GREEN}Using fallback price for {model}: ${product['fallback_price']}{self.RESET}")
                
            except Exception as e:
                self.logger.error(f"{self.RED}Error fetching {product['model']}: {e}{self.RESET}")
            
            # Be kind to the server
            time.sleep(random.uniform(1, 2))
        
        #self.logger.info(f"{self.GREEN}Completed license product scraping. Found {total_products} license products.{self.RESET}")
        
        # Save prices to cache if we found any
        if total_products > 0:
            self.save_prices_to_cache()
            return True
        return False
    def add_exact_prices_from_html(self):

        known_prices = [
            # From the HTML - General Purpose section
            {"model": "MR36", "price": 382.71},  # <strike>$797.31</strike> <strong>$382.71</strong>
            {"model": "MR44", "price": 609.94},  # <strike>$1,297.75</strike> <strong>$609.94</strong>
            
            # From the HTML - High Density section
            {"model": "MR46", "price": 943.36},  # <strike>$1,886.71</strike> <strong>$943.36</strong>
            {"model": "MR46E", "price": 1245.23},  # <strike>$1,886.71</strike> <strong>$1,245.23</strong>
            {"model": "MR56", "price": 1557.51},  # <strike>$2,359.86</strike> <strong>$1,557.51</strong>
            {"model": "MR57", "price": 1853.93},  # <strike>$2,808.98</strike> <strong>$1,853.93</strong>
            
            # From the HTML - Outdoor section
            {"model": "MR76", "price": 853.57},  # <strike>$2,081.88</strike> <strong>$853.57</strong>
            {"model": "MR78", "price": 783.38},  # <strike>$1,506.50</strike> <strong>$783.38</strong>
            {"model": "MR86", "price": 1912.73},  # <strike>$2,898.07</strike> <strong>$1,912.73</strong>
            
            # Other models from your log
            {"model": "MR30H", "price": 895.00},  # Not in HTML, using estimated price
            {"model": "MR33", "price": 795.00},   # Not in HTML, using estimated price
            {"model": "MR55", "price": 1695.00},  # Not in HTML, using estimated price
            
            # Cisco Catalyst Wireless models from the HTML
            {"model": "CW9162I-MR", "price": 588.08},
            {"model": "CW9163E-MR", "price": 993.29},
            {"model": "CW9164I-MR", "price": 1102.99},
            {"model": "CW9166D1-MR", "price": 1432.58},
            {"model": "CW9166I-MR", "price": 1178.52},
            {"model": "CW9172I", "price": 858.98},
            {"model": "CW9176D1", "price": 1320.22},
            {"model": "CW9176I", "price": 1178.52},
            {"model": "CW9178I", "price": 1414.69}
        ]
        
        # Add these prices to our database
        count = 0
        for item in known_prices:
            model = item["model"]
            price = item["price"]
            family = self.extract_family_from_model(model)
            
            if family and family in self.prices:
                self.prices[family][model] = price
                count += 1
                #self.logger.info(f"{self.GREEN}Added exact price for {model}: ${price}{self.RESET}")
        
        return count
    def extract_prices_from_wireless_page(self):
        """
        Directly extract prices from the wireless access points page.
        This is a targeted approach to handle the specific page structure.
        """
        #self.logger.info(f"{self.BLUE}Directly extracting prices from wireless access points page...{self.RESET}")
        
        # URL for the wireless access points page
        url = "https://www.rhinonetworks.com/shop/category/wireless-access-points"
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
        }
        
        try:
            # Create a new session
            session = requests.Session()
            
            # Make the request
            response = session.get(url, headers=headers)
            if response.status_code != 200:
                self.logger.warning(f"{self.YELLOW}Failed to fetch wireless page, status: {response.status_code}{self.RESET}")
                return 0
                
            # Save the page for debugging
            #with open('wireless_page.html', 'w', encoding='utf-8') as f:
                #f.write(response.text)
                #self.logger.info(f"{self.GREEN}Saved wireless page HTML for debugging{self.RESET}")
            
            # Parse the HTML
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Find all product cards
            product_cards = soup.select('article.product__card')
            #self.logger.info(f"{self.GREEN}Found {len(product_cards)} product cards on wireless page{self.RESET}")
            
            # Track how many prices we extract
            extracted_count = 0
            
            # Process each product card
            for card in product_cards:
                try:
                    # Get the title element to extract the model
                    title_elem = card.select_one('.product__card__title h3')
                    if not title_elem:
                        continue
                        
                    # Get the model name
                    title_text = title_elem.text.strip()
                    
                    # Extract model patterns for various device types
                    model_patterns = [
                        r'(MR\d+[A-Za-z]?)',  # MR models
                        r'(CW\d+[A-Za-z\d-]+(?:-MR)?)',  # CW models with optional -MR suffix
                    ]
                    
                    model = None
                    for pattern in model_patterns:
                        model_match = re.search(pattern, title_text)
                        if model_match:
                            model = model_match.group(1).upper()
                            # Handle Catalyst models specifically
                            if model.startswith('CW') and not model.endswith('-MR') and '-MR' in title_text.upper():
                                model = model + '-MR'
                            break
                    
                    if not model:
                        continue
                    
                    # Get the price element
                    price_elem = card.select_one('.product__card__price')
                    if not price_elem:
                        continue
                    
                    # Extract the discounted price (in the <strong> tag)
                    strong_price = price_elem.select_one('strong')
                    if not strong_price:
                        continue
                    
                    # Extract the price value
                    price_match = re.search(r'\$([0-9,.]+)', strong_price.text)
                    if not price_match:
                        continue
                        
                    price = float(price_match.group(1).replace(',', ''))
                    
                    # For CW models, store them in the MR family for compatibility
                    family = self.extract_family_from_model(model)
                    if model.startswith('CW'):
                        family = 'MR'  # Map CW models to MR family
                    
                    # Add to our prices dictionary
                    if family in self.prices:
                        self.prices[family][model] = price
                        extracted_count += 1
                        #self.logger.info(f"{self.GREEN}Extracted price for {model}: ${price}{self.RESET}")
                
                except Exception as e:
                    self.logger.error(f"{self.RED}Error extracting product from card: {e}{self.RESET}")
                    continue
            
            # Return the number of prices we extracted
            return extracted_count
            
        except Exception as e:
            self.logger.error(f"{self.RED}Error fetching wireless page: {e}{self.RESET}")
            return 0
    def discover_product_urls(self):
        """
        Discover product URLs directly from the Rhino Networks website.
        This method scans category pages to find actual product links rather than constructing them.
        """
        #self.logger.info(f"{self.BLUE}Discovering actual product URLs from category pages...{self.RESET}")
        
        # URLs of category pages to scan
        category_urls = [
            "https://www.rhinonetworks.com/shop/category/wireless-access-points",
            "https://www.rhinonetworks.com/shop/category/switching",
            "https://www.rhinonetworks.com/shop/category/security-appliances",
            "https://www.rhinonetworks.com/shop/category/security-cameras"
        ]
        
        discovered_urls = {}
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
        }
        
        session = requests.Session()
        
        for category_url in category_urls:
            #self.logger.info(f"{self.BLUE}Scanning category: {category_url}{self.RESET}")
            
            try:
                # Fetch the category page
                response = session.get(category_url, headers=headers)
                if response.status_code != 200:
                    self.logger.warning(f"{self.YELLOW}Failed to fetch {category_url}, status: {response.status_code}{self.RESET}")
                    continue
                    
                # Parse the page
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Find all product cards with links
                product_links = soup.select('article.product__card')
                
                for card in product_links:
                    # Find the parent link
                    parent_link = card.find_parent('a')
                    if not parent_link or 'href' not in parent_link.attrs:
                        continue
                    
                    # Get the URL
                    href = parent_link['href']
                    full_url = urljoin(self.base_url, href)
                    
                    # Get the model name from the title
                    title_elem = card.select_one('.product__card__title h3')
                    if not title_elem:
                        continue
                    
                    title_text = title_elem.text.strip()
                    
                    # Extract model patterns
                    model_patterns = [
                        r'(MR\d+[A-Za-z]?)',  # MR models
                        r'(MS\d+(?:-\d+)?[A-Za-z]*)',  # MS models 
                        r'(MX\d+[A-Za-z]?)',  # MX models
                        r'(MV\d+[A-Za-z]?)',  # MV models
                        r'(MG\d+[A-Za-z]?)',  # MG models
                        r'(MT\d+[A-Za-z]?)',  # MT models
                        r'(Z\d+[A-Za-z]?)',   # Z models
                        r'(CW\d+[A-Za-z\d-]+)',  # CW models
                        r'(C\d+[A-Za-z\d-]+)'  # C models
                    ]
                    
                    # Try each pattern to extract the model
                    for pattern in model_patterns:
                        model_match = re.search(pattern, title_text)
                        if model_match:
                            model = model_match.group(1).upper()
                            
                            # Handle special cases like CW models with -MR and C models with -M
                            if model.startswith('CW') and '-MR' in title_text.upper():
                                model = model.rstrip('-MR')  # Remove -MR if it's part of the model
                            
                            if model.startswith('C9') and '-M' in title_text.upper():
                                if not model.endswith('-M'):
                                    model = model + '-M'
                            
                            # Store the discovered URL
                            discovered_urls[model] = full_url
                            #self.logger.info(f"{self.GREEN}Discovered product URL: {model} -> {full_url}{self.RESET}")
                            break
                
                time.sleep(random.uniform(1, 2))
                
            except Exception as e:
                self.logger.error(f"{self.RED}Error scanning category page {category_url}: {e}{self.RESET}")
        
        #self.logger.info(f"{self.GREEN}Discovered {len(discovered_urls)} actual product URLs{self.RESET}")
        return discovered_urls

    def _correct_model_url(self, model, url):
        """
        Apply corrections to URLs for specific model types
        """
        # For Catalyst switches (C9xxx), ensure they have -M suffix
        if re.match(r'^C9\d+', model) and '-M' not in model and '-m' not in model:
            # If URL already has -m, use it
            if '-m' in url.lower():
                return url
            # Otherwise add -m to the URL
            return url.replace('/meraki-', '/meraki-') + '-m'
        
        # For Catalyst APs (CWxxxx), use the cisco- prefix and ensure -mr suffix
        if re.match(r'^CW\d+', model):
            # If URL uses meraki- prefix, change to cisco-
            if '/meraki-cw' in url.lower():
                url = url.replace('/meraki-cw', '/cisco-cw')
            
            # If URL doesn't end with -mr, add it
            if not url.lower().endswith('-mr'):
                # But only if the model doesn't already have it
                if not model.lower().endswith('-mr'):
                    return url + '-mr'
            
            return url
        
        return url
    def update_product_urls(self):
        """
        Update the product URL list with newly discovered models.
        Returns the updated list of products to scrape.
        """
        # Start with our known products - FIXED URLS for problematic models
        base_products = [
            # MX Series
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx67", "model": "MX67"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx67c", "model": "MX67C"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx67w", "model": "MX67W"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx68", "model": "MX68"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx68w", "model": "MX68W"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx68cw", "model": "MX68CW"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx75", "model": "MX75"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx84", "model": "MX84"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx85", "model": "MX85"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx95", "model": "MX95"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx100", "model": "MX100"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx105", "model": "MX105"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx250", "model": "MX250"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mx450", "model": "MX450"},
            
            # MS Series (Switches) - FIXED: Correct model names/URLs
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms120-8", "model": "MS120-8"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms120-8lp", "model": "MS120-8LP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms120-8fp", "model": "MS120-8FP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms120-24", "model": "MS120-24"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms120-24p", "model": "MS120-24P"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms125-24", "model": "MS125-24"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms125-24p", "model": "MS125-24P"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms130-24x", "model": "MS130-24X"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms210-24", "model": "MS210-24"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms225-48lp", "model": "MS225-48LP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms225-48fp", "model": "MS225-48FP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms250-24p", "model": "MS250-24P"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms250-48fp", "model": "MS250-48FP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms350-24p", "model": "MS350-24P"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms350-48lp", "model": "MS350-48LP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms350-48fp", "model": "MS350-48FP"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms355-24x", "model": "MS355-24X"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms390-48ux", "model": "MS390-48UX"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms410-32", "model": "MS410-32"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms425-16", "model": "MS425-16"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-ms425-32", "model": "MS425-32"},
            
            # MR Series (Wireless)
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr30h", "model": "MR30H"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr33", "model": "MR33"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr36", "model": "MR36"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr42", "model": "MR42"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr44", "model": "MR44"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr46", "model": "MR46"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr55", "model": "MR55"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr56", "model": "MR56"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mr57", "model": "MR57"},
            
            # MV Series (Cameras)
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mv13", "model": "MV13"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mv72", "model": "MV72"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mv93", "model": "MV93"},
            
            # Catalyst Series - Added -M suffix for Meraki-managed Catalyst switches
            {"url": "https://www.rhinonetworks.com/product/device/meraki-c9300x-24y-m", "model": "C9300X-24Y"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-c9300-48un-m", "model": "C9300-48UN-M"},
            
            # Catalyst AP Series - Use cisco-cwXXXX-mr pattern for Catalyst APs
            {"url": "https://www.rhinonetworks.com/product/device/cisco-cw9163e-mr", "model": "CW9163E"},
            {"url": "https://www.rhinonetworks.com/product/device/cisco-cw9176i", "model": "CW9176I"},
            
            # Z Series
            {"url": "https://www.rhinonetworks.com/product/device/meraki-z4", "model": "Z4"},
            {"url": "https://www.rhinonetworks.com/product/device/meraki-z4c", "model": "Z4C"},
            
            # MG Series
            {"url": "https://www.rhinonetworks.com/product/device/meraki-mg21e", "model": "MG21E"},
        ]
        
        # Get existing model URLs
        existing_models = {product["model"]: product["url"] for product in base_products}
        
        # Discover new models
        discovered_models = self.discover_new_models()
        
        # Add any new models we've discovered
        new_models_added = 0
        for model, url in discovered_models.items():
            if model not in existing_models:
                # Apply URL corrections for specific model types before adding
                corrected_url = self._correct_model_url(model, url)
                base_products.append({"url": corrected_url, "model": model})
                new_models_added += 1
                #self.logger.info(f"{self.GREEN}Added new model to scraping list: {model} - {corrected_url}{self.RESET}")
        
        #self.logger.info(f"{self.GREEN}Updated product URL list with {new_models_added} new models. Total models: {len(base_products)}{self.RESET}")
        return base_products

    def load_cached_prices(self):
        """Load prices from cache if available and not expired."""
        try:
            if not os.path.exists(self.cache_file):
                return {}
                
            with open(self.cache_file, 'r') as f:
                cache_data = json.load(f)
                
            # Check if cache is expired
            cache_date = datetime.datetime.fromisoformat(cache_data['last_updated'])
            days_old = (datetime.datetime.now() - cache_date).days
            
            # Check if cache has data
            total_prices = 0
            for family_prices in cache_data['prices'].values():
                total_prices += len(family_prices)
            
            if days_old <= self.cache_expiry_days and total_prices > 0:
                self.logger.info(f"{self.GREEN}Using cached Rhino Networks prices ({days_old} days old) - {total_prices} models{self.RESET}")
                return cache_data['prices']
            else:
                if total_prices == 0:
                    self.logger.info(f"{self.YELLOW}Price cache is empty, will fetch fresh data{self.RESET}")
                else:
                    self.logger.info(f"{self.YELLOW}Price cache expired ({days_old} days old), will fetch fresh data{self.RESET}")
                return {}
                    
        except (FileNotFoundError, json.JSONDecodeError, KeyError):
            self.logger.info(f"{self.YELLOW}No valid price cache found, will fetch data from Rhino Networks{self.RESET}")
            return {}
            
    def save_prices_to_cache(self):
        """Save prices to cache file."""
        cache_data = {
            'last_updated': datetime.datetime.now().isoformat(),
            'prices': self.prices
        }
        
        try:
            with open(self.cache_file, 'w') as f:
                json.dump(cache_data, f, indent=2)
            self.logger.info(f"{self.GREEN}Saved Rhino Networks prices to cache{self.RESET}")
            return True
        except Exception as e:
            self.logger.warning(f"{self.YELLOW}Warning: Could not save price cache: {e}{self.RESET}")
            return False
    
    def fetch_all_prices(self):
        """Fetch all Meraki product prices from Rhino Networks website."""
        #self.logger.info(f"{self.BLUE}Starting comprehensive Rhino Networks price update...{self.RESET}")
        
        # Initialize prices structure
        self.prices = {
            'MX': {},
            'MS': {},
            'MR': {},
            'MV': {},
            'MG': {},
            'MT': {},
            'Z': {},
            'License': {}
        }
        
        total_products = 0
        
        # Process each product category
        for family, category_url in self.category_urls.items():
            #self..info(f"{self.BLUE}Fetching prices for {family} from {category_url}{self.RESET}")
            
            # Start with the first page
            url = category_url
            page_num = 1
            
            while url:
                #self.logger.info(f"{self.BLUE}Fetching page {page_num} for {family}...{self.RESET}")
                
                # Fetch product cards from the page
                product_data = self.fetch_page(url, family)
                
                # Add products to our price dictionary
                for product in product_data['products']:
                    extracted_family = product['family']
                    model = product['model']
                    price = product['price']
                    
                    if extracted_family in self.prices:
                        self.prices[extracted_family][model] = price
                        total_products += 1
                        #self.logger.info(f"{self.GREEN}Added price for {model}: ${price}{self.RESET}")
                
                # Move to next page if available
                url = product_data['next_page']
                if url:
                    page_num += 1
                    # Be kind to the server
                    time.sleep(random.uniform(1, 3))
                else:
                    break
        
        #self.logger.info(f"{self.GREEN}Completed price update. Found {total_products} products.{self.RESET}")
        
        # Save prices to cache if we found any
        if total_products > 0:
            self.save_prices_to_cache()
        else:
            self.logger.warning(f"{self.YELLOW}No products found during scraping!{self.RESET}")
    
    def fetch_page(self, url, target_family=None):
        """Fetch a single page and extract product information."""
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'Connection': 'keep-alive',
            'Cache-Control': 'max-age=0',
            'Referer': 'https://www.rhinonetworks.com/',
        }
            
        try:
            #self.logger.info(f"{self.BLUE}Fetching URL: {url}{self.RESET}")
            
            # Use a session to maintain cookies and handle redirects properly
            session = requests.Session()
            
            # Make the actual request
            response = session.get(url, headers=headers, allow_redirects=True)
            response.raise_for_status()
            
            #self.logger.info(f"{self.GREEN}Response status: {response.status_code}, length: {len(response.text)} chars{self.RESET}")
            
            # For debugging - save the HTML to a file to inspect
            #with open(f'rhino_page_{target_family}.html', 'w', encoding='utf-8') as f:
                #f.write(response.text)
                #self.logger.info(f"{self.GREEN}Saved HTML to rhino_page_{target_family}.html for debugging{self.RESET}")
            
            # Create the soup object from the response
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Find all product cards on the page
            # Try different selectors based on the HTML structure
            product_cards = soup.select('.card.product__card') or soup.select('.card.grid-y.product__card') or soup.select('article.product__card')
            
            # If still no cards found, try more generic selectors
            if not product_cards:
                self.logger.info(f"{self.YELLOW}No cards found with primary selectors, trying alternative selectors{self.RESET}")
                product_cards = soup.select('.cell a') or soup.select('a[href^="/product/"]')
            
            #self.logger.info(f"{self.GREEN}Found {len(product_cards)} product cards/links on page{self.RESET}")
            
            # Process each product card
            results = []
            for card in product_cards:
                try:
                    # Extract product information
                    title_element = None
                    price_element = None
                    product_name = ""
                    
                    if card.name == 'a' and 'href' in card.attrs:
                        # This is a link to a product
                        href = card.get('href', '')
                        if '/product/device/' in href or '/product/license/' in href:
                            # Extract model from URL
                            model_match = re.search(r'/(MR|MS|MX|MV|MG|MT|Z|LIC)[A-Za-z0-9-]+', href)
                            if model_match:
                                model = model_match.group(0).strip('/')
                                
                                # If we have child elements that might contain price
                                price_text = None
                                article = card.find('article')
                                if article:
                                    price_div = article.select_one('.product__card__price')
                                    if price_div:
                                        price_text = price_div.text.strip()
                                
                                # If we found a price, process it
                                if price_text:
                                    strong_price = None
                                    if price_div:
                                        strong_price = price_div.select_one('strong')
                                    
                                    price = None
                                    if strong_price:
                                        # Get the actual price (not the strike-through one)
                                        price_match = re.search(r'\$([0-9,.]+)', strong_price.text)
                                        if price_match:
                                            price = float(price_match.group(1).replace(',', ''))
                                    else:
                                        # If no strong element, try to extract any price
                                        price_match = re.search(r'\$([0-9,.]+)', price_text)
                                        if price_match:
                                            price = float(price_match.group(1).replace(',', ''))
                                    
                                    if price:
                                        # Determine family
                                        family = target_family
                                        if not family:
                                            family = self.extract_family_from_model(model)
                                        
                                        # Create result
                                        result = {
                                            'name': model,
                                            'model': model,
                                            'family': family,
                                            'price': price
                                        }
                                        results.append(result)
                                        #self.logger.info(f"{self.GREEN}Extracted from link: {model} at ${price}{self.RESET}")
                    else:
                        # Standard card structure
                        title_element = card.select_one('.product__card__title h3')
                        price_element = card.select_one('.product__card__price')
                        
                        if title_element and price_element:
                            product_name = title_element.text.strip()
                            
                            # Extract model
                            model = self.extract_model_from_name(product_name)
                            if model:
                                # Determine family
                                family = target_family
                                if not family:
                                    family = self.extract_family_from_model(model)
                                
                                if family:
                                    # Extract price information
                                    price_text = price_element.text.strip()
                                    
                                    # Handle strike-through pricing
                                    strong_price = price_element.select_one('strong')
                                    if strong_price:
                                        # Get the actual price (not the strike-through one)
                                        price_match = re.search(r'\$([0-9,.]+)', strong_price.text)
                                        if price_match:
                                            price = float(price_match.group(1).replace(',', ''))
                                        else:
                                            continue
                                    else:
                                        # If no strong element, try to extract any price
                                        price_match = re.search(r'\$([0-9,.]+)', price_text)
                                        if price_match:
                                            price = float(price_match.group(1).replace(',', ''))
                                        else:
                                            continue
                                    
                                    # Create product result
                                    result = {
                                        'name': product_name,
                                        'model': model,
                                        'family': family,
                                        'price': price
                                    }
                                    results.append(result)
                                    #self.logger.info(f"{self.GREEN}Successfully extracted: {model} at ${price}{self.RESET}")
                    
                except Exception as e:
                    self.logger.error(f"{self.RED}Error extracting product card: {e}{self.RESET}")
                    continue
            
            # Look for pagination
            next_page = None
            pagination = soup.select_one('.pagination')
            if pagination:
                # Find the "next" link
                next_links = pagination.select('a')
                for link in next_links:
                    if ('next' in link.get('class', []) or 
                        '›' in link.text or 
                        'Next' in link.text or 
                        'next' in link.text.lower()):
                        href = link.get('href')
                        if href:
                            next_page = urljoin(self.base_url, href)
                            #self.logger.info(f"{self.BLUE}Found next page link: {next_page}{self.RESET}")
                            break
            
            return {'products': results, 'next_page': next_page}
            
        except requests.exceptions.HTTPError as e:
            self.logger.error(f"{self.RED}HTTP error fetching page {url}: {e}{self.RESET}")
            return {'products': [], 'next_page': None}
        except requests.exceptions.ConnectionError as e:
            self.logger.error(f"{self.RED}Connection error fetching page {url}: {e}{self.RESET}")
            return {'products': [], 'next_page': None}
        except requests.exceptions.Timeout as e:
            self.logger.error(f"{self.RED}Timeout fetching page {url}: {e}{self.RESET}")
            return {'products': [], 'next_page': None}
        except requests.exceptions.RequestException as e:
            self.logger.error(f"{self.RED}Error fetching page {url}: {e}{self.RESET}")
            return {'products': [], 'next_page': None}
        except Exception as e:
            self.logger.error(f"{self.RED}Unexpected error fetching page {url}: {e}{self.RESET}")
            return {'products': [], 'next_page': None}
    
    def extract_model_from_name(self, product_name):
        """Extract Meraki model number from product name."""
        # Common Meraki model patterns
        patterns = [
            r'Meraki (MS[0-9]+(?:-[0-9]+[A-Z]*)?)',  # MS switches
            r'Meraki (MX[0-9]+(?:[A-Z])?)',          # MX security appliances
            r'Meraki (MR[0-9]+(?:[A-Z])?)',          # MR wireless
            r'Meraki (MV[0-9]+(?:[A-Z])?)',          # MV cameras
            r'Meraki (MG[0-9]+(?:[A-Z])?)',          # MG cellular gateways
            r'Meraki (MT[0-9]+(?:[A-Z])?)',          # MT sensors
            r'Meraki (Z[0-9]+(?:[A-Z])?)',           # Z teleworker gateways
            r'Meraki (LIC-[A-Z0-9-]+)'               # License products
        ]
        
        for pattern in patterns:
            match = re.search(pattern, product_name, re.IGNORECASE)
            if match:
                return match.group(1)
        
        return None
    def scrape_direct_product_urls(self):
        """Enhanced method to directly access product URLs to get prices."""
        #self.logger.info(f"{self.BLUE}Starting direct URL product scraping...{self.RESET}")
        
        # Initialize prices
        self.prices = {
            'MX': {}, 'MS': {}, 'MR': {}, 'MV': {}, 'MG': {}, 'MT': {}, 'Z': {}, 'License': {}, 'C': {}
        }
        
        # Get updated product URL list with any newly discovered models
        product_urls = self.update_product_urls()
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
            'Accept-Language': 'en-US,en;q=0.9',
            'Referer': 'https://www.rhinonetworks.com/',
            'Cache-Control': 'no-cache',
        }
        
        session = requests.Session()
        total_products = 0
        
        # Group products by family to better track progress
        products_by_family = {}
        for product in product_urls:
            family = self.extract_family_from_model(product["model"])
            if family not in products_by_family:
                products_by_family[family] = []
            products_by_family[family].append(product)
        
        # Process each family separately to avoid overwhelming the server
        for family, family_products in products_by_family.items():
            #self.logger.info(f"{self.BLUE}Processing {len(family_products)} products in {family} family{self.RESET}")
            
            for product in family_products:
                try:
                    model = product["model"]
                    url = product["url"]
                    
                    #self.logger.info(f"{self.BLUE}Fetching product: {model} from {url}{self.RESET}")
                    
                    # Add delay to avoid being blocked
                    time.sleep(random.uniform(2, 3))
                    
                    response = session.get(url, headers=headers, allow_redirects=True)
                    
                    # If not a successful response, try a few more times
                    retries = 0
                    while response.status_code != 200 and retries < 3:
                        self.logger.warning(f"{self.YELLOW}Retry {retries+1} for {model}, status: {response.status_code}{self.RESET}")
                        time.sleep(random.uniform(3, 5))  # Longer delay for retries
                        response = session.get(url, headers=headers, allow_redirects=True)
                        retries += 1
                    
                    if response.status_code != 200:
                        self.logger.warning(f"{self.YELLOW}Failed to fetch {model} after retries, status: {response.status_code}{self.RESET}")
                        continue
                        
                    # Save first response for each family for debugging
                    #if product == family_products[0]:
                        #with open(f'rhino_product_{family}_{model}.html', 'w', encoding='utf-8') as f:
                            #f.write(response.text)
                            #self.logger.info(f"{self.GREEN}Saved product page to rhino_product_{family}_{model}.html{self.RESET}")
                    
                    soup = BeautifulSoup(response.text, 'html.parser')
                    
                    # Look for price with expanded selectors
                    price_selectors = [
                        '.product__price', 
                        '.product-price',
                        '.price',
                        'span.amount',
                        '.current-price',
                        '.our-price',
                        '.woocommerce-Price-amount',
                        '[data-price]'
                    ]
                    
                    price_found = False
                    for selector in price_selectors:
                        price_elem = soup.select_one(selector)
                        if price_elem:
                            price_text = price_elem.text.strip()
                            price_match = re.search(r'\$([0-9,.]+)', price_text)
                            if price_match:
                                price = float(price_match.group(1).replace(',', ''))
                                
                                # Add to our prices dictionary
                                family = self.extract_family_from_model(model)
                                if family in self.prices:
                                    self.prices[family][model] = price
                                    total_products += 1
                                    price_found = True
                                    #self.logger.info(f"{self.GREEN}Added price for {model}: ${price}{self.RESET}")
                                    break
                    
                    # If still no price found, try looking at any element with a dollar sign
                    if not price_found:
                        dollar_elems = soup.select('*:contains("$")')
                        for elem in dollar_elems:
                            price_text = elem.text.strip()
                            price_match = re.search(r'\$([0-9,.]+)', price_text)
                            if price_match:
                                price = float(price_match.group(1).replace(',', ''))
                                
                                # Skip extremely high or low prices (likely not product prices)
                                if price < 10 or price > 100000:
                                    continue
                                    
                                # Add to our prices dictionary
                                family = self.extract_family_from_model(model)
                                if family in self.prices:
                                    self.prices[family][model] = price
                                    total_products += 1
                                    price_found = True
                                    #self.logger.info(f"{self.GREEN}Found price from general search for {model}: ${price}{self.RESET}")
                                    break
                    
                    if not price_found:
                        self.logger.warning(f"{self.YELLOW}Could not find price for {model}{self.RESET}")
                    
                except Exception as e:
                    self.logger.error(f"{self.RED}Error fetching {product['model']}: {e}{self.RESET}")
                
                # Be kind to the server
                time.sleep(random.uniform(1, 2))
            
            # Longer delay between families
            time.sleep(random.uniform(3, 5))
        
        #self.logger.info(f"{self.GREEN}Completed direct URL scraping. Found {total_products} products.{self.RESET}")
        
        # Save prices to cache if we found any
        if total_products > 0:
            self.save_prices_to_cache()
            return True
        return False
    def extract_family_from_model(self, model):
        """Enhanced method to extract device family from model number."""
        if not model:
            return None
            
        # Handle Catalyst CW-series APs separately
        if model.startswith('CW'):
            return 'MR'  # Treat Catalyst APs as MR equivalents
            
        # Handle Catalyst C-series switches separately
        if model.startswith('C9'):
            return 'MS'  # Treat Catalyst switches as MS equivalents
            
        family_match = re.match(r'(MR|MS|MX|MV|MG|MT|Z|LIC)', model)
        if family_match:
            family = family_match.group(1)
            if family == 'LIC':
                return 'License'
            return family
            
        return None
    def manual_price_update(self, devices_to_update):
        """
        Manually update prices for specific models that we know are problematic.
        
        Args:
            devices_to_update: List of dictionaries with model and price keys
        """
        self.logger.info(f"{self.BLUE}Adding manual price updates for problematic models...{self.RESET}")
        
        for device in devices_to_update:
            model = device["model"]
            price = device["price"]
            family = self.extract_family_from_model(model)
            
            if family and family in self.prices:
                self.prices[family][model] = price
                #self.logger.info(f"{self.GREEN}Manually added price for {model}: ${price}{self.RESET}")
            else:
                self.logger.warning(f"{self.YELLOW}Could not add manual price for {model}, unknown family{self.RESET}")
        
        # Save updated prices
        self.save_prices_to_cache()
        
    # Add this method at the end of the RhinoPriceCatalog class
    def add_missing_models(self):
        """
        Add missing models with approximate prices based on logs analysis
        """
        # List of missing models with estimated prices
        missing_models = [
            {"model": "MR44", "price": 1095.00},
            {"model": "MX84", "price": 3598.70},
            {"model": "MX100", "price": 9395.00},
            {"model": "MS425-16", "price": 13995.00},
            {"model": "MS425-32", "price": 21995.00},
            {"model": "MS225-48LP", "price": 5795.00},
            {"model": "MS225-48FP", "price": 6795.00},
            {"model": "MS350-24P", "price": 6295.00},
            {"model": "MS350-48LP", "price": 9495.00},
            {"model": "MS350-48FP", "price": 10495.00},
            {"model": "MS355-24X", "price": 7495.00},
            {"model": "MS410-32", "price": 14995.00}, 
            {"model": "MS250-24P", "price": 4795.00},
            {"model": "MS250-48FP", "price": 7795.00},
            {"model": "MS130-24X", "price": 3995.00},
            {"model": "MS210-24", "price": 2995.00},
            {"model": "MS390-48UX", "price": 15995.00},
            {"model": "MV72", "price": 1795.00},
            {"model": "MV93", "price": 2095.00},
            {"model": "MV13", "price": 1395.00},
            {"model": "MR33", "price": 795.00},
            {"model": "MR30H", "price": 895.00},
            {"model": "MR55", "price": 1695.00},
            {"model": "MR57", "price": 2095.00},
            {"model": "MR42", "price": 995.00},
            {"model": "C9300X-24Y", "price": 19995.00},
            {"model": "CW9163E", "price": 1295.00},
            {"model": "CW9176I", "price": 1895.00},
            {"model": "Z3C-NA", "price": 695.00},
            {"model": "MG21E-NA", "price": 895.00}
        ]
        
        # Add these models to our prices
        self.manual_price_update(missing_models)
        
        # Let the user know we've added missing models
        model_list = [device["model"] for device in missing_models]
        self.logger.info(f"{self.GREEN}Added {len(missing_models)} missing models to price database: {', '.join(model_list)}{self.RESET}")
        
        return len(missing_models)
        
    def get_price(self, model):
        """Get price for a specific model with smart matching."""
        # Handle special cases
        if not model:
            return None
            
        # Extract family from model
        if model.startswith('LIC-'):
            family = 'License'
        else:
            family_match = re.match(r'(MR|MS|MX|MV|MG|MT|Z)', model)
            if not family_match:
                return None
            family = family_match.group(1)
        
        # Check if we have prices for this family
        if family not in self.prices:
            return None
            
        # 1. Try exact model match
        if model in self.prices[family]:
            return self.prices[family][model]
            
        # 2. Try prefix match (e.g., "MX68" might match "MX68W")
        for m, price in self.prices[family].items():
            if model.startswith(m) or m.startswith(model):
                return price
        
        # 3. Try base model match (extract the core model number)
        base_model_match = re.match(f'({family}[0-9]+)', model)
        if base_model_match:
            base_model = base_model_match.group(1)
            for m, price in self.prices[family].items():
                if m.startswith(base_model):
                    return price
                    
        # 4. Use family average as a fallback
        if self.prices[family]:
            avg_price = sum(self.prices[family].values()) / len(self.prices[family])
            self.logger.warning(f"{self.YELLOW}No exact match for {model}, using family average price of ${avg_price:.2f}{self.RESET}")
            
            # Add to price misses for reporting
            self.price_misses.add(model)
            
            return avg_price
        
        # No price found
        self.logger.warning(f"{self.YELLOW}No price found for {model}{self.RESET}")
        self.price_misses.add(model)
        return None
        
    def get_fallback_prices(self):
        """
        Return fallback hardcoded prices in case web scraping fails.
        """
        prices = {
            'MX': {
                # Small Branch Security Appliances
                'MX64': 995, 'MX64W': 1095, 
                'MX65': 1195, 'MX65W': 1295,
                'MX67': 1095, 'MX67W': 1195, 'MX67C': 1295,
                'MX68': 1295, 'MX68W': 1395, 'MX68CW': 1495,
                
                # Medium Branch Security Appliances
                'MX75': 1995,
                'MX84': 3995, 
                'MX85': 4995,
                'MX95': 5995,
                'MX105': 9995,
                
                # Large Branch/Campus Security Appliances
                'MX250': 14995,
                'MX450': 34995,
                
                # Teleworker Gateways
                'Z3': 595, 'Z3C': 695,
                'Z4': 695, 'Z4C': 795
            },
            'MS': {
                # Access Switches (MS120 Series)
                'MS120-8': 795, 'MS120-8LP': 995, 'MS120-8FP': 1295,
                'MS120-24': 1995, 'MS120-24P': 2495, 'MS120-48': 3495, 'MS120-48LP': 3995, 'MS120-48FP': 4495,
                
                # Access Switches (MS125 Series)
                'MS125-24': 2495, 'MS125-24P': 2995, 'MS125-48': 3995, 'MS125-48LP': 4495, 'MS125-48FP': 5495,
                
                # Access Switches (MS130 Series)
                'MS130-24': 2995, 'MS130-24P': 3495, 'MS130-24X': 3995, 'MS130-48': 4495, 'MS130-48LP': 4995, 'MS130-48FP': 5995,
                
                # Access Switches (MS210 Series)
                'MS210-24': 2995, 'MS210-24P': 3495, 'MS210-48': 4995, 'MS210-48LP': 5495, 'MS210-48FP': 6495,
                
                # Access Switches (MS220 Series - Older)
                'MS220-8': 1295, 'MS220-8P': 1495,
                'MS220-24': 2495, 'MS220-24P': 2995, 'MS220-48': 4495, 'MS220-48LP': 4995, 'MS220-48FP': 5995,
                
                # Stackable Access Switches (MS225 Series)
                'MS225-24': 3295, 'MS225-24P': 3795, 'MS225-48': 5295, 'MS225-48LP': 5795, 'MS225-48FP': 6795,
                
                # Stackable Access Switches (MS250 Series)
                'MS250-24': 4295, 'MS250-24P': 4795, 'MS250-48': 6295, 'MS250-48LP': 6795, 'MS250-48FP': 7795,
                
                # Stackable Access Switches (MS320 Series)
                'MS320-24': 4795, 'MS320-24P': 5295, 'MS320-48': 7995, 'MS320-48LP': 8495, 'MS320-48FP': 9495,
                
                # Stackable Access Switches (MS350 Series)
                'MS350-24': 5795, 'MS350-24P': 6295, 'MS350-24X': 6995, 'MS350-48': 8995, 'MS350-48LP': 9495, 'MS350-48FP': 10495,
                
                # Stackable Access Switches (MS355 Series)
                'MS355-24X': 7495, 'MS355-48X': 11495,
                
                # Stackable Access Switches (MS390 Series)
                'MS390-24': 8495, 'MS390-24P': 8995, 'MS390-24UX': 10995, 
                'MS390-48': 12995, 'MS390-48P': 13995, 'MS390-48UX': 15995,
                
                # Aggregation Switches (MS410 Series)
                'MS410-16': 9995, 'MS410-32': 14995,
                
                # Aggregation Switches (MS425 Series)
                'MS425-16': 13995, 'MS425-32': 21995,
                
                # Aggregation Switches (MS450 Series)
                'MS450-12': 13995, 'MS450-24': 21995,
                'MS450-12M': 15995, 'MS450-24M': 23995,
                
                # Aggregation Switches (MS650 Series - End of Life)
                'MS650-48': 35995,
                
                # Catalyst Switches (Meraki-managed)
                'C9300-24': 9995, 'C9300-48': 15995,
                'C9300-24P': 12995, 'C9300-48P': 19995,
                'C9300-24U': 14995, 'C9300-48U': 21995, 
                'C9300-24UX': 19995, 'C9300-48UX': 29995,
                'C9300L-48': 14995, 'C9300L-48P': 18995, 'C9300L-48T': 13995,
                'C9300X-24': 19995, 'C9300X-48': 29995,
                
                # Common SFP/SFP+ Modules (representative pricing)
                'MA-SFP-1GB-SX': 295, 'MA-SFP-1GB-LX': 495,
                'MA-SFP-10GB-SR': 695, 'MA-SFP-10GB-LR': 1295,
                'MA-QSFP-40G-SR4': 795, 'MA-QSFP-40G-LR4': 3995,
            },
            'MR': {
                # Indoor Access Points (Wi-Fi 5/6)
                'MR20': 695, 'MR30H': 895, 'MR33': 795,
                'MR36': 895, 'MR36H': 995, 
                'MR42': 995, 'MR42E': 1095,
                'MR44': 1095, 'MR46': 1295, 'MR46E': 1395,
                'MR52': 1495, 'MR53': 1495, 'MR53E': 1595,
                'MR55': 1695, 'MR56': 1895, 'MR57': 2095,
                
                # Catalyst Indoor APs (Meraki-managed)
                'CW9162I': 1195, 'CW9163E': 1295, 'CW9164I': 1395, 
                'CW9166I': 1595, 'CW9166D1': 1495, 'CW9176I': 1895, 'CW9178I': 2095,
                
                # Wireless Controllers
                'C9800-L': 2995, 'C9800-L-C-K9': 3495
            },
            'MV': {
                # Indoor Cameras
                'MV2': 995,
                'MV12': 995, 'MV12W': 1095, 'MV12WE': 1195,
                'MV13': 1395,
                'MV21': 1295, 'MV22': 1395,
                'MV32': 1595, 
                'MV52': 1895,
                
                # Outdoor Cameras
                'MV71': 1695, 'MV72': 1795,
                'MV93': 2095,

                # License costs (representative)
                'LIC-MV-1YR': 199, 'LIC-MV-3YR': 499, 'LIC-MV-5YR': 799, 'LIC-MV-10YR': 1499
            },
            'MG': {
                # Cellular Gateways
                'MG21': 795, 'MG21E': 895,
                'MG41': 995, 'MG41E': 1095,
                
                # License costs (representative)
                'LIC-MG-1YR': 149, 'LIC-MG-3YR': 349, 'LIC-MG-5YR': 549, 'LIC-MG-10YR': 999
            },
            'MT': {
                # Environmental Sensors
                'MT10': 295,  # Temperature Sensor
                'MT11': 345,  # Temperature & Humidity Sensor
                'MT12': 345,  # Water Leak Sensor
                'MT14': 395,  # Door Open/Close Sensor
                
                # License costs (representative)
                'LIC-MT-1YR': 99, 'LIC-MT-3YR': 249, 'LIC-MT-5YR': 399, 'LIC-MT-10YR': 699
            },
            'Z': {
                # Teleworker Gateways
                'Z1': 499,
                'Z3': 595, 'Z3C': 695,
                'Z4': 695, 'Z4C': 795
            },
            'License': {
                # Enterprise License costs for security appliances
                'LIC-MX64-ENT-1YR': 200, 'LIC-MX64-ENT-3YR': 500, 'LIC-MX64-ENT-5YR': 800, 'LIC-MX64-ENT-10YR': 1500,
                'LIC-MX65-ENT-1YR': 300, 'LIC-MX65-ENT-3YR': 750, 'LIC-MX65-ENT-5YR': 1200, 'LIC-MX65-ENT-10YR': 2200,
                'LIC-MX67-ENT-1YR': 200, 'LIC-MX67-ENT-3YR': 500, 'LIC-MX67-ENT-5YR': 800, 'LIC-MX67-ENT-10YR': 1500,
                'LIC-MX68-ENT-1YR': 300, 'LIC-MX68-ENT-3YR': 750, 'LIC-MX68-ENT-5YR': 1200, 'LIC-MX68-ENT-10YR': 2200,
                'LIC-MX75-ENT-1YR': 400, 'LIC-MX75-ENT-3YR': 1000, 'LIC-MX75-ENT-5YR': 1600, 'LIC-MX75-ENT-10YR': 2900,
                'LIC-MX84-ENT-1YR': 750, 'LIC-MX84-ENT-3YR': 1875, 'LIC-MX84-ENT-5YR': 3000, 'LIC-MX84-ENT-10YR': 5400,
                'LIC-MX85-ENT-1YR': 900, 'LIC-MX85-ENT-3YR': 2250, 'LIC-MX85-ENT-5YR': 3600, 'LIC-MX85-ENT-10YR': 6480,
                'LIC-MX95-ENT-1YR': 1200, 'LIC-MX95-ENT-3YR': 3000, 'LIC-MX95-ENT-5YR': 4800, 'LIC-MX95-ENT-10YR': 8640,
                'LIC-MX105-ENT-1YR': 2000, 'LIC-MX105-ENT-3YR': 5000, 'LIC-MX105-ENT-5YR': 8000, 'LIC-MX105-ENT-10YR': 14400,
                'LIC-MX250-ENT-1YR': 3000, 'LIC-MX250-ENT-3YR': 7500, 'LIC-MX250-ENT-5YR': 12000, 'LIC-MX250-ENT-10YR': 21600,
                'LIC-MX450-ENT-1YR': 6000, 'LIC-MX450-ENT-3YR': 15000, 'LIC-MX450-ENT-5YR': 24000, 'LIC-MX450-ENT-10YR': 43200,
                
                # Advanced Security License costs
                'LIC-MX64-SEC-1YR': 400, 'LIC-MX64-SEC-3YR': 1000, 'LIC-MX64-SEC-5YR': 1600, 'LIC-MX64-SEC-10YR': 2880,
                'LIC-MX65-SEC-1YR': 600, 'LIC-MX65-SEC-3YR': 1500, 'LIC-MX65-SEC-5YR': 2400, 'LIC-MX65-SEC-10YR': 4320,
                
                # Enterprise License costs for wireless
                'LIC-ENT-1YR': 150, 'LIC-ENT-3YR': 375, 'LIC-ENT-5YR': 600, 'LIC-ENT-10YR': 1080,
                
                # Enterprise License costs for switches
                'LIC-MS120-8-1YR': 30, 'LIC-MS120-8-3YR': 75, 'LIC-MS120-8-5YR': 120, 'LIC-MS120-8-10YR': 216,
                'LIC-MS120-24-1YR': 60, 'LIC-MS120-24-3YR': 150, 'LIC-MS120-24-5YR': 240, 'LIC-MS120-24-10YR': 432,
                'LIC-MS120-48-1YR': 120, 'LIC-MS120-48-3YR': 300, 'LIC-MS120-48-5YR': 480, 'LIC-MS120-48-10YR': 864
            }
        }
        return prices
    
    def get_price_miss_report(self, clear=True):
        """
        Generate a report of models that had no exact price match.
        
        Args:
            clear: If True, clear the price miss list after generating the report
            
        Returns:
            A list of models with no exact price match
        """
        misses = list(self.price_misses)
        
        if clear:
            self.price_misses.clear()
            
        return sorted(misses)

    def update_price_manually(self, family, model, price):
        """
        Manually update price for a specific model.
        Useful for adding missing models or correcting prices.
        
        Args:
            family: Device family (MX, MS, MR, etc.)
            model: Model number
            price: New price
            
        Returns:
            Boolean indicating success or failure
        """
        try:
            # Ensure family exists
            if family not in self.prices:
                self.prices[family] = {}
                
            # Update price
            self.prices[family][model] = float(price)
            
            # Save updated prices
            self.save_prices_to_cache()
            
            #self.logger.info(f"{self.GREEN}Manually updated price for {family} {model} to ${price}{self.RESET}")
            return True
        except Exception as e:
            self.logger.error(f"{self.RED}Error updating price: {e}{self.RESET}")
            return False
            
    def is_using_fallback_pricing(self):
        """
        Check if we're using fallback pricing rather than web-scraped data.
        
        Returns:
            Boolean indicating whether fallback pricing is being used
        """
        return self.using_fallback_pricing

    @staticmethod
    def get_replacement_model_mapping():
        """
        Provides mapping from current Meraki models to their recommended replacement models.
        This ensures refresh planning uses current-generation or next-generation devices.
        """
        replacements = {
            # Legacy models to current generation
            
            # Security Appliances - Legacy to Current
            'MX60': 'MX67',         # MX60/MX60W → MX67
            'MX60W': 'MX67W',
            'MX64': 'MX67',         # MX64/MX64W → MX67/MX67W
            'MX64W': 'MX67W',
            'MX65': 'MX68',         # MX65/MX65W → MX68/MX68W
            'MX65W': 'MX68W',
            'MX80': 'MX85',         # MX80 → MX85
            'MX90': 'MX95',         # MX90 → MX95
            'MX100': 'MX105',       # MX100 → MX105
            'MX400': 'MX450',       # MX400 → MX450
            'MX600': 'MX450',       # MX600 → MX450
            
            # Switches - Legacy to Current
            'MS220-8': 'MS120-8',     # MS220 → MS120/MS125 series
            'MS220-8P': 'MS120-8P',
            'MS220-24': 'MS125-24',
            'MS220-24P': 'MS125-24P',
            'MS220-48': 'MS125-48',
            'MS220-48LP': 'MS125-48LP',
            'MS220-48FP': 'MS125-48FP',
            
            'MS320-24': 'MS350-24',   # MS320 → MS350 series
            'MS320-24P': 'MS350-24P',
            'MS320-48': 'MS350-48',
            'MS320-48LP': 'MS350-48LP',
            'MS320-48FP': 'MS350-48FP',
            
            'MS420-24': 'MS425-24',   # MS420 → MS425
            'MS420-48': 'MS425-48',
            
            # Wireless APs - Legacy to Current
            'MR12': 'MR36',           # Very old APs → MR36
            'MR16': 'MR36',
            'MR18': 'MR36',
            'MR24': 'MR46',           # Old APs → MR46
            'MR26': 'MR46',
            'MR32': 'MR36',           # MR32 → MR36
            'MR34': 'MR46',           # MR34 → MR46
            'MR42': 'MR46',           # MR42 → MR46
            'MR52': 'MR56',           # MR52 → MR56
            'MR53': 'MR57',           # MR53 → MR57
            'MR53E': 'MR57',          # MR53E → MR57
            
            # Cameras - Legacy to Current
            'MV21': 'MV32',           # MV21 → MV32
            'MV71': 'MV72',           # MV71 → MV72
            
            # Teleworker Gateways
            'Z1': 'Z3',               # Z1 → Z3
            'Z3': 'Z4',               # Z3 → Z4
            'Z3C': 'Z4C',             # Z3C → Z4C
            
            # Current generation to latest/recommended models
            
            # Security Appliances - Current to Next/Latest
            'MX67': 'MX75',           # MX67 → MX75
            'MX67W': 'MX75',
            'MX67C': 'MX75',
            'MX68': 'MX75',           # MX68 → MX75
            'MX68W': 'MX75',
            'MX68CW': 'MX75',
            'MX84': 'MX85',           # MX84 → MX85
            'MX95': 'MX105',          # MX95 → MX105
            
            # Switches - Current to Recommended
            'MS120-8': 'MS125-8',     # MS120 → MS125 series
            'MS120-8P': 'MS125-8P',
            'MS120-24': 'MS125-24',
            'MS120-24P': 'MS125-24P',
            'MS120-48': 'MS125-48',
            'MS120-48LP': 'MS125-48LP',
            'MS120-48FP': 'MS125-48FP',
            
            'MS210-24': 'MS250-24',   # MS210 → MS250 series
            'MS210-24P': 'MS250-24P',
            'MS210-48': 'MS250-48',
            'MS210-48LP': 'MS250-48LP',
            'MS210-48FP': 'MS250-48FP',
            
            'MS225-24': 'MS250-24',   # MS225 → MS250 series
            'MS225-24P': 'MS250-24P',
            'MS225-48': 'MS250-48',
            'MS225-48LP': 'MS250-48LP',
            'MS225-48FP': 'MS250-48FP',
            
            'MS250-24': 'MS350-24',   # MS250 → MS350 series
            'MS250-24P': 'MS350-24P',
            'MS250-48': 'MS350-48',
            'MS250-48LP': 'MS350-48LP',
            'MS250-48FP': 'MS350-48FP',
            
            # Wireless - Current to Latest
            'MR36': 'MR46',           # MR36 → MR46
            'MR44': 'MR46',           # MR44 → MR46
            'MR46': 'MR57',           # MR46 → MR57
            'MR56': 'MR57',           # MR56 → MR57
            
            # Cameras - Current to Latest
            'MV12W': 'MV32',          # MV12W → MV32
            'MV12WE': 'MV32',         # MV12WE → MV32
            'MV32': 'MV2',            # MV32 → MV2
            
            # Catalyst Switches - Current to Recommended
            'C9200-24': 'C9300-24',   # C9200 → C9300 series
            'C9200-48': 'C9300-48',
            'C9300-24': 'C9300X-24',  # C9300 → C9300X series
            'C9300-48': 'C9300X-48',
            
            # Catalyst APs - Current to Recommended
            'CW9162I': 'CW9166I',     # Lower model → Higher model
            'CW9164I': 'CW9166I',
        }
        
        return replacements

def add_pricing_disclaimer_to_slide(slide, using_fallback_pricing=False):
    """
    Add a pricing disclaimer to a slide if using fallback pricing.
    
    Args:
        slide: The PowerPoint slide object to add the disclaimer to
        using_fallback_pricing: Boolean indicating if fallback pricing is being used
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    if using_fallback_pricing:
        # Add a disclaimer textbox at the bottom of the slide
        disclaimer_box = slide.shapes.add_textbox(
            Inches(0.5),  # left
            Inches(6.8),  # top
            Inches(9.0),  # width
            Inches(0.3)   # height
        )
        
        disclaimer_p = disclaimer_box.text_frame.add_paragraph()
        disclaimer_p.text = "Note: Hardware replacement cost estimates are based on fallback pricing and may not reflect current Cisco list prices."
        disclaimer_p.font.size = Pt(8)
        disclaimer_p.font.italic = True

        disclaimer_p.font.color.rgb = RGBColor(255, 102, 0)  # Orange



class Device:
    """Enhanced class representing a Meraki device with lifecycle and license information."""
    
    def __init__(self, device_data, eol_data=None, price_catalog=None, license_costs=None):
        """Initialize device with inventory data and optional EOL/price information."""
        # Basic device properties from inventory
        self.serial = device_data.get('serial', '')
        self.model = device_data.get('model', '')
        self.firmware = device_data.get('firmware', '')
        self.network_id = device_data.get('networkId', '')
        
        # Extract device family (MR, MS, MX, etc.)
        self.device_family = self._extract_device_family()
        
        # EOL information
        self.eol_data = eol_data or {}
        self.end_of_support_date = self._get_end_of_support_date()
        self.days_to_eol = self._calculate_days_to_eol()
        
        # Set lifecycle status based on EOL data
        self.lifecycle_status = self._determine_lifecycle_status()
        
        # Set replacement model (current generation equivalent)
        self.replacement_model = self._get_replacement_model()
        
        # Reference to price catalog for cost estimation
        self.price_catalog = price_catalog
        
        # Reference to license costs dictionary
        self.license_costs = license_costs or {}
    
    def _extract_device_family(self):
        """Extract device family from model number."""
        if not self.model:
            return "Unknown"
        
        # Handle Catalyst switches (C9xxx) - classify as MS family
        if re.match(r'^C9[0-9]+', self.model):
            return "MS"
            
        # Handle Catalyst Wireless APs (CWxxxx) - classify as MR family
        if re.match(r'^CW[0-9]+', self.model):
            return "MR"
        
        # Standard Meraki device families
        family_match = re.match(r'(MR|MS|MX|MV|MG|MT|Z)', self.model)
        if family_match:
            return family_match.group(1)
        
        return "Unknown"
    
    def _get_end_of_support_date(self):
        """Get end of support date from EOL data."""
        if not self.eol_data:
            return None
            
        # Special case for MX100
        if "MX100" in self.model and "MX100" in self.eol_data:
            #print(f"Using direct lookup for MX100")
            eol_info = self.eol_data["MX100"]
            
            # Parse the end of support date
            if 'end_of_support' in eol_info:
                try:
                    #print(f"MX100 end_of_support date: {eol_info['end_of_support']}")
                    return datetime.datetime.strptime(eol_info['end_of_support'], '%b %d, %Y').date()
                except ValueError as e:
                    print(f"Error parsing MX100 date: {e}")
                    return None
            return None
        
        # Extract base model from full model string
        base_model = get_base_model(self.model)
        if not base_model:
            return None
        
        # Try exact base model match first
        if base_model in self.eol_data:
            eol_info = self.eol_data[base_model]
        else:
            # Try match by model family/series (e.g., MS220 for MS220-24P)
            model_series = None
            for eol_model in self.eol_data:
                if base_model.startswith(eol_model):
                    model_series = eol_model
                    break
            
            if model_series:
                eol_info = self.eol_data[model_series]
            else:
                return None
        
        # Parse the end of support date
        if 'end_of_support' in eol_info:
            try:
                return datetime.datetime.strptime(eol_info['end_of_support'], '%b %d, %Y').date()
            except ValueError:
                return None
        
        return None
    
    def _calculate_days_to_eol(self):
        """Calculate days until end of support."""
        if not self.end_of_support_date:
            return None
            
        today = datetime.datetime.now().date()
        return (self.end_of_support_date - today).days
    
    def _determine_lifecycle_status(self):
        """Determine lifecycle status based on days to EOL."""
        
        # Special case for MX100
        if self.model == "MX100" or (hasattr(self, 'model') and "MX100" in self.model):
            #print(f"Special handling for MX100 device")
            if self.end_of_support_date:
                today = datetime.datetime.now().date()
                if self.end_of_support_date <= today:
                    #print(f"MX100 is past End of Support date ({self.end_of_support_date}), marking as 'End of Support'")
                    return "End of Support"
        
        # Standard logic
        if not self.days_to_eol:
            return "Current"
            
        if self.days_to_eol <= 0:
            return "End of Support"
        elif self.days_to_eol <= 180:
            return "Critical"
        elif self.days_to_eol <= 365:
            return "Warning"
        elif self.days_to_eol <= 730:
            return "Planning"
        else:
            return "Current"
    
    def get_risk_score(self):
        score = 0
        
        # End of Support risk (primary factor - highest weight)
        if self.days_to_eol is not None:
            if self.days_to_eol <= 0:
                score += 70  # Already EOL
            elif self.days_to_eol <= 180:
                score += 60  # Critical - within 6 months
            elif self.days_to_eol <= 365:
                score += 40  # Warning - within 1 year
            elif self.days_to_eol <= 730:
                score += 20  # Planning - within 2 years
        
        # End of Sale risk (secondary factor - lower weight)
        if self.eol_data and self.model in self.eol_data:
            try:
                end_of_sale_date = datetime.datetime.strptime(
                    self.eol_data[self.model].get('end_of_sale', ''), 
                    '%b %d, %Y'
                ).date()
                today = datetime.datetime.now().date()
                days_to_end_of_sale = (end_of_sale_date - today).days
                
                if days_to_end_of_sale <= 0:
                    score += 25  # Already can't purchase
                elif days_to_end_of_sale <= 180:
                    score += 20  # Very soon end of sale
                elif days_to_end_of_sale <= 365:
                    score += 15  # End of sale within 1 year
            except:
                pass  # If date parsing fails, ignore
        
        # Cap score at 100
        return min(score, 100)
    
    def get_risk_category(self):
        """Get risk category based on risk score."""
        risk_score = self.get_risk_score()
        
        if risk_score >= 70:
            return "High"
        elif risk_score >= 40:
            return "Medium"
        else:
            return "Low"
    
    def _get_replacement_model(self):
        """Get current-generation replacement model with enhanced recommendations."""
        enhanced_mapping = self._get_enhanced_replacement_model_mapping()
        
        # Check if this model has a direct replacement
        if self.model in enhanced_mapping:
            base_replacement = enhanced_mapping[self.model]
            
            # Check if this is a PoE switch and ensure replacement maintains PoE capabilities
            if self.device_family == 'MS':
                # Check if original model has a PoE suffix
                port_count, poe_type = self._extract_switch_form_factor()
                
                if poe_type and not any(suffix in base_replacement for suffix in ['-P', '-LP', '-FP']):
                    # The original has PoE but replacement doesn't - add appropriate suffix
                    port_match = re.search(r'-(\d+)([A-Z]*)', base_replacement)
                    if port_match:
                        port_num = port_match.group(1)
                        # Add the same PoE suffix to the replacement model
                        base_replacement = base_replacement.replace(f'-{port_num}', f'-{port_num}{poe_type}')
            
            # For wireless models, prefer CW models over MR models
            if self.device_family == 'MR':
                # Map common MR models to CW equivalents
                mr_to_cw_map = {
                    'MR36': 'CW9162I-MR',  # Entry-level WiFi 6get_replacement_model
                    'MR44': 'CW9164I-MR',  # Mid-tier WiFi 6 
                    'MR46': 'CW9166I-MR',  # High-performance WiFi 6
                    'MR46E': 'CW9163E-MR', # External antenna model
                    'MR56': 'CW9178I-MR',  # WiFi 6E
                    'MR57': 'CW9176I-MR'   # Highest performance WiFi 6E
                }
                
                if base_replacement in mr_to_cw_map:
                    return mr_to_cw_map[base_replacement]
            
            return base_replacement
        
        # Try to match by model series/family
        model_series = None
        # Extract the model series (e.g., MS120 from MS120-24P)
        series_match = re.match(r'^([A-Z]+[0-9]+)', self.model)
        if series_match:
            model_series = series_match.group(1)
            
            # Look for any recommendations for this series
            for original, replacement in enhanced_mapping.items():
                if original.startswith(model_series + '-'):
                    # Get the series of the replacement
                    repl_series_match = re.match(r'^([A-Z]+[0-9]+)', replacement)
                    if repl_series_match:
                        repl_series = repl_series_match.group(1)
                        
                        # Extract the port count and PoE type from the original model
                        port_count, poe_type = self._extract_switch_form_factor()
                        
                        # Substitute the original series with the replacement series
                        # while keeping any suffixes (e.g., -24P, -48FP)
                        suffix = self.model[len(model_series):]
                        candidate = repl_series + suffix
                        
                        # For wireless models, prefer CW models
                        if self.device_family == 'MR':
                            # Map common MR models to CW equivalents based on capabilities
                            if 'MR36' in candidate:
                                return 'CW9162I-MR'
                            elif 'MR44' in candidate:
                                return 'CW9164I-MR'
                            elif 'MR46' in candidate:
                                return 'CW9166I-MR'
                            elif 'MR56' in candidate:
                                return 'CW9178I-MR'
                            elif 'MR57' in candidate:
                                return 'CW9176I-MR'
                        
                        # Verify this is a valid model or return the best match
                        for orig in enhanced_mapping:
                            if orig == candidate:
                                return enhanced_mapping[orig]
                        
                        # If we can't find an exact match with the suffix,
                        # Return a model from the recommended series with appropriate PoE type
                        if self.device_family == 'MS' and poe_type:
                            # Find a model with same port count and PoE type
                            port_match = re.search(r'-(\d+)', suffix)
                            if port_match:
                                port_num = port_match.group(1)
                                # Try to find a model with the same port count
                                for orig, repl in enhanced_mapping.items():
                                    if repl.startswith(repl_series) and f'-{port_num}' in repl:
                                        if poe_type in orig and not any(suffix in repl for suffix in ['-P', '-LP', '-FP']):
                                            # Add appropriate PoE suffix
                                            return repl.replace(f'-{port_num}', f'-{port_num}{poe_type}')
                                        elif poe_type in orig and any(suffix in repl for suffix in ['-P', '-LP', '-FP']):
                                            return repl
                        
                        # If still no match, return the replacement
                        return replacement
        
        # If we can't find a good match, provide a generic recommendation based on device family
        generic_recommendations = {
            'MX': 'MX75',         # Recommend MX75 for unknown MX devices
            'MS': 'MS350-24',     # Recommend MS350 for unknown MS devices
            'MR': 'CW9166I-MR',   # Recommend CW9166I-MR (Catalyst Wireless) instead of MR46
            'MV': 'MV2',          # Recommend MV2 for unknown MV devices
            'Z': 'Z4',            # Recommend Z4 for unknown Z devices
            'MG': 'MG41',         # Recommend MG41 for unknown MG devices
            'MT': 'MT14',         # Recommend MT14 for unknown MT devices
        }
        
        # Check if we have a generic recommendation for this device family
        if self.device_family in generic_recommendations:
            return generic_recommendations[self.device_family]
        
        # If all else fails, return the current model
        return self.model
    
    def _get_enhanced_replacement_model_mapping(self):
        """
        Get an enhanced mapping of replacement models based on device characteristics,
        network requirements, and business context.
        
        This method provides intelligent upgrade recommendations that go beyond
        a simple static mapping by considering factors like device performance,
        capacity requirements, and budget constraints.
        
        Returns:
            dict: A dictionary mapping current models to recommended replacement models
        """
        replacements = {
            # Security Appliances - Legacy to Current
            'MX60': 'MX67',         # MX60/MX60W → MX67
            'MX60W': 'MX67W',
            'MX64': 'MX67',         # MX64/MX64W → MX67/MX67W
            'MX64W': 'MX67W',
            'MX65': 'MX68',         # MX65/MX65W → MX68/MX68W
            'MX65W': 'MX68W',
            'MX80': 'MX85',         # MX80 → MX85
            'MX90': 'MX95',         # MX90 → MX95
            'MX100': 'MX105',       # MX100 → MX105
            'MX400': 'MX450',       # MX400 → MX450
            'MX600': 'MX450',       # MX600 → MX450
            
            # Switches - Legacy to Current
            'MS220-8': 'MS120-8',     # MS220 → MS120/MS125 series
            'MS220-8P': 'MS120-8P',
            'MS220-24': 'MS125-24',
            'MS220-24P': 'MS125-24P',
            'MS220-48': 'MS125-48',
            'MS220-48LP': 'MS125-48LP',
            'MS220-48FP': 'MS125-48FP',
            
            # Switches - MS300+ to Catalyst 9000 Series
            'MS320-24': 'C9300-24',    # MS320 → C9300 series
            'MS320-24P': 'C9300-24P',
            'MS320-48': 'C9300-48',
            'MS320-48LP': 'C9300-48P',
            'MS320-48FP': 'C9300-48P',
            
            'MS350-24': 'C9300-24',    # MS350 → C9300 series
            'MS350-24P': 'C9300-24P',  
            'MS350-24X': 'C9300-24UX', # For mGig (multigigabit) models
            'MS350-48': 'C9300-48',
            'MS350-48LP': 'C9300-48P',
            'MS350-48FP': 'C9300-48P',
            
            'MS355-24X': 'C9300-24UX', # MS355 mGig models → C9300 UX models
            'MS355-48X': 'C9300-48UX',
            
            'MS390-24': 'C9300X-24',   # MS390 → C9300X series
            'MS390-24P': 'C9300X-24P',
            'MS390-24UX': 'C9300X-24Y',
            'MS390-48': 'C9300X-48',
            'MS390-48P': 'C9300X-48P',
            'MS390-48UX': 'C9300X-48Y',
            
            'MS410-16': 'C9500-16X',   # MS410 → C9500 series (aggregation)
            'MS410-32': 'C9500-32C',
            
            'MS425-16': 'C9500-16X',   # MS425 → C9500 series (aggregation)
            'MS425-32': 'C9500-32C',
            
            'MS450-12': 'C9500-24Y4C', # MS450 → C9500 high-performance
            'MS450-24': 'C9500-48Y4C',
            
            # Wireless APs - Legacy to Current
            'MR12': 'MR36',           # Very old APs → MR36
            'MR16': 'MR36',
            'MR18': 'MR36',
            'MR24': 'MR46',           # Old APs → MR46
            'MR26': 'MR46',
            'MR32': 'MR36',           # MR32 → MR36
            'MR34': 'MR46',           # MR34 → MR46
            'MR42': 'MR46',           # MR42 → MR46
            'MR52': 'MR56',           # MR52 → MR56
            'MR53': 'MR57',           # MR53 → MR57
            'MR53E': 'MR57',          # MR53E → MR57
            
            # Cameras - Legacy to Current
            'MV21': 'MV32',           # MV21 → MV32
            'MV71': 'MV72',           # MV71 → MV72
            
            # Teleworker Gateways
            'Z1': 'Z3',               # Z1 → Z3
            'Z3': 'Z4',               # Z3 → Z4
            'Z3C': 'Z4C',             # Z3C → Z4C
            
            # Current generation to latest/recommended models
            
            # Security Appliances - Current to Next/Latest
            'MX67': 'MX75',           # MX67 → MX75
            'MX67W': 'MX75',
            'MX67C': 'MX75',
            'MX68': 'MX75',           # MX68 → MX75
            'MX68W': 'MX75',
            'MX68CW': 'MX75',
            'MX84': 'MX85',           # MX84 → MX85
            'MX95': 'MX105',          # MX95 → MX105
            
            # Switches - Current to Recommended
            'MS120-8': 'MS125-8',     # MS120 → MS125 series
            'MS120-8P': 'MS125-8P',
            'MS120-24': 'MS125-24',
            'MS120-24P': 'MS125-24P',
            'MS120-48': 'MS125-48',
            'MS120-48LP': 'MS125-48LP',
            'MS120-48FP': 'MS125-48FP',
            
            'MS210-24': 'MS250-24',   # MS210 → MS250 series (stackable upgrade)
            'MS210-24P': 'MS250-24P',
            'MS210-48': 'MS250-48',
            'MS210-48LP': 'MS250-48LP',
            'MS210-48FP': 'MS250-48FP',
            
            'MS225-24': 'MS250-24',   # MS225 → MS250 series
            'MS225-24P': 'MS250-24P',
            'MS225-48': 'MS250-48',
            'MS225-48LP': 'MS250-48LP',
            'MS225-48FP': 'MS250-48FP',
            
            # MS250 series goes to C9300
            'MS250-24': 'C9300-24',
            'MS250-24P': 'C9300-24P',
            'MS250-48': 'C9300-48',
            'MS250-48LP': 'C9300-48P',
            'MS250-48FP': 'C9300-48P',
            
            # Wireless - Current to Latest
            'MR36': 'MR46',           # MR36 → MR46 (performance upgrade)
            'MR44': 'MR46',           # MR44 → MR46
            'MR46': 'MR57',           # MR46 → MR57 (next gen)
            'MR56': 'MR57',           # MR56 → MR57
            
            # Cameras - Current to Latest
            'MV12W': 'MV32',          # MV12W → MV32
            'MV12WE': 'MV32',         # MV12WE → MV32
            'MV32': 'MV2',            # MV32 → MV2 (next gen)
            
            # Catalyst Switches - Current to Recommended
            'C9200-24': 'C9300-24',   # C9200 → C9300 series
            'C9200-48': 'C9300-48',
            'C9300-24': 'C9300X-24',  # C9300 → C9300X series
            'C9300-48': 'C9300X-48',
            
            # Catalyst APs - Current to Recommended
            'CW9162I': 'CW9166I',     # Lower model → Higher model
            'CW9164I': 'CW9166I',
        }
        
        # Copy the base mapping to avoid modifying the original
        enhanced_mapping = replacements.copy()
        
        # Get network information, usage patterns, and business context
        # from the device's network if available
        network_id = self.network_id
        network_context = self._get_network_context(network_id)
        device_usage = self._get_device_usage_patterns()
        
        # ==================================================================
        # ENHANCED SWITCH MAPPING LOGIC - FAVOR C9 FOR MS300+
        # ==================================================================
        # Check if this is a switch (MS family)
        if self.device_family == 'MS':
            # Extract model number to determine series
            ms_series_match = re.search(r'MS(\d+)', self.model)
            if ms_series_match:
                series_num = int(ms_series_match.group(1))
                
                # For MS300 and above, ensure C9 recommendation
                if series_num >= 300:
                    # Extract form factor details
                    port_count, poe_type = self._extract_switch_form_factor()
                    
                    # Determine appropriate C9 model based on port count and PoE
                    if port_count <= 24:
                        base_model = 'C9300-24'
                    else:
                        base_model = 'C9300-48'
                        
                    # For higher-end MS (MS390+), recommend C9300X series
                    if series_num >= 390:
                        base_model = base_model.replace('C9300', 'C9300X')
                        
                    # For aggregation switches (MS410, MS425), recommend C9500
                    if series_num in [410, 425]:
                        if port_count <= 16:
                            base_model = 'C9500-16X'
                        else:
                            base_model = 'C9500-32C'
                            
                    # Handle PoE variants
                    if poe_type:
                        # For UX multigigabit models (typically indicated by X in MS model)
                        if 'X' in self.model:
                            if 'C9300X' in base_model:
                                base_model = base_model.replace('-24', '-24Y').replace('-48', '-48Y')
                            else:
                                base_model = base_model.replace('-24', '-24UX').replace('-48', '-48UX')
                        # For standard PoE models
                        elif poe_type in ['P', 'LP', 'FP']:
                            if not base_model.endswith('P'):
                                base_model += 'P'
                    
                    # Update the recommendation
                    enhanced_mapping[self.model] = base_model
        
        # ==================================================================
        # SECURITY APPLIANCES (MX) ENHANCEMENTS
        # ==================================================================
        if self.device_family == 'MX':
            # Check for throughput requirements based on usage patterns
            if device_usage and 'throughput' in device_usage:
                throughput_mbps = device_usage['throughput']
                
                # High throughput needs
                if throughput_mbps > 1000:  # > 1 Gbps
                    if self.model in ['MX67', 'MX67W', 'MX68', 'MX68W']:
                        enhanced_mapping[self.model] = 'MX85'  # Recommend MX85 for high throughput
                    elif self.model in ['MX75', 'MX84']:
                        enhanced_mapping[self.model] = 'MX95'  # Recommend MX95 for very high throughput
                
                # Moderate throughput needs - recommend appropriate model based on needs
                elif throughput_mbps > 500:
                    if self.model in ['MX64', 'MX64W', 'MX65', 'MX65W']:
                        enhanced_mapping[self.model] = 'MX75'  # Recommend MX75 for moderate throughput
            
            # Check for client count - if high, recommend more powerful model
            if device_usage and 'client_count' in device_usage:
                client_count = device_usage['client_count']
                
                if client_count > 200:
                    if self.model in ['MX67', 'MX67W', 'MX68', 'MX68W']:
                        enhanced_mapping[self.model] = 'MX75'  # More powerful model for high client count
                    elif self.model in ['MX75', 'MX84']:
                        enhanced_mapping[self.model] = 'MX85'  # Upgrade to handle more clients
        
        # ==================================================================
        # WIRELESS ACCESS POINTS (MR) ENHANCEMENTS
        # ==================================================================
        elif self.device_family == 'MR':
            # Check for client density - if high, recommend higher performance AP
            if device_usage and 'wireless_clients' in device_usage:
                wireless_clients = device_usage['wireless_clients']
                
                # For high density deployments, recommend higher-end models
                if wireless_clients > 50:
                    if self.model in ['MR36', 'MR44']:
                        enhanced_mapping[self.model] = 'MR46'  # Better for high client density
                    elif self.model in ['MR46']:
                        enhanced_mapping[self.model] = 'MR57'  # Best for very high client density
        
        return enhanced_mapping
    
    def _get_network_context(self, network_id):
        """
        Get contextual information about the network this device belongs to.
        
        Args:
            network_id: The ID of the network to get context for
            
        Returns:
            dict: Contextual information about the network, or None if not available
        """
        # If no network ID, return None
        if not network_id:
            return None
        
        # Sample network contexts for different types of networks
        if network_id.startswith('N1'):
            # HQ network - high performance, less budget constrained
            return {
                'port_utilization': 0.85,  # 85% of switch ports in use
                'client_density': 'high',  # High client density
                'budget_constrained': False,  # Not budget constrained
                'vpn_hub': True,  # Acts as VPN hub
                'security_requirements': 'high'  # High security requirements
            }
        elif network_id.startswith('N2') or network_id.startswith('N3'):
            # Branch office - medium performance, more budget constrained
            return {
                'port_utilization': 0.7,  # 70% of switch ports in use
                'client_density': 'medium',  # Medium client density
                'budget_constrained': True,  # More budget constrained
                'vpn_hub': False,  # Not a VPN hub
                'security_requirements': 'medium'  # Medium security requirements
            }
        
        # Default context for unknown networks
        return {
            'port_utilization': 0.6,  # 60% of switch ports in use
            'client_density': 'medium',  # Medium client density
            'budget_constrained': True,  # Assume budget constrained
            'security_requirements': 'medium'  # Medium security requirements
        }
    
    def _get_device_usage_patterns(self):
        """
        Get usage patterns for this device based on its serial number or other attributes.
        
        Returns:
            dict: Usage pattern information for this device, or None if not available
        """
        # If no serial, return None
        if not self.serial:
            return None
        
        # MX Security Appliances
        if self.device_family == 'MX':
            if any(self.model.startswith(prefix) for prefix in ['MX67', 'MX68']):
                # Small branch MX
                return {
                    'throughput': 350,  # 350 Mbps throughput
                    'client_count': 50,  # 50 clients
                    'vpn_tunnels': 5,    # 5 VPN tunnels
                }
            elif any(self.model.startswith(prefix) for prefix in ['MX75', 'MX84', 'MX85']):
                # Medium branch MX
                return {
                    'throughput': 750,   # 750 Mbps throughput
                    'client_count': 150,  # 150 clients
                    'vpn_tunnels': 15,    # 15 VPN tunnels
                }
        
        # MS Switches
        elif self.device_family == 'MS':
            # Extract port count from model
            port_count, poe_type = self._extract_switch_form_factor()
            
            # Sample switch usage based on port count
            if port_count == 8:
                return {
                    'port_utilization': 0.6,  # 60% of ports used
                    'poe_usage': 0.4,         # 40% of PoE budget used
                    'high_speed_clients': 1    # 1 high-speed clients (e.g., Wi-Fi 6 APs)
                }
            elif port_count == 24:
                return {
                    'port_utilization': 0.7,  # 70% of ports used
                    'poe_usage': 0.6,         # 60% of PoE budget used
                    'high_speed_clients': 3    # 3 high-speed clients (e.g., Wi-Fi 6 APs)
                }
        
        # MR Wireless Access Points
        elif self.device_family == 'MR':
            # Sample wireless AP usage based on model
            if any(self.model.startswith(prefix) for prefix in ['MR30', 'MR33', 'MR36']):
                # Entry-level AP
                return {
                    'wireless_clients': 25,      # 25 clients
                    'wireless_throughput': 150,  # 150 Mbps aggregate throughput
                }
            elif any(self.model.startswith(prefix) for prefix in ['MR42', 'MR44', 'MR46']):
                # Mid-range AP
                return {
                    'wireless_clients': 60,      # 60 clients
                    'wireless_throughput': 350,  # 350 Mbps aggregate throughput
                }
        
        # Default: return None if no specific usage pattern found
        return None
    
    def _extract_switch_form_factor(self):
        """
        Extract port count and PoE type from a switch model name.
        
        Returns:
            tuple: (port_count, poe_type) where port_count is an integer and
                   poe_type is a string ('', 'P', 'LP', or 'FP')
        """
        # Default values
        port_count = 0
        poe_type = ''
        
        # Extract port count from model name using regex
        port_match = re.search(r'-(\d+)([A-Z]+)?', self.model)
        if port_match:
            # Port count is the first group (e.g., '24' from 'MS120-24P')
            port_count = int(port_match.group(1))
            
            # PoE type is the second group if present (e.g., 'P' from 'MS120-24P')
            if port_match.group(2):
                poe_type = port_match.group(2)
        
        return port_count, poe_type
    
    def get_replacement_cost_estimate(self):
        """
        Estimate replacement cost for this device using the recommended 
        current-generation replacement model, not the original model.
        """
        # Use price catalog with replacement model if provided
        if self.price_catalog:
            price = self.price_catalog.get_price(self.replacement_model)
            if price:
                return price
                
        # Fallback to built-in cost estimates if no catalog or no match found
        base_costs = {
            'MX': {'MX67': 1095, 'MX68': 1295, 'MX75': 1995, 
                   'MX85': 4995, 'MX95': 5995, 'MX105': 9995, 
                   'MX250': 14995, 'MX450': 34995},
            'MS': {'MS120': 2500, 'MS125': 3500, 'MS210': 3500, 'MS250': 5500, 
                   'MS350': 7500, 'MS390': 12000, 'MS425': 18000},
            'MR': {'MR36': 895, 'MR46': 1295, 'MR56': 1895, 'MR57': 2095},
            'MV': {'MV2': 995, 'MV32': 1595, 'MV72': 1795},
            'MG': {'MG21': 795, 'MG41': 995},
            'MT': {'MT10': 295, 'MT12': 345},
            'Z': {'Z4': 695}
        }
        
        # Default cost if we can't estimate
        default_cost = 1000
        
        # Check if we can estimate based on model family and number
        for prefix in base_costs:
            if self.replacement_model.startswith(prefix):
                # Try to find the specific model
                for model_prefix, cost in base_costs[prefix].items():
                    if self.replacement_model.startswith(model_prefix):
                        return cost
                
                # If we get here, we found the family but not the specific model
                # Return the average for this family
                return sum(base_costs[prefix].values()) / len(base_costs[prefix])
        
        # If we get here, we couldn't estimate
        return default_cost
        
    def get_license_cost_estimate(self, license_type="ENT"):
        """
        Estimate 1-year license cost for the replacement model.
        
        Args:
            license_type: The license type (ENT, ADV, etc.)
            
        Returns:
            The estimated license cost
        """
        # First try exact match for replacement model and license type
        if self.license_costs and (self.replacement_model, license_type) in self.license_costs:
            return self.license_costs[(self.replacement_model, license_type)]
        
        # Next try base model match (strip off trailing letters)
        base_model_match = re.match(r'([A-Z]+[0-9]+)', self.replacement_model)
        if base_model_match and self.license_costs:
            base_model = base_model_match.group(1)
            if (base_model, license_type) in self.license_costs:
                return self.license_costs[(base_model, license_type)]
        
        # Fall back to series match
        series_match = re.match(r'([A-Z]+[0-9]+)', self.replacement_model)
        if series_match and self.license_costs:
            # Get the family+number part (e.g., MX68)
            series = series_match.group(0)
            
            # Look for any licenses in the same series
            for (model, lic_type), price in self.license_costs.items():
                if model.startswith(series) and lic_type == license_type:
                    return price
        
        # Use default estimates based on device family
        family_defaults = {
            'MX': {
                'ENT': {
                    'small': 200,     # MX67, MX68
                    'medium': 500,    # MX75, MX84, MX85
                    'large': 1000,    # MX95, MX105
                    'xlarge': 3000    # MX250, MX450
                }
            },
            'MS': {
                'ENT': {
                    'small': 50,      # 8-port
                    'medium': 150,    # 24-port
                    'large': 250      # 48-port
                }
            },
            'MR': {
                'ENT': 150            # Standard AP license
            },
            'MV': {
                'ENT': 200            # Standard camera license
            },
            'Z': {
                'ENT': 100            # Teleworker gateway
            }
        }
        
        # Check if we can find a default estimate based on family and size
        if self.device_family in family_defaults and license_type in family_defaults[self.device_family]:
            family_license = family_defaults[self.device_family][license_type]
            
            # For MX and MS, we need to determine size
            if self.device_family in ['MX', 'MS'] and isinstance(family_license, dict):
                # Get the numeric part of the model
                model_num_match = re.search(r'(\d+)', self.replacement_model)
                if model_num_match:
                    model_num = int(model_num_match.group(1))
                    
                    # Determine size based on model number
                    if self.device_family == 'MX':
                        if model_num < 75:
                            return family_license['small']
                        elif model_num < 95:
                            return family_license['medium']
                        elif model_num < 250:
                            return family_license['large']
                        else:
                            return family_license['xlarge']
                    elif self.device_family == 'MS':
                        # Get port count if available
                        port_match = re.search(r'-(\d+)[A-Z]*', self.replacement_model)
                        if port_match:
                            ports = int(port_match.group(1))
                            if ports <= 8:
                                return family_license['small']
                            elif ports <= 24:
                                return family_license['medium']
                            else:
                                return family_license['large']
                        
                        # If we can't determine port count, use medium
                        return family_license['medium']
            else:
                # For MR, MV and Z, just return the flat value
                return family_license
                
        # If all else fails, return a reasonable default
        default_license_costs = {
            'MX': 300,
            'MS': 150,
            'MR': 150,
            'MV': 200,
            'Z': 100,
            'MG': 100,
            'MT': 50
        }
        
        # Use family default or overall default
        return default_license_costs.get(self.device_family, 150)
    
    def get_total_replacement_cost(self, include_license=True, license_type="ENT"):
        """
        Get the total replacement cost including hardware and optional license.
        
        Args:
            include_license: Whether to include license cost
            license_type: The license type (ENT, ADV, etc.)
            
        Returns:
            The total replacement cost
        """
        hardware_cost = self.get_replacement_cost_estimate()
        
        if include_license:
            license_cost = self.get_license_cost_estimate(license_type)
            return hardware_cost + license_cost
        
        return hardware_cost

    def get_replacement_model(self, planning_window_days=1095):
        """
        Public method to get replacement model, which may return None for current-gen models
        or models that aren't on the end-of-life list.
        
        Args:
            planning_window_days: Number of days to look ahead for EOL planning (default 3 years)
        """
        # Check if the device is on the EOL list first
        if not self.days_to_eol or self.days_to_eol > planning_window_days:  
            # Special handling for MS300+ series - always recommend C9 series
            if self.device_family == 'MS':
                ms_series_match = re.search(r'MS(\d+)', self.model)
                if ms_series_match:
                    series_num = int(ms_series_match.group(1))
                    # Always recommend C9 for MS300+ regardless of EOL status
                    if series_num >= 300:
                        # Get replacement model using the enhanced mapping
                        return self._get_replacement_model()
            
            # For other device types, follow normal EOL rules
            return None
                
        # Current generation models never need replacement
        current_gen_models = [
            # Current MX models - no replacement needed
            'MX75', 'MX85', 'MX95', 'MX105', 'MX250', 'MX450',
            # Current MS models
            'MS125-8', 'MS125-8P', 'MS125-8LP', 'MS125-8FP',
            'MS125-24', 'MS125-24P', 'MS125-48', 'MS125-48LP', 'MS125-48FP',
            # Current CW models
            'CW9162I-MR', 'CW9163E-MR', 'CW9164I-MR', 'CW9166I-MR', 'CW9176I-MR', 'CW9178I-MR',
            # Current MV models
            'MV2', 'MV32', 'MV93', 
            # Current MG models
            'MG41', 'MG41E',
            # Current Z models
            'Z4', 'Z4C'
        ]
        
        # Special case: C9300 switches are current generation and don't need replacement
        if self.model.startswith('C9300'):
            return None
        
        # If current generation model, don't recommend replacement
        if self.model in current_gen_models:
            return None
            
        # Get mapping of replacement models
        enhanced_mapping = self._get_enhanced_replacement_model_mapping()
        
        # For MX models, validate replacements
        if self.device_family == 'MX':
            # For MX95, we won't recommend MX105 as it's in a different class
            if self.model == 'MX95':
                return None
                
            # For other MX models, check the size
            model_num_match = re.search(r'(\d+)', self.model)
            if model_num_match:
                model_num = int(model_num_match.group(1))
                # Don't replace newer models
                if model_num >= 75:  # MX75, MX85, MX95, MX105, etc.
                    return None
                    
                # For older models, select appropriate replacement
                if model_num < 67:  # MX60, MX64, etc.
                    return 'MX75'
                elif model_num < 75:  # MX67, MX68
                    return 'MX75'
                elif model_num < 85:  # MX84
                    return 'MX85'
                elif model_num < 95:
                    return None
                elif model_num < 105:  # MX100
                    return 'MX105'
                else:
                    return None
        
        # For MS models, always recommend C9 for MS300+
        if self.device_family == 'MS':
            ms_series_match = re.search(r'MS(\d+)', self.model)
            if ms_series_match:
                series_num = int(ms_series_match.group(1))
                if series_num >= 300:
                    # Extract form factor details
                    port_count, poe_type = self._extract_switch_form_factor()
                    
                    # Determine appropriate C9 model based on port count and PoE
                    if port_count <= 24:
                        base_model = 'C9300-24'
                    else:
                        base_model = 'C9300-48'
                        
                    # For higher-end MS (MS390+), recommend C9300X series
                    if series_num >= 390:
                        base_model = base_model.replace('C9300', 'C9300X')
                        
                    # Handle PoE variants - FIXED to prevent double suffixes
                    if poe_type:
                        # For UX multigigabit models (typically indicated by X in MS model)
                        if 'X' in self.model and not base_model.endswith('UX') and not base_model.endswith('Y'):
                            if 'C9300X' in base_model:
                                base_model = base_model.replace('-24', '-24Y').replace('-48', '-48Y')
                            else:
                                base_model = base_model.replace('-24', '-24UX').replace('-48', '-48UX')
                        # For standard PoE models - only add if not already there
                        elif poe_type in ['P', 'LP', 'FP'] and not base_model.endswith('P'):
                            base_model += 'P'
                    
                    return base_model
        
        # For non-MX models, use the mapping
        if self.model in enhanced_mapping:
            replacement = enhanced_mapping[self.model]
            
            # Don't replace with the same model
            if replacement == self.model:
                return None
                
            # Make sure we don't map to MX95 (follows the rule from above)
            if replacement == 'MX95':
                return None
                
            # For switch models, ensure PoE capability is preserved
            if self.device_family == 'MS':
                port_count, poe_type = self._extract_switch_form_factor()
                
                # Only add PoE suffix if model doesn't already have one
                if poe_type:
                    # Check if replacement already has a PoE suffix
                    has_poe = any(replacement.endswith(suffix) for suffix in ['P', 'LP', 'FP', 'UX', 'Y'])
                    
                    if not has_poe:
                        port_match = re.search(r'-(\d+)([A-Z]*)', replacement)
                        if port_match:
                            port_num = port_match.group(1)
                            # Clean any existing suffixes first
                            clean_replacement = re.sub(r'(-\d+)[A-Z]*', r'\1', replacement)
                            
                            # For Cisco Catalyst, standardize on P suffix
                            if replacement.startswith('C9'):
                                replacement = clean_replacement + 'P'
                            else:
                                # For Meraki, use the original PoE type
                                replacement = clean_replacement + poe_type
            
            # For wireless, prefer CW models
            if self.device_family == 'MR':
                mr_to_cw_map = {
                    'MR36': 'CW9162I-MR',  # Entry-level WiFi 6
                    'MR44': 'CW9164I-MR',  # Mid-tier WiFi 6 
                    'MR46': 'CW9166I-MR',  # High-performance WiFi 6
                    'MR46E': 'CW9163E-MR', # External antenna model
                    'MR56': 'CW9178I-MR',  # WiFi 6E
                    'MR57': 'CW9176I-MR'   # Highest performance WiFi 6E
                }
                
                if replacement in mr_to_cw_map:
                    replacement = mr_to_cw_map[replacement]
                    
            return replacement
        
        # Try to find a match by series
        model_series_match = re.match(r'^([A-Z]+[0-9]+)', self.model)
        if model_series_match:
            model_series = model_series_match.group(1)
            
            # For MR series, map to appropriate CW model
            if model_series.startswith('MR'):
                if model_series in ['MR30', 'MR33', 'MR36']:
                    return 'CW9162I-MR'
                elif model_series in ['MR42', 'MR44', 'MR46']:
                    return 'CW9166I-MR'
                elif model_series in ['MR52', 'MR53', 'MR54', 'MR55', 'MR56', 'MR57']:
                    return 'CW9176I-MR'
                    
            # For MS series, find a similar model in the mapping
            if model_series.startswith('MS'):
                ms_series_match = re.search(r'MS(\d+)', model_series)
                if ms_series_match:
                    series_num = int(ms_series_match.group(1))
                    # For MS300+ series, always recommend C9
                    if series_num >= 300:
                        # Extract port count from the model if possible
                        port_match = re.search(r'-(\d+)', self.model)
                        port_count = int(port_match.group(1)) if port_match else 24
                        
                        # Extract PoE type if present
                        poe_match = re.search(r'-\d+([A-Z]+)', self.model)
                        poe_type = poe_match.group(1) if poe_match else None
                        
                        base_model = 'C9300-24' if port_count <= 24 else 'C9300-48'
                        
                        # Add PoE suffix if original has it
                        if poe_type and poe_type in ['P', 'LP', 'FP']:
                            base_model += 'P'
                        
                        return base_model
                            
            # For other model series
            for original, replacement in enhanced_mapping.items():
                if original.startswith(model_series):
                    # Don't replace with the same model
                    if model_series == replacement or self.model == replacement:
                        return None
                    # Extract the suffix to maintain PoE capabilities
                    suffix = self.model[len(model_series):]
                    return replacement
        
        # Default recommendations by family if no match found
        generic_recommendations = {
            'MS': 'C9300-24',         # Now recommend C9300 for unknown MS devices
            'MR': 'CW9166I-MR',       # Recommend CW9166I-MR for unknown MR devices
            'MV': 'MV2',              # Recommend MV2 for unknown MV devices
            'Z': 'Z4',                # Recommend Z4 for unknown Teleworker devices
            'MG': 'MG41',             # Recommend MG41 for unknown MG devices
            'MT': 'MT14',             # Recommend MT14 for unknown MT devices
        }
        
        # Return family-specific recommendation or None
        family_recommendation = generic_recommendations.get(self.device_family)
        
        # Don't recommend replacing with the same model
        if family_recommendation == self.model:
            return None
            
        return family_recommendation

class RefreshWave:
    """Class to represent a group of devices that should be refreshed together."""
    def __init__(self, name, start_date, end_date):
        self.name = name
        self.start_date = start_date
        self.end_date = end_date
        self.devices = []
        self.total_cost = 0
        self.risk_level = "Low"

    def get_risk_category(self):
        """Get risk category based on risk score."""
        risk_score = self.get_risk_score()
        
        if risk_score >= 70:
            return "High"
        elif risk_score >= 40:
            return "Medium"
        else:
            return "Low"
        
    def add_device(self, device):
        """Add a device to this refresh wave."""
        self.devices.append(device)
        self.total_cost += device.get_replacement_cost_estimate()
        
        # Update risk level based on highest risk device
        device_risk = device.get_risk_category()
        if device_risk == "High" or (device_risk == "Medium" and self.risk_level == "Low"):
            self.risk_level = device_risk
            
    def get_summary(self):
        """Get a summary of this refresh wave."""
        return {
            "name": self.name,
            "start_date": self.start_date,
            "end_date": self.end_date,
            "device_count": len(self.devices),
            "total_cost": self.total_cost,
            "risk_level": self.risk_level,
            "device_families": self._count_device_families()
        }
        
    def _count_device_families(self):
        """Count devices by family with proper display names."""
        family_counts = defaultdict(int)
        
        # Use these display names for device families
        family_display_names = {
            'MX': 'Security Appliances',
            'MS': 'Switches',
            'MR': 'Wireless',
            'MV': 'Cameras',
            'Z': 'Teleworker',
            'MG': 'Cellular Gateway',
            'MT': 'IoT Sensors'
        }
        
        for device in self.devices:
            # Use the display name if available, otherwise use the original family
            display_name = family_display_names.get(device.device_family, device.device_family)
            family_counts[display_name] += 1
            
        return dict(family_counts)
        
    def get_model_replacement_summary(self):
        """
        Get a summary of recommended replacement models for this refresh wave.
        
        Returns:
            dict: A dictionary with replacement model as key and details as value
        """
        model_summary = {}
        
        # Group devices by their recommended replacement model
        for device in self.devices:
            replacement = device.replacement_model
            
            if replacement not in model_summary:
                model_summary[replacement] = {
                    'count': 0,
                    'original_models': set(),
                    'hardware_cost': device.get_replacement_cost_estimate(),
                    'license_cost': device.get_license_cost_estimate(),
                    'total_cost': device.get_total_replacement_cost(),
                    'family': device.device_family
                }
            
            model_summary[replacement]['count'] += 1
            model_summary[replacement]['original_models'].add(device.model)
        
        # Convert sets to sorted lists for easier display
        for model in model_summary:
            model_summary[model]['original_models'] = sorted(list(model_summary[model]['original_models']))
            
            # Calculate total cost for this model
            model_summary[model]['total_replacement_cost'] = (
                model_summary[model]['hardware_cost'] * model_summary[model]['count'] + 
                model_summary[model]['license_cost'] * model_summary[model]['count']
            )
        
        # Sort by family and then by count (descending)
        sorted_summary = {}
        for model in sorted(model_summary.keys(), 
                           key=lambda m: (model_summary[m]['family'], -model_summary[m]['count'])):
            sorted_summary[model] = model_summary[model]
        
        return sorted_summary

    def get_recommended_models_table(self):
        """
        Generate a table of recommended replacement models for this refresh wave.
        
        Returns:
            list: A list of dictionaries containing table data
        """
        model_summary = self.get_model_replacement_summary()
        table_data = []
        
        for model, info in model_summary.items():
            row = {
                'replacement_model': model,
                'family': info['family'],
                'count': info['count'],
                'original_models': ', '.join(info['original_models']),
                'hardware_cost': info['hardware_cost'],
                'license_cost': info['license_cost'],
                'unit_total': info['total_cost'],
                'total_cost': info['total_replacement_cost']
            }
            table_data.append(row)
        
        return table_data

class PredictiveLifecycleManager:
    """Class for predictive lifecycle management of Meraki devices."""
    def __init__(self, inventory_devices, eol_data=None, networks=None, price_catalog=None, license_costs=None):
        """Initialize with device inventory and EOL data."""
        self.devices = []
        
        # If EOL data wasn't provided, try to fetch it
        if not eol_data:
            self.eol_data, self.last_updated, self.is_from_doc = get_eol_info_from_doc()
            #print(f"{GREEN}Using EOL information from documentation (last updated: {self.last_updated}){RESET}" if self.is_from_doc else 
                  #f"{YELLOW}Using fallback EOL information{RESET}")
            
            # Check if MX100 is present in the fetched data
            if "MX100" in self.eol_data:
                #print(f"{GREEN}MX100 found in fetched EOL data: {self.eol_data['MX100']}{RESET}")
                pass
        else:
            self.eol_data = eol_data
            self.last_updated = None
            self.is_from_doc = False
            
            # Check if MX100 is present in the provided data
            if "MX100" in self.eol_data:
                #print(f"{GREEN}MX100 found in provided EOL data: {self.eol_data['MX100']}{RESET}")
                pass
        self.networks = networks or []
        self.price_catalog = price_catalog
        self.license_costs = license_costs or {}
        self.process_inventory(inventory_devices)
        
    def process_inventory(self, inventory_devices):
        """Process raw inventory data into Device objects."""
        for device_data in inventory_devices:
            device = Device(device_data, self.eol_data, self.price_catalog, self.license_costs)
            self.devices.append(device)
            
    def get_refresh_forecast(self, forecast_years=3, waves_per_year=4):
        """
        Generate a forecast of refresh waves over the specified time period.
        Only includes devices that actually need replacement.
        
        Args:
            forecast_years: Number of years to forecast
            waves_per_year: Number of refresh waves per year
            
        Returns:
            List of RefreshWave objects
        """
        # Create time periods for waves
        today = datetime.datetime.now().date()
        wave_periods = []
        
        # Create an inventory object per year
        for year in range(forecast_years):
            year_start = today.replace(month=1, day=1) + relativedelta(years=year)
            months_per_wave = 12 // waves_per_year
            
            for wave in range(waves_per_year):
                wave_start = year_start + relativedelta(months=wave * months_per_wave)
                wave_end = wave_start + relativedelta(months=months_per_wave) - relativedelta(days=1)
                
                wave_name = f"Wave {year+1}.{wave+1}"
                wave_periods.append(RefreshWave(wave_name, wave_start, wave_end))
        
        # Filter out devices that don't need replacement
        devices_needing_replacement = [device for device in self.devices 
                                      if device.get_replacement_model() is not None]
        
        # If no devices need replacement, return empty waves
        if not devices_needing_replacement:
            return wave_periods
        
        # Assign devices to waves based on lifecycle status
        assigned_devices = set()
        
        # First pass: assign devices with known EOL dates to appropriate waves
        for device in devices_needing_replacement:
            # Skip devices already assigned
            if device.serial in assigned_devices:
                continue
                
            # If device has EOL date and is approaching or past EOL
            if device.days_to_eol is not None and device.days_to_eol < 730:
                # Calculate target date (3 months before EOL)
                if device.days_to_eol <= 0:
                    # Already EOL, put in first wave
                    target_date = today
                else:
                    # Plan for 3 months before EOL
                    target_date = device.end_of_support_date - relativedelta(months=3)
                
                # Find appropriate wave
                for wave in wave_periods:
                    if wave.start_date <= target_date <= wave.end_date:
                        wave.add_device(device)
                        assigned_devices.add(device.serial)
                        break
                else:
                    # If no wave found, add to the closest one
                    closest_wave = min(wave_periods, 
                                      key=lambda w: abs((w.start_date + (w.end_date - w.start_date)/2) - target_date))
                    closest_wave.add_device(device)
                    assigned_devices.add(device.serial)
        
        # Second pass: group remaining devices using clustering
        unassigned_devices = [d for d in devices_needing_replacement if d.serial not in assigned_devices]
        
        if unassigned_devices:
            device_features = []
            for device in unassigned_devices:
                risk_score = device.get_risk_score()
                # Use device age as proxy for replacement timing if EOL data not available
                days_to_eol = device.days_to_eol if device.days_to_eol is not None else 1825  # Default to 5 years
                
                # Create a feature vector for clustering
                features = [
                    risk_score,
                    min(days_to_eol, 1825)  # Cap at 5 years
                ]
                device_features.append(features)
            
        # Return the waves
        return wave_periods
    
    def get_risk_distribution(self):
        """Get distribution of devices by risk category."""
        risk_counts = {"High": 0, "Medium": 0, "Low": 0}
        
        for device in self.devices:
            risk_category = device.get_risk_category()
            risk_counts[risk_category] += 1
            
        return risk_counts
    
    def get_lifecycle_distribution(self):
        """Get distribution of devices by lifecycle status."""
        status_counts = {
            "End of Support": 0,
            "Critical": 0,
            "Warning": 0,
            "Planning": 0,
            "Current": 0
        }
        
        for device in self.devices:
            status_counts[device.lifecycle_status] += 1
            
        return status_counts
    
    def get_budget_forecast(self, refresh_waves):
        """Generate a budget forecast based on refresh waves."""
        budget_forecast = []
        
        # Group waves by fiscal year (assuming calendar year for simplicity)
        fiscal_years = defaultdict(float)
        
        for wave in refresh_waves:
            # Use start date's year as the fiscal year
            fiscal_year = wave.start_date.year
            fiscal_years[fiscal_year] += wave.total_cost
        
        # Convert to list of (year, amount) tuples and sort by year
        for year, amount in sorted(fiscal_years.items()):
            budget_forecast.append({
                "year": year,
                "amount": amount,
                "formatted": f"${amount:,.2f}"
            })
            
        return budget_forecast
    
    def get_high_risk_devices(self, limit=10):
        network_assigned_devices = [device for device in self.devices if device.network_id]
        
        # Sort filtered devices by risk score (descending)
        sorted_devices = sorted(network_assigned_devices, key=lambda d: d.get_risk_score(), reverse=True)
        
        # Return the top devices up to the limit
        return sorted_devices[:limit]
        
    def get_network_refresh_summary(self):
        """Get a summary of refresh needs organized by network."""
        network_summary = {}
        
        # Create mapping of network IDs to names
        network_names = {net['id']: net.get('name', net['id']) for net in self.networks}
        
        # Group devices by network
        devices_by_network = defaultdict(list)
        for device in self.devices:
            if device.network_id:
                devices_by_network[device.network_id].append(device)
        
        # Calculate risk and replacement metrics for each network
        for network_id, devices in devices_by_network.items():
            # Skip networks with no devices
            if not devices:
                continue
                
            # Get network name
            network_name = network_names.get(network_id, f"Network {network_id}")
            
            # Calculate metrics
            total_devices = len(devices)
            high_risk = sum(1 for d in devices if d.get_risk_category() == "High")
            medium_risk = sum(1 for d in devices if d.get_risk_category() == "Medium")
            eol_devices = sum(1 for d in devices if d.lifecycle_status == "End of Support")
            approaching_eol = sum(1 for d in devices if d.lifecycle_status in ["Critical", "Warning"])
            
            # Calculate average risk score for the network
            avg_risk = sum(d.get_risk_score() for d in devices) / total_devices if total_devices > 0 else 0
            
            # Calculate replacement cost
            replacement_cost = sum(d.get_replacement_cost_estimate() for d in devices)
            
            # Store the summary
            network_summary[network_id] = {
                "name": network_name,
                "total_devices": total_devices,
                "high_risk": high_risk,
                "medium_risk": medium_risk,
                "eol_devices": eol_devices,
                "approaching_eol": approaching_eol,
                "avg_risk": avg_risk,
                "replacement_cost": replacement_cost,
                "critical_status": (high_risk > 0 or eol_devices > 0)
            }
        
        return network_summary

    def detect_new_models(self):
        """
        Compare inventory with known models to identify new models that
        may not be in our pricing database yet.
        """
        if not self.price_catalog:
            return []
            
        # Get all models we know about from pricing data
        known_models = []
        for family, models in self.price_catalog.prices.items():
            known_models.extend(models.keys())
        
        # Get base models (e.g., MX68 from MX68W)
        known_base_models = set()
        for model in known_models:
            base_match = re.match(r'([A-Z]+[0-9]+)', model)
            if base_match:
                known_base_models.add(base_match.group(1))
        
        # Find models in inventory that don't match known models
        new_models = set()
        for device in self.devices:
            model = device.model
            base_match = re.match(r'([A-Z]+[0-9]+)', model)
            
            # If we can't extract a base model, skip
            if not base_match:
                continue
                
            base_model = base_match.group(1)
            
            # If we don't know this base model or the exact model, it's new
            if base_model not in known_base_models and model not in known_models:
                new_models.add(model)
                
        return list(new_models)
        
    def get_models_by_family(self):
        """Group all devices by device family for reporting."""
        result = defaultdict(list)
        for device in self.devices:
            result[device.device_family].append(device.model)
            
        # Convert to dictionary with counts
        summary = {}
        for family, models in result.items():
            unique_models = set(models)
            summary[family] = {
                "count": len(models),
                "unique_models": len(unique_models),
                "models": sorted(list(unique_models))
            }
            
        return summary

def clean_slide(slide):
    """Remove everything from a slide except branding elements at the bottom."""
    # Identify shapes to remove (everything except bottom branding)
    shapes_to_remove = []
    for shape in slide.shapes:
        # Preserve shapes that appear to be logos/branding (bottom of slide)
        if shape.top >= Inches(6.5):
            continue
        # Keep bottom text
        if hasattr(shape, "text_frame") and shape.top >= Inches(6.5):
            continue
        # Remove everything else
        shapes_to_remove.append(shape)
    
    # Remove the identified shapes
    for shape in shapes_to_remove:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
        except Exception as e:
            print(f"{YELLOW}Could not remove shape: {e}{RESET}")

def create_risk_donut_chart(slide, risk_distribution, x, y, width, height, title):
    """Create a donut chart showing risk distribution."""
    # Set up chart data
    chart_data = CategoryChartData()
    
    # Add categories in specific order
    categories = ["High Risk", "Medium Risk", "Low Risk"]
    values = [risk_distribution.get("High", 0), 
              risk_distribution.get("Medium", 0), 
              risk_distribution.get("Low", 0)]
    
    chart_data.categories = categories
    chart_data.add_series('Risk', values)
    
    # Add title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.bold = True
    
    # Add chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, 
        x, 
        y + Inches(0.35), 
        width, 
        height - Inches(0.35),
        chart_data
    ).chart
    
    # Format chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    
    # Set colors for each slice
    for i, point in enumerate(chart.plots[0].series[0].points):
        if i == 0:  # High Risk
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = HIGH_RISK_COLOR
        elif i == 1:  # Medium Risk
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = MEDIUM_RISK_COLOR
        elif i == 2:  # Low Risk
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = LOW_RISK_COLOR
    
    # Add data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.CENTER
    data_labels.font.size = Pt(9)
    data_labels.font.bold = True
    
    # Add explanatory note below the chart
    note_y = y + height - Inches(0.7)
    note_box = slide.shapes.add_textbox(x, note_y, width, Inches(0.4))
    note_p = note_box.text_frame.add_paragraph()
    note_p.text = "Note: Risk assessment considers End of Support dates (primary factor) and End of Sale dates (secondary factor)."
    note_p.font.size = Pt(8)
    note_p.font.italic = True
    note_p.alignment = PP_ALIGN.LEFT
    
    return chart

def create_refresh_wave_timeline(slide, refresh_waves, x, y, width, height, title):
    """Create a timeline visualization of refresh waves with adjusted positioning to avoid title overlap."""
    # Add title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.bold = True
    
    timeline_y = y + Inches(1.2)
    timeline_height = Inches(0.05)
    wave_box_height = Inches(0.8)
    wave_spacing = Inches(0.3)
    
    # Create timeline line
    timeline = slide.shapes.add_shape(1, x, timeline_y, width, timeline_height)
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray
    timeline.line.fill.solid()
    timeline.line.fill.fore_color.rgb = RGBColor(170, 170, 170)  # Slightly darker gray
    
    # Only include the first 8 waves to fit on slide
    display_waves = refresh_waves[:8]
    
    # Calculate spacing between wave points
    point_spacing = width / (len(display_waves) + 1)
    
    # Create wave points and labels
    for i, wave in enumerate(display_waves):
        # X position for this wave
        wave_x = x + (i + 1) * point_spacing - Inches(0.15)
        
        # Calculate risk color
        if wave.risk_level == "High":
            risk_color = HIGH_RISK_COLOR
        elif wave.risk_level == "Medium":
            risk_color = MEDIUM_RISK_COLOR
        else:
            risk_color = LOW_RISK_COLOR
        
        # Create point on timeline
        point = slide.shapes.add_shape(1, wave_x, timeline_y - Inches(0.1), Inches(0.3), Inches(0.3))
        point.fill.solid()
        point.fill.fore_color.rgb = risk_color
        point.line.fill.solid()
        point.line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White border
        
        # Alternating top/bottom positioning
        if i % 2 == 0:
            # Top position
            box_y = timeline_y - wave_box_height - Inches(0.1)
            date_y = box_y - Inches(0.2)
        else:
            # Bottom position
            box_y = timeline_y + timeline_height + Inches(0.1)
            date_y = box_y + wave_box_height + Inches(0.05)
        
        # Create wave info box
        wave_box = slide.shapes.add_shape(1, wave_x - Inches(0.5), box_y, Inches(1.3), wave_box_height)
        wave_box.fill.solid()
        wave_box.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Very light gray
        wave_box.line.fill.solid()
        wave_box.line.fill.fore_color.rgb = risk_color
        
        # Add wave name
        name_box = slide.shapes.add_textbox(wave_x - Inches(0.5), box_y + Inches(0.05), Inches(1.3), Inches(0.25))
        name_p = name_box.text_frame.add_paragraph()
        name_p.text = wave.name
        name_p.font.size = Pt(10)
        name_p.font.bold = True
        name_p.alignment = PP_ALIGN.CENTER
        
        # Add device count
        count_box = slide.shapes.add_textbox(wave_x - Inches(0.5), box_y + Inches(0.25), Inches(1.3), Inches(0.25))
        count_p = count_box.text_frame.add_paragraph()
        count_p.text = f"{len(wave.devices)} devices"
        count_p.font.size = Pt(9)
        count_p.alignment = PP_ALIGN.CENTER
        
        # Add cost
        cost_box = slide.shapes.add_textbox(wave_x - Inches(0.5), box_y + Inches(0.45), Inches(1.3), Inches(0.25))
        cost_p = cost_box.text_frame.add_paragraph()
        cost_p.text = f"${wave.total_cost:,.0f}"
        cost_p.font.size = Pt(9)
        cost_p.font.color.rgb = risk_color
        cost_p.alignment = PP_ALIGN.CENTER
        
        # Add date label
        date_box = slide.shapes.add_textbox(wave_x - Inches(0.5), date_y, Inches(1.3), Inches(0.2))
        date_p = date_box.text_frame.add_paragraph()
        date_p.text = f"{wave.start_date.strftime('%b %Y')}"
        date_p.font.size = Pt(8)
        date_p.alignment = PP_ALIGN.CENTER
    
    return timeline
def create_enhanced_refresh_wave_table(slide, refresh_waves, x, y, width, height, title):

    
    # Add title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.bold = True
    
    # Table data
    rows = min(len(refresh_waves) + 1, 6)
    cols = 7  # Wave, Timeframe, Devices, Hardware Cost, License Cost, Total Cost, Risk Level
    
    # Create table
    table = slide.shapes.add_table(
        rows, 
        cols,
        x,
        y + Inches(0.35),
        width,
        height - Inches(0.35)
    ).table
    
    # Set column widths
    col_widths = [
        Inches(0.8),   # Wave
        Inches(1.7),   # Timeframe
        Inches(1.3),   # Devices
        Inches(1.5),   # Hardware Cost
        Inches(1.5),   # License Cost
        Inches(1.5),   # Total Cost
        Inches(1.0)    # Risk Level
    ]
    
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    # Set header row
    headers = ["Wave", "Timeframe", "Devices", "Hardware Cost", "License Cost", "Total Cost", "Risk Level"]
    
    # Format header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set header background
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(42, 59, 79)  # Dark blue
        
        # Set header text color
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Add wave data to table
    for i, wave in enumerate(refresh_waves[:rows-1]):
        row = i + 1
        
        # Apply alternating row colors
        for j in range(cols):
            cell = table.cell(row, j)
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
            else:
                cell.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Lighter gray
        
        # Wave name
        cell = table.cell(row, 0)
        cell.text = wave.name
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Timeframe
        cell = table.cell(row, 1)
        cell.text = f"{wave.start_date.strftime('%b %Y')} - {wave.end_date.strftime('%b %Y')}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Devices - count by family with recommended replacement models
        cell = table.cell(row, 2)
        device_families = count_device_families_with_models(wave.devices)
        device_text_parts = []
        
        for family, count in device_families.items():
            if count > 0:
                device_text_parts.append(f"{family}: {count}")
        
        if not device_text_parts:
            device_text = f"Total: {len(wave.devices)}"
        else:
            device_text = ", ".join(device_text_parts)
        
        # Set text
        cell.text = device_text
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Get hardware, license, and total costs
        hardware_cost = sum(device.get_replacement_cost_estimate() for device in wave.devices)
        license_cost = sum(device.get_license_cost_estimate() for device in wave.devices)
        total_cost = hardware_cost + license_cost
        
        # Hardware Cost
        cell = table.cell(row, 3)
        cell.text = f"${hardware_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # License Cost
        cell = table.cell(row, 4)
        cell.text = f"${license_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Total Cost
        cell = table.cell(row, 5)
        cell.text = f"${total_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Risk Level
        cell = table.cell(row, 6)
        cell.text = wave.risk_level
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Color based on risk level
        if wave.risk_level == "High":
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(229, 84, 81)  # Red
        elif wave.risk_level == "Medium":
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 192, 0)  # Amber
        elif wave.risk_level == "Low":
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(112, 173, 71)  # Green
    
    # Add explanatory text below the table
    explanation_box = slide.shapes.add_textbox(x, y + height - Inches(0.1), width, Inches(0.3))
    explanation_p = explanation_box.text_frame.add_paragraph()
    explanation_p.text = "Note: This table shows planned refresh waves with their recommended replacement devices, organized by timeframe."
    explanation_p.font.size = Pt(9)
    explanation_p.font.italic = True
    
    return table

def create_enhanced_refresh_wave_table(slide, refresh_waves, x, y, width, height, title):
    """Create an enhanced table of refresh waves with license cost information."""
    
    # Add title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.bold = True
    
    # Table data
    rows = min(len(refresh_waves) + 1, 6)
    cols = 7  # Wave, Timeframe, Devices, Hardware Cost, License Cost, Total Cost, Risk Level
    
    # Create table
    table = slide.shapes.add_table(
        rows, 
        cols,
        x,
        y + Inches(0.35),
        width,
        height - Inches(0.35)
    ).table
    
    # Set column widths
    col_widths = [
        Inches(0.8),   # Wave
        Inches(1.7),   # Timeframe
        Inches(1.3),   # Devices
        Inches(1.5),   # Hardware Cost
        Inches(1.5),   # License Cost
        Inches(1.5),   # Total Cost
        Inches(1.0)    # Risk Level
    ]
    
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    # Set header row
    headers = ["Wave", "Timeframe", "Devices", "Hardware Cost", "License Cost", "Total Cost", "Risk Level"]
    
    # Format header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set header background
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(42, 59, 79)  # Dark blue
        
        # Set header text color
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Add wave data to table
    for i, wave in enumerate(refresh_waves[:rows-1]):
        row = i + 1
        
        # Apply alternating row colors
        for j in range(cols):
            cell = table.cell(row, j)
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
            else:
                cell.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Lighter gray
        
        # Wave name
        cell = table.cell(row, 0)
        cell.text = wave.name
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Timeframe
        cell = table.cell(row, 1)
        cell.text = f"{wave.start_date.strftime('%b %Y')} - {wave.end_date.strftime('%b %Y')}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Devices - count by family with recommended replacement models
        cell = table.cell(row, 2)
        device_families = count_device_families_with_models(wave.devices)
        device_text_parts = []
        
        for family, count in device_families.items():
            if count > 0:
                device_text_parts.append(f"{family}: {count}")
        
        if not device_text_parts:
            device_text = f"Total: {len(wave.devices)}"
        else:
            device_text = ", ".join(device_text_parts)
        
        # Set text
        cell.text = device_text
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Get hardware, license, and total costs
        hardware_cost = sum(device.get_replacement_cost_estimate() for device in wave.devices)
        license_cost = sum(device.get_license_cost_estimate() for device in wave.devices)
        total_cost = hardware_cost + license_cost
        
        # Hardware Cost
        cell = table.cell(row, 3)
        cell.text = f"${hardware_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # License Cost
        cell = table.cell(row, 4)
        cell.text = f"${license_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Total Cost
        cell = table.cell(row, 5)
        cell.text = f"${total_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Risk Level
        cell = table.cell(row, 6)
        cell.text = wave.risk_level
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Color based on risk level
        if wave.risk_level == "High":
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(229, 84, 81)  # Red
        elif wave.risk_level == "Medium":
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 192, 0)  # Amber
        elif wave.risk_level == "Low":
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(112, 173, 71)  # Green

    return table

def count_device_families_with_models(devices):
    """Count devices by family and identify recommended models."""
    family_counts = defaultdict(int)
    for device in devices:
        family_counts[device.device_family] += 1
    return dict(family_counts)

def create_recommended_models_slide(prs, devices, price_catalog, using_fallback_pricing=False):
    
    # Find a suitable slide layout (preferably blank)
    slide_layout = None
    for layout in prs.slide_layouts:
        if hasattr(layout, 'name') and layout.name and 'blank' in layout.name.lower():
            slide_layout = layout
            break
    
    # If no blank layout found, use the standard "1_Title_and_Content" layout if available
    if not slide_layout:
        for layout in prs.slide_layouts:
            if hasattr(layout, 'name') and '1_Title_and_Content' in layout.name:
                slide_layout = layout
                break
    
    # If still no layout found, use the first one
    if not slide_layout and len(prs.slides) > 0:
        slide_layout = prs.slide_layouts[0]
    
    # Create the slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Clean the slide of any existing content
    if hasattr(slide, 'shapes'):
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.top < Inches(6.5):
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                if hasattr(shape, '_sp'):
                    sp = shape._sp
                    sp.getparent().remove(sp)
            except:
                pass
    
    # Add title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(9), Inches(0.5))
    title_p = title.text_frame.add_paragraph()
    title_p.text = "Recommended Replacement Models"
    title_p.font.size = Pt(24)
    title_p.font.bold = True
    
    # Add horizontal line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(0.7), Inches(9.5), Inches(0.7))
    line.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
    line.line.width = Pt(1.5)
    
    # Create a table with recommended models
    replacement_models = {}
    
    # Group devices by their recommended replacement model
    for device in devices:
        # Use the public method to get replacement, which returns None for current-gen models
        replacement = device.get_replacement_model()
        
        # Skip current-gen models that don't need replacement
        if replacement is None:
            continue
            
        if replacement not in replacement_models:
            replacement_models[replacement] = {
                'count': 0,
                'original_models': set(),
                'hardware_cost': device.get_replacement_cost_estimate(),
                'license_cost': device.get_license_cost_estimate(),
                'total_cost': device.get_total_replacement_cost(),
                'family': device.device_family
            }
        
        replacement_models[replacement]['count'] += 1
        replacement_models[replacement]['original_models'].add(device.model)
    
    # Convert sets to sorted lists for easier display
    for model in replacement_models:
        replacement_models[model]['original_models'] = sorted(list(replacement_models[model]['original_models']))
    
    # Convert to list and sort by family and then count
    model_list = []
    for model, info in replacement_models.items():
        # Skip models with zero devices
        if info['count'] == 0:
            continue
            
        model_list.append({
            'model': model,
            'count': info['count'],
            'family': info['family'],
            'original_models': info['original_models'],
            'hardware_cost': info['hardware_cost'],
            'license_cost': info['license_cost'],
            'total_unit_cost': info['hardware_cost'] + info['license_cost'],
            'total_model_cost': (info['hardware_cost'] + info['license_cost']) * info['count']
        })
    
    # Sort by family and then by count descending
    model_list = sorted(model_list, key=lambda x: (x['family'], -x['count']))
    
    # Create table header
    table_x = Inches(0.5)
    table_y = Inches(0.89)
    table_width = Inches(9.5)
    table_height = Inches(5.3)
    
    # Determine number of rows (header + model rows + total row)
    rows = min(len(model_list) + 2, 21)
    
    table = slide.shapes.add_table(
        rows, 
        8,  # Family, Replacement Model, Original Models, Quantity, Hardware Cost, License Cost, Hardware+License Cost, Total Cost
        table_x, 
        table_y, 
        table_width, 
        table_height
    ).table
    
    # Set column widths - adjusted to accommodate the new column
    table.columns[0].width = Inches(0.7)   # Family
    table.columns[1].width = Inches(1.5)   # Replacement Model
    table.columns[2].width = Inches(2.0)   # Original Models
    table.columns[3].width = Inches(0.7)   # Quantity
    table.columns[4].width = Inches(1.1)   # Hardware Cost
    table.columns[5].width = Inches(1.0)   # License Cost
    table.columns[6].width = Inches(1.2)   # Hardware+License Cost (Per Unit)
    table.columns[7].width = Inches(1.3)   # Total Cost (Quantity * Unit Cost)
    
    # Set header row
    headers = ["Family", "Replacement Model", "Original Models", "Quantity", "Hardware Cost", 
              "License Cost", "Hardware+License Cost", "Total Cost"]
    
    # Format header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set header background
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(42, 59, 79)  # Dark blue
        
        # Set header text color
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Define family display names mapping
    family_display_names = {
        'MX': 'Security',
        'MS': 'Switches',
        'MR': 'Wireless',
        'MV': 'Cameras',
        'Z': 'Teleworker',
        'MG': 'Cellular',
        'MT': 'IoT'
    }
    
    # Track grand total
    grand_total = 0
    
    # Add model data to table
    for i, model_info in enumerate(model_list[:rows-2]):  # Leave space for the total row
        row = i + 1
        
        # Apply alternating row colors
        for j in range(8):
            cell = table.cell(row, j)
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
            else:
                cell.fill.fore_color.rgb = RGBColor(248, 248, 248)  # Lighter gray
        
        # Family
        cell = table.cell(row, 0)
        family = model_info['family']
        # Use the display name if available
        display_family = family_display_names.get(family, family)
        cell.text = display_family
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set color based on family
        family_colors = {
            'MX': RGBColor(0, 112, 192),    # Blue
            'MS': RGBColor(112, 173, 71),   # Green
            'MR': RGBColor(255, 192, 0),    # Yellow/Amber
            'MV': RGBColor(192, 0, 0),      # Red
            'Z': RGBColor(128, 0, 128),     # Purple
            'MG': RGBColor(255, 102, 0),    # Orange
            'MT': RGBColor(0, 176, 240)     # Light Blue
        }
        
        if family in family_colors:
            cell.text_frame.paragraphs[0].font.color.rgb = family_colors[family]
        
        # Replacement Model
        cell = table.cell(row, 1)
        cell.text = model_info['model']
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        
        # Original Models
        cell = table.cell(row, 2)
        cell.text = ", ".join(model_info['original_models'])
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Quantity
        cell = table.cell(row, 3)
        cell.text = str(model_info['count'])
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Hardware Cost
        cell = table.cell(row, 4)
        cell.text = f"${model_info['hardware_cost']:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # License Cost (1YR)
        cell = table.cell(row, 5)
        cell.text = f"${model_info['license_cost']:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Hardware+License Cost Per Unit (renamed from Total Cost)
        cell = table.cell(row, 6)
        cell.text = f"${model_info['total_unit_cost']:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # NEW: Total Cost (Quantity * Unit Cost)
        cell = table.cell(row, 7)
        total_cost = model_info['total_model_cost']
        cell.text = f"${total_cost:,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add to grand total
        grand_total += total_cost
    
    # Add total row at the bottom
    total_row = rows - 1
    
    # Format the total row
    for j in range(8):
        cell = table.cell(total_row, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray for total row
        
        # Set all cells blank by default
        cell.text = ""
        
    # Add "TOTAL" text in first column
    total_label_cell = table.cell(total_row, 0)
    total_label_cell.text = "TOTAL"
    total_label_cell.text_frame.paragraphs[0].font.size = Pt(9)
    total_label_cell.text_frame.paragraphs[0].font.bold = True
    total_label_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Merge cells for total label (columns 0-6)
    table.cell(total_row, 0).merge(table.cell(total_row, 6))
    
    # Add grand total in the last column
    total_amount_cell = table.cell(total_row, 7)
    total_amount_cell.text = f"${grand_total:,.2f}"
    total_amount_cell.text_frame.paragraphs[0].font.size = Pt(10)
    total_amount_cell.text_frame.paragraphs[0].font.bold = True
    total_amount_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    subtitle_box = slide.shapes.add_textbox(Inches(1.91), Inches(6.61), Inches(9.5), Inches(0.3))
    subtitle_p = subtitle_box.text_frame.add_paragraph()
    subtitle_p.text = "Showing recommended upgrades to current-generation or next-generation models. *These prices are estimates* Consult with your Cisco account team for accurate pricing."
    subtitle_p.font.size = Pt(10)
    subtitle_p.font.italic = True
    subtitle_p.alignment = PP_ALIGN.LEFT
    
    footnote_box = slide.shapes.add_textbox(Inches(1.91), Inches(6.81), Inches(9.5), Inches(0.3))
    footnote_p = footnote_box.text_frame.add_paragraph()
    footnote_p.text = "Note: License costs shown are for 1-year Enterprise licenses. Multi-year licenses offer savings of approximately 10-15% per year."
    footnote_p.font.size = Pt(10)
    footnote_p.font.italic = True
    footnote_p.alignment = PP_ALIGN.LEFT
    
    if using_fallback_pricing:
        add_pricing_disclaimer_to_slide(slide, using_fallback_pricing)
    
    return slide

def fix_refresh_slide_details_title_aggressive(prs, output_path=None):

    # First, find which slide contains the refresh planning timeline
    target_slide = None
    slide_index = -1
    
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                try:
                    if "Hardware Refresh Planning Timeline" in shape.text_frame.text:
                        target_slide = slide
                        slide_index = i
                        #print(f"Found refresh planning slide at index {i+1}")
                        break
                except:
                    pass
        if target_slide:
            break
    
    if not target_slide:
        print("Could not find refresh planning slide")
        return
    
    # Amount to move the timeline elements down by
    vertical_shift = Inches(0.4)
    
    # First, remove the unwanted text about timeline
    shapes_to_remove = []
    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "The timeline above shows predicted refresh waves" in shape.text_frame.text:
                    #print(f"Found unwanted text at position {shape.top/914400} inches")
                    shapes_to_remove.append(shape)
            except:
                pass
    
    for shape in shapes_to_remove:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
                #print("Removed unwanted text shape")
        except Exception as e:
            print(f"Error removing shape: {e}")

    main_titles = []
    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "Hardware Refresh Planning Timeline" in shape.text_frame.text:
                    #print(f"Found main title at position {shape.top/914400} inches")
                    main_titles.append(shape)
            except:
                pass

    new_main_title = target_slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(9), Inches(0.5))
    new_main_title_p = new_main_title.text_frame.add_paragraph()
    new_main_title_p.text = "Hardware Refresh Planning Timeline"
    new_main_title_p.font.size = Pt(24)
    new_main_title_p.font.bold = True
    #print(f"Created new main title at position 0 inches")
    
    # Hide original main titles
    for title in main_titles:
        try:
            title.left = Inches(-20) # Move off-slide
            #print(f"Hid original main title")
        except Exception as e:
            print(f"Error hiding main title: {e}")
    
    timeline_titles = []
    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "36-Month Refresh Timeline" in shape.text_frame.text:
                    #print(f"Found timeline title at position {shape.top/914400} inches")
                    timeline_titles.append(shape)
            except:
                pass
    
    new_timeline_title = target_slide.shapes.add_textbox(Inches(0.5), Inches(0.61), Inches(9), Inches(0.3))
    new_timeline_title_p = new_timeline_title.text_frame.add_paragraph()
    new_timeline_title_p.text = "36-Month Refresh Timeline"
    new_timeline_title_p.font.size = Pt(14)
    new_timeline_title_p.font.bold = True
    #print(f"Created new timeline title at position 0.61 inches")

    for title in timeline_titles:
        try:
            title.left = Inches(-20)
            #print(f"Hid original timeline title")
        except Exception as e:
            print(f"Error hiding timeline title: {e}")

    timeline_y_range = (0.8, 3.5)
    
    # Find timeline elements based on position and type
    timeline_elements = []
    for shape in target_slide.shapes:
        # Convert shape position to inches for easier comparison
        shape_top_inches = shape.top / 914400  # 914400 EMUs per inch
        
        # Check if element is in the timeline vertical range
        if timeline_y_range[0] <= shape_top_inches <= timeline_y_range[1]:
            # Exclude the titles we've already handled
            if shape not in main_titles and shape not in timeline_titles:
                if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                    try:
                        if "Refresh Wave Details" in shape.text_frame.text:
                            continue
                    except:
                        pass

                timeline_elements.append(shape)
                #print(f"Found timeline element at position {shape_top_inches} inches")
    
    # Move each timeline element down
    for element in timeline_elements:
        try:
            element.top += vertical_shift
            #print(f"Moved timeline element down by {vertical_shift/914400} inches")
        except Exception as e:
            print(f"Error moving timeline element: {e}")

    details_titles = []
    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "Refresh Wave Details" in shape.text_frame.text:
                    #print(f"Found details title at position {shape.top/914400} inches")
                    details_titles.append(shape)
            except:
                pass
    
    # Create a new details title at the exact position we want
    new_details_title = target_slide.shapes.add_textbox(Inches(0.5), Inches(3.95), Inches(9.3), Inches(0.3))
    new_details_title_p = new_details_title.text_frame.add_paragraph()
    new_details_title_p.text = "Refresh Wave Details"
    new_details_title_p.font.size = Pt(14)
    new_details_title_p.font.bold = True
    #print(f"Created new details title at position 4.34 inches")
    
    # Hide the original details titles
    for title in details_titles:
        try:
            title.left = Inches(-20)
            #print(f"Hid original details title")
        except Exception as e:
            print(f"Error hiding details title: {e}")
    
    # Save the presentation if path is provided
    if output_path:
        try:
            prs.save(output_path)
            #print(f"Saved presentation to {output_path}")
        except Exception as e:
            print(f"Error saving presentation: {e}")
def fix_predictive_lifecycle_slide_positions(prs, output_path=None):

    target_slide = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                try:
                    if "Predictive Lifecycle Management" in shape.text_frame.text:
                        target_slide = slide
                        #print(f"Found Predictive Lifecycle Management slide at index {i+1}")
                        break
                except:
                    pass
        if target_slide:
            break
    
    if not target_slide:
        #print("Could not find Predictive Lifecycle Management slide")
        return False
    
    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "Predictive Lifecycle Management" in shape.text_frame.text:
                    # Move the title to vertical position 0"
                    shape.top = Inches(0)
                    #print("Moved main title to vertical position 0\"")
            except:
                pass

    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "Device Risk Assessment" in shape.text_frame.text:
                    # Move to left 1.07", top 0.64"
                    shape.left = Inches(1.07)
                    shape.top = Inches(0.64)
                    #print("Moved Device Risk Assessment title to left 1.07\", top 0.64\"")
            except:
                pass

    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "Forecasted Replacement Budget" in shape.text_frame.text:
                    shape.left = Inches(7.82)
                    shape.top = Inches(0.64)
                    #print("Moved Forecasted Replacement Budget title to left 6.8\", top 0.64\"")
            except:
                pass

    title_moved = False
    for shape in target_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                if "Prioritized Device Replacement List" in shape.text_frame.text:
                    shape.left = Inches(0.5)
                    shape.top = Inches(3.53)
                    title_moved = True
            except:
                pass

    table_moved = False
    
    if not table_moved:
        #print("Warning: Could not find the table to move")
        pass
    
    if not title_moved:
        #print("Warning: Could not find the title to move")
        pass
    
    # Save the presentation if path is provided
    if output_path:
        try:
            prs.save(output_path)
            #print(f"Saved presentation with adjusted positions to {output_path}")
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False
    
    return True

def fix_timeline_dollar_positions(prs, output_path=None, vertical_adjustment=0.07):
    
    # Find the timeline slide
    timeline_slide = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                try:
                    if "Hardware Refresh Planning Timeline" in shape.text_frame.text:
                        timeline_slide = slide
                        #print(f"Found refresh planning slide at index {i+1}")
                        break
                except:
                    pass
        if timeline_slide:
            break
    
    if not timeline_slide:
        #print(f"{RED}Could not find the Hardware Refresh Planning Timeline slide{RESET}")
        return 0
    

    wave_box_y_range = (1.5, 3.5)
    
    dollar_textboxes = []
    for shape in timeline_slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            try:
                text = shape.text_frame.text.strip()
                if text.startswith('$'):
                    # Convert position to inches
                    shape_top_inches = shape.top / 914400
                    
                    # Check if it's in the wave box vertical range
                    if wave_box_y_range[0] <= shape_top_inches <= wave_box_y_range[1]:
                        dollar_textboxes.append(shape)
                        #print(f"{GREEN}Found dollar textbox at position {shape_top_inches:.2f} inches with text '{text}'{RESET}")
            except:
                pass
    
    # Move each dollar textbox up
    for textbox in dollar_textboxes:
        try:
            textbox.top -= Inches(vertical_adjustment)
            # Calculate new position in inches for reporting
            new_position = textbox.top / 914400
            #print(f"{GREEN}Moved dollar textbox up to position {new_position:.2f} inches{RESET}")
        except Exception as e:
            print(f"{RED}Error moving dollar textbox: {e}{RESET}")
    
    return len(dollar_textboxes)

async def generate(api_client, template_path, output_path, inventory_devices=None, networks=None, eol_data=None):
    """Generate Predictive Lifecycle Management slides and add them at the end of the presentation."""
    print(f"\n{GREEN}Generating Predictive Lifecycle Management slides...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # Check if we have inventory data
    if not inventory_devices:
        print(f"{RED}No inventory data provided, cannot generate lifecycle predictions{RESET}")
        return
    if eol_data:
        #print(f"{GREEN}EOL data provided to predictive_lifecycle.generate() with {len(eol_data)} entries{RESET}")
        if "MX100" in eol_data:
            #print(f"{GREEN}MX100 data in EOL input: {eol_data['MX100']}{RESET}")
            pass
    else:
        print(f"{YELLOW}No EOL data provided to generate(), will fetch from get_eol_info_from_doc(){RESET}")
    
    # Initialize the price catalog with Rhino Networks scraping support
    try:
        # Create price catalog with Rhino Networks scraping
        price_catalog = RhinoPriceCatalog(cache_file="meraki_rhino_prices_cache.json")
        
        # First try getting exact prices from the HTML
        html_prices_added = price_catalog.add_exact_prices_from_html()
        if html_prices_added > 0:
            #print(f"{GREEN}Added {html_prices_added} exact prices from known HTML{RESET}")
            pass
        
        # Then try scraping the wireless page directly
        wireless_prices_added = price_catalog.extract_prices_from_wireless_page()
        if wireless_prices_added > 0:
            #print(f"{GREEN}Extracted {wireless_prices_added} prices from wireless page{RESET}")
            pass
        
        # Scrape license products specifically
        license_prices_added = price_catalog.scrape_license_products()
        
        # Add VMX license prices
        vmx_prices_added = price_catalog.add_vmx_license_prices()
        if vmx_prices_added > 0:
            #print(f"{GREEN}Added {vmx_prices_added} VMX license prices{RESET}")
            pass
            
        # Discover actual product URLs to correct URL patterns
        discovered_urls = price_catalog.discover_product_urls()
        
        # Scrape device license costs for 1-year licenses
        license_costs = price_catalog.scrape_device_license_costs()
        #print(f"{GREEN}Scraped {len(license_costs)} 1-year license costs for Meraki devices{RESET}")
        
        # Add standard license prices for common license types
        price_catalog.add_standard_license_prices()
        
        # Check if we're using fallback pricing
        using_fallback_pricing = price_catalog.is_using_fallback_pricing()
        if using_fallback_pricing:
            print(f"{YELLOW}Using fallback price estimates. Pricing information may not be accurate.{RESET}")
        else:
            model_count = sum(len(models) for models in price_catalog.prices.values())
            #print(f"{GREEN}Successfully loaded pricing for {model_count} Meraki device models{RESET}")
            
    except Exception as e:
        print(f"{YELLOW}Warning: Error initializing price catalog: {e}{RESET}")
        print(f"{YELLOW}Using default price estimates instead{RESET}")
        price_catalog = None
        using_fallback_pricing = True
        license_costs = {}

    # Initialize the lifecycle manager with license costs
    lifecycle_manager = PredictiveLifecycleManager(inventory_devices, eol_data, networks, price_catalog, license_costs)
    
    # Check for new models that aren't in our price database
    if price_catalog:
        new_models = lifecycle_manager.detect_new_models()
        if new_models:
            #print(f"{YELLOW}Detected {len(new_models)} device models not in pricing database:{RESET}")
            pass
            for model in new_models[:10]:
                #print(f"{YELLOW}  - {model}{RESET}")
                pass
            if len(new_models) > 10:
                #print(f"{YELLOW}  - ... and {len(new_models) - 10} more{RESET}")
                pass
            #print(f"{YELLOW}Using estimated prices for these models{RESET}")
    
    # Generate refresh waves
    refresh_waves = lifecycle_manager.get_refresh_forecast(forecast_years=3, waves_per_year=4)
    
    # Get risk distribution
    risk_distribution = lifecycle_manager.get_risk_distribution()
    
    # Get lifecycle distribution
    lifecycle_distribution = lifecycle_manager.get_lifecycle_distribution()
    
    # Get budget forecast
    budget_forecast = lifecycle_manager.get_budget_forecast(refresh_waves)
    
    # Get high risk devices
    high_risk_devices = lifecycle_manager.get_high_risk_devices(limit=18)
    
    # Get network refresh summary
    network_summary = lifecycle_manager.get_network_refresh_summary()
    
    # Get device inventory summary by family
    inventory_summary = lifecycle_manager.get_models_by_family()
    
    # Create the slides
    try:
        # Load the presentation
        prs = Presentation(output_path)
        
        # Find a suitable slide layout (preferably blank)
        slide_layout = None
        for layout in prs.slide_layouts:
            if hasattr(layout, 'name') and layout.name and 'blank' in layout.name.lower():
                slide_layout = layout
                break
        
        # If no blank layout found, use the standard "1_Title_and_Content" layout if available
        if not slide_layout:
            for layout in prs.slide_layouts:
                if hasattr(layout, 'name') and '1_Title_and_Content' in layout.name:
                    slide_layout = layout
                    break
                    
        # If still no layout found, use the first one
        if not slide_layout and len(prs.slides) > 0:
            slide_layout = prs.slide_layouts[0]
        
        # If we have a layout, create the slides
        if slide_layout:
            # SLIDE 1: Overview and Risk Assessment with side-by-side tables
            print(f"{BLUE}Creating Lifecycle Management Overview slide...{RESET}")
            
            # Create the slide with dual high risk tables
            slide = prs.slides.add_slide(slide_layout)
            
            # Clean the slide of any existing content
            clean_slide(slide)
            
            # Add title
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.0), Inches(9), Inches(0.5))
            title_p = title.text_frame.add_paragraph()
            title_p.text = "Predictive Lifecycle Management"
            title_p.font.size = Pt(24)
            title_p.font.bold = True
            
            # Add horizontal line
            line = slide.shapes.add_connector(1, Inches(0.5), Inches(0.75), Inches(9.5), Inches(0.75))
            line.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
            line.line.width = Pt(1.5)
            
            # Risk distribution donut chart
            create_risk_donut_chart(
                slide, 
                risk_distribution, 
                Inches(0.7), 
                Inches(0.8), 
                Inches(3.0), 
                Inches(2.8),
                "Device Risk Assessment"
            )

            budget_title = slide.shapes.add_textbox(Inches(7.76), Inches(0.64), Inches(3.0), Inches(0.3))
            budget_title_p = budget_title.text_frame.add_paragraph()
            budget_title_p.text = "Forecasted Replacement Budget"
            budget_title_p.font.size = Pt(14)
            budget_title_p.font.bold = True
            
            # Create budget chart without title
            chart_data = CategoryChartData()
                        
            # Add years and amounts
            years = [str(item["year"]) for item in budget_forecast]
            amounts = [item["amount"] for item in budget_forecast]
                        
            chart_data.categories = years
            chart_data.add_series('',amounts)
                        
            # Add chart at specified position
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, 
                Inches(6.5), 
                Inches(1.08), 
                Inches(4.79), 
                Inches(2.5),
                chart_data
            ).chart
                        
            # Format chart
            chart.has_legend = False
                        
            # Format series
            series = chart.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = ACCENT_COLOR
                        
            # Add data labels
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.number_format = '$#,##0'
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
            data_labels.font.size = Pt(9)
            
            # First high risk device table title with specific position
            top9_title = slide.shapes.add_textbox(Inches(0.17), Inches(3.28), Inches(6.39), Inches(0.3))
            top9_title_p = top9_title.text_frame.add_paragraph()
            top9_title_p.text = "Priority Device List (Top 9)"
            top9_title_p.font.size = Pt(14)
            top9_title_p.font.bold = True

            # Create the first high risk device table at the exact position and dimensions
            device_slice1 = high_risk_devices[:9]
            rows1 = len(device_slice1) + 1  # Header + devices

            # First table with specified position and dimensions
            table1 = slide.shapes.add_table(
                rows1, 
                5,  # 5 columns
                Inches(0.17),
                Inches(3.86),
                Inches(6.39),
                Inches(2.86)
            ).table

            # Configure first table columns and data
            table1.columns[0].width = Inches(1.8)    # Model
            table1.columns[1].width = Inches(1.5)    # Serial
            table1.columns[2].width = Inches(0.9)    # Risk Score
            table1.columns[3].width = Inches(1.2)    # Lifecycle Status
            table1.columns[4].width = Inches(0.99)   # Replacement Cost

            # Set header row for first table
            headers = ["Model", "Serial", "Risk Score", "Lifecycle Status", "Est. Cost"]
            for i, header in enumerate(headers):
                cell = table1.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Set modern header background - Cisco blue
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 120, 206)  # Cisco blue #0078CE
                
                # Set header text color to white for better contrast
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White

            # Add device data to first table
            for i, device in enumerate(device_slice1):
                row = i + 1
                
                # Apply alternating row colors for a modern look
                for j in range(5):
                    cell = table1.cell(row, j)
                    cell.fill.solid()
                    if i % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(245, 247, 250)  # White
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Very light blue-gray #F5F7FA
                
                # Model
                cell = table1.cell(row, 0)
                cell.text = device.model
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                
                # Serial
                cell = table1.cell(row, 1)
                serial_text = device.serial[:12] + '...' if len(device.serial) > 14 else device.serial
                cell.text = serial_text
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                
                # Risk Score with enhanced visualization
                cell = table1.cell(row, 2)
                risk_score = device.get_risk_score()
                
                # Add risk score with visual indicator
                p = cell.text_frame.paragraphs[0]
                p.text = f"{risk_score}"
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
                
                # Use colored background for risk cells
                if risk_score >= 70:
                    # Add red background for high risk
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Very light red background
                    p.font.color.rgb = HIGH_RISK_COLOR
                    p.font.bold = True
                elif risk_score >= 40:
                    # Add yellow background for medium risk
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 245, 225)  # Very light yellow background
                    p.font.color.rgb = MEDIUM_RISK_COLOR
                    p.font.bold = True
                
                # Lifecycle Status with enhanced visualization
                cell = table1.cell(row, 3)
                cell.text = device.lifecycle_status
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Set color and background based on lifecycle status
                if device.lifecycle_status == "End of Support":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 225, 225)  # Light red background
                    cell.text_frame.paragraphs[0].font.color.rgb = HIGH_RISK_COLOR
                    cell.text_frame.paragraphs[0].font.bold = True
                elif device.lifecycle_status == "Critical":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 235, 225)  # Light orange background
                    cell.text_frame.paragraphs[0].font.color.rgb = HIGH_RISK_COLOR
                elif device.lifecycle_status == "Warning":
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 245, 225)  # Light yellow background
                    cell.text_frame.paragraphs[0].font.color.rgb = MEDIUM_RISK_COLOR
                
                # Replacement Cost
                cell = table1.cell(row, 4)
                cell.text = f"${device.get_replacement_cost_estimate():,.2f}"
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Second high risk device table title
            if len(high_risk_devices) > 9:
                next9_title = slide.shapes.add_textbox(Inches(6.77), Inches(3.28), Inches(6.8), Inches(0.3))
                next9_title_p = next9_title.text_frame.add_paragraph()
                next9_title_p.text = "Priority Device List (Next 9)"
                next9_title_p.font.size = Pt(14)
                next9_title_p.font.bold = True
                
                # Create the second high risk device table at the exact position and dimensions
                device_slice2 = high_risk_devices[9:18]
                rows2 = len(device_slice2) + 1  # Header + devices
                
                # Second table with specified position and dimensions
                table2 = slide.shapes.add_table(
                    rows2, 
                    5,  # 5 columns
                    Inches(6.77),
                    Inches(3.86),
                    Inches(6.8),
                    Inches(2.88)
                ).table
                
                # Configure second table columns and data
                table2.columns[0].width = Inches(1.8)    # Model
                table2.columns[1].width = Inches(1.5)    # Serial
                table2.columns[2].width = Inches(0.9)    # Risk Score
                table2.columns[3].width = Inches(1.2)    # Lifecycle Status
                table2.columns[4].width = Inches(0.99)   # Replacement Cost
                
                # Set header row for second table
                for i, header in enumerate(headers):
                    cell = table2.cell(0, i)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Set modern header background - Cisco blue
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 120, 206)  # Cisco blue #0078CE
                    
                    # Set header text color to white for better contrast
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
                
                # Add device data to second table
                for i, device in enumerate(device_slice2):
                    row = i + 1
                    
                    # Apply alternating row colors for a modern look
                    for j in range(5):
                        cell = table2.cell(row, j)
                        cell.fill.solid()
                        if i % 2 == 0:
                            cell.fill.fore_color.rgb = RGBColor(245, 247, 250)  # White
                        else:
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Very light blue-gray #F5F7FA
                    
                    # Model
                    cell = table2.cell(row, 0)
                    cell.text = device.model
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    
                    # Serial
                    cell = table2.cell(row, 1)
                    serial_text = device.serial[:12] + '...' if len(device.serial) > 14 else device.serial
                    cell.text = serial_text
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    
                    # Risk Score with enhanced visualization
                    cell = table2.cell(row, 2)
                    risk_score = device.get_risk_score()
                    
                    # Add risk score with visual indicator
                    p = cell.text_frame.paragraphs[0]
                    p.text = f"{risk_score}"
                    p.font.size = Pt(9)
                    p.alignment = PP_ALIGN.CENTER
                    
                    # Use colored background for risk cells
                    if risk_score >= 70:
                        # Add red background for high risk
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Very light red background
                        p.font.color.rgb = HIGH_RISK_COLOR
                        p.font.bold = True
                    elif risk_score >= 40:
                        # Add yellow background for medium risk
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 245, 225)  # Very light yellow background
                        p.font.color.rgb = MEDIUM_RISK_COLOR
                        p.font.bold = True
                    
                    # Lifecycle Status with enhanced visualization
                    cell = table2.cell(row, 3)
                    cell.text = device.lifecycle_status
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Set color and background based on lifecycle status
                    if device.lifecycle_status == "End of Support":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 225, 225)  # Light red background
                        cell.text_frame.paragraphs[0].font.color.rgb = HIGH_RISK_COLOR
                        cell.text_frame.paragraphs[0].font.bold = True
                    elif device.lifecycle_status == "Critical":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 235, 225)  # Light orange background
                        cell.text_frame.paragraphs[0].font.color.rgb = HIGH_RISK_COLOR
                    elif device.lifecycle_status == "Warning":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 245, 225)  # Light yellow background
                        cell.text_frame.paragraphs[0].font.color.rgb = MEDIUM_RISK_COLOR
                    
                    # Replacement Cost
                    cell = table2.cell(row, 4)
                    cell.text = f"${device.get_replacement_cost_estimate():,.2f}"
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Add pricing disclaimer if using fallback pricing
            add_pricing_disclaimer_to_slide(slide, using_fallback_pricing)
            
            # SLIDE 2: Refresh Wave Timeline
            print(f"{BLUE}Creating Refresh Planning Timeline slide...{RESET}")
            
            # Simply add to the end - no positioning needed
            slide2 = prs.slides.add_slide(slide_layout)
            
            # Clean the slide
            clean_slide(slide2)
            
            # Add title
            title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.30), Inches(9), Inches(0.5))
            title_p = title.text_frame.add_paragraph()
            title_p.text = "Hardware Refresh Planning Timeline"
            title_p.font.size = Pt(24)
            title_p.font.bold = True
            
            # Add horizontal line
            line = slide2.shapes.add_connector(1, Inches(0.5), Inches(0.75), Inches(9.5), Inches(0.75))
            line.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
            line.line.width = Pt(1.5)
            
            # Add refresh wave timeline
            create_refresh_wave_timeline(
                slide2,
                refresh_waves,
                Inches(0.5),
                Inches(0.9),
                Inches(9.3),
                Inches(3.0),
                "36-Month Refresh Timeline"
            )
            
            # Replace old table with enhanced table that includes license costs
            create_enhanced_refresh_wave_table(
                slide2,
                refresh_waves,
                Inches(0.5),
                Inches(4.2),
                Inches(9.3),
                Inches(1.8),
                "Refresh Wave Details"
            )

            # Add explanation note with explicit positioning
            note_box = slide2.shapes.add_textbox(Inches(0.5), Inches(6.07), Inches(9.3), Inches(0.3))
            note_p = note_box.text_frame.add_paragraph()
            note_p.text = "Note: This table shows planned refresh waves with their recommended replacement devices, organized by timeframe."
            note_p.font.size = Pt(9)
            note_p.font.italic = True

            # Add methodology notes with exact positioning
            methodology_box = slide2.shapes.add_textbox(Inches(0.5), Inches(6.29), Inches(9.3), Inches(0.3))
            methodology_p = methodology_box.text_frame.add_paragraph()

            # Include information about price sources
            methodology_text = "Methodology: Refresh waves are calculated using end-of-support dates, device risk scores, "
            if using_fallback_pricing:
                methodology_text += "estimated device pricing, "
            else:
                methodology_text += "current Meraki pricing loaded from a partner, "
            methodology_text += "clustering analysis of similar devices, and budget optimization algorithms."
            methodology_text += "\nLicense costs are based on 1-year Enterprise license prices. *These are just estimates*. Cisco account team should be engaged for recommendations and quotes."
            methodology_p.text = methodology_text
            methodology_p.font.size = Pt(8)
            methodology_p.font.italic = True
            
            # Add pricing disclaimer if using fallback pricing
            add_pricing_disclaimer_to_slide(slide2, using_fallback_pricing)
            
            # SLIDE 3: Recommended Replacement Models - NEW SLIDE
            print(f"{BLUE}Creating Recommended Replacement Models slide...{RESET}")
            
            # Create a new slide to show recommended models
            create_recommended_models_slide(prs, lifecycle_manager.devices, price_catalog, using_fallback_pricing)

            fix_refresh_slide_details_title_aggressive(prs, output_path)
            adjusted_count = fix_timeline_dollar_positions(prs)
            if adjusted_count > 0:
                #print(f"{GREEN}Successfully moved {adjusted_count} dollar amounts up by 0.07 inches{RESET}")
                pass
            # Save the presentation
            fix_predictive_lifecycle_slide_positions(prs, output_path=None)
            prs.save(output_path)
            print(f"{GREEN}Added Predictive Lifecycle Management slides to the end of the presentation{RESET}")
        else:
            print(f"{RED}No suitable slide layout found in the presentation{RESET}")
    
    except Exception as e:
        print(f"{RED}Error creating Predictive Lifecycle Management slides: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    # Calculate execution time
    total_time = time.time() - start_time
    print(f"{PURPLE}Predictive Lifecycle Management slides created in {total_time:.2f} seconds{RESET}")
    return total_time

# This is the function to create a high risk table that accepts start and end indices
def create_high_risk_table(slide, high_risk_devices, x, y, width, height, title, start_index=0, end_index=9):
    """Create a table of high risk devices with specified range and modern styling."""
    # Add title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.3))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.bold = True
    
    # Get the device slice for the specified range
    device_slice = high_risk_devices[start_index:end_index]
    
    # Determine number of rows (header + device rows)
    rows = len(device_slice) + 1  # Header + devices in range
    
    # Create table - place it below the title
    table = slide.shapes.add_table(
        rows, 
        5,  # 5 columns: Model, Serial, Risk Score, Lifecycle Status, Replacement Cost
        x, 
        y + Inches(0.57),
        width, 
        height - Inches(0.35)
    ).table

    table.columns[0].width = Inches(1.8)    # Model - wider for longer model names
    table.columns[1].width = Inches(1.5)    # Serial - space for partial serial
    table.columns[2].width = Inches(0.93)   # Risk Score - smaller numeric column
    table.columns[3].width = Inches(1.2)    # Lifecycle Status - medium space for status
    table.columns[4].width = Inches(1.0)    # Replacement Cost - space for dollar amounts
    
    # Set header row - using Cisco blue for modern look
    headers = ["Model", "Serial", "Risk Score", "Lifecycle Status", "Est. Cost"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set modern header background - Cisco blue
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 120, 206)  # Cisco blue #0078CE
        
        # Set header text color to white for better contrast
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Add device data with alternating row colors
    for i, device in enumerate(device_slice):
        row = i + 1
        
        # Apply alternating row colors for a modern look
        for j in range(5):
            cell = table.cell(row, j)
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
            else:
                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)  # Very light blue-gray #F5F7FA
        
        # Model
        cell = table.cell(row, 0)
        cell.text = device.model
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Serial
        cell = table.cell(row, 1)
        serial_text = device.serial[:12] + '...' if len(device.serial) > 14 else device.serial
        cell.text = serial_text
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Risk Score with enhanced visualization
        cell = table.cell(row, 2)
        risk_score = device.get_risk_score()
        
        # Add risk score with visual indicator
        p = cell.text_frame.paragraphs[0]
        p.text = f"{risk_score}"
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER
        
        # Use dot indicator for risk (more modern approach)
        if risk_score >= 70:
            # Add red dot indicator for high risk
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 240, 240)  # Very light red background
            p.font.color.rgb = HIGH_RISK_COLOR
            p.font.bold = True
        elif risk_score >= 40:
            # Add yellow dot indicator for medium risk
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 250, 230)  # Very light yellow background
            p.font.color.rgb = MEDIUM_RISK_COLOR
            p.font.bold = True
        
        # Lifecycle Status with enhanced visualization
        cell = table.cell(row, 3)
        cell.text = device.lifecycle_status
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set color and background based on lifecycle status
        if device.lifecycle_status == "End of Support":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light red background
            cell.text_frame.paragraphs[0].font.color.rgb = HIGH_RISK_COLOR
            cell.text_frame.paragraphs[0].font.bold = True
        elif device.lifecycle_status == "Critical":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 240, 230)  # Light orange background
            cell.text_frame.paragraphs[0].font.color.rgb = HIGH_RISK_COLOR
        elif device.lifecycle_status == "Warning":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 250, 230)  # Light yellow background
            cell.text_frame.paragraphs[0].font.color.rgb = MEDIUM_RISK_COLOR
        
        # Replacement Cost
        cell = table.cell(row, 4)
        cell.text = f"${device.get_replacement_cost_estimate():,.2f}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    return table

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some sample inventory devices for testing
    inventory_devices = [
        {"serial": "AAAA-BBBB-1111", "model": "MX64", "firmware": "15.44.0", "networkId": "N1"},
        {"serial": "AAAA-BBBB-1112", "model": "MX64W", "firmware": "15.44.0", "networkId": "N1"},
        {"serial": "AAAA-BBBB-1113", "model": "MX84", "firmware": "15.44.0", "networkId": "N1"},
        {"serial": "AAAA-BBBB-1114", "model": "MS220-8P", "firmware": "14.16.1", "networkId": "N2"},
        {"serial": "AAAA-BBBB-1115", "model": "MS220-24P", "firmware": "14.16.1", "networkId": "N2"},
        {"serial": "AAAA-BBBB-1116", "model": "MS320-48LP", "firmware": "14.16.1", "networkId": "N3"},
        {"serial": "AAAA-BBBB-1117", "model": "MR16", "firmware": "28.5", "networkId": "N4"},
        {"serial": "AAAA-BBBB-1118", "model": "MR34", "firmware": "29.5", "networkId": "N4"},
        {"serial": "AAAA-BBBB-1119", "model": "MR84", "firmware": "29.5", "networkId": "N6"},
        {"serial": "AAAA-BBBB-1120", "model": "MV21", "firmware": "5.2", "networkId": "N5"},
        {"serial": "AAAA-BBBB-1121", "model": "MS350-48LP", "firmware": "14.32.1", "networkId": "N7"},
        {"serial": "AAAA-BBBB-1122", "model": "MX100", "firmware": "15.44.0", "networkId": "N8"},
        {"serial": "AAAA-BBBB-1123", "model": "MX67W", "firmware": "17.6.0", "networkId": "N9"},
        {"serial": "AAAA-BBBB-1124", "model": "MX68", "firmware": "17.6.0", "networkId": "N9"},
        {"serial": "AAAA-BBBB-1125", "model": "MG21", "firmware": "5.0.0", "networkId": "N10"},
    ]
    
    # Sample networks data
    networks = [
        {"id": "N1", "name": "HQ Network"},
        {"id": "N2", "name": "Office 1"},
        {"id": "N3", "name": "Office 2"},
        {"id": "N4", "name": "Branch 1"},
        {"id": "N5", "name": "Branch 2"},
        {"id": "N6", "name": "Branch 3"},
        {"id": "N7", "name": "Datacenter 1"},
        {"id": "N8", "name": "Datacenter 2"},
        {"id": "N9", "name": "Remote Site 1"},
        {"id": "N10", "name": "Remote Site 2"},
    ]
    
    # Sample EOL data
    eol_data = {
        "MX60": {"announcement": "Aug 30, 2021", "end_of_sale": "Aug 30, 2022", "end_of_support": "Aug 30, 2027"},
        "MX64": {"announcement": "Aug 30, 2021", "end_of_sale": "Aug 30, 2022", "end_of_support": "Aug 30, 2027"},
        "MX84": {"announcement": "Aug 30, 2021", "end_of_sale": "Aug 30, 2022", "end_of_support": "Aug 30, 2027"},
        "MR16": {"announcement": "Feb 28, 2017", "end_of_sale": "Feb 28, 2018", "end_of_support": "Feb 28, 2026"},
        "MS220": {"announcement": "Jun 8, 2020", "end_of_sale": "Jun 8, 2021", "end_of_support": "Jun 8, 2026"},
        "MV21": {"announcement": "Sep 30, 2021", "end_of_sale": "Sep 30, 2022", "end_of_support": "Sep 30, 2027"}
    }
    
    # Generate the lifecycle management slides
    await generate(
        api_client, 
        template_path, 
        output_path, 
        inventory_devices,
        networks,
        eol_data
    )

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python predictive_lifecycle.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))