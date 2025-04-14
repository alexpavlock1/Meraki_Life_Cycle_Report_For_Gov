import os
import sys
import asyncio
import time
import re
import requests
import json
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from collections import defaultdict
import random

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Colors for the firmware compliance slide
GOOD_COLOR = RGBColor(108, 184, 108)  # Green
WARNING_COLOR = RGBColor(248, 196, 71)  # Yellow/Amber
CRITICAL_COLOR = RGBColor(227, 119, 84)  # Red/Orange
TITLE_COLOR = RGBColor(39, 110, 55)  # Dark green for subtitle

# Product type mapping - For MG, MV, MT devices
PRODUCT_MAPPING = {
    'cellularGateway': 'MG',
    'camera': 'MV',
    'sensor': 'MT'
}

# Import the AdaptiveRateLimiter from the clients module if available
try:
    from clients import AdaptiveRateLimiter, rate_limited_api_call, get_api_key, logger
    #print(f"{GREEN}Successfully imported rate limiting from clients module{RESET}")
    HAS_RATE_LIMITER = True
except ImportError:
    #print(f"{YELLOW}Could not import AdaptiveRateLimiter from clients, using built-in version{RESET}")
    HAS_RATE_LIMITER = False
    
    # Optimized AdaptiveRateLimiter implementation
    # Set up a local logger for the built-in version
    import logging
    local_logger = logging.getLogger('meraki_firmware_compliance')
    local_logger.setLevel(logging.INFO)
    if not local_logger.handlers:
        file_handler = logging.FileHandler('logs/meraki_api_detailed.log')
        file_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        local_logger.addHandler(file_handler)
        local_logger.propagate = False
    
    class AdaptiveRateLimiter:
        """Highly aggressive rate limiter for maximum Meraki API performance."""
        
        def __init__(self, initial_limit=50, min_limit=30, max_limit=60):
            self.current_limit = initial_limit
            self.min_limit = min_limit
            self.max_limit = max_limit
            self.semaphore = asyncio.Semaphore(self.current_limit)
            self.error_count = 0
            self.success_count = 0
            self.consecutive_successes = 0
            self.last_adjustment_time = time.time()
        
        async def wait(self):
            """Acquire the semaphore before making an API request."""
            await self.semaphore.acquire()
            return True
        
        def release(self):
            """Release the semaphore after making an API request."""
            self.semaphore.release()
        
        def decrease_limit(self):
            """Decrease the concurrency limit after a 429 error."""
            if self.current_limit > self.min_limit:
                # Decrease by up to 2 depending on error severity
                decrease_amount = 2 if self.error_count > 3 else 1
                self.current_limit = max(self.min_limit, self.current_limit - decrease_amount)
                self.semaphore = asyncio.Semaphore(self.current_limit)
                local_logger.info(f"Rate limit hit! Decreasing concurrency limit to {self.current_limit}")
            self.last_adjustment_time = time.time()
        
        def increase_limit(self):
            """More aggressively increase the concurrency limit after consecutive successes."""
            if self.current_limit < self.max_limit:
                # Much faster ramp-up based on consecutive successes
                if self.consecutive_successes > 500:
                    increase_amount = 5  # Jump by 5 with many successes
                elif self.consecutive_successes > 200:
                    increase_amount = 3  # Jump by 3 with moderate successes
                else:
                    increase_amount = 1  # Default increase
                self.current_limit = min(self.max_limit, self.current_limit + increase_amount)
                self.semaphore = asyncio.Semaphore(self.current_limit)
                local_logger.info(f"Increased concurrency limit to {self.current_limit}")
                # If we're ramping up rapidly without errors, be more aggressive next time
                if increase_amount > 1:
                    self.consecutive_successes -= 100  # Reduce but don't reset consecutive successes
            self.last_adjustment_time = time.time()
            
        def check_and_adjust(self):
            """More responsive adjustment based on recent performance."""
            current_time = time.time()
            
            # Check more frequently for quicker rate limit adjustments
            if current_time - self.last_adjustment_time < 1:  # Check every second instead of every 5 seconds
                return
                
            # Increase rate limit faster with fewer successes
            if self.success_count > 20 and self.error_count == 0:  # Only need 20 successes instead of 50
                self.consecutive_successes += self.success_count
                self.increase_limit()
                self.success_count = 0
            # If we're already at max and still succeeding, keep tracking consecutive successes
            elif self.success_count > 0 and self.error_count == 0 and self.current_limit == self.max_limit:
                self.consecutive_successes += self.success_count
                self.success_count = 0
            # If we've had errors, decrease the limit
            elif self.error_count > 0:
                self.decrease_limit()
                self.error_count = 0
                self.success_count = 0
                self.consecutive_successes = 0
    
    # Fallback rate_limited_api_call function
    async def rate_limited_api_call(api_func, rate_limiter, max_retries=4, *args, **kwargs):
        """
        Execute an API call with simpler error handling to avoid stalls.
        
        Args:
            api_func: The API function to call
            rate_limiter: AdaptiveRateLimiter instance
            max_retries: Maximum number of retries on failure
            *args, **kwargs: Arguments to pass to the API function
            
        Returns:
            API call result or raises an exception after retries
        """
        retries = 0
        last_error = None
        
        while retries <= max_retries:
            try:
                # Extremely minimal jitter to avoid synchronized requests
                jitter = random.uniform(0.001, 0.01)  # Ultra-minimal jitter (1-10ms)
                await asyncio.sleep(jitter)
                
                async with rate_limiter.semaphore:
                    # Set a timeout on the API call to prevent hanging indefinitely
                    # Increased timeout for first attempt, reduce for retries
                    timeout = 20.0 if retries == 0 else 15.0
                    result = await asyncio.wait_for(api_func(*args, **kwargs), timeout=timeout)
                    # Simple success counter, no complex logic
                    rate_limiter.success_count += 1
                    return result
                    
            except asyncio.TimeoutError:
                # Handle timeout errors
                #print(f"{RED}Timeout in API call {api_func.__name__}{RESET}")
                retries += 1
                if retries <= max_retries:
                    # Almost no backoff for timeouts to maximize throughput
                    await asyncio.sleep(0.05)  # Minimal timeout backoff
                else:
                    raise Exception(f"API timeout after {max_retries} retries")
                    
            except Exception as e:
                error_str = str(e)
                is_rate_limit = "429" in error_str
                
                if is_rate_limit:
                    # Rate limit error - increment counter and retry
                    rate_limiter.error_count += 1
                    retries += 1
                    if retries <= max_retries:
                        # Minimal backoff with decay for ultra-fast recovery
                        # Wait almost no time to maximize throughput
                        wait_time = 0.2 * (0.5 ** retries) * retries
                        #print(f"{YELLOW}Rate limit hit, retrying in {wait_time:.1f}s (attempt {retries}/{max_retries}){RESET}")
                        await asyncio.sleep(wait_time)
                    else:
                        raise Exception(f"Rate limit exceeded after {max_retries} retries")
                else:
                    # Other error - pass through but add context if available
                    if 'errors' in error_str or 'Error' in error_str:
                        # Add retry information to help diagnose specific API errors
                        raise type(e)(f"{error_str} (during API call to {api_func.__name__})") from e
                    else:
                        raise
    
    # Fallback API key function
    def get_api_key():
        """Retrieve Meraki API key from environment variable."""
        api_key = os.environ.get("MERAKI_API_KEY")
        if not api_key:
            raise ValueError("MERAKI_API_KEY environment variable is not set")
        return api_key

async def get_network_firmware_upgrades(aiomeraki, network_id, rate_limiter=None):
    """Get firmware upgrade information for a network with rate limiting."""
    try:
        # Use the rate_limited_api_call function for better rate limiting
        return await rate_limited_api_call(
            aiomeraki.networks.getNetworkFirmwareUpgrades,
            rate_limiter,
            networkId=network_id
        )
    except Exception as e:
        print(f"{RED}Error getting firmware upgrades for network {network_id}: {e}{RESET}")
        return None

def get_latest_stable_firmware(available_versions):
    """Find the latest stable firmware version from available versions."""
    stable_versions = [v for v in available_versions if v.get('releaseType') == 'stable']
    if not stable_versions:
        return None
    
    # Sort by release date (newest first)
    sorted_versions = sorted(stable_versions, key=lambda x: x.get('releaseDate', ''), reverse=True)
    return sorted_versions[0] if sorted_versions else None

def get_firmware_major_version(firmware_str):
    """Extract major version from firmware string (first two version numbers)."""
    if not firmware_str:
        return None
        
    # Extract version pattern (like 18.211.5 from various formats)
    # Try different patterns to handle various firmware string formats
    patterns = [
        r'(\d+\.\d+)(?:\.\d+)*',  # Standard format like 18.211.5
        r'[A-Z]+ (\d+\.\d+)(?:\.\d+)*',  # Format like "MV 18.211.5"
        r'[a-z]+-(\d+\.\d+)(?:\.\d+)*'  # Format like "firmware-18.211.5"
    ]
    
    for pattern in patterns:
        match = re.search(pattern, firmware_str)
        if match:
            return match.group(1)
    
    return None

def extract_firmware_version(firmware_str):
    """Extract the full version from a firmware string."""
    if not firmware_str:
        return None
    
    # Extract version pattern (all numbers including patches)
    patterns = [
        r'(\d+\.\d+(?:\.\d+)*)',  # Standard format like 18.211.5
        r'[A-Z]+ (\d+\.\d+(?:\.\d+)*)',  # Format like "MV 18.211.5"
        r'[a-z]+-(\d+\.\d+(?:\.\d+)*)'  # Format like "firmware-18.211.5"
    ]
    
    for pattern in patterns:
        match = re.search(pattern, firmware_str)
        if match:
            return match.group(1)
    
    return None

def categorize_firmware_status(current_firmware, latest_stable_firmware):
    """
    Categorize firmware status as Good, Warning, or Critical.
    
    Good: Running latest major version with latest patch
    Warning: Running latest major version but old patch
    Critical: Running older major version
    """
    if not current_firmware or not latest_stable_firmware:
        return "Critical"  # Default to critical if we can't determine
        
    current_major = get_firmware_major_version(current_firmware)
    latest_major = get_firmware_major_version(latest_stable_firmware)
    
    current_full = extract_firmware_version(current_firmware)
    latest_full = extract_firmware_version(latest_stable_firmware)
    
    if not current_major or not latest_major:
        return "Critical"  # Default to critical if we can't parse versions
    
    if current_major == latest_major:
        if current_full == latest_full:
            return "Good"  # Same major version and patch
        else:
            return "Warning"  # Same major version but different patch
    else:
        return "Critical"  # Different major version

async def analyze_network_firmware(aiomeraki, networks, rate_limiter):
    """Analyze firmware status for all networks."""
    #print(f"{BLUE}Analyzing firmware status for {len(networks)} networks...{RESET}")
    
    firmware_stats = {
        'MG': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'Versions': defaultdict(int)},
        'MV': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'Versions': defaultdict(int)},
        'MT': {'Good': 0, 'Warning': 0, 'Critical': 0, 'Total': 0, 'Versions': defaultdict(int)}
    }
    
    latest_firmware = {
        'MG': None,
        'MV': None,
        'MT': None
    }
    
    # Find the latest stable firmware versions first
    for product_type in PRODUCT_MAPPING.values():
        latest_firmware[product_type] = {
            'version': None,
            'releaseDate': None
        }
    
    # Process networks in massive batches for extreme speed
    chunk_size = 30  # Very large batch size for massive parallelism
    network_chunks = [networks[i:i + chunk_size] for i in range(0, len(networks), chunk_size)]
    
    # First pass: Find the latest stable firmware for each product type
    #print(f"{BLUE}First pass: Finding latest stable firmware versions...{RESET}")
    for i, chunk in enumerate(network_chunks):
        progress = (i + 1) / len(network_chunks) * 100
        if i % 10 == 0:  # Show progress every 10 chunks
            #print(f"{BLUE}Processing chunk {i+1}/{len(network_chunks)} ({progress:.1f}%)...{RESET}")
            pass
        # No delay needed with our aggressive parallelism
        # Process each network in the chunk
        tasks = []
        for network in chunk:
            task = get_network_firmware_upgrades(aiomeraki, network['id'], rate_limiter)
            tasks.append(task)
            
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        for network, result in zip(chunk, results):
            # Skip if the result is an exception or None
            if isinstance(result, Exception) or result is None:
                continue
                
            firmware_data = result
            if 'products' not in firmware_data:
                continue
                
            products = firmware_data.get('products', {})
            
            for product_key, product_category in PRODUCT_MAPPING.items():
                product_data = products.get(product_key)
                if not product_data or 'availableVersions' not in product_data:
                    continue
                
                # Find the latest stable version from available versions
                available_versions = product_data.get('availableVersions', [])
                latest_stable = get_latest_stable_firmware(available_versions)
                
                if latest_stable:
                    version_str = latest_stable.get('shortName', '') or latest_stable.get('firmware', '')
                    release_date = latest_stable.get('releaseDate', '')
                    
                    # Update our record of the latest firmware if it's newer
                    if not latest_firmware[product_category]['version'] or \
                       (release_date and (not latest_firmware[product_category]['releaseDate'] or 
                                          release_date > latest_firmware[product_category]['releaseDate'])):
                        latest_firmware[product_category]['version'] = version_str
                        latest_firmware[product_category]['releaseDate'] = release_date
                        #print(f"{GREEN}Found newer {product_category} firmware: {version_str} (released: {release_date}){RESET}")
        
        # Check and adjust rate limiter periodically
        if hasattr(rate_limiter, 'check_and_adjust'):
            rate_limiter.check_and_adjust()
    
    # Second pass: Categorize networks based on firmware status
    #print(f"{BLUE}Second pass: Categorizing networks based on firmware status...{RESET}")
    for i, chunk in enumerate(network_chunks):
        progress = (i + 1) / len(network_chunks) * 100
        if i % 10 == 0:  # Show progress every 10 chunks
            #print(f"{BLUE}Categorizing chunk {i+1}/{len(network_chunks)} ({progress:.1f}%)...{RESET}")
            pass
            
        # No delay needed with our aggressive parallelism
        # Process each network in the chunk
        tasks = []
        for network in chunk:
            task = get_network_firmware_upgrades(aiomeraki, network['id'], rate_limiter)
            tasks.append(task)
            
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        for network, result in zip(chunk, results):
            # Skip if the result is an exception or None
            if isinstance(result, Exception) or result is None:
                continue
                
            firmware_data = result
            if 'products' not in firmware_data:
                continue
                
            products = firmware_data.get('products', {})
            
            for product_key, product_category in PRODUCT_MAPPING.items():
                product_data = products.get(product_key)
                if not product_data or 'currentVersion' not in product_data:
                    continue
                
                # We have a network with this product type
                firmware_stats[product_category]['Total'] += 1
                
                # Get current firmware version
                current_version = product_data.get('currentVersion', {})
                current_firmware = current_version.get('shortName', '') or current_version.get('firmware', '')
                
                # Record this firmware version
                if current_firmware:
                    firmware_stats[product_category]['Versions'][current_firmware] += 1
                
                # Get the latest stable firmware we determined in the first pass
                latest_stable_firmware = latest_firmware[product_category]['version']
                
                if latest_stable_firmware:
                    # Categorize firmware status
                    status = categorize_firmware_status(current_firmware, latest_stable_firmware)
                    firmware_stats[product_category][status] += 1
                else:
                    # If we couldn't determine the latest stable firmware, default to critical
                    firmware_stats[product_category]['Critical'] += 1
        
        # Check and adjust rate limiter periodically
        if hasattr(rate_limiter, 'check_and_adjust'):
            rate_limiter.check_and_adjust()
    
    # Print summary of findings
    #print(f"{GREEN}Firmware analysis completed:{RESET}")
    for product, stats in firmware_stats.items():
        total = stats['Total']
        if total > 0:
            #print(f"{BLUE}{product} networks: {total}{RESET}")
            #print(f"  - Good: {stats['Good']} ({stats['Good']/total*100:.1f}%)")
            #print(f"  - Warning: {stats['Warning']} ({stats['Warning']/total*100:.1f}%)")
            #print(f"  - Critical: {stats['Critical']} ({stats['Critical']/total*100:.1f}%)")
            
            
            # Print top 5 versions
            top_versions = sorted(stats['Versions'].items(), key=lambda x: x[1], reverse=True)[:5]
            #print(f"  - Top versions:")
            for version, count in top_versions:
                #print(f"    * {version}: {count} networks")
                pass
    
    # Strip out just the version strings from the latest_firmware dict
    latest_versions = {k: v['version'] for k, v in latest_firmware.items()}
    
    return firmware_stats, latest_versions

def draw_percentage_circle(slide, x, y, radius, percentage, color, line_width=2):
    """Draw a circle with percentage in the middle."""
    # Create the circle shape - using OVAL (type 9) for proper circle
    circle = slide.shapes.add_shape(9, x - radius, y - radius, radius * 2, radius * 2)
    circle.fill.background()
    circle.line.color.rgb = color
    circle.line.width = Pt(line_width)
    
    # Add percentage text
    text_box = slide.shapes.add_textbox(x - radius, y - Inches(0.5), radius * 2, radius)
    text_frame = text_box.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = f"{percentage}%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

async def generate(api_client, template_path, output_path, networks=None, inventory_devices=None):
    """Generate the Firmware Compliance slide for MG, MV, MT."""
    print(f"\n{GREEN}Generating MG/MV/MT Firmware Compliance slide (Slide 9)...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # If networks list is not provided, would need to fetch it
    if not networks:
        print(f"{RED}No networks data provided{RESET}")
        return
    
    #print(f"{BLUE}Using network data for {len(networks)} networks{RESET}")
    
    # Import meraki API for firmware data
    try:
        # Get API key - use imported function or fallback
        api_key = get_api_key()
        
        # Create a rate limiter with very aggressive settings for maximum speed
        rate_limiter = AdaptiveRateLimiter(initial_limit=50, min_limit=30, max_limit=60)
        #print(f"{GREEN}Using adaptive rate limiter with initial concurrency limit of {rate_limiter.current_limit}{RESET}")
        
        # Set up Meraki client
        import meraki.aio
        async with meraki.aio.AsyncDashboardAPI(
            api_key=api_key,
            suppress_logging=True,
            maximum_retries=3,
            base_url="https://api.gov-meraki.com/api/v1"
        ) as aiomeraki:
            # Analyze firmware for all networks
            firmware_stats, latest_firmware = await analyze_network_firmware(aiomeraki, networks, rate_limiter)
    except ImportError as e:
        print(f"{YELLOW}Could not import required modules: {e}. Using mock data for testing.{RESET}")
        # Use mock data for testing if API access isn't available
        firmware_stats = {
            'MG': {
                'Good': 42, 'Warning': 158, 'Critical': 72, 'Total': 272,
                'Versions': {
                    '1.15.0': 42,
                    '1.14.2': 158,
                    '1.12.0': 40,
                    '1.11.0': 20,
                    '1.10.0': 12
                }
            },
            'MV': {
                'Good': 210, 'Warning': 135, 'Critical': 85, 'Total': 430,
                'Versions': {
                    '4.18.0': 210,
                    '4.17.0': 135,
                    '4.0.0': 35,
                    '3.12.0': 30,
                    '3.10.0': 20
                }
            },
            'MT': {
                'Good': 90, 'Warning': 65, 'Critical': 20, 'Total': 175,
                'Versions': {
                    '1.5.0': 90,
                    '1.4.0': 65,
                    '1.2.0': 20
                }
            }
        }
        
        latest_firmware = {
            'MG': '1.15.0',
            'MV': '4.18.0',
            'MT': '1.5.0'
        }
    except Exception as e:
        print(f"{RED}Error analyzing firmware: {e}{RESET}")
        import traceback
        traceback.print_exc()
        return
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    #print(f"{BLUE}Updating PowerPoint with MG/MV/MT firmware compliance data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)
        
        # Use slide 9 (index 8)
        # If the slide doesn't exist, add it
        if len(prs.slides) < 9:
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides[8]
        
        # Clear existing shapes except for title
        title_shape = None
        shapes_to_remove = []

        # Look for existing title - any title at the top of the slide is preserved
        for shape in slide.shapes:
            # Identify likely title shapes based on position and properties
            if hasattr(shape, "text_frame") and shape.top < Inches(1.0):
                title_shape = shape
                continue
            
            # Mark non-title shapes for removal
            shapes_to_remove.append(shape)

        # Remove all other shapes
        for shape in shapes_to_remove:
            try:
                if hasattr(shape, '_sp'):
                    sp = shape._sp
                    sp.getparent().remove(sp)
            except Exception as e:
                print(f"{RED}Error removing shape: {e}{RESET}")

        if title_shape:
            #print(f"{BLUE}Using existing title from template{RESET}")
            pass
        else:
            # Create a title if none exists
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "Firmware Compliance - MG/MV/MT"
            title_p.font.size = Pt(44)
            title_p.font.bold = True
            #print(f"{YELLOW}Added new title: 'Firmware Compliance - MG/MV/MT'{RESET}")
        
        # Add horizontal line across the full width of the slide
        line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(11.0), Inches(1.2))
        line.line.color.rgb = TITLE_COLOR  # Dark green
        line.line.width = Pt(2)
        
        # Add "By Network" subtitle
        subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.6))
        subtitle_p = subtitle.text_frame.add_paragraph()
        subtitle_p.text = "By Network"
        subtitle_p.font.size = Pt(32)  # Adjusted font size
        subtitle_p.font.color.rgb = TITLE_COLOR  # Dark green
        
        # Define spacing constants
        row_spacing = 1.6
        
        # Row positions
        y_positions = [Inches(2.7), Inches(2.7 + row_spacing), Inches(2.7 + 2*row_spacing)]
        category_labels = ["\"Good\"", "\"Warning\"", "\"Critical\""]
        descriptions = [
            "FW beyond\n180 days from\nEOST",
            "FW within 180\ndays of EOST",
            "FW past EOST\nDate"
        ]
        
        # Add category labels and descriptions
        for i, (y, label, desc) in enumerate(zip(y_positions, category_labels, descriptions)):
            # Add category label
            label_box = slide.shapes.add_textbox(Inches(0.5), y - Inches(0.3), Inches(1.5), Inches(0.5))
            label_p = label_box.text_frame.add_paragraph()
            label_p.text = label
            label_p.font.size = Pt(22)  # Adjusted font size
            label_p.font.bold = True
            
            # Add description
            desc_box = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.2), Inches(1.5), Inches(1.0))
            desc_p = desc_box.text_frame.add_paragraph()
            desc_p.text = desc
            desc_p.font.size = Pt(12)
        
        # Define product column positions
        product_cols = {
            'MG': Inches(3.5),
            'MV': Inches(7.25),
            'MT': Inches(10.75)
        }
        
        # Define positions array from column dict for easier iteration
        product_positions = [product_cols['MG'], product_cols['MV'], product_cols['MT']]
        products = ['MG', 'MV', 'MT']
        
        # Add vertical dividers between product columns
        divider1 = slide.shapes.add_connector(1, Inches(5.75), Inches(1.8), Inches(5.75), Inches(7.0))
        divider1.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
        divider1.line.width = Pt(1)
        
        divider2 = slide.shapes.add_connector(1, Inches(9.25), Inches(1.8), Inches(9.25), Inches(7.0))
        divider2.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
        divider2.line.width = Pt(1)
        
        # Add product type headers
        for i, (pos, product) in enumerate(zip(product_positions, products)):
            # Position headers below the line but not too far
            header = slide.shapes.add_textbox(pos - Inches(0.6), Inches(1.3), Inches(1.2), Inches(0.5))
            p = header.text_frame.add_paragraph()
            p.text = product
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        
        # Add percentage circles and stats for each product and category
        colors = [GOOD_COLOR, WARNING_COLOR, CRITICAL_COLOR]
        
        for row, (y, category, color) in enumerate(zip(y_positions, ['Good', 'Warning', 'Critical'], colors)):
            for col, (x, product) in enumerate(zip(product_positions, products)):
                # Calculate percentage
                total = firmware_stats[product]['Total']
                if total > 0:
                    count = firmware_stats[product][category]
                    percentage = int(round(count / total * 100))
                else:
                    count = 0
                    percentage = 0
                
                # Draw percentage circle
                circle_x = x
                circle_y = y    
                draw_percentage_circle(slide, circle_x, circle_y, Inches(0.6), percentage, color, 2)
                
                # Add network count below the circle
                count_box = slide.shapes.add_textbox(x - Inches(1.0), y + Inches(0.4), Inches(2.0), Inches(0.4))
                count_p = count_box.text_frame.add_paragraph()
                count_p.text = f"{count:,}/{total:,} Networks"
                count_p.font.size = Pt(12)
                count_p.alignment = PP_ALIGN.CENTER
                count_p.font.color.rgb = color
                
                # Define positions for version listings in each column
                version_x = x + Inches(0.7)
                version_y = y - Inches(0.9)
                
                # Filter versions for this category based on their relationship to the latest firmware
                filtered_versions = {}
                
                # Always show the latest stable version for the "Good" category
                if category == 'Good' and latest_firmware[product]:
                    latest_version = latest_firmware[product]
                    total_version_count = 0
                    
                    # See if any networks are actually running this version
                    for version, v_count in firmware_stats[product]['Versions'].items():
                        if version == latest_version:
                            total_version_count = v_count
                    
                    # Add the version to our display list - even if count is 0
                    filtered_versions[latest_version] = total_version_count
                else:
                    # For Warning and Critical, just show actual versions in use
                    for version, v_count in firmware_stats[product]['Versions'].items():
                        status = categorize_firmware_status(version, latest_firmware[product])
                        if status == category:
                            filtered_versions[version] = v_count
                
                # Sort and display top 5 versions for this category
                top_versions = sorted(filtered_versions.items(), key=lambda x: x[1], reverse=True)[:5]
                if top_versions:
                    # Add each line to the slide - with proper spacing
                    for i, (version, v_count) in enumerate(top_versions):
                        version_box = slide.shapes.add_textbox(version_x, version_y + Inches(i * 0.22), Inches(2.2), Inches(0.25))
                        version_p = version_box.text_frame.add_paragraph()
                        version_p.text = f"{version} = {v_count}"
                        version_p.font.size = Pt(11)
                        version_p.font.color.rgb = color
        
        # Save the presentation
        prs.save(output_path)
        print(f"{GREEN}Updated MG/MV/MT Firmware Compliance slide (Slide 9){RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    print(f"{PURPLE}MG/MV/MT Firmware Compliance slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some sample networks for testing
    sample_networks = [
        {"id": "N1", "name": "Network 1"},
        {"id": "N2", "name": "Network 2"},
        {"id": "N3", "name": "Network 3"},
    ]
    
    await generate(api_client, template_path, output_path, networks=sample_networks)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python slide9.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))