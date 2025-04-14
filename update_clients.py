import os
import sys
import json
import traceback
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ANSI color codes for terminal output
BLUE = '\033[94m'
GREEN = '\033[92m'
RED = '\033[91m'
YELLOW = '\033[93m'
RESET = '\033[0m'

def update_title_slide(prs, org_names):
    """Update the title slide with organization names."""
    try:
        if len(prs.slides) > 0:
            title_slide = prs.slides[0]
            
            # Search for the title shape with the text to replace
            title_shape = None
            for shape in title_slide.shapes:
                if hasattr(shape, 'text') and 'Company - Meraki Bi-Weekly Life Cycle Report' in shape.text:
                    title_shape = shape
                    break
            
            # If we found the title shape, update it
            if title_shape:
                # Create organization names string
                if len(org_names) == 1:
                    org_names_str = next(iter(org_names.values()))
                else:
                    org_names_str = " & ".join(org_names.values())
                
                # Limit length to avoid overflow
                if len(org_names_str) > 60:
                    org_names_str = org_names_str[:57] + "..."
                
                # Get current date
                import datetime
                current_date = datetime.datetime.now().strftime("%B %d, %Y")
                
                # Store original text and formatting
                original_text = title_shape.text
                
                # We need to preserve the formatting of all paragraphs and runs
                original_paragraphs = []
                for p in title_shape.text_frame.paragraphs:
                    p_info = {'text': p.text, 'runs': []}
                    for run in p.runs:
                        run_info = {
                            'text': run.text,
                            'font': {
                                'bold': getattr(run.font, 'bold', None),
                                'italic': getattr(run.font, 'italic', None),
                                'size': getattr(run.font, 'size', None),
                                'name': getattr(run.font, 'name', None)
                            }
                        }
                        p_info['runs'].append(run_info)
                    original_paragraphs.append(p_info)
                
                # Create the new text
                new_text = original_text.replace(
                    "Company - Meraki Bi-Weekly Life Cycle Report", 
                    f"{org_names_str} - Meraki Bi-Weekly Life Cycle Report"
                )
                
                # Also update the date if it's in the title
                if "March 22, 2025" in new_text:
                    new_text = new_text.replace("March 22, 2025", current_date)
                
                # Clear the text frame
                title_shape.text = ""
                
                # Add the new text
                paragraph = title_shape.text_frame.paragraphs[0]
                paragraph.text = new_text
                
                # Explicitly set font size to 40pt
                for run in paragraph.runs:
                    run.font.size = Pt(40)
                
                #print(f"{GREEN}Updated title slide with organization name(s): {org_names_str}{RESET}")
                return True
            else:
                print(f"{YELLOW}Could not find title text to update on slide 1{RESET}")
        else:
            print(f"{YELLOW}No slides found in the presentation{RESET}")
        
        return False
    
    except Exception as e:
        print(f"{RED}Error updating title slide: {e}{RESET}")
        traceback.print_exc()
        return False

def update_client_days_header(slide, days):
    """Update the clients section header with the correct number of days."""
    try:
        # Find and remove any existing days header
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text.strip()
                if "Clients (for last" in text and "day" in text:
                    shapes_to_remove.append(shape)
                    #print(f"Found client days header to remove: '{text}'")
        
        # Remove the identified shapes
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
                #print(f"{GREEN}Removed existing client days header{RESET}")
            except Exception as e:
                print(f"{YELLOW}Couldn't remove existing header: {e}{RESET}")
        
        # Correct format for days (singular/plural)
        days_text = "day" if days == 1 else "days"
        new_text = f" (for last {days} {days_text})"
        
        # Add new textbox with exact positioning and formatting
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # Exact positioning (2.63" horizontal, 3.57" vertical)
        textbox = slide.shapes.add_textbox(
            Inches(2.63), Inches(3.57), Inches(5), Inches(0.4)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Set the text
        p = tf.paragraphs[0]
        p.text = new_text
        p.alignment = PP_ALIGN.LEFT

        from pptx.enum.dml import MSO_THEME_COLOR
        
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = False
        p.font.color.rgb = RGBColor(39, 118, 38)  # Dark green color for "Clients" text
        
        #print(f"{GREEN}Created client header with correct formatting: '{new_text}'{RESET}")
        return True
    
    except Exception as e:
        print(f"{RED}Error updating client days header: {e}{RESET}")
        traceback.print_exc()
        return False

def update_dashboard_slide(stats, template_path, output_path, days=14, org_names=None):
    """Update dashboard statistics in PowerPoint."""
    try:
        #print(f"{BLUE}Opening template: {template_path}{RESET}")
        #print(f"{BLUE}Will save to: {output_path}{RESET}")
        #print(f"{BLUE}Stats to update: {stats}{RESET}")
        #print(f"{BLUE}Days parameter: {days}{RESET}")
        if org_names:
            #print(f"{BLUE}Organization names: {org_names}{RESET}")
            pass
        
        # Open the template
        prs = Presentation(template_path)
        #print(f"{GREEN}Successfully opened template with {len(prs.slides)} slides{RESET}")
        
        # Update title slide if org_names are provided
        if org_names:
            update_title_slide(prs, org_names)
        
        # Use slide 2 (index 1)
        if len(prs.slides) <= 1:
            print(f"{RED}Error: Template doesn't have a second slide{RESET}")
            return 0
            
        slide = prs.slides[1]
        
        # Update the client days header with the days parameter
        update_client_days_header(slide, days)
        
        #print("Debug: Looking for shapes with specific text content")
        
        # Dictionary to store target shapes
        target_shapes = {
            'networks': None,
            'inventory': None,
            'active_nodes': None,
            'unique_clients_total': None,
            'unique_clients_daily': None,
            'non_unique_clients_total': None,
            'non_unique_clients_daily': None
        }
        
        # Collect all shapes at all levels
        all_shapes = []
        
        # Process all shapes (including in groups)
        def collect_shapes_recursively(shape_container, all_shapes_list):
            """Collect all shapes from a container (slide or group) recursively."""
            for shape in shape_container.shapes:
                all_shapes_list.append(shape)
                
                # If it's a group shape, look inside it
                if hasattr(shape, 'shapes'):
                    collect_shapes_recursively(shape, all_shapes_list)
        
        # Collect all shapes from the slide
        collect_shapes_recursively(slide, all_shapes)
        
        # First pass: identify the existing shapes by text content
        for shape in all_shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                    
                shape_type = type(shape).__name__
                shape_id = getattr(shape, 'shape_id', 'unknown')
                shape_name = getattr(shape, 'name', 'unnamed')
                width = getattr(shape, "width", 0)
                height = getattr(shape, "height", 0)
                
                #print(f"Shape: Type={shape_type}, ID={shape_id}, Name={shape_name}")
                #print(f"  - Text: '{text}'")
                #print(f"  - Size: width={width}, height={height}")
                
                # Look for network stats shapes
                if "networks" in text.lower():
                    target_shapes['networks'] = shape
                    #print(f"  - Identified as Networks shape")
                    
                elif "inventory" in text.lower():
                    target_shapes['inventory'] = shape
                    #print(f"  - Identified as Total Inventory shape")
                    
                elif "active" in text.lower() and "node" in text.lower():
                    target_shapes['active_nodes'] = shape
                    #print(f"  - Identified as Total Active Nodes shape")
                
                # Look for unique client shapes
                elif "unique clients total" in text.lower() and "non" not in text.lower():
                    target_shapes['unique_clients_total'] = shape
                    #print(f"  - Identified as Unique clients total shape")
                
                elif "avg unique clients" in text.lower() and "non" not in text.lower():
                    target_shapes['unique_clients_daily'] = shape
                    #print(f"  - Identified as Unique clients daily shape")
                
                # Look for non-unique client shapes
                elif "non-unique clients total" in text.lower():
                    target_shapes['non_unique_clients_total'] = shape
                    #print(f"  - Identified as Non-unique clients total shape")
                
                elif "non-unique clients" in text.lower() and "per day" in text.lower():
                    target_shapes['non_unique_clients_daily'] = shape 
                    #print(f"  - Identified as Non-unique clients daily shape")
        
        # Helper function to carefully update text while preserving all formatting
        def update_shape_value(shape, new_value):
            """Update a shape's value while preserving its formatting including color."""
            if not shape or not hasattr(shape, "text_frame") or not shape.text_frame:
                return False
                
            try:
                # Format the new value with commas
                new_value_formatted = f"{new_value:,}"
                
                # Get the first paragraph (which contains the value)
                if not shape.text_frame.paragraphs:
                    return False
                
                paragraph = shape.text_frame.paragraphs[0]
                
                # Get the original text
                original_text = paragraph.text
                
                # Store the original formatting of all runs
                original_runs = []
                for run in paragraph.runs:
                    # Store all attributes we can
                    run_info = {
                        'text': run.text,
                        'bold': getattr(run.font, 'bold', None),
                        'italic': getattr(run.font, 'italic', None),
                        'size': getattr(run.font, 'size', None)
                    }
                    
                    # Carefully extract color information
                    if hasattr(run.font, 'color'):
                        # Check if it has RGB color
                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                            run_info['rgb'] = run.font.color.rgb
                        # Check if it has a theme color
                        if hasattr(run.font.color, 'theme_color') and run.font.color.theme_color is not None:
                            run_info['theme_color'] = run.font.color.theme_color
                    
                    original_runs.append(run_info)
                
                # Create the new text by replacing the value part
                # Find digits at the start of text
                digits_match = re.match(r'^[\d,]+', original_text)
                if digits_match:
                    # Replace just the digits
                    new_text = new_value_formatted + original_text[digits_match.end():]
                else:
                    # If no digits found, handle it differently
                    parts = original_text.split(' ', 1)
                    if len(parts) > 1:
                        new_text = new_value_formatted + ' ' + parts[1]
                    else:
                        new_text = new_value_formatted
                
                # Store remaining paragraphs
                remaining_paragraphs = []
                for i in range(1, len(shape.text_frame.paragraphs)):
                    p = shape.text_frame.paragraphs[i]
                    remaining_paragraphs.append(p.text)
                
                # Clear the text frame
                for i in range(len(shape.text_frame.paragraphs)):
                    if i == 0:
                        # Clear first paragraph
                        shape.text_frame.paragraphs[0].clear()
                    else:
                        # Remove extra paragraphs
                        try:
                            shape.text_frame.paragraphs[-1]._p.getparent().remove(shape.text_frame.paragraphs[-1]._p)
                        except:
                            pass
                
                # Add the new text to first paragraph
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.text = new_text
                
                # Apply the original formatting to the new text
                if paragraph.runs and original_runs:
                    run = paragraph.runs[0]
                    
                    # Apply first run's formatting
                    if original_runs[0]['bold'] is not None:
                        run.font.bold = original_runs[0]['bold']
                    if original_runs[0]['italic'] is not None:
                        run.font.italic = original_runs[0]['italic']
                    if original_runs[0]['size'] is not None:
                        run.font.size = original_runs[0]['size']
                        
                    # Apply color - try different methods
                    if 'rgb' in original_runs[0] and original_runs[0]['rgb'] is not None:
                        # Direct RGB assignment
                        run.font.color.rgb = original_runs[0]['rgb']
                    elif 'theme_color' in original_runs[0] and original_runs[0]['theme_color'] is not None:
                        # Theme color assignment
                        run.font.color.theme_color = original_runs[0]['theme_color']
                
                # Add back remaining paragraphs
                for para_text in remaining_paragraphs:
                    para = shape.text_frame.add_paragraph()
                    para.text = para_text
                
                return True
            except Exception as e:
                print(f"{RED}Error updating shape value: {e}{RESET}")
                traceback.print_exc()
                return False
        
        # Shapes updated counter
        shapes_updated = 0
        
        # Helper function to create a new shape if not found
        def create_shape_if_missing(shape_key, value, label, x, y, width, height, color_rgb):
            if not target_shapes[shape_key]:
                #print(f"{YELLOW}Shape for {shape_key} not found, creating a new one{RESET}")
                
                # Create a new textbox
                left = Inches(x)
                top = Inches(y)
                box_width = Inches(width)
                box_height = Inches(height)
                
                textbox = slide.shapes.add_textbox(left, top, box_width, box_height)
                tf = textbox.text_frame
                tf.word_wrap = True
                
                # Add value with commas
                value_p = tf.add_paragraph()
                value_p.text = f"{value:,}"
                value_p.font.size = Pt(36)
                value_p.font.bold = True
                value_p.font.color.rgb = RGBColor(*color_rgb)
                
                # Add label
                label_p = tf.add_paragraph()
                label_p.text = label
                label_p.font.size = Pt(14)
                
                # Store the new shape
                target_shapes[shape_key] = textbox
                return True
            
            return False
                
        # Update Networks
        if target_shapes['networks']:
            if update_shape_value(target_shapes['networks'], stats['total_networks']):
                #print(f"Updated Networks value to {stats['total_networks']:,}")
                shapes_updated += 1
        else:
            if create_shape_if_missing('networks', stats['total_networks'], 'Networks', 0.5, 2.0, 3.0, 1.5, (0, 150, 0)):
                shapes_updated += 1
        
        # Update Inventory
        if target_shapes['inventory']:
            if update_shape_value(target_shapes['inventory'], stats['total_inventory']):
                #print(f"Updated Total Inventory value to {stats['total_inventory']:,}")
                shapes_updated += 1
        else:
            if create_shape_if_missing('inventory', stats['total_inventory'], 'Total Inventory', 4.0, 2.0, 3.0, 1.5, (0, 0, 150)):
                shapes_updated += 1
        
        # Update Active Nodes
        if target_shapes['active_nodes']:
            if update_shape_value(target_shapes['active_nodes'], stats['total_active_nodes']):
                #print(f"Updated Total Active Nodes value to {stats['total_active_nodes']:,}")
                shapes_updated += 1
        else:
            if create_shape_if_missing('active_nodes', stats['total_active_nodes'], 'Total Active Nodes', 7.5, 2.0, 3.0, 1.5, (150, 0, 0)):
                shapes_updated += 1
        
        # Update client statistics
        if 'total_unique_clients' in stats:
            # Update Unique clients total
            if target_shapes['unique_clients_total']:
                if update_shape_value(target_shapes['unique_clients_total'], stats['total_unique_clients']):
                    #print(f"Updated Unique clients total to {stats['total_unique_clients']:,}")
                    shapes_updated += 1
            else:
                if create_shape_if_missing('unique_clients_total', stats['total_unique_clients'], 'Unique clients total', 0.5, 4.0, 2.5, 1.5, (0, 120, 120)):
                    shapes_updated += 1
            
            # Update Unique clients daily
            if target_shapes['unique_clients_daily']:
                if update_shape_value(target_shapes['unique_clients_daily'], stats['avg_unique_clients_per_day']):
                    #print(f"Updated Unique clients daily to {stats['avg_unique_clients_per_day']:,}")
                    shapes_updated += 1
            else:
                if create_shape_if_missing('unique_clients_daily', stats['avg_unique_clients_per_day'], 'Avg unique clients per day', 3.0, 4.0, 2.5, 1.5, (0, 120, 120)):
                    shapes_updated += 1
            
            # Update Non-unique clients total
            if target_shapes['non_unique_clients_total']:
                if update_shape_value(target_shapes['non_unique_clients_total'], stats['total_non_unique_clients']):
                    #print(f"Updated Non-unique clients total to {stats['total_non_unique_clients']:,}")
                    shapes_updated += 1
            else:
                if create_shape_if_missing('non_unique_clients_total', stats['total_non_unique_clients'], 'Non-unique clients total', 0.5, 5.5, 2.5, 1.5, (150, 75, 0)):
                    shapes_updated += 1
            
            # Update Non-unique clients daily
            if target_shapes['non_unique_clients_daily']:
                if update_shape_value(target_shapes['non_unique_clients_daily'], stats['avg_non_unique_clients_per_day']):
                    #print(f"Updated Non-unique clients daily to {stats['avg_non_unique_clients_per_day']:,}")
                    shapes_updated += 1
            else:
                if create_shape_if_missing('non_unique_clients_daily', stats['avg_non_unique_clients_per_day'], 'Non-unique clients per day', 3.0, 5.5, 2.5, 1.5, (150, 75, 0)):
                    shapes_updated += 1
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Successfully updated {shapes_updated} target shapes{RESET}")
        
        return shapes_updated
    
    except Exception as e:
        print(f"{RED}Error in update_dashboard_slide: {e}{RESET}")
        traceback.print_exc()
        return 0

if __name__ == "__main__":
    # Read arguments when run as script
    if len(sys.argv) < 4:
        print("Usage: python update_clients.py <template_path> <output_path> <stats_json> [days] [org_names_json]")
        sys.exit(1)
    
    template_path = sys.argv[1]
    output_path = sys.argv[2]
    stats = json.loads(sys.argv[3])
    
    # Default days value
    days = 14
    
    # Parse days parameter if provided
    if len(sys.argv) > 4:
        try:
            days = int(sys.argv[4])
            #print(f"{BLUE}Using days parameter: {days}{RESET}")
        except ValueError:
            print(f"{YELLOW}Invalid days parameter, using default: {days}{RESET}")
    
    # Parse org_names parameter if provided
    org_names = None
    if len(sys.argv) > 5:
        try:
            org_names = json.loads(sys.argv[5])
            #print(f"{BLUE}Using organization names: {org_names}{RESET}")
        except:
            print(f"{YELLOW}Invalid org_names parameter, skipping title slide update{RESET}")
    
    update_dashboard_slide(stats, template_path, output_path, days, org_names)