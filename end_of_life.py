import os
import sys
import asyncio
import time
import re
import requests
from bs4 import BeautifulSoup
import datetime
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from copy import deepcopy

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Colors for the EOL charts and labels
GOOD_COLOR = RGBColor(108, 184, 108)  # Green
WARNING_COLOR = RGBColor(248, 196, 71)  # Yellow/Amber
CRITICAL_COLOR = RGBColor(227, 119, 84)  # Red/Orange

# Fallback End of Life data if documentation cannot be accessed. Updated 4-4-25
# Format: {model_prefix: {'announcement': date, 'end_of_sale': date, 'end_of_support': date}}
EOL_FALLBACK_DATA = {
    # MX Devices
    "MX60": {"announcement": "Jul 10, 2015", "end_of_sale": "Oct 24, 2015", "end_of_support": "Oct 24, 2022"},
    "MX64": {"announcement": "Jan 26, 2022", "end_of_sale": "July 26, 2022", "end_of_support": "July 26, 2027"},
    "MX65": {"announcement": "Nov 20, 2018", "end_of_sale": "May 28, 2019", "end_of_support": "May 28, 2026"},
    "MX80": {"announcement": "Jan 26, 2016", "end_of_sale": "Aug 30, 2016", "end_of_support": "Aug 30, 2023"},
    "MX84": {"announcement": "Aug 10, 2021", "end_of_sale": "Oct 31, 2021", "end_of_support": "Oct 31, 2026"},
    "MX90": {"announcement": "Nov 5, 2013", "end_of_sale": "Apr 26, 2014", "end_of_support": "Apr 26, 2021"},
    "MX100": {"announcement": "Aug 10, 2021", "end_of_sale": "Feb 1, 2022", "end_of_support": "Feb 1, 2027"},
    "MX400": {"announcement": "Feb 28, 2018", "end_of_sale": "May 20, 2018", "end_of_support": "May 20, 2025"},
    "MX600": {"announcement": "Feb 28, 2018", "end_of_sale": "May 20, 2018", "end_of_support": "May 20, 2025"},
    
    # MS Devices
    "MS220": {"announcement": "Mar 16, 2017", "end_of_sale": "Jul 29, 2017", "end_of_support": "Jul 29, 2024"},
    "MS22": {"announcement": "Nov 5, 2013", "end_of_sale": "Apr 26, 2014", "end_of_support": "Apr 26, 2021"},
    "MS320": {"announcement": "Dec 8, 2016", "end_of_sale": "Mar 31, 2017", "end_of_support": "Mar 31, 2024"},
    "MS42": {"announcement": "Nov 5, 2013", "end_of_sale": "Apr 26, 2014", "end_of_support": "Apr 26, 2021"},
    "MS420": {"announcement": "Aug 1, 2016", "end_of_sale": "Oct 31, 2016", "end_of_support": "Oct 31, 2023"},
    
    # MR Devices (Access Points)
    "MR12": {"announcement": "Jul 27, 2015", "end_of_sale": "Oct 24, 2015", "end_of_support": "Oct 24, 2022"},
    "MR16": {"announcement": "Feb 27, 2014", "end_of_sale": "May 31, 2014", "end_of_support": "May 31, 2021"},
    "MR18": {"announcement": "Dec 8, 2016", "end_of_sale": "Feb 13, 2017", "end_of_support": "Mar 31, 2024"},
    "MR24": {"announcement": "Feb 27, 2014", "end_of_sale": "May 31, 2014", "end_of_support": "May 31, 2021"},
    "MR26": {"announcement": "Feb 9, 2016", "end_of_sale": "May 9, 2016", "end_of_support": "May 9, 2023"},
    "MR32": {"announcement": "Jan 18, 2017", "end_of_sale": "Apr 30, 2017", "end_of_support": "Jul 31, 2024"},
    "MR34": {"announcement": "Aug 1, 2016", "end_of_sale": "Oct 31, 2016", "end_of_support": "Oct 31, 2023"},
    "MR62": {"announcement": "Aug 15, 2017", "end_of_sale": "Nov 15, 2017", "end_of_support": "Nov 15, 2024"},
    "MR66": {"announcement": "Jun 7, 2017", "end_of_sale": "Jun 9, 2017", "end_of_support": "Jun 9, 2024"},
    "MR72": {"announcement": "Mar 7, 2017", "end_of_sale": "Apr 30, 2017", "end_of_support": "Apr 30, 2024"},
    "MR84": {"announcement": "Jan 27, 2021", "end_of_sale": "May 7, 2021", "end_of_support": "Jul 21, 2026"},
    
    # MV Devices (Cameras)
    "MV21": {"announcement": "Mar 19, 2019", "end_of_sale": "Jun 19, 2019", "end_of_support": "Jun 19, 2026"},
    "MV71": {"announcement": "Mar 19, 2019", "end_of_sale": "Jun 19, 2019", "end_of_support": "Jun 19, 2026"},
    
    # MG Devices (Cellular Gateway)
    "MG21": {"announcement": "Mar 18, 2024", "end_of_sale": "Mar 18, 2025", "end_of_support": "Sep 18, 2029"},
    
    # Z Series (Teleworker Gateway)
    "Z1": {"announcement": "April 27, 2018", "end_of_sale": "July 27, 2018", "end_of_support": "July 27, 2025"},
    "Z3": {"announcement": "Mar 4, 2024", "end_of_sale": "Sep 4, 2024", "end_of_support": "Sep 4, 2029"},
    
    # Additional models from the documentation
    "Z3C-HW-NA": {"announcement": "Mar 4, 2024", "end_of_sale": "Sep 4, 2024", "end_of_support": "Sep 4, 2029"},
    "Z3C-HW-WW": {"announcement": "Aug 11, 2023", "end_of_sale": "Feb 11, 2024", "end_of_support": "Feb 11, 2029"},
    "MX60W": {"announcement": "Jul 10, 2015", "end_of_sale": "Oct 24, 2015", "end_of_support": "Oct 24, 2022"},
    "MX64W": {"announcement": "Jan 26, 2022", "end_of_sale": "July 26, 2022", "end_of_support": "July 26, 2027"},
    "MX70": {"announcement": "Jan 18, 2012", "end_of_sale": "Mar 31, 2012", "end_of_support": "Mar 31, 2017"},
    "MX50": {"announcement": "Jul 15, 2011", "end_of_sale": "Sep 1, 2011", "end_of_support": "Sep 1, 2016"},
    "MG21E-HW-NA": {"announcement": "Mar 18, 2024", "end_of_sale": "Sep 18, 2024", "end_of_support": "Sep 18, 2029"},
    "MG21E-HW-WW": {"announcement": "Mar 18, 2024", "end_of_sale": "Jul 10, 2024", "end_of_support": "Sep 18, 2029"},
    "MR11": {"announcement": "May 29, 2012", "end_of_sale": "Aug 30, 2012", "end_of_support": "Aug 30, 2017"},
    "MR14": {"announcement": "May 29, 2012", "end_of_sale": "Aug 30, 2012", "end_of_support": "Aug 30, 2017"},
    "MR20": {"announcement": "Dec 19, 2022", "end_of_sale": "Jun 1, 2023", "end_of_support": "Jun 13, 2028"},
    "MR30H": {"announcement": "Mar 25, 2022", "end_of_sale": "May 31, 2022", "end_of_support": "Jul 26, 2027"},
    "MR33": {"announcement": "Jan 27, 2021", "end_of_sale": "Jul 14, 2022", "end_of_support": "Jul 21, 2026"},
    "MR42": {"announcement": "Jan 27, 2021", "end_of_sale": "Jul 14, 2022", "end_of_support": "Jul 21, 2026"},
    "MR42E": {"announcement": "Jan 27, 2021", "end_of_sale": "Apr 22, 2022", "end_of_support": "Jul 21, 2026"},
    "MR45": {"announcement": "Jan 27, 2021", "end_of_sale": "Jul 21, 2021", "end_of_support": "Jul 21, 2026"},
    "MR52": {"announcement": "Jan 27, 2021", "end_of_sale": "Apr 7, 2022", "end_of_support": "Jul 21, 2026"},
    "MR53": {"announcement": "Jan 27, 2021", "end_of_sale": "May 7, 2021", "end_of_support": "Jul 21, 2026"},
    "MR53E": {"announcement": "Jan 27, 2021", "end_of_sale": "Apr 7, 2022", "end_of_support": "Jul 21, 2026"},
    "MR55": {"announcement": "Feb 15, 2022", "end_of_sale": "Apr 7, 2022", "end_of_support": "Aug 1, 2027"},
    "MR56-HW": {"announcement": "Feb 07, 2025", "end_of_sale": "Aug 07, 2025", "end_of_support": "Aug 7, 2030"},
    "MR58": {"announcement": "Jul 26, 2012", "end_of_sale": "Oct 30, 2012", "end_of_support": "Oct 30, 2017"},
    "MR70": {"announcement": "Aug 18, 2023", "end_of_sale": "Feb 19, 2024", "end_of_support": "Feb 19, 2029"},
    "MR74": {"announcement": "Jan 27, 2021", "end_of_sale": "Jul 21, 2021", "end_of_support": "Jul 21, 2026"},
    "MS120-8FP-HW": {"announcement": "Mar 28, 2024", "end_of_sale": "Feb 20, 2025", "end_of_support": "Mar 28, 2030"},
    "MS120 FAMILY": {"announcement": "Mar 28, 2024", "end_of_sale": "Mar 28, 2025", "end_of_support": "Mar 28, 2030"},
    "MS125 FAMILY": {"announcement": "Mar 28, 2024", "end_of_sale": "Mar 28, 2025", "end_of_support": "Mar 28, 2030"},
    "MS220-8": {"announcement": "Jan 9, 2018", "end_of_sale": "Sep 21, 2018", "end_of_support": "Sep 21, 2025"},
    "MS250 FAMILY": {"announcement": "Aug 28, 2024", "end_of_sale": "Aug 8, 2025", "end_of_support": "Aug 8, 2030"},
    "MS350 FAMILY": {"announcement": "Aug 28, 2024", "end_of_sale": "Aug 8, 2025", "end_of_support": "Aug 8, 2030"},
    "MS355 FAMILY": {"announcement": "Aug 28, 2024", "end_of_sale": "Aug 8, 2025", "end_of_support": "Aug 8, 2030"},
    "MS390 FAMILY": {"announcement": "Apr 4, 2024", "end_of_sale": "Mar 28, 2025", "end_of_support": "Mar 28, 2032"},
    "MS410 FAMILY": {"announcement": "Mar 28, 2024", "end_of_sale": "Sep 28, 2024", "end_of_support": "Sep 28, 2029"},
    "MS425 FAMILY": {"announcement": "Mar 28, 2024", "end_of_sale": "Jun 24, 2024", "end_of_support": "Sep 28, 2029"},
    "vMX100": {"announcement": "Sep 29, 2020", "end_of_sale": "Dec 22, 2020", "end_of_support": "Dec 22, 2027"},
    
    # GR/GX/GS Series
    "GR10 (US)": {"announcement": "Oct 21, 2022", "end_of_sale": "Apr 24, 2023", "end_of_support": "Apr 24, 2025"},
    "GR12-HW-US": {"announcement": "Apr 29, 2024", "end_of_sale": "Jan 9, 2025", "end_of_support": "Apr 29, 2027"},
    "GR60 (US)": {"announcement": "Oct 21, 2022", "end_of_sale": "Apr 24, 2023", "end_of_support": "Apr 24, 2025"},
    "GR62-HW-US": {"announcement": "Apr 29, 2024", "end_of_sale": "Apr 29, 2025", "end_of_support": "Apr 29, 2027"},
    "GX20-HW-US": {"announcement": "Jan 22, 2024", "end_of_sale": "Jun 20, 2024", "end_of_support": "Jun 20, 2026"},
    "GX50-HW-US": {"announcement": "Apr 29, 2024", "end_of_sale": "Apr 29, 2025", "end_of_support": "Apr 29, 2027"},
    "GS110-8 (US)": {"announcement": "Feb 24, 2023", "end_of_sale": "Aug 24, 2023", "end_of_support": "Aug 24, 2025"},
    "GS110-8P (US)": {"announcement": "Feb 24, 2023", "end_of_sale": "Jun 1, 2023", "end_of_support": "Aug 24, 2025"},
    "GS110-24 (US)": {"announcement": "Feb 24, 2023", "end_of_sale": "Aug 24, 2023", "end_of_support": "Aug 24, 2025"},
    "GS110-24P (US)": {"announcement": "Feb 24, 2023", "end_of_sale": "Jun 1, 2023", "end_of_support": "Aug 24, 2025"},
    "GS110-48 (US)": {"announcement": "Feb 24, 2023", "end_of_sale": "Jun 1, 2023", "end_of_support": "Aug 24, 2025"},
    "GS110-48P (US)": {"announcement": "Feb 24, 2023", "end_of_sale": "Jun 1, 2023", "end_of_support": "Aug 24, 2025"},
}

# Last updated date - fallback value
EOL_LAST_UPDATED = "April 4th, 2025"

def get_blank_layout(prs):
    """Find the most suitable blank layout in the presentation."""
    # Try to find by name first
    for layout in prs.slide_layouts:
        try:
            if hasattr(layout, 'name') and layout.name and layout.name.lower() in ['blank', 'blank slide', 'empty']:
                #print(f"{BLUE}Found blank layout by name: '{layout.name}'{RESET}")
                return layout
        except:
            pass
    
    # Look for layout with fewest placeholders
    least_placeholders = float('inf')
    best_layout = None
    
    for layout in prs.slide_layouts:
        placeholder_count = 0
        try:
            # Count placeholders
            for shape in layout.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    placeholder_count += 1
            
            if placeholder_count < least_placeholders:
                least_placeholders = placeholder_count
                best_layout = layout
        except:
            pass
    
    if best_layout:
        # print(f"{BLUE}Found layout with minimum {least_placeholders} placeholders{RESET}")
        return best_layout
    
    # Fallback to the last layout (often blank) or the standard blank (index 6)
    if len(prs.slide_layouts) > 6:
        # print(f"{YELLOW}Using default blank layout at index 6{RESET}")
        return prs.slide_layouts[6]
    else:
        # print(f"{YELLOW}Using last layout as blank layout{RESET}")
        return prs.slide_layouts[-1]

def create_clean_slide(prs, layout):
    """Create a new slide with only branding/logos kept."""
    # Create new slide
    new_slide = prs.slides.add_slide(layout)
    
    # Clean slide
    clean_slide(new_slide)
    
    return new_slide

def clean_slide(slide):
    """Remove everything from a slide except branding elements at the bottom."""
    # Identify shapes to remove (everything except bottom branding)
    shapes_to_remove = []
    for shape in slide.shapes:
        # Preserve shapes that appear to be logos/branding (bottom of slide)
        if shape.top >= Inches(6.5):
            continue
        # Keep bottom text
        if hasattr(shape, "text_frame") and shape.top >= Inches(6.5):
            continue
        # Remove everything else
        shapes_to_remove.append(shape)
    
    # Remove the identified shapes
    for shape in shapes_to_remove:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
        except Exception as e:
            # print(f"{YELLOW}Could not remove shape: {e}{RESET}")
            pass
    # Check for any remaining connectors at the bottom that might be green lines
    bottom_connectors = []
    for shape in slide.shapes:
        # If it's a connector/line at the bottom part of the slide
        if hasattr(shape, 'shape_type') and shape.shape_type == 6 and shape.top > Inches(5.0):
            bottom_connectors.append(shape)
    
    # Remove any bottom connectors
    for shape in bottom_connectors:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
        except Exception as e:
            print(f"{YELLOW}Could not remove bottom connector: {e}{RESET}")

def add_slide_content(slide, title, models, last_updated_date, is_from_doc):
    """Add standard content to a slide."""
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.27), Inches(9), Inches(0.8))
    title_p = title_shape.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    
    # Try to set font to Inter if available
    try:
        title_p.font.name = "Inter"
    except:
        # If Inter font isn't available, fall back to Arial or system default
        try:
            title_p.font.name = "Arial"
        except:
            pass  # Use system default if nothing else works
            
    title_p.alignment = PP_ALIGN.LEFT
    
    # Add horizontal line under the title
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(11.0), Inches(1.2))
    line.line.color.rgb = RGBColor(39, 110, 55)  # Dark green
    line.line.width = Pt(2)
    
    # Add last updated date with data source indicator
    update_text = f"EOL information last updated {last_updated_date}"
    if not is_from_doc:
        update_text += " (using fallback data)"
        
    update_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
    update_tf = update_box.text_frame
    update_p = update_tf.add_paragraph()
    update_p.text = update_text
    update_p.font.size = Pt(10)
    update_p.font.italic = True
    
    # Add table with model data
    add_model_table(slide, models)

def add_model_table(slide, models):
    """Add a table with model information to the slide with modernized styling."""
    if not models:
        # No models to display
        no_data_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(10.0), Inches(1.0))
        no_data_p = no_data_box.text_frame.add_paragraph()
        no_data_p.text = "No device models found in inventory."
        no_data_p.font.size = Pt(14)
        no_data_p.alignment = PP_ALIGN.CENTER
        return
    
    # Create the table
    rows = len(models) + 1  # Add header row
    cols = 5  # Model, Count, Announcement, EOS, EOSS
    
    # MODIFIED: Reduce row height and add maximum table height
    fixed_row_height = Inches(0.30)  # Reduced from 0.35 to 0.30
    
    # Maximum height to ensure the table doesn't go too low
    max_table_height = Inches(5.0)  # Leave at least 1 inch at the bottom for the logo
    
    # Calculate total height, but cap it if needed
    calculated_height = fixed_row_height * rows
    total_height = min(calculated_height, max_table_height)
    
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5),        # X position
        Inches(1.4),        # Y position
        Inches(10.0),       # Width
        total_height        # Height
    ).table
    
    # Set fixed column widths
    table.columns[0].width = Inches(2.2)    # Model
    table.columns[1].width = Inches(0.8)    # Count
    table.columns[2].width = Inches(2.3)    # Announcement
    table.columns[3].width = Inches(2.3)    # End of Sale
    table.columns[4].width = Inches(2.4)    # End of Support
    
    # Set fixed row heights for consistency
    for i in range(rows):
        table.rows[i].height = fixed_row_height
    
    # Set headers with modernized styling
    headers = ["Model", "Count", "Announcement Date", "End of Sale Date", "End of Support Date"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Set modern header background - Cisco blue
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RGBColor(0, 120, 206)
        
        # Set header text color to white for better contrast
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Add model data with modernized styling
    for i, model_data in enumerate(models, 1):
        # Apply alternating row colors for a modern look
        for j in range(5):
            cell = table.cell(i, j)
            cell_fill = cell.fill
            cell_fill.solid()
            if i % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
            else:
                cell_fill.fore_color.rgb = RGBColor(245, 247, 250)  # Very light blue-gray #F5F7FA
        
        # Model
        cell = table.cell(i, 0)
        cell.text = model_data['model']
        if model_data['status'] == 'EOL':
            # For EOL models, make cell background light red in addition to red text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light red background
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)  # Dark red text for EOL models
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Count
        cell = table.cell(i, 1)
        cell.text = str(model_data['count'])
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Announcement date
        cell = table.cell(i, 2)
        cell.text = model_data['announcement']
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # End of Sale date
        cell = table.cell(i, 3)
        cell.text = model_data['end_of_sale']
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Highlight if End of Sale date is within 1 year (assuming it's not N/A)
        if model_data['end_of_sale'] != 'N/A':
            try:
                eos_date = datetime.datetime.strptime(model_data['end_of_sale'], "%b %d, %Y")
                current_date = datetime.datetime.now()
                days_to_eos = (eos_date - current_date).days
                
                if days_to_eos <= 365 and days_to_eos > 0:  # Within a year
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 245, 225)  # Light yellow background
                    cell.text_frame.paragraphs[0].font.color.rgb = WARNING_COLOR
                    cell.text_frame.paragraphs[0].font.bold = True
                elif days_to_eos <= 0:  # Already past EOS
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light red background
                    cell.text_frame.paragraphs[0].font.color.rgb = CRITICAL_COLOR
                    cell.text_frame.paragraphs[0].font.bold = True
            except:
                pass  # If date parsing fails, don't highlight
        
        # End of Support date
        cell = table.cell(i, 4)
        cell.text = model_data['end_of_support'] if model_data['end_of_support'] is not None else 'N/A'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Highlight if End of Support date is within 1 year (assuming it's not N/A)
        if model_data['end_of_support'] is not None and model_data['end_of_support'] != 'N/A':
            try:
                eoss_date = datetime.datetime.strptime(model_data['end_of_support'], "%b %d, %Y")
                current_date = datetime.datetime.now()
                days_to_eoss = (eoss_date - current_date).days
                
                if days_to_eoss <= 365 and days_to_eoss > 0:  # Within a year
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 245, 225)  # Light yellow background
                    cell.text_frame.paragraphs[0].font.color.rgb = WARNING_COLOR
                    cell.text_frame.paragraphs[0].font.bold = True
                elif days_to_eoss <= 0:  # Already past EOSS
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light red background
                    cell.text_frame.paragraphs[0].font.color.rgb = CRITICAL_COLOR
                    cell.text_frame.paragraphs[0].font.bold = True
            except:
                pass  # If date parsing fails, don't highlight

def add_pie_chart(slide, data, title, x, y, width, height, total_devices, status_descriptions=None):

    # Add chart title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.5))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(18)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Create chart data with detailed categories including counts and percentages
    chart_data = CategoryChartData()
    
    # Define categories and values in the correct order with detailed labels
    categories = []
    values = []
    
    # Create detailed category labels that include counts and percentages
    for status in ['Good', 'Warning', 'Critical']:
        count = data.get(status, 0)
        if count > 0 or (title == "End of Support Status" and status == "Warning"):
            # Create detailed label with count and percentage
            percentage = (count / total_devices) * 100
            label = f"{status} ({count}/{total_devices}, {percentage:.1f}%)"
            categories.append(label)
            values.append(count)
    
    # Add categories and series to chart data
    chart_data.categories = categories
    chart_data.add_series('Status', values)
    
    # Add chart to slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, 
        x, 
        y + Inches(0.5), 
        width, 
        Inches(2.5),
        chart_data
    ).chart
    
    # Basic chart formatting
    chart.has_title = False
    chart.has_legend = True
    
    # Position legend on the right
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)  # Smaller font to fit detailed text
    

    plot = chart.plots[0]
    plot.has_data_labels = False  
    
    # Apply colors to chart segments
    for i, point in enumerate(plot.series[0].points):
        status = categories[i].split()[0]  # Extract status from detailed label
        if status == 'Good':
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = GOOD_COLOR
        elif status == 'Warning':
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = WARNING_COLOR
        elif status == 'Critical':
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = CRITICAL_COLOR
    
    # Define legend positioning based on chart title
    if title == "End of Sale Status":
        # End of Sale legend positions
        legend_x = Inches(0.5)    
        text_x = Inches(0.7)      
    else:  # End of Support Status
        # End of Support legend positions
        legend_x = Inches(6.0)    
        text_x = Inches(6.2)      
    
    # Updated vertical positions for all legends
    green_y = Inches(5.68)        
    yellow_y = Inches(5.99)       
    red_y = Inches(6.33)          
    
    # Add status descriptions
    if status_descriptions:
        # Always add all three statuses for consistency
        # 1. Good (Green)
        square = slide.shapes.add_shape(1, legend_x, green_y, Inches(0.2), Inches(0.2))
        square.fill.solid()
        square.fill.fore_color.rgb = GOOD_COLOR
        square.line.color.rgb = RGBColor(0, 0, 0)
        
        # Position text with proper vertical alignment
        text_box = slide.shapes.add_textbox(text_x, green_y - Inches(0.10), Inches(4.5), Inches(0.2))
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_p = text_frame.add_paragraph()
        text_p.text = f"Good: {status_descriptions['Good']}"
        text_p.font.size = Pt(10)
        text_p.alignment = PP_ALIGN.LEFT
        
        # 2. Warning (Yellow)
        square = slide.shapes.add_shape(1, legend_x, yellow_y, Inches(0.2), Inches(0.2))
        square.fill.solid()
        square.fill.fore_color.rgb = WARNING_COLOR
        square.line.color.rgb = RGBColor(0, 0, 0)
        
        # Position text with proper vertical alignment
        text_box = slide.shapes.add_textbox(text_x, yellow_y - Inches(0.10), Inches(4.5), Inches(0.2))
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_p = text_frame.add_paragraph()
        text_p.text = f"Warning: {status_descriptions['Warning']}"
        text_p.font.size = Pt(10)
        text_p.alignment = PP_ALIGN.LEFT
        
        # 3. Critical (Original Orange)
        square = slide.shapes.add_shape(1, legend_x, red_y, Inches(0.2), Inches(0.2))
        square.fill.solid()
        square.fill.fore_color.rgb = CRITICAL_COLOR  # Orange
        square.line.color.rgb = RGBColor(0, 0, 0)
        
        # Position text with proper vertical alignment
        text_box = slide.shapes.add_textbox(text_x, red_y - Inches(0.10), Inches(4.5), Inches(0.2))
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_p = text_frame.add_paragraph()
        text_p.text = f"Critical: {status_descriptions['Critical']}"
        text_p.font.size = Pt(10)
        text_p.alignment = PP_ALIGN.LEFT

def get_eol_info_from_doc():

    try:
        # Attempt to fetch documentation
        #print(f"{BLUE}Attempting to fetch EOL information from documentation{RESET}")
        
        # URLs to try - Meraki sometimes changes their documentation paths
        urls_to_try = [
            "https://documentation.meraki.com/General_Administration/Other_Topics/Meraki_End-of-Life_(EOL)_Products_and_Dates",
            "https://documentation.meraki.com/General_Administration/Other_Topics/Cisco_Meraki_End-of-Life_(EOL)_Products_and_Dates"
        ]
        
        # Add User-Agent header to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml',
            'Accept-Language': 'en-US,en;q=0.9'
        }
        
        response = None
        html_content = None
        
        # Try each URL with a retry mechanism
        for url in urls_to_try:
            retry_count = 0
            max_retries = 3
            
            while retry_count < max_retries:
                try:
                    #print(f"{BLUE}Trying URL: {url} (Attempt {retry_count + 1}/{max_retries}){RESET}")
                    # Make the request with a timeout and headers
                    response = requests.get(url, timeout=15, headers=headers)
                    
                    if response.status_code == 200:
                        html_content = response.text
                        #print(f"{GREEN}Successfully retrieved content from {url}{RESET}")
                        break
                    else:
                        print(f"{YELLOW}Got status code {response.status_code} from {url}{RESET}")
                        
                except requests.RequestException as e:
                    print(f"{YELLOW}Request failed: {e}{RESET}")
                
                retry_count += 1
                time.sleep(1)  # Brief pause before retry
            
            if html_content:
                break
        
        if not html_content:
            print(f"{RED}Failed to retrieve documentation from any URL{RESET}")
            return EOL_FALLBACK_DATA, EOL_LAST_UPDATED, False
        
        # TARGETED APPROACH: Extract date from meta tags and schema.org data
        last_updated = None
        
        # Look for meta tag with article:modified_time
        meta_match = re.search(r'<meta\s+property="article:modified_time"\s+content="([^"]+)"', html_content)
        if meta_match:
            iso_date = meta_match.group(1)
            # Convert ISO date to readable format
            try:
                import datetime
                dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                #print(f"{GREEN}Found last updated date in meta tag: '{last_updated}'{RESET}")
            except Exception as e:
                # If datetime conversion fails, use the raw date
                print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                last_updated = iso_date
        
        # If not found in meta tags, look for dateModified in JSON-LD
        if not last_updated:
            schema_match = re.search(r'"dateModified":"([^"]+)"', html_content)
            if schema_match:
                iso_date = schema_match.group(1)
                # Convert ISO date to readable format
                try:
                    import datetime
                    dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                    last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                    #print(f"{GREEN}Found last updated date in schema.org data: '{last_updated}'{RESET}")
                except Exception as e:
                    # If datetime conversion fails, use the raw date
                    print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                    last_updated = iso_date
        
        # If still not found, try the original methods
        if not last_updated:
            # Parse the HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # First, try to find the date in header/metadata elements
            date_elements = soup.select('.doc-updated, .last-updated, .page-metadata, .page-info')
            for element in date_elements:
                element_text = element.get_text()
                #print(f"{BLUE}Checking metadata element: {element_text}{RESET}")
                date_match = re.search(r'(?:last\s+updated|updated)(?:\s+on)?:?\s*(\w+\s+\d+,\s+\d{4})', element_text, re.IGNORECASE)
                if date_match:
                    last_updated = date_match.group(1)
                    #print(f"{GREEN}Found last updated date in metadata: {last_updated}{RESET}")
                    break
            
            # If not found in dedicated elements, look in the page text
            if not last_updated:
                page_text = soup.get_text()
                date_patterns = [
                    r'Last updated:?\s*(\w+\s+\d+,\s+\d{4})',
                    r'Updated:?\s*(\w+\s+\d+,\s+\d{4})',
                    r'Last modified:?\s*(\w+\s+\d+,\s+\d{4})',
                    r'\*\*\*Last updated\*\*\*\s*(\w+\s+\d+,\s+\d{4})',
                    r'updated\s+(\w+\s+\d+,\s+\d{4})'
                ]
                
                for pattern in date_patterns:
                    date_match = re.search(pattern, page_text, re.IGNORECASE)
                    if date_match:
                        last_updated = date_match.group(1)
                        #print(f"{GREEN}Found last updated date in text: {last_updated}{RESET}")
                        break
        
        # If not found, use our known fallback date
        if not last_updated:
            last_updated = EOL_LAST_UPDATED
            print(f"{YELLOW}Couldn't find last updated date, using fallback date: {last_updated}{RESET}")
        
        # Now continue with the existing code to parse the EOL data
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Initialize collection for EOL data
        eol_data = {}
        
        # Look for tables with EOL information
        #print(f"{BLUE}Scanning tables for EOL information...{RESET}")
        
        tables = soup.find_all('table')
        #print(f"{BLUE}Found {len(tables)} tables in the document{RESET}")
        
        for table_idx, table in enumerate(tables):
            # Get table content sample for debugging
            table_text = table.get_text().lower()
            sample_text = table_text[:100] + "..." if len(table_text) > 100 else table_text
            #print(f"{BLUE}Table {table_idx + 1} content sample: {sample_text}{RESET}")
            
            # Try to parse the table regardless of keywords - we'll check if it contains product info
            rows = table.find_all('tr')
            if not rows:
                #print(f"{YELLOW}Table {table_idx + 1} has no rows{RESET}")
                continue
                
            # Get headers from first row
            header_cells = rows[0].find_all(['th', 'td'])
            if not header_cells:
                #print(f"{YELLOW}Table {table_idx + 1} first row has no cells{RESET}")
                continue
                
            headers = [cell.get_text().strip().lower() for cell in header_cells]
            #print(f"{BLUE}Table {table_idx + 1} headers: {headers}{RESET}")
            
            # Find relevant columns
            product_col = None
            announcement_col = None
            end_of_sale_col = None
            end_of_support_col = None
            
            # Check if any header contains product/model terms
            for i, header in enumerate(headers):
                #print(f"Checking header {i}: '{header}'")
                
                # Product/model column
                if any(term in header for term in ['product', 'model', 'device', 'hardware']):
                    product_col = i
                    #print(f"{GREEN}Found product column at index {i}: '{header}'{RESET}")
                
                # Announcement column
                elif any(term in header for term in ['announcement', 'announced', 'eol announcement']):
                    announcement_col = i
                    #print(f"{GREEN}Found announcement column at index {i}: '{header}'{RESET}")
                
                # End of sale column
                elif any(term in header for term in ['end of sale', 'end-of-sale', 'eos', 'sales', 'order', 'last order']):
                    end_of_sale_col = i
                    #print(f"{GREEN}Found end of sale column at index {i}: '{header}'{RESET}")
                
                # End of support column
                elif any(term in header for term in ['end of support', 'end-of-support', 'eoss', 'support', 'service', 'last date']):
                    end_of_support_col = i
                    #print(f"{GREEN}Found end of support column at index {i}: '{header}'{RESET}")
            
            # If we found at least product column and one date column
            if product_col is not None and any(col is not None for col in [announcement_col, end_of_sale_col, end_of_support_col]):
                #print(f"{GREEN}Table {table_idx + 1} has product info and at least one date column{RESET}")
                
                # Process each row (skip header) 
                for row_idx, row in enumerate(rows[1:], 1):
                    cells = row.find_all(['td', 'th'])
                    
                    # Skip rows without enough cells
                    if len(cells) <= max(filter(None, [product_col, announcement_col, end_of_sale_col, end_of_support_col])):
                        #print(f"{YELLOW}Row {row_idx} has insufficient cells - skipping{RESET}")
                        continue
                    
                    # Get product text
                    product_text = cells[product_col].get_text().strip().upper()  # Convert to uppercase for consistency
                    
                    # Skip empty product cells
                    if not product_text:
                        #print(f"{YELLOW}Row {row_idx} has empty product text - skipping{RESET}")
                        continue
                        
                    #print(f"{BLUE}Row {row_idx}, Product text: {product_text}{RESET}")
                    
                    # Debug: Log the product text for MS120/MS390 entries
                    if 'MS120' in product_text or 'MS390' in product_text:
                        # print(f"{YELLOW}DEBUG PARSING: Product text: '{product_text}'{RESET}")
                        pass
                    
                    # First try to match FAMILY entries specifically (handle both regular space and &nbsp;)
                    family_match = re.search(r'((?:VMX|MR|MS|MX|MV|MG|MT|Z|CW)\d+(?:-\d+(?:[A-Z]+)?)?)\s*FAMILY', product_text.replace('\xa0', ' '), re.IGNORECASE)
                    
                    if family_match:
                        # Found a family entry, use the full "XXX FAMILY" format
                        base_model = family_match.group(1).upper()
                        model_match = type('obj', (object,), {'group': lambda self, x: f"{base_model} FAMILY"})()
                        if 'MS120' in product_text or 'MS390' in product_text:
                            # print(f"{GREEN}DEBUG PARSING: Found FAMILY entry: {base_model} FAMILY{RESET}")
                            pass
                    else:
                        # Look for individual model entries
                        model_match = re.search(r'((?:VMX|MR|MS|MX|MV|MG|MT|Z|CW)\d+(?:-\d+(?:LP|FP|UP|X)?)?|(?:VMX|MR|MS|MX|MV|MG|MT|Z|CW)\d+(?:[A-Z])?)', product_text, re.IGNORECASE)
                        if model_match and ('MS120' in product_text or 'MS390' in product_text):
                            # print(f"{BLUE}DEBUG PARSING: Found individual entry: {model_match.group(1)}{RESET}")
                            pass

                    if model_match:
                        # Convert to uppercase for consistency and normalize
                        model = model_match.group(1).upper().strip()
                        
                        # FIXED: Handle non-breaking spaces and keep "FAMILY" suffix for family entries
                        # Replace non-breaking spaces with regular spaces first
                        model = model.replace('\xa0', ' ')
                        
                        # Keep the FAMILY suffix for family entries - this is important for matching!
                        
                        # Debug output to check what we're extracting
                        if any(debug_key in product_text for debug_key in ['MS120', 'MS390', 'FAMILY']):
                            #print(f"Extracted model from doc: '{model}' from product text: '{product_text}'")
                            pass
                            
                        # Initialize data for this model
                        if model not in eol_data:
                            eol_data[model] = {
                                'announcement': None,
                                'end_of_sale': None,
                                'end_of_support': None
                            }
                        else:
                            # Check if we're trying to overwrite a FAMILY entry with a specific model
                            if 'FAMILY' in model:
                                # This is a family entry, always use it (it's more authoritative)
                                print(f"{GREEN}DEBUG: Processing FAMILY entry: {model}{RESET}")
                                pass  
                            else:
                                # This is a specific model - check if we already have a family entry for it
                                base_family = re.match(r'([A-Z]+\d+)', model)
                                if base_family:
                                    family_key = f"{base_family.group(1)} FAMILY"
                                    if family_key in eol_data:
                                        # We already have family data, skip this specific model
                                        # print(f"{BLUE}DEBUG: Skipping specific model {model} because {family_key} already exists{RESET}")
                                        continue
                                    else:
                                        # print(f"{YELLOW}DEBUG: No family entry found for {family_key}, processing individual model {model}{RESET}")
                                        pass
                        
                        # Extract dates - handle different date formats
                        if announcement_col is not None:
                            announcement_text = cells[announcement_col].get_text().strip()
                            date_match = re.search(r'(\w+\s+\d+,?\s+\d{4})', announcement_text, re.IGNORECASE)
                            if date_match:
                                # Normalize date format
                                date = date_match.group(1).replace(',', '').title()
                                date_parts = date.split()
                                if len(date_parts) == 3:
                                    month, day, year = date_parts
                                    # Ensure month is abbreviated to 3 letters
                                    month = month[:3].title()
                                    normalized_date = f"{month} {day}, {year}"
                                    eol_data[model]['announcement'] = normalized_date
                                    #print(f"{GREEN}Announcement date for {model}: {normalized_date}{RESET}")
                        
                        if end_of_sale_col is not None:
                            end_of_sale_text = cells[end_of_sale_col].get_text().strip()
                            date_match = re.search(r'(\w+\s+\d+,?\s+\d{4})', end_of_sale_text, re.IGNORECASE)
                            if date_match:
                                # Normalize date format
                                date = date_match.group(1).replace(',', '').title()
                                date_parts = date.split()
                                if len(date_parts) == 3:
                                    month, day, year = date_parts
                                    # Ensure month is abbreviated to 3 letters
                                    month = month[:3].title()
                                    normalized_date = f"{month} {day}, {year}"
                                    eol_data[model]['end_of_sale'] = normalized_date
                                    #print(f"{GREEN}End of Sale date for {model}: {normalized_date}{RESET}")
                        
                        if end_of_support_col is not None:
                            end_of_support_text = cells[end_of_support_col].get_text().strip()
                            date_match = re.search(r'(\w+\s+\d+,?\s+\d{4})', end_of_support_text, re.IGNORECASE)
                            if date_match:
                                # Normalize date format
                                date = date_match.group(1).replace(',', '').title()
                                date_parts = date.split()
                                if len(date_parts) == 3:
                                    month, day, year = date_parts
                                    # Ensure month is abbreviated to 3 letters
                                    month = month[:3].title()
                                    normalized_date = f"{month} {day}, {year}"
                                    eol_data[model]['end_of_support'] = normalized_date
                                    #print(f"{GREEN}End of Support date for {model}: {normalized_date}{RESET}")
                    else:
                        #print(f"{YELLOW}Could not extract model from product text: {product_text}{RESET}")
                        pass
            else:
                print(f"{YELLOW}Table {table_idx + 1} doesn't have required columns for EOL information{RESET}")
        
        # If we found useful data, validate and merge with fallback when appropriate
        if eol_data:
            #print(f"{GREEN}Successfully parsed EOL information from documentation{RESET}")
            #print(f"Found EOL info for {len(eol_data)} models")
            
            # IMPROVED: Validate critical family entries and supplement with fallback if needed
            validated_data = eol_data.copy()
            
            # Check for potentially problematic entries (dates in the past that might be errors)
            import datetime
            current_date = datetime.datetime.now()
            
            for model, data in eol_data.items():
                eos_date = data.get('end_of_sale')
                if eos_date:
                    try:
                        eos_datetime = datetime.datetime.strptime(eos_date, "%b %d, %Y")
                        # If EOS date is more than 2 years in the past, it might be an error
                        days_past = (current_date - eos_datetime).days
                        if days_past > 730:  # More than 2 years past
                            # Check if fallback has a more reasonable date
                            fallback_key = f"{model} FAMILY"
                            if fallback_key in EOL_FALLBACK_DATA:
                                fallback_eos = EOL_FALLBACK_DATA[fallback_key].get('end_of_sale')
                                if fallback_eos:
                                    try:
                                        fallback_datetime = datetime.datetime.strptime(fallback_eos, "%b %d, %Y")
                                        if fallback_datetime > current_date:  # Fallback is in the future
                                            #print(f"{YELLOW}Warning: {model} doc date ({eos_date}) is old, fallback has newer date ({fallback_eos}){RESET}")
                                            pass
                                    except:
                                        pass
                    except:
                        pass
            
            return validated_data, last_updated, True
        else:
            print(f"{YELLOW}Could not parse EOL information from documentation, using fallback{RESET}")
            return EOL_FALLBACK_DATA, EOL_LAST_UPDATED, False
            
    except Exception as e:
        print(f"{RED}Error fetching/parsing documentation: {e}{RESET}")
        import traceback
        traceback.print_exc()
        
        # Use fallback values
        print(f"{YELLOW}Using fallback EOL information{RESET}")
        return EOL_FALLBACK_DATA, EOL_LAST_UPDATED, False

def get_base_model(model):
    """
    Extract the base model from a full model string.
    Enhanced to handle model variants like MS220-8P.
    
    Args:
        model: The model string to process
        
    Returns:
        str or None: The base model or None if no match
    """
    if not model:
        return None
        
    # Debug important models
    debug_model = False
    if 'MS220-8P' in model or 'MX100' in model:
        debug_model = True
        #print(f"Extracting base model from: {model}")
    
    # Try different patterns from most specific to least specific
    
    # First, try to match models like MS220-8 (with number after hyphen)
    hyphen_match = re.match(r'((?:MR|MS|MX|MV|MG|MT|Z|CW)\d+-\d+)', model)
    if hyphen_match:
        base_model = hyphen_match.group(1)
        if debug_model:
            #print(f"Extracted hyphenated base model: {base_model}")
            pass
        return base_model
    
    # Next, try base models without hyphens (MX100, MR36, etc.)
    base_match = re.match(r'((?:MR|MS|MX|MV|MG|MT|Z|CW)\d+)', model)
    if base_match:
        base_model = base_match.group(1)
        if debug_model:
            #print(f"Extracted simple base model: {base_model}")
            pass
        return base_model
    
    # Try specific pattern for Catalyst models
    catalyst_match = re.match(r'(C\d+(?:\d+)?(?:-\d+)?)', model)
    if catalyst_match:
        base_model = catalyst_match.group(1)
        if debug_model:
            #print(f"Extracted Catalyst base model: {base_model}")
            pass
        return base_model
    
    # If all else fails, return None
    if debug_model:
        #print(f"Could not extract base model from: {model}")
        pass
    return None

def is_model_eol(model, eol_data):
    """
    Check if a model is in the EOL list and return its EOL info.
    Enhanced to prefer more specific model matches.
    
    Args:
        model: The device model to check
        eol_data: Dictionary of EOL information
        
    Returns:
        dict or None: EOL information for the model or None if not EOL
    """
    # Debug logging for important models
    if 'MS120' in model or 'MS390' in model or 'MS220-8P' in model or 'MX100' in model:
        # print(f"{YELLOW}DEBUG: Checking EOL status for model: {model}{RESET}")
        pass
        
    # Extract base model from full model string
    base_model = get_base_model(model)
    if not base_model:
        return None
        
    if 'MS120' in model or 'MS390' in model or 'MS220-8P' in model or 'MX100' in model:
        # print(f"{YELLOW}DEBUG: Base model extracted: {base_model}{RESET}")
        pass
    
    # 1. First, check for FAMILY entries (prioritize over specific models)
    family_key = f"{base_model} FAMILY"
    if 'MS120' in model or 'MS390' in model:
        # print(f"{YELLOW}DEBUG: Looking for family key: {family_key} in eol_data{RESET}")
        # print(f"{YELLOW}DEBUG: Available keys containing {base_model}: {[k for k in eol_data.keys() if base_model in k]}{RESET}")
        pass
    
    if family_key in eol_data:
        if 'MS120' in model or 'MS390' in model or 'MS220-8P' in model or 'MX100' in model:
            # print(f"{YELLOW}DEBUG: Found family match: {family_key}{RESET}")
            pass
        return eol_data[family_key]
    
    # 2. Check for exact match with the full model
    if model in eol_data:
        if 'MS120' in model or 'MS390' in model or 'MS220-8P' in model or 'MX100' in model:
            # print(f"{YELLOW}DEBUG: Exact match found for {model} in EOL data{RESET}")
            pass
        return eol_data[model]
    
    # 3. For models like MS220-8P, try to match with MS220-8
    # Extract the base part without the final character if it's a P, LP, FP, etc.
    model_without_suffix = None
    if re.search(r'-\d+[A-Z]+$', model):
        # For models like MS120-48LP, MS120-24P, etc.
        match = re.match(r'(.*-\d+)[A-Z]+$', model)
        if match:
            model_without_suffix = match.group(1)
            
            if 'MS220-8P' in model:
                #print(f"Checking for model without suffix: {model_without_suffix}")
                pass
                
            if model_without_suffix in eol_data:
                if 'MS220-8P' in model:
                    #print(f"Found match for {model} using {model_without_suffix}")
                    pass
                return eol_data[model_without_suffix]
    
    # 4. Check if the base model is in the EOL data
    if base_model in eol_data:
        if 'MS120' in model or 'MS390' in model or 'MS220-8P' in model or 'MX100' in model:
            # print(f"{YELLOW}DEBUG: Base model match found: {base_model}{RESET}")
            pass
        return eol_data[base_model]
    
    # 5. Find all potential matches and use improved priority logic
    potential_matches = []
    
    # First, explicitly check for FAMILY entries based on device series
    # Extract device series (e.g., MS390, MS120) from the model
    device_series_match = re.match(r'^([A-Z]+\d+)', model)
    if device_series_match:
        device_series = device_series_match.group(1)
        family_key = f"{device_series} FAMILY"
        if family_key in eol_data:
            potential_matches.append(family_key)
    
    # Then gather all other potential matches
    for eol_model in eol_data:
        # Skip if already added as family match
        if eol_model in potential_matches:
            continue
            
        # Check if the model starts with this EOL model
        if model.startswith(eol_model) or (model_without_suffix and model_without_suffix.startswith(eol_model)):
            potential_matches.append(eol_model)
        # Also check if the base model starts with this EOL model
        elif base_model and base_model.startswith(eol_model):
            potential_matches.append(eol_model)
    
    if ('MS120' in model or 'MS390' in model or 'MS220-8P' in model) and potential_matches:
        # print(f"{YELLOW}DEBUG: Potential matches for {model}: {potential_matches}{RESET}")
        pass
    
    if potential_matches:
        # IMPROVED: Separate family entries from specific model variants
        family_matches = []
        specific_matches = []
        
        for match in potential_matches:
            # FAMILY entries are those that explicitly end with "FAMILY"
            if match.endswith(' FAMILY'):
                family_matches.append(match)
            else:
                specific_matches.append(match)
        
        # IMPROVED: Prefer FAMILY entries for better EOL date consistency
        # FAMILY entries typically have the most up-to-date and relevant dates
        if family_matches:
            # Sort family matches by length (shorter first, as they're more general)
            family_matches.sort(key=len)
            best_family_match = family_matches[0]
            
            # Debug output for important models
            if any(debug_model in model for debug_model in ['MS120', 'MS390', 'MS250', 'MS350']):
                # print(f"{YELLOW}DEBUG: Using family match '{best_family_match}' for model '{model}'{RESET}")
                pass
            
            return eol_data[best_family_match]
        
        # If no family matches, use specific matches
        elif specific_matches:
            # First, try to find models with hyphen and number (like MS220-8) over just series (like MS220)
            hyphen_models = [m for m in specific_matches if '-' in m]
            if hyphen_models:
                # Sort by length (descending) to get most specific match
                hyphen_models.sort(key=len, reverse=True)
                
                # For MS220-8P, prefer MS220-8 over MS220
                if 'MS220-8P' in model:
                    ms220_8_match = next((m for m in hyphen_models if m == 'MS220-8'), None)
                    if ms220_8_match:
                        #print(f"Using specific match {ms220_8_match} for {model}")
                        return eol_data[ms220_8_match]
                
                best_match = hyphen_models[0]
                if 'MS220-8P' in model or 'MX100' in model:
                    #print(f"Using best hyphen match: {best_match}")
                    pass
                return eol_data[best_match]
            
            # If no hyphen models, sort remaining by length (descending)
            specific_matches.sort(key=len, reverse=True)
            best_match = specific_matches[0]
            
            if 'MS220-8P' in model or 'MX100' in model:
                #print(f"Using best general match: {best_match}")
                pass
            return eol_data[best_match]
    
    # No match found
    if 'MS220-8P' in model or 'MX100' in model:
        #print(f"No match found for {model} in EOL data")
        pass
    return None

def categorize_eol_status(eol_date, current_date):
    """
    Categorize the EOL status based on the date:
    - Green/Good: More than 2 years away or no date (None)
    - Yellow/Warning: Within 2 years
    - Red/Critical: Within 1 year
    
    Args:
        eol_date: String date in format "MMM DD, YYYY" or None
        current_date: datetime.datetime object of current date
        
    Returns:
        str: "Good", "Warning", or "Critical"
    """
    if not eol_date:
        return "Good"  # No EOL date means it's a newer device
    
    try:
        # Convert date string to datetime object
        eol_datetime = datetime.datetime.strptime(eol_date, "%b %d, %Y")
        
        # Calculate time difference in years
        time_diff = (eol_datetime - current_date).days / 365.25
        
        if time_diff <= 1:
            return "Critical"  # Within 1 year
        elif time_diff <= 2:
            return "Warning"   # Within 2 years
        else:
            return "Good"      # More than 2 years away
    except Exception as e:
        print(f"{YELLOW}Error parsing date {eol_date}: {e}{RESET}")
        return "Good"  # If there's any error parsing the date, assume Good

async def generate(api_client, template_path, output_path, inventory_devices=None, networks=None):
    """Generate the End of Life Products slide."""
    # Silenced to reduce terminal output
    # print(f"\n{GREEN}Generating End of Life Products slide (Slide 11)...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # If inventory_devices is provided, use it
    if not inventory_devices:
        print(f"{RED}No inventory data provided{RESET}")
        return
    
    #print(f"{BLUE}Using inventory data with {len(inventory_devices)} devices{RESET}")
    
    # Documentation URL to add to slide notes
    eol_doc_url = "https://documentation.meraki.com/General_Administration/Other_Topics/Meraki_End-of-Life_(EOL)_Products_and_Dates"
    
    # Get EOL information from documentation (or use hardcoded fallback)
    eol_data, last_updated_date, is_from_doc = get_eol_info_from_doc()
    
    # Log the source of EOL information
    if is_from_doc:
        print(f"{GREEN}Using EOL information from documentation (last updated: {last_updated_date}){RESET}")
        print(f"{BLUE}Found EOL data for {len(eol_data)} models from documentation{RESET}")
        # Log some sample entries to verify parsing
        family_entries = [k for k in eol_data.keys() if 'FAMILY' in k]
        if family_entries:
            print(f"{BLUE}FAMILY entries found: {family_entries[:5]}{RESET}")
        
        # Debug: Show specific MS120 and MS390 entries that were parsed
        ms_entries = {k: v for k, v in eol_data.items() if k.startswith('MS120') or k.startswith('MS390')}
        # if ms_entries:
        #     print(f"{YELLOW}DEBUG: MS120/MS390 entries from documentation:{RESET}")
        #     for model, dates in ms_entries.items():
        #         print(f"  {model}: {dates}")
        #     print(f"{RESET}")
            
        # Debug: Show all FAMILY entries
        family_entries = {k: v for k, v in eol_data.items() if 'FAMILY' in k}
        # if family_entries:
        #     print(f"{GREEN}DEBUG: All FAMILY entries in final data:{RESET}")
        #     for model, dates in family_entries.items():
        #         print(f"  {model}: {dates}")
        #     print(f"{RESET}")
        # else:
        #     print(f"{RED}DEBUG: NO FAMILY entries found in final data!{RESET}")
            
        # Debug: Check if we're getting None values
        none_entries = {k: v for k, v in eol_data.items() if any(val is None for val in v.values())}
        if none_entries:
            # print(f"{RED}DEBUG: Entries with None values (causing predictive lifecycle error):{RESET}")
            # for model, dates in none_entries.items():
            #     print(f"  {model}: {dates}")
            # print(f"{RESET}")
            
            # Clean up None values to prevent crashes in other modules
            # print(f"{YELLOW}Cleaning up entries with None values...{RESET}")
            for model_key in list(none_entries.keys()):
                if model_key in eol_data:
                    del eol_data[model_key]
                    # print(f"  Removed {model_key} with None values")
            # print(f"{GREEN}Cleanup complete. {len(none_entries)} entries removed.{RESET}")
    else:
        print(f"{YELLOW}Using fallback EOL information - documentation unavailable{RESET}")
    
    # Process device data
    process_start_time = time.time()
    #print(f"{PURPLE}[{time.strftime('%H:%M:%S')}] Processing device data against EOL information...{RESET}")
    
    # Get current date for calculations
    current_date = datetime.datetime.now()
    
    # Track devices by product type and EOL status
    device_types = {
        'MX': [],  # Security appliances
        'MS': [],  # Switches
        'MR': [],  # Wireless APs (including Catalyst Wireless)
        'MV': [],  # Cameras
        'MG': [],  # Cellular gateways
        'MT': [],  # IoT sensors
        'Z': [],   # Teleworker gateways
        'Other': []  # Others
    }
    
    # Note: We no longer need a separate 'CW' entry since we'll map CW models to 'MR'
    
    # Track EOL devices
    eol_devices = defaultdict(list)
    
    # Track device counts by EOL status for pie charts
    eos_status = {'Good': 0, 'Warning': 0, 'Critical': 0}
    eoss_status = {'Good': 0, 'Warning': 0, 'Critical': 0}
    all_devices = []  # Track all analyzed devices
    
    # Check each device against EOL data
    for device in inventory_devices:
        model = device.get('model', 'unknown')
        base_model = get_base_model(model)
        
        # Skip if we can't determine the base model
        if not base_model:
            continue
        
        # Add to all_devices list
        all_devices.append(device)
        
        # Determine device type
        device_type = 'Other'
        
        # Special case for CW models - treat them as MR (wireless)
        if base_model.startswith('CW'):
            device_type = 'MR'  # Map Catalyst Wireless to MR category
        else:
            # Check other device types
            for prefix in device_types.keys():
                if base_model.startswith(prefix):
                    device_type = prefix
                    break
        
        # Add to device type list
        device_types[device_type].append(device)
        
        # Check if device is EOL
        eol_info = is_model_eol(model, eol_data)
        if eol_info:
            eol_devices[base_model].append((model, eol_info))
            
            # Get EOL dates
            eos_date = eol_info.get('end_of_sale')
            eoss_date = eol_info.get('end_of_support')
            
            # Categorize based on dates
            eos_category = categorize_eol_status(eos_date, current_date)
            eoss_category = categorize_eol_status(eoss_date, current_date)
            
            # Update counters
            eos_status[eos_category] += 1
            eoss_status[eoss_category] += 1
        else:
            # No EOL info means it's a newer device, categorize as Good
            eos_status['Good'] += 1
            eoss_status['Good'] += 1
    
    # Count statistics
    total_devices = len(all_devices)
    total_eol_devices = sum(len(devices) for devices in eol_devices.values())
    eol_percentage = (total_eol_devices / total_devices * 100) if total_devices > 0 else 0
    
    # Print statistics for verification
    #print(f"{BLUE}EOL Device Statistics:{RESET}")
    #print(f"Total devices analyzed: {total_devices}")
    #print(f"Total EOL devices found: {total_eol_devices} ({eol_percentage:.1f}%)")
    
    # Print EOS status
    #print(f"{BLUE}End of Sale Status:{RESET}")
    #print(f"Good: {eos_status['Good']} devices ({eos_status['Good']/total_devices*100:.1f}%)")
    #print(f"Warning: {eos_status['Warning']} devices ({eos_status['Warning']/total_devices*100:.1f}%)")
    #print(f"Critical: {eos_status['Critical']} devices ({eos_status['Critical']/total_devices*100:.1f}%)")
    
    # Print EOSS status
    #print(f"{BLUE}End of Support Status:{RESET}")
    #print(f"Good: {eoss_status['Good']} devices ({eoss_status['Good']/total_devices*100:.1f}%)")
    #print(f"Warning: {eoss_status['Warning']} devices ({eoss_status['Warning']/total_devices*100:.1f}%)")
    #print(f"Critical: {eoss_status['Critical']} devices ({eoss_status['Critical']/total_devices*100:.1f}%)")
    
    # Group EOL devices by product type for reporting
    eol_by_type = defaultdict(list)
    for base_model, devices in eol_devices.items():
        device_type = 'Other'
        
        # Special case for CW models - treat them as MR (wireless)
        if base_model.startswith('CW'):
            device_type = 'MR'  # Map Catalyst Wireless to MR category
        else:
            # Check other device types
            for prefix in device_types.keys():
                if base_model.startswith(prefix):
                    device_type = prefix
                    break
        
        for model, eol_info in devices:
            eol_by_type[device_type].append((base_model, eol_info))
    
    process_time = time.time() - process_start_time
    # print(f"{BLUE}Device data processing completed in {process_time:.2f} seconds{RESET}")
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    #print(f"{BLUE}Updating PowerPoint with EOL data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)
        
        # Use slide 11 (index 10) - updated from slide 10 (index 9)
        # If the slide doesn't exist, add it
        if len(prs.slides) < 11:
            # Add a blank slide using our clean slide functionality
            blank_layout = get_blank_layout(prs)
            slide = create_clean_slide(prs, blank_layout)
        else:
            slide = prs.slides[10]  # Updated from index 9 to 10
            # Clean the existing slide
            clean_slide(slide)
        
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.27), Inches(9), Inches(0.8))
        title_p = title_shape.text_frame.add_paragraph()
        title_p.text = "End of Life Products"
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # Try to set font to Inter if available
        try:
            title_p.font.name = "Inter"
        except:
            # If Inter font isn't available, fall back to Arial or system default
            try:
                title_p.font.name = "Arial"
            except:
                pass  # Use system default if nothing else works
                
        title_p.alignment = PP_ALIGN.LEFT
        
        # Add horizontal line under the title
        line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(11.0), Inches(1.2))
        line.line.color.rgb = RGBColor(39, 110, 55)  # Dark green
        line.line.width = Pt(2)
        
        # Add last updated date with data source indicator
        update_text = f"EOL information last updated {last_updated_date}"
        if not is_from_doc:
            update_text += " (using fallback data)"
            
        update_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
        update_tf = update_box.text_frame
        update_p = update_tf.add_paragraph()
        update_p.text = update_text
        update_p.font.size = Pt(10)
        update_p.font.italic = True
        
        # Add summary statistics
        summary_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(10), Inches(0.4))
        summary_tf = summary_box.text_frame
        summary_p = summary_tf.add_paragraph()
        summary_p.text = f"EOL Devices: {total_eol_devices} of {total_devices} devices ({eol_percentage:.1f}%)"
        summary_p.font.size = Pt(14)
        summary_p.font.bold = True

        if eol_percentage >= 50:
            summary_p.font.color.rgb = RGBColor(192, 0, 0)  # Red for critical (≥50%)
        elif eol_percentage >= 20:
            summary_p.font.color.rgb = RGBColor(227, 119, 84)  # Orange/Amber for warning (20-50%)
        else:
            summary_p.font.color.rgb = RGBColor(108, 184, 108)  # Green for good (<20%)
        
        # Add descriptions for the status categories
        status_descriptions = {
            'Good': 'More than 2 years until date or not EOL',
            'Warning': 'Within 2 years of date',
            'Critical': 'Within 1 year of date'
        }
        
        # Add pie charts for End of Sale and End of Support
        add_pie_chart(
            slide,
            eos_status,
            "End of Sale Status",
            Inches(0.5),
            Inches(2.0),
            Inches(4.5),
            3.5,
            total_devices,
            status_descriptions
        )
        
        add_pie_chart(
            slide,
            eoss_status,
            "End of Support Status",
            Inches(6.0),
            Inches(2.0),
            Inches(4.5),
            3.5,
            total_devices,
            status_descriptions
        )
        
        # Add a section for detailed EOL device listing if space permits
        if total_eol_devices > 0 and total_eol_devices <= 10:  # Only show details if few devices
            details_title = slide.shapes.add_textbox(Inches(0.65), Inches(5.8), Inches(10), Inches(0.4))
            details_tf = details_title.text_frame
            details_p = details_tf.add_paragraph()
            details_p.text = "EOL Device Details"
            details_p.font.size = Pt(14)
            details_p.font.bold = True
            
            # Add a simple table with device details
            rows = 1 + min(9, total_eol_devices)  # header + devices (limit to 9)
            cols = 4  # Model, Announcement, EOS, EOSS
            
            if rows > 1:
                table = slide.shapes.add_table(
                    rows, cols,
                    Inches(0.65), Inches(6.2),
                    Inches(9.7), Inches(0.3 * rows)
                ).table
                
                # Set headers with modernized styling
                headers = ["Model", "Announcement", "End of Sale", "End of Support"]
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Set modern header background - Cisco blue
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 120, 206)  # Cisco blue #0078CE
                    
                    # Set header text color to white for better contrast
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
                
                # Add device data with modernized styling
                row_idx = 1
                for device_type in sorted(eol_by_type.keys()):
                    for base_model, eol_info in sorted(eol_by_type[device_type], key=lambda x: x[0]):
                        if row_idx >= rows:
                            break
                            
                        # Apply alternating row colors
                        for j in range(cols):
                            cell = table.cell(row_idx, j)
                            cell.fill.solid()
                            if row_idx % 2 == 0:
                                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                            else:
                                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)  # Very light blue-gray #F5F7FA
                        
                        # Model
                        cell = table.cell(row_idx, 0)
                        cell.text = base_model
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)  # Dark red for EOL models
                        cell.text_frame.paragraphs[0].font.bold = True
                        
                        # Customize model cell with light red background
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light red background
                        
                        # Announcement date
                        cell = table.cell(row_idx, 1)
                        cell.text = eol_info.get('announcement', 'N/A')
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        
                        # End of Sale date
                        cell = table.cell(row_idx, 2)
                        cell.text = eol_info.get('end_of_sale', 'N/A')
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        
                        # End of Support date
                        cell = table.cell(row_idx, 3)
                        cell.text = eol_info.get('end_of_support', 'N/A')
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        
                        row_idx += 1
                    
                    if row_idx >= rows:
                        break
        
        # Set compatibility mode before saving
        if hasattr(prs, 'core_properties'):
            prs.core_properties.revision = 1
            if hasattr(prs.core_properties, 'category'):
                prs.core_properties.category = 'Meraki Dashboard Report'
        
        # Add URL to slide notes (visible only to the presenter)
        if hasattr(slide, 'notes_slide'):
            notes = slide.notes_slide
        else:
            notes = slide.notes_slide = prs.notes_master.clone_master_slide()
        
        # Clear any existing notes
        for shape in notes.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()
        
        # Add the URL to the slide notes
        notes_text_frame = notes.notes_text_frame
        note_p = notes_text_frame.add_paragraph()
        note_p.text = f"Source: {eol_doc_url}"
        note_p.font.size = Pt(12)
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Updated End of Life Products slide (Slide 11){RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    # print(f"{PURPLE}End of Life Products slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def generate_detail_slide(api_client, template_path, output_path, inventory_devices=None, networks=None):
    """Generate the Model Details slide showing EOL dates for all models."""
    # print(f"\n{GREEN}Generating Model Details slide (Slide 12+)...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # If inventory_devices is provided, use it
    if not inventory_devices:
        print(f"{RED}No inventory data provided{RESET}")
        return
    
    #print(f"{BLUE}Using inventory data with {len(inventory_devices)} devices{RESET}")
    
    # Documentation URL to add to slide notes
    eol_doc_url = "https://documentation.meraki.com/General_Administration/Other_Topics/Meraki_End-of-Life_(EOL)_Products_and_Dates"
    
    # Get EOL information from documentation (or use hardcoded fallback)
    eol_data, last_updated_date, is_from_doc = get_eol_info_from_doc()
    
    # Log the source of EOL information
    if is_from_doc:
        #print(f"{GREEN}Using EOL information from documentation (last updated: {last_updated_date}){RESET}")
        pass
    else:
        print(f"{YELLOW}Using fallback EOL information - documentation unavailable{RESET}")
    
    # Process device data
    process_start_time = time.time()
    #print(f"{PURPLE}[{time.strftime('%H:%M:%S')}] Processing device data for details slide...{RESET}")
    
    # Create a dictionary to store unique models and their EOL data
    models_data = {}
    
    # Process each device to extract unique models and their EOL info
    for device in inventory_devices:
        model = device.get('model', 'unknown')
        base_model = get_base_model(model)
        
        # Skip if we can't determine the base model
        if not base_model:
            continue
            
        # Skip if we already processed this model
        if base_model in models_data:
            models_data[base_model]['count'] += 1
            continue
            
        # Check for EOL info
        eol_info = is_model_eol(model, eol_data)
        
        # Store model data
        models_data[base_model] = {
            'model': base_model,
            'count': 1,  # Initialize count to 1 since we found one
            'announcement': eol_info.get('announcement', 'N/A') if eol_info else 'N/A',
            'end_of_sale': eol_info.get('end_of_sale', 'N/A') if eol_info else 'N/A',
            'end_of_support': eol_info.get('end_of_support', 'N/A') if eol_info else 'N/A',
            'status': 'EOL' if eol_info else 'Current'
        }
    
    # Sort models by EOL status (EOL first) then by model name
    sorted_models = sorted(
        models_data.values(),
        key=lambda x: (0 if x['status'] == 'EOL' else 1, x['model'])
    )
    
    process_time = time.time() - process_start_time
    #print(f"{BLUE}Device data processing completed in {process_time:.2f} seconds{RESET}")
    #print(f"{BLUE}Found {len(sorted_models)} unique device models{RESET}")
    
    # Constants for table layout
    MODELS_PER_SLIDE = 15  # Maximum number of models per slide (excluding header row)
    TOTAL_MODELS = len(sorted_models)
    TOTAL_SLIDES_NEEDED = (TOTAL_MODELS + MODELS_PER_SLIDE - 1) // MODELS_PER_SLIDE  # Ceiling division
    
    #print(f"{BLUE}Need {TOTAL_SLIDES_NEEDED} slides to display all {TOTAL_MODELS} models{RESET}")
    
    try:
        # Load the presentation
        prs = Presentation(output_path)
        
        # Find slide 12 index (updated from 11)
        slide_index_12 = 11  # 0-based index for slide 12
        
        # Get a layout to use (preferably a blank one)
        base_layout = get_blank_layout(prs)
        
        # Create slide 12 directly without adding any intermediate slides
        if len(prs.slides) <= slide_index_12:
            #print(f"{YELLOW}Adding slide 12 as it doesn't exist yet{RESET}")
            
            # Add slide 12 directly
            slide_12 = create_clean_slide(prs, base_layout)
        else:
            # If slide 12 exists, clean it
            slide_12 = prs.slides[slide_index_12]
            clean_slide(slide_12)
        
        # Update slide 12 with first batch of models
        slide_title = "Device Models and EOL Dates"
        if TOTAL_SLIDES_NEEDED > 1:
            slide_title += f" (Page 1 of {TOTAL_SLIDES_NEEDED})"
        
        models_for_first_slide = sorted_models[:MODELS_PER_SLIDE]
        
        # Add content to slide 12
        add_slide_content(slide_12, slide_title, models_for_first_slide, last_updated_date, is_from_doc)
        
        # Create additional slides if needed
        if TOTAL_SLIDES_NEEDED > 1:
            for page_num in range(2, TOTAL_SLIDES_NEEDED + 1):
                # Create a clean new slide
                new_slide = create_clean_slide(prs, base_layout)
                
                # Get models for this slide
                start_idx = (page_num - 1) * MODELS_PER_SLIDE
                end_idx = min(start_idx + MODELS_PER_SLIDE, TOTAL_MODELS)
                models_for_this_slide = sorted_models[start_idx:end_idx]
                
                # Create title for this slide
                slide_title = f"Device Models and EOL Dates (Page {page_num} of {TOTAL_SLIDES_NEEDED})"
                
                # Add content to the new slide
                add_slide_content(new_slide, slide_title, models_for_this_slide, last_updated_date, is_from_doc)
                
                # Add URL to slide notes (visible only to presenter)
                if hasattr(new_slide, 'notes_slide'):
                    notes = new_slide.notes_slide
                else:
                    notes = new_slide.notes_slide = prs.notes_master.clone_master_slide()
                
                # Clear any existing notes
                for shape in notes.shapes:
                    if shape.has_text_frame:
                        shape.text_frame.clear()
                
                # Add the URL to the slide notes
                notes_text_frame = notes.notes_text_frame
                note_p = notes_text_frame.add_paragraph()
                note_p.text = f"Source: {eol_doc_url}"
                note_p.font.size = Pt(12)
                
                #print(f"{GREEN}Created slide for page {page_num} of {TOTAL_SLIDES_NEEDED}{RESET}")
        
        # Add URL to slide 12 notes as well
        if hasattr(slide_12, 'notes_slide'):
            notes = slide_12.notes_slide
        else:
            notes = slide_12.notes_slide = prs.notes_master.clone_master_slide()
        
        # Clear any existing notes
        for shape in notes.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()
        
        # Add the URL to the slide notes
        notes_text_frame = notes.notes_text_frame
        note_p = notes_text_frame.add_paragraph()
        note_p.text = f"Source: {eol_doc_url}"
        note_p.font.size = Pt(12)
        
        # Save the presentation
        prs.save(output_path)
        # print(f"{GREEN}Created {TOTAL_SLIDES_NEEDED} Device Models slides{RESET}")
        
    except Exception as e:
        print(f"{RED}Error creating model detail slides: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    # Calculate total execution time
    total_time = time.time() - start_time
    # print(f"{PURPLE}Model Details slides generation completed in {total_time:.2f} seconds{RESET}")
    
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing both slides
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some dummy inventory data for testing
    inventory_devices = [
        {"model": "MX64", "firmware": "15.44.0", "networkId": "N1"},
        {"model": "MX84-HW", "firmware": "15.44.0", "networkId": "N1"},
        {"model": "MS220-8P", "firmware": "14.16.1", "networkId": "N2"},
        {"model": "MS220-24P", "firmware": "14.16.1", "networkId": "N2"},
        {"model": "MS320-48LP", "firmware": "14.16.1", "networkId": "N3"},
        {"model": "MR16", "firmware": "28.5", "networkId": "N4"},
        {"model": "MR34", "firmware": "29.5", "networkId": "N4"},
        {"model": "MV21", "firmware": "5.2", "networkId": "N5"},
        {"model": "MR84", "firmware": "29.5", "networkId": "N6"},
        {"model": "MS350-48LP", "firmware": "14.32.1", "networkId": "N7"}
    ]
    
    # Generate both slides
    await generate(api_client, template_path, output_path, inventory_devices)
    await generate_detail_slide(api_client, template_path, output_path, inventory_devices)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python slide11_12.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))
