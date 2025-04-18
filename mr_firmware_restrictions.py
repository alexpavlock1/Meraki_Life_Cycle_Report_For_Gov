import os
import sys
import asyncio
import time
import re
import requests
from bs4 import BeautifulSoup
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Fallback firmware restrictions - used only if documentation cannot be accessed
# IMPORTANT: Only include models that are actually restricted (not "Current")
MR_FIRMWARE_RESTRICTIONS = {
    "26": [
        "MR12", "MR16", "MR24", "MR62", "MR66", 
        "MR18", "MR26", "MR34",
        "MR32", "MR72"
    ],
    "30": [
        "MR42", "MR52", "MR53", "MR84",
        "MR30H", "MR33", "MR74",
        "MR42E", "MR53E", "MR20", "MR70"
    ],
    "27.5": ["MR44"],
    "28.6": ["MR57"]
}

# Models known to support "Current" firmware - fallback list
# Note: CW models now need to be explicitly listed from documentation
MR_UNRESTRICTED_MODELS = [
    "MR36", "MR46", "MR45", "MR55", "MR56", "MR76", "MR86"
]

# Last updated date - fallback value
RESTRICTIONS_LAST_UPDATED = "Mar 11, 2025"

def get_firmware_restrictions_from_doc():
    """
    Fetch MR access point maximum runnable firmware from documentation.
    
    Returns:
        tuple: (firmware_restrictions dict, unrestricted_models list, last_updated string, is_from_doc bool)
    """
    try:
        # Attempt to fetch documentation
        # print(f"{BLUE}Attempting to fetch MR firmware information from documentation{RESET}")
        
        # Default fallback data
        fallback_restrictions = MR_FIRMWARE_RESTRICTIONS
        fallback_unrestricted = MR_UNRESTRICTED_MODELS
        
        # Use the correct URL for firmware information - try both pages
        doc_urls = [
            "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions",
            "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/AP_Firmware_Versions"
        ]
        
        # Add User-Agent header to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Initialize collections for firmware data
        firmware_restrictions = {}  # model -> max firmware version
        unrestricted_models = []    # models that can run "Current" firmware
        last_updated = None
        
        # Try each URL in turn
        for doc_url in doc_urls:
            try:
                # print(f"{BLUE}Checking URL: {doc_url}{RESET}")
                # Make the request with a timeout and headers
                response = requests.get(doc_url, timeout=15, headers=headers)
                response.raise_for_status()
                
                # Get the raw HTML content
                html_content = response.text
                
                # TARGETED APPROACH: Extract date from meta tags and schema.org data
                if not last_updated:
                    # Look for meta tag with article:modified_time
                    meta_match = re.search(r'<meta\s+property="article:modified_time"\s+content="([^"]+)"', html_content)
                    if meta_match:
                        iso_date = meta_match.group(1)
                        # Convert ISO date to readable format
                        try:
                            import datetime
                            dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                            last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                            # print(f"{GREEN}Found last updated date in meta tag: '{last_updated}'{RESET}")
                        except Exception as e:
                            # If datetime conversion fails, use the raw date
                            # print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                            last_updated = iso_date
                
                # If not found in meta tags, look for dateModified in JSON-LD
                if not last_updated:
                    schema_match = re.search(r'"dateModified":"([^"]+)"', html_content)
                    if schema_match:
                        iso_date = schema_match.group(1)
                        # Convert ISO date to readable format
                        try:
                            import datetime
                            dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                            last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                            # print(f"{GREEN}Found last updated date in schema.org data: '{last_updated}'{RESET}")
                        except Exception as e:
                            # If datetime conversion fails, use the raw date
                            # print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                            last_updated = iso_date
                
                # Parse the HTML with BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # If we still didn't find the date with meta tags, try the original methods
                if not last_updated:
                    # First, try to find the date in header/metadata elements
                    date_elements = soup.select('.doc-updated, .last-updated, .page-metadata')
                    for element in date_elements:
                        element_text = element.get_text()
                        date_match = re.search(r'(?:last\s+updated|updated)(?:\s+on)?:?\s*(\w+\s+\d+,\s+\d{4})', element_text, re.IGNORECASE)
                        if date_match:
                            last_updated = date_match.group(1)
                            # print(f"{GREEN}Found last updated date in metadata: {last_updated}{RESET}")
                            break
                    
                    # If not found in dedicated elements, look in the page text
                    if not last_updated:
                        page_text = soup.get_text()
                        date_patterns = [
                            r'Last updated:?\s*(\w+\s+\d+,\s+\d{4})',
                            r'Updated:?\s*(\w+\s+\d+,\s+\d{4})',
                            r'Last modified:?\s*(\w+\s+\d+,\s+\d{4})',
                            r'\*\*\*Last updated\*\*\*\s*(\w+\s+\d+,\s+\d{4})'
                        ]
                        
                        for pattern in date_patterns:
                            date_match = re.search(pattern, page_text, re.IGNORECASE)
                            if date_match:
                                last_updated = date_match.group(1)
                                # print(f"{GREEN}Found last updated date in text: {last_updated}{RESET}")
                                break
                
                # APPROACH #1: Look for tables with firmware information
                # print(f"{BLUE}Scanning tables for MR firmware information...{RESET}")
                
                tables = soup.find_all('table')
                
                for table in tables:
                    # Check if this table might contain MR firmware information
                    table_text = table.get_text().lower()
                    if ('mr' in table_text and 'firmware' in table_text) or ('access point' in table_text and 'firmware' in table_text):
                        # print(f"{BLUE}Found table with MR and firmware mentions{RESET}")
                        pass
                        
                        # Check table headers to understand structure
                        headers = []
                        rows = table.find_all('tr')
                        
                        if rows:
                            header_cells = rows[0].find_all(['th', 'td'])
                            headers = [cell.get_text().strip().lower() for cell in header_cells]
                            # print(f"{BLUE}Table headers: {headers}{RESET}")
                        
                        # Find the relevant columns
                        product_col = None
                        max_firmware_col = None
                        
                        for i, header in enumerate(headers):
                            if any(term in header for term in ['product', 'model', 'access point', 'device']):
                                product_col = i
                            if any(term in header for term in ['maximum', 'max', 'firmware restriction']):
                                max_firmware_col = i
                            elif 'firmware' in header and any(term in header for term in ['maximum', 'max', 'restriction']):
                                max_firmware_col = i
                        
                        # If we couldn't identify columns but "maximum runnable firmware" is in headers
                        if product_col is None and max_firmware_col is None:
                            if 'maximum runnable firmware' in headers:
                                max_firmware_col = headers.index('maximum runnable firmware')
                                # Look for product/model column - often the first column
                                if 'product' in headers:
                                    product_col = headers.index('product')
                                else:
                                    product_col = 0  # Assume first column is product/model
                        
                        # If we identified the needed columns, extract the data
                        if product_col is not None and max_firmware_col is not None:
                            # print(f"{GREEN}Found table with product (col {product_col}) and max firmware (col {max_firmware_col}) columns{RESET}")
                            pass
                            
                            for row in rows[1:]:
                                cells = row.find_all(['td', 'th'])
                                
                                if len(cells) > max(product_col, max_firmware_col):
                                    product_text = cells[product_col].get_text().strip()
                                    max_firmware_text = cells[max_firmware_col].get_text().strip().lower()
                                    
                                    # Extract the base model (e.g., MR33 from MR33-HW or CW9162I)
                                    mr_models = re.findall(r'(MR\d+\w*|CW\d+\w*)', product_text)
                                    
                                    for model in mr_models:
                                        # Check if this model has a firmware restriction or can run "Current"
                                        if any(term in max_firmware_text for term in ['current', 'latest', 'newest', 'unrestricted']):
                                            if model not in unrestricted_models:
                                                unrestricted_models.append(model)
                                                # print(f"{GREEN}Found unrestricted model: {model} (can run Current firmware){RESET}")
                                        else:
                                            version_match = re.search(r'(\d+(?:\.\d+)?)', max_firmware_text)
                                            if version_match:
                                                version = version_match.group(1)
                                                if version not in firmware_restrictions:
                                                    firmware_restrictions[version] = []
                                                
                                                if model not in firmware_restrictions[version]:
                                                    firmware_restrictions[version].append(model)
                                                    # print(f"{GREEN}Found restriction: {model} -> MR {version}{RESET}")
                
                # APPROACH #2: Look for AP models and firmware mentions in text
                # print(f"{BLUE}Looking for MR firmware information in page text...{RESET}")
                
                # Look for specific patterns like "MR models that can support the latest firmware"
                page_text = soup.get_text()
                
                # Pattern for unrestricted models (both MR and CW)
                unrestricted_pattern = re.compile(r'(MR\d+\w*|CW\d+\w*).*?(?:can|will).*?(?:run|support).*?(?:current|latest|newest)', re.IGNORECASE)
                for match in unrestricted_pattern.finditer(page_text):
                    model = match.group(1)
                    if model not in unrestricted_models:
                        unrestricted_models.append(model)
                        # print(f"{GREEN}Found unrestricted model (text): {model} (can run Current firmware){RESET}")
                
                # Pattern for restricted models (both MR and CW)
                restricted_pattern = re.compile(r'(MR\d+\w*|CW\d+\w*).*?(?:restricted|limited|maximum).*?(?:firmware|version).*?(\d+(?:\.\d+)?)', re.IGNORECASE)
                for match in restricted_pattern.finditer(page_text):
                    model = match.group(1)
                    version = match.group(2)
                    
                    if version not in firmware_restrictions:
                        firmware_restrictions[version] = []
                    
                    if model not in firmware_restrictions[version]:
                        firmware_restrictions[version].append(model)
                        # print(f"{GREEN}Found restriction (text): {model} -> MR {version}{RESET}")
                
                # APPROACH #3: Look for paragraphs/sections about firmware restrictions
                # print(f"{BLUE}Scanning sections for MR firmware mentions...{RESET}")
                
                # Find sections that might discuss firmware restrictions
                for header in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5']):
                    header_text = header.get_text().lower()
                    if ('mr' in header_text or 'access point' in header_text) and any(term in header_text for term in ['firmware', 'version', 'restriction']):
                        # print(f"{GREEN}Found relevant section: {header.get_text().strip()}{RESET}")
                        pass
                        
                        # Check paragraphs following this header
                        next_elem = header.find_next_sibling()
                        while next_elem and next_elem.name not in ['h1', 'h2', 'h3', 'h4', 'h5']:
                            if next_elem.name in ['p', 'ul', 'ol', 'div']:
                                elem_text = next_elem.get_text()
                                
                                # Look for MR and CW models mentioned
                                mr_models = re.findall(r'(MR\d+\w*|CW\d+\w*)', elem_text)
                                
                                for model in mr_models:
                                    # Check if this paragraph mentions current/latest firmware
                                    if re.search(r'(?:current|latest|newest).*?(?:firmware|version)', elem_text, re.IGNORECASE):
                                        if model not in unrestricted_models:
                                            unrestricted_models.append(model)
                                            # print(f"{GREEN}Found unrestricted model (section): {model}{RESET}")
                                    
                                    # Check if this paragraph mentions a specific version restriction
                                    version_match = re.search(r'(?:restricted|limited|maximum).*?(?:firmware|version).*?(\d+(?:\.\d+)?)', elem_text, re.IGNORECASE)
                                    if version_match:
                                        version = version_match.group(1)
                                        
                                        if version not in firmware_restrictions:
                                            firmware_restrictions[version] = []
                                        
                                        if model not in firmware_restrictions[version]:
                                            firmware_restrictions[version].append(model)
                                            # print(f"{GREEN}Found restriction (section): {model} -> MR {version}{RESET}")
                            
                            next_elem = next_elem.find_next_sibling()

                if firmware_restrictions and unrestricted_models:
                    break
                    
            except Exception as url_error:
                # print(f"{YELLOW}Error processing URL {doc_url}: {url_error}{RESET}")
                pass
        
        if firmware_restrictions or unrestricted_models:

            # print(f"{GREEN}Successfully parsed MR firmware information from documentation{RESET}")
            
            if firmware_restrictions:
                # print(f"Found {len(firmware_restrictions)} firmware restrictions:")
                # for version, models in sorted(firmware_restrictions.items(), key=lambda x: float(x[0]) if x[0].replace('.','').isdigit() else 0, reverse=True):
                #     print(f"  - MR {version}: {len(models)} models - {', '.join(sorted(models))}")
                pass
            
            if unrestricted_models:
                # print(f"Found {len(unrestricted_models)} unrestricted models that can run Current firmware:")
                # print(f"  - {', '.join(sorted(unrestricted_models))}")
                pass
            
            return firmware_restrictions, unrestricted_models, last_updated, True
        else:
            # print(f"{YELLOW}Could not parse firmware information from documentation, using fallback{RESET}")
            return fallback_restrictions, fallback_unrestricted, None, False
            
    except Exception as e:
        # print(f"{RED}Error fetching/parsing documentation: {e}{RESET}")
        # traceback.print_exc()
        
        # Use fallback values but no fallback date
        # print(f"{YELLOW}Using fallback firmware information{RESET}")
        return MR_FIRMWARE_RESTRICTIONS, MR_UNRESTRICTED_MODELS, None, False


# Helper function to safely check if a line has a specific RGB color
def has_rgb_color(shape, target_rgb):
    """Check if shape has a line with the target RGB color, safely handling None cases."""
    if not hasattr(shape, 'line'):
        return False
    
    # Check if shape has a line
    line = shape.line
    if line is None:
        return False
    
    # Check if line has a color
    if not hasattr(line, 'color') or line.color is None:
        return False
        
    # Check if color has rgb attribute
    if not hasattr(line.color, 'rgb') or line.color.rgb is None:
        return False
    
    # Finally check the RGB value
    return line.color.rgb == target_rgb

# Helper function to extract base model
def get_base_model(model):
    """Extract the base model (e.g., MR33 from MR33-HW)."""
    base_match = re.match(r'(MR\d+\w*|CW\d+\w*)', model)
    return base_match.group(1) if base_match else None

# Helper function to check if model has firmware restriction
def get_model_firmware_version(model, firmware_restrictions, unrestricted_models):
    """
    Determine if a model has a firmware restriction, and if so, which version.
    
    Args:
        model: The full model string (e.g., MR33-HW)
        firmware_restrictions: Dict of firmware versions and their restricted models
        unrestricted_models: List of models that can run Current firmware
        
    Returns:
        str or None: The firmware version restriction or None if unrestricted
    """
    # Extract the base model
    base_model = get_base_model(model)
    
    if not base_model:
        return None  # Not a recognizable model
    
    # Process Cisco Wireless models similar to MR models
    # No longer automatically treating CW models as unrestricted
    
    # Check if model is explicitly listed as unrestricted
    for um in unrestricted_models:
        if base_model == um or base_model.startswith(um):
            return None
    
    # Check if model has a firmware restriction
    for version, models in firmware_restrictions.items():
        for rm in models:
            if base_model == rm or base_model.startswith(rm):
                return version
    
    # If not found in either list, treat as unrestricted
    return None

async def generate(api_client, template_path, output_path, inventory_devices=None):
    """Generate the MR Firmware Restrictions slide."""
    print(f"\n{GREEN}Generating MR Firmware Restrictions slide (Slide 5)...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    # If inventory_devices is provided, use it
    # Otherwise, would need to fetch it (not implemented here)
    if not inventory_devices:
        print(f"{RED}No inventory data provided{RESET}")
        return
    
    # print(f"{BLUE}Using inventory data provided from slide 1{RESET}")
    
    # Get firmware restrictions from documentation (or use hardcoded fallback)
    firmware_restrictions, unrestricted_models, last_updated_date, is_from_doc = get_firmware_restrictions_from_doc()
    
    # Log the source of firmware restrictions
    if is_from_doc:
        # print(f"{GREEN}Using MR firmware information from documentation (last updated: {last_updated_date}){RESET}")
        pass
    else:
        # print(f"{YELLOW}Using fallback MR firmware information - documentation unavailable{RESET}")
        pass
    
    # Process MR device data
    process_start_time = time.time()
    print(f"{PURPLE}[{time.strftime('%H:%M:%S')}] Processing MR device data...{RESET}")
    
    # Filter only MR devices and Cisco Wireless devices
    mr_devices = [device for device in inventory_devices 
                 if device.get('model', '').startswith('MR') or
                    device.get('model', '').startswith('CW')]
    
    # Count devices by firmware version and model
    restricted_devices = {}
    unrestricted_devices = {}
    total_mr_devices = len(mr_devices)
    
    # Group devices by their firmware restriction and model
    for device in mr_devices:
        model = device.get('model', 'unknown')
        
        # Check if model has a firmware restriction
        restricted_version = get_model_firmware_version(model, firmware_restrictions, unrestricted_models)
        
        if restricted_version:
            # This model has a firmware restriction
            if restricted_version not in restricted_devices:
                restricted_devices[restricted_version] = {}
            
            if model not in restricted_devices[restricted_version]:
                restricted_devices[restricted_version][model] = 0
            
            restricted_devices[restricted_version][model] += 1
        else:
            # This model doesn't have a specific restriction (can run "Current")
            if model not in unrestricted_devices:
                unrestricted_devices[model] = 0
            
            unrestricted_devices[model] += 1
    
    # Print statistics for verification
    # print(f"{BLUE}MR Device Statistics:{RESET}")
    # print(f"Total MR devices found: {total_mr_devices}")
    # 
    # for version in restricted_devices:
    #     device_count = sum(restricted_devices[version].values())
    #     print(f"MR {version}: {device_count} devices")
    #     # Print device models
    #     for model, count in sorted(restricted_devices[version].items()):
    #         print(f"  - {model}: {count}")
    # 
    # unrestricted_count = sum(unrestricted_devices.values())
    # print(f"Not Firmware Restricted: {unrestricted_count} devices")
    # # Print unrestricted device models
    # for model, count in sorted(unrestricted_devices.items()):
    #     print(f"  - {model}: {count}")
    
    process_time = time.time() - process_start_time
    print(f"{BLUE}MR data processing completed in {process_time:.2f} seconds{RESET}")
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    print(f"{BLUE}Updating PowerPoint with MR data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)

        # If the slide doesn't exist, add it
        if len(prs.slides) < 5:
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides[4]
        
        # Clear existing shapes except for title
        title_shape = None
        teal_line = None
        black_line = None
        
        # Look for existing title and lines
        for shape in slide.shapes:
            # Find title
            if hasattr(shape, "text_frame") and "MR Firmware Restrictions" in shape.text_frame.text:
                title_shape = shape
                continue
                
            # Find teal horizontal line
            if has_rgb_color(shape, RGBColor(80, 200, 192)):
                teal_line = shape
                continue
                
            # Find black horizontal line
            if has_rgb_color(shape, RGBColor(0, 0, 0)):
                black_line = shape
                continue
        
        # Create title if it doesn't exist
        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "MR Firmware Restrictions"
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            # print(f"{YELLOW}Added new title: 'MR Firmware Restrictions'{RESET}")
        else:
            # print(f"{BLUE}Found existing textbox title: 'MR Firmware Restrictions'{RESET}")
            pass
        
        # Remove all shapes except title and lines
        shapes_to_preserve = [title_shape, teal_line, black_line]
        shapes_to_remove = []
        
        for shape in slide.shapes:
            if shape not in shapes_to_preserve and shape is not None:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                if hasattr(shape, '_sp'):
                    sp = shape._sp
                    sp.getparent().remove(sp)
            except Exception as e:
                # print(f"{RED}Error removing shape: {e}{RESET}")
                pass
        
        # print(f"{BLUE}Removing {len(shapes_to_remove)} shapes while preserving title and lines{RESET}")
        
        # Check if we need to add horizontal lines
        if teal_line is None:
            # print(f"{YELLOW}No teal horizontal line found, this will be added by the template{RESET}")
            pass
        
        if black_line is None:
            # print(f"{YELLOW}No black horizontal line found, this will be added by the template{RESET}")
            pass
        
        # Add last updated date with data source indicator
        update_text = f"Firmware restriction last updated {last_updated_date}"
        if not is_from_doc:
            update_text += " (using fallback data)"
            
        update_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
        update_tf = update_box.text_frame
        update_p = update_tf.add_paragraph()
        update_p.text = update_text
        update_p.font.size = Pt(10)
        update_p.font.italic = True
        
        # Add an explanatory note to clarify what "firmware restrictions" means
        explanation_text = "Note: These values represent the maximum firmware versions these devices can run."
        explanation_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(8), Inches(0.3))
        explanation_tf = explanation_box.text_frame
        explanation_p = explanation_tf.add_paragraph()
        explanation_p.text = explanation_text
        explanation_p.font.size = Pt(10)
        explanation_p.font.italic = True
        
        # Setup style settings
        header_size = Pt(16)
        item_size = Pt(12)
        
        # Current Y position for content
        current_y = Inches(1.9)
        
        # Define column positions
        left_col_x = Inches(0.5)
        right_col_x = Inches(4.75)
        
        # Left Column - Not Firmware Restricted
        if unrestricted_devices:
            header = slide.shapes.add_textbox(left_col_x, current_y, Inches(4), Inches(0.3))
            tf = header.text_frame
            p = tf.add_paragraph()
            p.text = "Not Firmware Restricted"
            p.font.size = header_size
            p.font.bold = True
            
            left_content_y = current_y + Inches(0.5)
            
            # Split devices into Cisco Wireless and Meraki AP categories
            cw_models = {model: count for model, count in unrestricted_devices.items() 
                       if model.startswith('CW')}
            mr_models = {model: count for model, count in unrestricted_devices.items() 
                       if model.startswith('MR')}
            
            # Add Cisco Wireless devices if any
            if cw_models:
                # Add Cisco Wireless header
                cw_header = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                tf = cw_header.text_frame
                p = tf.add_paragraph()
                p.text = "Cisco Wireless Access Points:"
                p.font.size = item_size
                p.font.bold = True
                
                left_content_y += Inches(0.3)
                
                # Group CW models
                cw_groups = {}
                for model, count in cw_models.items():
                    base = re.match(r'(CW\d+)', model)
                    base_model = base.group(1) if base else model
                    
                    if base_model not in cw_groups:
                        cw_groups[base_model] = []
                    
                    cw_groups[base_model].append((model, count))
                
                # Process each CW group
                for base_model, models in sorted(cw_groups.items()):
                    current_line = ""
                    all_lines = []
                    
                    for model, count in sorted(models):
                        model_text = f"{model} ({count})"
                        
                        # Check if adding this would make the line too long
                        if current_line and len(current_line) + len(model_text) + 2 > 40:
                            all_lines.append(current_line)
                            current_line = model_text
                        else:
                            if current_line:
                                current_line += ", " + model_text
                            else:
                                current_line = model_text
                    
                    if current_line:
                        all_lines.append(current_line)
                    
                    # Add each line to the slide
                    for line in all_lines:
                        item = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                        tf = item.text_frame
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = item_size
                        
                        left_content_y += Inches(0.25)
                
                left_content_y += Inches(0.2)
            
            if mr_models:
                # Add Meraki AP header
                mr_header = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                tf = mr_header.text_frame
                p = tf.add_paragraph()
                p.text = "Meraki Access Points:"
                p.font.size = item_size
                p.font.bold = True
                
                left_content_y += Inches(0.3)
                
                # Group MR models for better organization
                mr_groups = {}
                for model, count in mr_models.items():
                    base = re.match(r'(MR\d+)', model)
                    base_model = base.group(1) if base else model
                    
                    if base_model not in mr_groups:
                        mr_groups[base_model] = []
                    
                    mr_groups[base_model].append((model, count))
                
                # Process each MR group
                for base_model, models in sorted(mr_groups.items()):
                    current_line = ""
                    all_lines = []
                    
                    for model, count in sorted(models):
                        model_text = f"{model} ({count})"
                        
                        # Check if adding this would make the line too long (shorter to fit column)
                        if current_line and len(current_line) + len(model_text) + 2 > 30:
                            all_lines.append(current_line)
                            current_line = model_text
                        else:
                            if current_line:
                                current_line += ", " + model_text
                            else:
                                current_line = model_text
                    
                    if current_line:
                        all_lines.append(current_line)
                    
                    # Add each line to the slide
                    for line in all_lines:
                        item = slide.shapes.add_textbox(left_col_x + Inches(0.15), left_content_y, Inches(3.5), Inches(0.25))
                        tf = item.text_frame
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = item_size
                        
                        left_content_y += Inches(0.25)
        
        sorted_versions = sorted(restricted_devices.keys(), 
                                key=lambda x: float(x) if x.replace('.','').isdigit() else 0, 
                                reverse=True)
        
        if sorted_versions:
            right_content_y = current_y
            
            # Process each version in the right column
            for version_index, version in enumerate(sorted_versions):
                # Add firmware version header
                header = slide.shapes.add_textbox(right_col_x, right_content_y, Inches(4), Inches(0.3))
                tf = header.text_frame
                p = tf.add_paragraph()
                p.text = f"MR {version}"
                p.font.size = header_size
                p.font.bold = True
                
                right_content_y += Inches(0.4)
                
                # Add subtitle
                subtitle = slide.shapes.add_textbox(right_col_x + Inches(0.15), right_content_y, Inches(4), Inches(0.25))
                tf = subtitle.text_frame
                p = tf.add_paragraph()
                p.text = "Meraki Access Points:"
                p.font.size = item_size
                p.font.bold = True
                
                right_content_y += Inches(0.3)
                
                # Group models by base model
                model_groups = {}
                for model, count in sorted(restricted_devices[version].items()):
                    # Match both MR and CW models
                    base = re.match(r'(MR\d+|CW\d+)', model)
                    base_model = base.group(1) if base else model
                    
                    if base_model not in model_groups:
                        model_groups[base_model] = []
                    
                    model_groups[base_model].append((model, count))
                
                # Add each model group
                for base_model, models in sorted(model_groups.items()):
                    line_text = ""
                    for model, count in sorted(models):
                        if line_text:
                            line_text += ", "
                        line_text += f"{model} ({count})"
                    
                    # Add the line
                    item = slide.shapes.add_textbox(right_col_x + Inches(0.15), right_content_y, Inches(4), Inches(0.25))
                    tf = item.text_frame
                    p = tf.add_paragraph()
                    p.text = line_text
                    p.font.size = item_size
                    
                    right_content_y += Inches(0.25)
                
                right_content_y += Inches(0.3)
        
        total_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(3), Inches(0.4))
        tf = total_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Total MR Devices: {total_mr_devices}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT
        
        # Add documentation URL to slide notes (visible only to the presenter)
        documentation_url = "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions#MR"
        
        if hasattr(slide, 'notes_slide'):
            notes = slide.notes_slide
        else:
            notes = slide.notes_slide = prs.notes_master.clone_master_slide()
        
        # Clear any existing notes
        for shape in notes.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()
        
        # Add the URL to the slide notes
        notes_text_frame = notes.notes_text_frame
        note_p = notes_text_frame.add_paragraph()
        note_p.text = f"Source: {documentation_url}"
        note_p.font.size = Pt(12)
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Updated MR slide (Slide 5) with proper firmware categorization{RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        import traceback
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    print(f"{PURPLE}MR Firmware Restrictions slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Would need to fetch inventory devices in a real scenario
    # For testing, create some dummy data
    inventory_devices = [
        # MR 26 firmware models
        {"model": "MR18", "firmware": "26.8.1", "networkId": "N1"},
        {"model": "MR24", "firmware": "26.8.1", "networkId": "N2"},
        {"model": "MR32", "firmware": "26.8.1", "networkId": "N3"},
        
        # MR 30 firmware models
        {"model": "MR33", "firmware": "30.0.1", "networkId": "N4"},
        {"model": "MR74", "firmware": "30.0.1", "networkId": "N5"},
        {"model": "MR42", "firmware": "30.0.1", "networkId": "N6"},
        
        # Unrestricted models (Current)
        {"model": "MR36", "firmware": "29.0.1", "networkId": "N7"},
        {"model": "MR46", "firmware": "29.0.1", "networkId": "N8"},
        {"model": "MR55", "firmware": "29.0.1", "networkId": "N9"},
        {"model": "CW9166I", "firmware": "29.0.1", "networkId": "N10"}
    ]
    
    await generate(api_client, template_path, output_path, inventory_devices=inventory_devices)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python mr_firmware_restrictions.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path))
