import os
import sys
import asyncio
import time
import re
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# Colors for the checkboxes
CHECK_COLOR = RGBColor(108, 184, 108)  # Green
CROSS_COLOR = RGBColor(212, 212, 212)  # Light Gray

def get_blank_layout(prs):
    """Find the most suitable blank layout in the presentation."""
    # Try to find by name first
    for layout in prs.slide_layouts:
        try:
            if hasattr(layout, 'name') and layout.name and layout.name.lower() in ['blank', 'blank slide', 'empty']:
                #print(f"{BLUE}Found blank layout by name: '{layout.name}'{RESET}")
                return layout
        except:
            pass
    
    # Look for layout with fewest placeholders
    least_placeholders = float('inf')
    best_layout = None
    
    for layout in prs.slide_layouts:
        placeholder_count = 0
        try:
            # Count placeholders
            for shape in layout.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    placeholder_count += 1
            
            if placeholder_count < least_placeholders:
                least_placeholders = placeholder_count
                best_layout = layout
        except:
            pass
    
    if best_layout:
        #print(f"{BLUE}Found layout with minimum {least_placeholders} placeholders{RESET}")
        return best_layout
    
    # Fallback to the last layout (often blank) or the standard blank (index 6)
    if len(prs.slide_layouts) > 6:
        #print(f"{YELLOW}Using default blank layout at index 6{RESET}")
        return prs.slide_layouts[6]
    else:
        #print(f"{YELLOW}Using last layout as blank layout{RESET}")
        return prs.slide_layouts[-1]

def create_clean_slide(prs, layout):
    """Create a new slide with only branding/logos kept."""
    # Create new slide
    new_slide = prs.slides.add_slide(layout)
    
    # Clean it
    clean_slide(new_slide)
    
    return new_slide

def clean_slide(slide):
    """Remove everything from a slide except branding elements at the bottom."""
    # Identify shapes to remove (everything except bottom branding)
    shapes_to_remove = []
    for shape in slide.shapes:
        # Preserve shapes that appear to be logos/branding (bottom of slide)
        if shape.top >= Inches(6.5):
            continue
        # Keep bottom text
        if hasattr(shape, "text_frame") and shape.top >= Inches(6.5):
            continue
        # Remove everything else
        shapes_to_remove.append(shape)
    
    # Remove the identified shapes
    for shape in shapes_to_remove:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
        except Exception as e:
            print(f"{YELLOW}Could not remove shape: {e}{RESET}")
    
    # Check for any remaining connectors at the bottom that might be green lines
    bottom_connectors = []
    for shape in slide.shapes:
        # If it's a connector/line at the bottom part of the slide
        if hasattr(shape, 'shape_type') and shape.shape_type == 6 and shape.top > Inches(5.0):
            bottom_connectors.append(shape)
    
    # Remove any bottom connectors
    for shape in bottom_connectors:
        try:
            if hasattr(shape, '_sp'):
                sp = shape._sp
                sp.getparent().remove(sp)
        except Exception as e:
            print(f"{YELLOW}Could not remove bottom connector: {e}{RESET}")

def create_checkbox(slide, x, y, is_checked=False):
    """Create a checkbox with or without a checkmark."""
    # Create checkbox
    checkbox = slide.shapes.add_shape(1, x, y, Inches(0.2), Inches(0.2))
    checkbox.fill.solid()
    checkbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    checkbox.line.color.rgb = RGBColor(128, 128, 128)  # Gray border
    
    # If checked, add a checkmark (adjusted for smaller size)
    if is_checked:
        # Create checkmark (simplified version using a thick line)
        check = slide.shapes.add_shape(5, x + Inches(0.03), y + Inches(0.03), Inches(0.14), Inches(0.14))
        check.fill.solid()
        check.fill.fore_color.rgb = CHECK_COLOR
        check.line.fill.solid()
        check.line.fill.fore_color.rgb = CHECK_COLOR
    
    return checkbox

def determine_product_availability(inventory_devices, manual_config=None):
    """
    Determine which Meraki products are available in the organization.
    """
    products = {
        'MX': False,
        'MS': False,
        'MR': False,
        'MG': False,
        'MV': False,
        'MT': False,
        'Secure Connect': False,
        'Umbrella Secure Internet Gateway': False,
        'Thousand Eyes': False,
        'Spaces': False,
        'XDR': False
    }
    
    # Determine product availability from inventory devices
    if inventory_devices:
        products['MX'] = any(device.get('model', '').upper().startswith('MX') for device in inventory_devices)
        products['MS'] = any(device.get('model', '').upper().startswith('MS') for device in inventory_devices)
        products['MR'] = any(device.get('model', '').upper().startswith('MR') or device.get('model', '').upper().startswith('CW') for device in inventory_devices)
        products['MG'] = any(device.get('model', '').upper().startswith('MG') for device in inventory_devices)
        products['MV'] = any(device.get('model', '').upper().startswith('MV') for device in inventory_devices)
        products['MT'] = any(device.get('model', '').upper().startswith('MT') for device in inventory_devices)
    
    # Override with manual configurations if provided
    if manual_config:
        for product, value in manual_config.items():
            if product in products:
                products[product] = value
    
    return products

def create_adoption_slide(prs, products):
    """
    Create a slide showing Meraki product adoption.
    """
    # Get blank layout
    layout = get_blank_layout(prs)
    
    # Create clean slide
    slide = create_clean_slide(prs, layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.8))
    title_p = title_shape.text_frame.add_paragraph()
    title_p.text = "Meraki Product Adoption"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    
    # Try to set font to Inter if available
    try:
        title_p.font.name = "Inter"
    except:
        # If Inter font isn't available, fall back to Arial or system default
        try:
            title_p.font.name = "Arial"
        except:
            pass  # Use system default if nothing else works
            
    title_p.alignment = PP_ALIGN.LEFT
    
    # Add horizontal line under the title
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(11.0), Inches(1.2))
    line.line.color.rgb = RGBColor(39, 110, 55)  # Dark green
    line.line.width = Pt(2)
    
    # Add date of report
    date_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
    date_tf = date_box.text_frame
    date_p = date_tf.add_paragraph()
    date_p.text = f"Report generated on {datetime.datetime.now().strftime('%B %d, %Y')}"
    date_p.font.size = Pt(10)
    date_p.font.italic = True
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(10), Inches(0.4))
    desc_tf = desc_box.text_frame
    desc_p = desc_tf.add_paragraph()
    desc_p.text = "Products and services this organization has adopted and opportunities to further your Meraki footprint."
    desc_p.font.size = Pt(14)
    
    # Add note about manual verification for certain products in RED
    note_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.9), Inches(10), Inches(0.4))
    note_tf = note_box.text_frame
    note_p = note_tf.add_paragraph()
    note_p.text = "Note: Software and Services require manual verification. Optional flags can be set to denote certain product adoption when executing the program. See read me for flags."
    note_p.font.size = Pt(10)
    note_p.font.italic = True
    note_p.font.color.rgb = RGBColor(192, 0, 0)  # RED
    
    # Product descriptions
    product_descriptions = {
        'MX': "Security appliances that protect your network with integrated threat prevention, content filtering, and SD-WAN capabilities. Reduces attack surface while simplifying management and cutting MPLS costs.",
        
        'MS': "Intelligent switches that simplify deployment with virtual stacking, streamline troubleshooting with automatic alerts, and provide deep visibility across your entire network from a single dashboard.",
        
        'MR': "High-performance wireless access points with integrated location analytics, automatic RF optimization, and seamless roaming. Delivers enterprise-grade connectivity with consumer-level simplicity. Meraki access points also have integrated security to provide L3-L7 firewalling capabilities.",
        
        'MG': "Cellular gateways providing reliable backup connectivity and or primary internet for branch locations. Ensures business continuity with automatic failover. MG can extend your network where ethernet is not an option bringing connectivity through cellular to your unreachable sites.",
        
        'MV': "Smart security cameras with built-in storage and analytics that eliminate NVR infrastructure. Provides actionable business insights while enhancing physical security through motion detection and intelligent alerts.",
        
        'MT': "Environmental sensors that monitor temperature, humidity, door status, and water leaks. Protects critical infrastructure with real-time alerts, preventing costly downtime and damage to sensitive equipment.",
        
        'Secure Connect': "Zero-trust solution that secures branches and remote users without traditional VPN complexity while also providing decentralization or VPN Hubs. Provides seamless connectivity for remote workers while enforcing granular security policies based on user, device, and application context.",
        
        'Umbrella Secure Internet Gateway': "Cloud-delivered security service that blocks threats at the DNS layer before they reach your network. Stops malware, ransomware, and phishing attacks while providing secure internet access from any location.",
        
        'Thousand Eyes': "End-to-end visibility solution that identifies performance issues across your entire digital supply chain. Reduces troubleshooting time from hours to minutes by pinpointing exactly where problems occur—in your network, ISP, or cloud services.",
        
        'Spaces': "Smart workplace platform that uses existing Meraki infrastructure to provide real-time occupancy insights. Optimizes space utilization, improves workplace experience, and delivers actionable data for real estate decisions.",
        
        'XDR': "Extended detection and response security that correlates threats across your network, endpoints, and cloud. Enhances security posture by automating threat detection, investigation, and response for faster remediation of sophisticated attacks."
    }
    
    # Column headers
    hardware_title = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(4), Inches(0.3))
    hardware_title_p = hardware_title.text_frame.add_paragraph()
    hardware_title_p.text = "Hardware Products"
    hardware_title_p.font.size = Pt(18)
    hardware_title_p.font.bold = True
    
    software_title = slide.shapes.add_textbox(Inches(6.2), Inches(2.1), Inches(4), Inches(0.3))
    software_title_p = software_title.text_frame.add_paragraph()
    software_title_p.text = "Software & Services"
    software_title_p.font.size = Pt(18)
    software_title_p.font.bold = True
    
    # Create product rows with tighter spacing
    # Hardware products
    hardware_products = [
        ('MX', products.get('MX', False)),
        ('MS', products.get('MS', False)),
        ('MR', products.get('MR', False)),
        ('MG', products.get('MG', False)),
        ('MV', products.get('MV', False)),
        ('MT', products.get('MT', False))
    ]
    
    # Software products
    software_products = [
        ('Secure Connect', products.get('Secure Connect', False)),
        ('Umbrella Secure Internet Gateway', products.get('Umbrella Secure Internet Gateway', False)),
        ('Thousand Eyes', products.get('Thousand Eyes', False)),
        ('Spaces', products.get('Spaces', False)),
        ('XDR', products.get('XDR', False))
    ]
    
    # Set constants for better positioning
    TEXT_HEIGHT = 0.25      # Height of text 
    CHECKBOX_SIZE = 0.2     # Size of checkbox
    VERTICAL_ADJUSTMENT = 0.1
    
    # Different row spacing for hardware and software
    HW_ROW_SPACING = 0.6    # Spacing for hardware
    SW_ROW_SPACING = 0.9    # Spacing for software
    
    # Adjusted height for description text boxes
    DESC_HEIGHT = 0.45
    
    # Starting positions
    HW_START_Y = 2.7        # Hardware items start Y position
    SW_START_Y = 2.9        # Software items start Y position
    
    # Add hardware products (left column) - using original spacing
    for i, (product, is_deployed) in enumerate(hardware_products):
        # Calculate row position
        row_base = Inches(HW_START_Y + (i * HW_ROW_SPACING))
        
        # Calculate vertical center of row for better alignment
        row_center = row_base + Inches(TEXT_HEIGHT/2)
        
        # Position checkbox with additional adjustment for PowerPoint's rendering
        checkbox_top = row_center - Inches(CHECKBOX_SIZE/2) + Inches(VERTICAL_ADJUSTMENT)
        create_checkbox(slide, Inches(1.2), checkbox_top, is_deployed)
        
        # Add product name
        text_box = slide.shapes.add_textbox(Inches(1.5), row_base, Inches(2), Inches(TEXT_HEIGHT))
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text
        text_p = text_frame.add_paragraph()
        text_p.text = product
        text_p.font.size = Pt(14)
        text_p.alignment = PP_ALIGN.LEFT
        
        # Add status text
        status_text = "Deployed" if is_deployed else "Not Deployed"
        status_color = CHECK_COLOR if is_deployed else CROSS_COLOR
        
        status_box = slide.shapes.add_textbox(Inches(2.5), row_base, Inches(2), Inches(TEXT_HEIGHT))
        status_frame = status_box.text_frame
        status_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text
        status_p = status_frame.add_paragraph()
        status_p.text = status_text
        status_p.font.size = Pt(12)
        status_p.font.color.rgb = status_color
        
        # Add description for products that aren't deployed
        if not is_deployed and product in product_descriptions:
            # Create a text box that can fit longer text
            desc_box = slide.shapes.add_textbox(Inches(1.5), row_base + Inches(0.2), 
                                               Inches(3.7), Inches(DESC_HEIGHT))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True  # Enable word wrap
            desc_p = desc_frame.add_paragraph()
            desc_p.text = product_descriptions[product]
            desc_p.font.size = Pt(8)
            desc_p.font.italic = True
            desc_p.font.color.rgb = RGBColor(100, 100, 100)  # Dark gray
            
            try:
                desc_p.line_spacing = 0.85  # Tighter line spacing
            except:
                pass  # Not all PowerPoint versions support this
    
    # Add software products
    for i, (product, is_deployed) in enumerate(software_products):
        # Calculate row position
        row_base = Inches(SW_START_Y + (i * SW_ROW_SPACING))
        
        # Calculate vertical center of row
        row_center = row_base + Inches(TEXT_HEIGHT/2)
        
        # Position checkbox with additional adjustment for PowerPoint's rendering
        checkbox_top = row_center - Inches(CHECKBOX_SIZE/2) + Inches(VERTICAL_ADJUSTMENT)
        create_checkbox(slide, Inches(6.2), checkbox_top, is_deployed)
        
        # Add product name
        text_box = slide.shapes.add_textbox(Inches(6.5), row_base, Inches(3.5), Inches(TEXT_HEIGHT))
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text
        text_p = text_frame.add_paragraph()
        text_p.text = product
        text_p.font.size = Pt(14)
        text_p.alignment = PP_ALIGN.LEFT
        
        # Add status text
        status_text = "Deployed" if is_deployed else "Not Deployed"
        status_color = CHECK_COLOR if is_deployed else CROSS_COLOR
        
        # Position status text more to the right to accommodate longer product names
        status_box = slide.shapes.add_textbox(Inches(10.0), row_base, Inches(1.5), Inches(TEXT_HEIGHT))
        status_frame = status_box.text_frame
        status_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text
        status_p = status_frame.add_paragraph()
        status_p.text = status_text
        status_p.font.size = Pt(12)
        status_p.font.color.rgb = status_color
        
        # Add description for products that aren't deployed
        if not is_deployed and product in product_descriptions:
            # Create a text box that can fit longer text
            desc_box = slide.shapes.add_textbox(Inches(6.5), row_base + Inches(0.2), 
                                               Inches(3.7), Inches(DESC_HEIGHT))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_p = desc_frame.add_paragraph()
            desc_p.text = product_descriptions[product]
            desc_p.font.size = Pt(8)
            desc_p.font.italic = True
            desc_p.font.color.rgb = RGBColor(100, 100, 100)  # Dark gray
            
            try:
                desc_p.line_spacing = 0.85
            except:
                pass
    
    return slide

async def generate(api_client, template_path, output_path, inventory_devices=None, networks=None, manual_config=None):
    """
    Generate the Meraki Product Adoption slide.
    """
    print(f"\n{GREEN}Generating Meraki Product Adoption slide...{RESET}")
    
    # Start timer
    start_time = time.time()
    
    try:
        # Determine product availability
        products = determine_product_availability(inventory_devices, manual_config)
        
        # Print product availability
        #print(f"{BLUE}Product Availability:{RESET}")
        for product, available in products.items():
            status = f"{GREEN}Available{RESET}" if available else f"{YELLOW}Not Available{RESET}"
            #print(f"  {product}: {status}")
        
        # Load the presentation
        #print(f"{BLUE}Loading presentation from {output_path}{RESET}")
        prs = Presentation(output_path)
        
        # Get current slide count for reporting
        slide_count_before = len(prs.slides)
        
        # Create the adoption slide (will be added at the end)
        slide = create_adoption_slide(prs, products)
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Added Meraki Product Adoption slide to {output_path} (slide {slide_count_before + 1}){RESET}")
        
        # Calculate execution time
        total_time = time.time() - start_time
        print(f"{PURPLE}Generated Meraki Product Adoption slide in {total_time:.2f} seconds{RESET}")
        
        return total_time
        
    except Exception as e:
        print(f"{RED}Error generating Meraki Product Adoption slide: {e}{RESET}")
        import traceback
        traceback.print_exc()
        return 0

async def main_async(org_ids, template_path=None, output_path=None, manual_config=None):
    """
    Standalone async entry point for testing.
    """
    # Default paths
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Create some dummy inventory data for testing
    inventory_devices = [
        {"model": "MX64", "firmware": "15.44.0", "networkId": "N1"},
        {"model": "MS220-8P", "firmware": "14.16.1", "networkId": "N2"},
        {"model": "MR34", "firmware": "29.5", "networkId": "N4"},
        {"model": "MV12", "firmware": "4.13", "networkId": "N5"},
    ]
    
    # Default manual configuration (all False)
    if manual_config is None:
        manual_config = {
            'Secure Connect': False,
            'Umbrella Secure Internet Gateway': False,
            'Thousand Eyes': False,
            'Spaces': False,
            'XDR': False
        }
    
    # Generate slide
    await generate(api_client, template_path, output_path, inventory_devices, None, manual_config)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python adoption.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    # Example manual configuration (for testing)
    test_manual_config = {
        'Secure Connect': True,
        'Umbrella Secure Internet Gateway': False,
        'Thousand Eyes': True,
        'Spaces': False,
        'XDR': False
    }
    
    # Run the generation
    asyncio.run(main_async(["dummy_org"], template_path, output_path, test_manual_config))