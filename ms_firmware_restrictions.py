import os
import sys
import asyncio
import time
import re
import requests
from bs4 import BeautifulSoup
import datetime
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import traceback

# ANSI color codes for terminal output
BLUE = '\033[94m'      # General information highlights
PURPLE = '\033[95m'    # Timer information
YELLOW = '\033[93m'    # Warnings
RED = '\033[91m'       # Errors
GREEN = '\033[92m'     # Success
RESET = '\033[0m'      # Reset to default color

# MS firmware version restrictions - ONLY include models that are actually restricted
# These will only be used as fallback if documentation cannot be accessed
MS_FIRMWARE_RESTRICTIONS = {
    "15.23": ["MS420"],
    "11.29": ["MS390"],
    "11.22": ["MS125"],
    "9.36": ["MS120"],
    "9": ["MS210", "MS250"],
    "8": ["MS410"],  
    "5": ["MS220", "MS320", "MS350", "MS355", "MS425"]
}

def extract_last_updated_date(soup):
    """
    Function to extract the last updated date from Meraki documentation.
    
    Args:
        soup: BeautifulSoup object of the parsed HTML
        
    Returns:
        str or None: The extracted date or None if not found
    """
    last_updated = None
    
    # APPROACH 1: Look for the exact "Last updated" text with asterisks (Markdown style)
    raw_html = str(soup)
    markdown_patterns = [
        r'\*\*Last updated\*\*\s*(Mar\s+\d+,?\s+2025)',
        r'\*\*Last updated\*\*\s*(March\s+\d+,?\s+2025)'
    ]
    
    for pattern in markdown_patterns:
        match = re.search(pattern, raw_html)
        if match:
            last_updated = match.group(1)
            # print(f"{GREEN}Found last updated date in Markdown: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 2: Look in the metadata section at the top of the page
    # Find any element that might contain metadata information
    meta_sections = soup.select('.page-metadata, .doc-metadata, .metadata, header p, .last-updated')
    for section in meta_sections:
        section_text = section.get_text()
        # Look for variations of "Last updated" followed by a date
        date_match = re.search(r'(?:Last\s+updated|Updated)(?:\s+on)?:?\s*(\w+\s+\d+,?\s+\d{4})', section_text, re.IGNORECASE)
        if date_match:
            last_updated = date_match.group(1)
            # print(f"{GREEN}Found last updated date in metadata: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 3: Look in the first few elements after the title
    # Targets the common pattern where the date appears right after the title
    title = soup.find('h1')
    if title:
        # Check the next few siblings of the title
        next_elements = []
        current = title.next_sibling
        for _ in range(5):
            if current:
                if hasattr(current, 'get_text'):
                    next_elements.append(current)
                current = current.next_sibling
        
        parent = title.parent
        if parent:
            next_elements.extend(parent.find_all(recursive=False))
        
        # Check these elements for the date
        for elem in next_elements:
            if hasattr(elem, 'get_text'):
                elem_text = elem.get_text()
                date_match = re.search(r'(?:Last\s+updated|Updated)(?:\s+on)?:?\s*(\w+\s+\d+,?\s+\d{4})', elem_text, re.IGNORECASE)
                if date_match:
                    last_updated = date_match.group(1)
                    # print(f"{GREEN}Found last updated date near title: '{last_updated}'{RESET}")
                    return last_updated
    
    specific_date_patterns = [
        r'(Mar\s+11,?\s+2025)',
        r'(March\s+11,?\s+2025)'
    ]
    
    for pattern in specific_date_patterns:
        match = re.search(pattern, raw_html)
        if match:
            last_updated = match.group(1)
            #print(f"{GREEN}Found specific date: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 5: Look for any date in the first part of the page
    # Find all text nodes in the first part of the document
    body = soup.find('body')
    if body:
        # Get the first ~20% of the HTML content
        first_part = str(body)[:int(len(str(body))*0.2)]
        # Look for any date with Mar/March 2025
        date_matches = re.findall(r'((?:Mar|March)\s+\d+,?\s+2025)', first_part)
        if date_matches:
            last_updated = date_matches[0]
            #print(f"{GREEN}Found date in first part of page: '{last_updated}'{RESET}")
            return last_updated
    
    # APPROACH 6: Look for "Last updated" line in any text node
    for tag in soup.find_all(['p', 'div', 'span']):
        if tag.string:
            text = tag.string.strip()
            if "Last updated" in text and "2025" in text:
                date_match = re.search(r'Last updated:?\s*(\w+\s+\d+,?\s+\d{4})', text, re.IGNORECASE)
                if date_match:
                    last_updated = date_match.group(1)
                    #print(f"{GREEN}Found last updated in clean text node: '{last_updated}'{RESET}")
                    return last_updated
    
    # If we've gone through all approaches and still haven't found a date
    #print(f"{YELLOW}Could not find the last updated date using any extraction method{RESET}")
    return None

def debug_date_extraction(soup):
    """
    A debugging function to print extensive information about the page structure
    to help identify where the last updated date might be located.
    
    Args:
        soup: BeautifulSoup object of the parsed HTML
    """
    #print(f"{BLUE}DEBUGGING DATE EXTRACTION{RESET}")
    
    # 1. Check for any text containing "Last updated" or the specific date
    raw_html = str(soup)
    
    # Look for "Last updated" mentions
    last_updated_matches = re.findall(r'(.{0,20}Last updated.{0,20})', raw_html)
    #print(f"{BLUE}Found {len(last_updated_matches)} mentions of 'Last updated' in HTML:{RESET}")
    for i, match in enumerate(last_updated_matches[:5]):  # Show first 5 matches
        #print(f"  {i+1}: {match}")
        pass
    
    # Look for the specific date
    date_matches = re.findall(r'(.{0,20}Mar 11,? 2025.{0,20})', raw_html)
    #print(f"{BLUE}Found {len(date_matches)} mentions of 'Mar 11, 2025' in HTML:{RESET}")
    for i, match in enumerate(date_matches[:5]):  # Show first 5 matches
        #print(f"  {i+1}: {match}")
        pass
    
    # 2. Examine the page structure near the title
    #print(f"{BLUE}Examining page structure near title:{RESET}")
    title = soup.find('h1')
    if title:
        #print(f"  Title: {title.get_text().strip()}")
        
        # Check elements right after the title
        #print(f"{BLUE}Elements after title:{RESET}")
        current = title.next_sibling
        for i in range(5):
            if current:
                if hasattr(current, 'name'):
                    #print(f"  Next {i+1} ({current.name}): {current.get_text().strip()[:50]}...")
                    pass
                else:
                    #print(f"  Next {i+1} (NavigableString): {str(current).strip()[:50]}...")
                    pass
                current = current.next_sibling
            else:
                print(f"  Next {i+1}: None")
    
    # 3. Look for any structured metadata
    #print(f"{BLUE}Searching for metadata elements:{RESET}")
    meta_selectors = ['.page-metadata', '.doc-metadata', '.metadata', '.last-updated', '.updated']
    for selector in meta_selectors:
        elements = soup.select(selector)
        #print(f"  Selector '{selector}': {len(elements)} elements found")
        for elem in elements:
            print(f"    Content: {elem.get_text().strip()[:50]}...")
    
    # 4. Check for specific page structure that might contain the date
    #print(f"{BLUE}Looking for specific page structures:{RESET}")
    
    # Check for a div with class "content" (common in documentation sites)
    content_div = soup.find('div', class_='content')
    if content_div:
        # Look for paragraphs in the first part of the content
        first_paragraphs = content_div.find_all('p', limit=3)
        #print(f"  Found {len(first_paragraphs)} paragraphs in content div:")
        for i, p in enumerate(first_paragraphs):
            #print(f"    Paragraph {i+1}: {p.get_text().strip()[:50]}...")
            pass
    
    # 5. Try to identify any text that appears to be a date in March 2025
    #print(f"{BLUE}Looking for any March 2025 dates in text nodes:{RESET}")
    date_pattern = re.compile(r'Mar(?:ch)?\s+\d+,?\s+2025')
    date_elements = []
    
    for tag in soup.find_all(['p', 'div', 'span']):
        if tag.string and date_pattern.search(tag.string):
            date_elements.append((tag.name, tag.string.strip()))
    
    #print(f"  Found {len(date_elements)} elements with March 2025 dates:")
    for i, (tag_name, content) in enumerate(date_elements[:5]):
        #print(f"    {i+1} ({tag_name}): {content[:50]}...")
        pass
    
    # 6. Check if date might be in a footer or header
    #print(f"{BLUE}Checking header and footer:{RESET}")
    for section_name in ['header', 'footer']:
        section = soup.find(section_name)
        if section:
            section_text = section.get_text().strip()
            #print(f"  {section_name.capitalize()} contains:")
            #print(f"    {section_text[:100]}...")
            
            # Check if date pattern is in this section
            if date_pattern.search(section_text):
                #print(f"    ✓ Contains a March 2025 date!")
                pass
            else:
                print(f"    ✗ No March 2025 date found.")
        else:
            print(f"  No {section_name} element found.")

    #print(f"{BLUE}First 1000 characters of HTML:{RESET}")
    #print(raw_html[:1000])

def get_firmware_restrictions_from_doc():
    """
    Fetch MS switch maximum runnable firmware versions from documentation.
    This represents the newest/highest firmware version each model can run.
    
    Returns:
        tuple: (max_firmware_versions dict, unrestricted_models list, last_updated date string, is_from_doc boolean)
    """
    try:
        # Attempt to fetch documentation
        #print(f"{BLUE}Attempting to fetch MS firmware information from documentation{RESET}")
        
        # Default fallback data for restrictions and models only (not date)
        fallback_restrictions = MS_FIRMWARE_RESTRICTIONS
        fallback_unrestricted = ["MS130", "MS225", "MS450", "CW9"]
        
        # Use the correct URL for firmware information
        doc_url = "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions"
        
        # Add User-Agent header to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Make the request with a timeout and headers
        response = requests.get(doc_url, timeout=15, headers=headers)
        response.raise_for_status()
        
        # Get the raw HTML content
        html_content = response.text
        
        # TARGETED APPROACH: Extract date from meta tags and schema.org data
        last_updated = None
        
        # Look for meta tag with article:modified_time
        meta_match = re.search(r'<meta\s+property="article:modified_time"\s+content="([^"]+)"', html_content)
        if meta_match:
            iso_date = meta_match.group(1)
            # Convert ISO date to readable format
            try:
                import datetime
                dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                #rint(f"{GREEN}Found last updated date in meta tag: '{last_updated}'{RESET}")
            except Exception as e:
                # If datetime conversion fails, use the raw date
                print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                last_updated = iso_date
        
        # If not found in meta tags, look for dateModified in JSON-LD
        if not last_updated:
            schema_match = re.search(r'"dateModified":"([^"]+)"', html_content)
            if schema_match:
                iso_date = schema_match.group(1)
                # Convert ISO date to readable format
                try:
                    import datetime
                    dt = datetime.datetime.fromisoformat(iso_date.replace('Z', '+00:00'))
                    last_updated = dt.strftime('%b %d, %Y')  # Format as "Mar 11, 2025"
                    #print(f"{GREEN}Found last updated date in schema.org data: '{last_updated}'{RESET}")
                except Exception as e:
                    # If datetime conversion fails, use the raw date
                    print(f"{YELLOW}Error converting date: {e}, using raw date{RESET}")
                    last_updated = iso_date
                    
        # Now parse the HTML with BeautifulSoup for firmware restrictions
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Initialize collections for firmware data
        firmware_restrictions = {}  # model -> max firmware version
        unrestricted_models = []    # models that can run "Current" firmware
        
        # APPROACH #1: Look for tables with firmware information
        #print(f"{BLUE}Scanning tables for MS firmware information...{RESET}")
        
        tables = soup.find_all('table')
        
        for table in tables:
            # Check if this table might contain MS firmware information
            table_text = table.get_text().lower()
            if ('ms' in table_text and 'firmware' in table_text) or ('switch' in table_text and 'firmware' in table_text):
                #print(f"{BLUE}Found table with MS and firmware mentions{RESET}")
                
                # Check table headers to understand structure
                headers = []
                rows = table.find_all('tr')
                
                if rows:
                    header_cells = rows[0].find_all(['th', 'td'])
                    headers = [cell.get_text().strip().lower() for cell in header_cells]
                    #print(f"{BLUE}Table headers: {headers}{RESET}")
                
                # Find the relevant columns
                product_col = None
                max_firmware_col = None
                
                for i, header in enumerate(headers):
                    if any(term in header for term in ['product', 'model', 'switch', 'device']):
                        product_col = i
                    if any(term in header for term in ['maximum', 'max', 'firmware restriction']):
                        max_firmware_col = i
                    elif 'firmware' in header and any(term in header for term in ['maximum', 'max', 'restriction']):
                        max_firmware_col = i
                
                # If we couldn't identify columns but "maximum runnable firmware" is in headers
                if product_col is None and max_firmware_col is None:
                    if 'maximum runnable firmware' in headers:
                        max_firmware_col = headers.index('maximum runnable firmware')
                        # Look for product/model column - often the first column
                        if 'product' in headers:
                            product_col = headers.index('product')
                        else:
                            product_col = 0  # Assume first column is product/model
                
                # If we identified the needed columns, extract the data
                if product_col is not None and max_firmware_col is not None:
                    #print(f"{GREEN}Found table with product (col {product_col}) and max firmware (col {max_firmware_col}) columns{RESET}")
                    
                    for row in rows[1:]:  # Skip header row
                        cells = row.find_all(['td', 'th'])
                        
                        if len(cells) > max(product_col, max_firmware_col):
                            product_text = cells[product_col].get_text().strip()
                            max_firmware_text = cells[max_firmware_col].get_text().strip().lower()
                            
                            # Extract the base model (e.g., MS225 from MS225-24P)
                            ms_models = re.findall(r'(MS\d+)', product_text)
                            
                            for model in ms_models:
                                # Check if this model has a firmware restriction or can run "Current"
                                if any(term in max_firmware_text for term in ['current', 'latest', 'newest', 'unrestricted']):
                                    if model not in unrestricted_models:
                                        unrestricted_models.append(model)
                                        #print(f"{GREEN}Found unrestricted model: {model} (can run Current firmware){RESET}")
                                else:
                                    # Extract version number
                                    version_match = re.search(r'(\d+(?:\.\d+)?)', max_firmware_text)
                                    if version_match:
                                        version = version_match.group(1)
                                        if version not in firmware_restrictions:
                                            firmware_restrictions[version] = []
                                        
                                        if model not in firmware_restrictions[version]:
                                            firmware_restrictions[version].append(model)
                                            #print(f"{GREEN}Found restriction: {model} -> MS {version}{RESET}")
        
        # APPROACH #2: Look for MS models and firmware mentions outside tables
        if not firmware_restrictions and not unrestricted_models:
            #print(f"{BLUE}Looking for MS firmware information in page text...{RESET}")
            
            # Look for MS models followed by firmware info
            page_text = soup.get_text()
            model_firmware_pattern = re.compile(r'(MS\d+).*?(?:maximum|restricted to|cannot run beyond).*?(?:firmware|version).*?(?:(current|latest)|(?:MS)?\s*(\d+(?:\.\d+)?))', re.IGNORECASE)
            
            for match in model_firmware_pattern.finditer(page_text):
                model = match.group(1)  # The MS model
                is_current = match.group(2)  # "current" or "latest" if matched
                version = match.group(3)  # Version number if matched
                
                if is_current:
                    # This model can run current firmware
                    if model not in unrestricted_models:
                        unrestricted_models.append(model)
                        #print(f"{GREEN}Found unrestricted model (text): {model} (can run Current firmware){RESET}")
                elif version:
                    # This model has a firmware restriction
                    if version not in firmware_restrictions:
                        firmware_restrictions[version] = []
                    
                    if model not in firmware_restrictions[version]:
                        firmware_restrictions[version].append(model)
                        #print(f"{GREEN}Found restriction (text): {model} -> MS {version}{RESET}")
        
        if firmware_restrictions or unrestricted_models:
            # Print summary of findings
            # print(f"{GREEN}Successfully parsed MS firmware information from documentation{RESET}")
            
            if firmware_restrictions:
                # print(f"Found {len(firmware_restrictions)} firmware restrictions:")
                # for version, models in sorted(firmware_restrictions.items(), key=lambda x: float(x[0]), reverse=True):
                #     print(f"  - MS {version}: {len(models)} models - {', '.join(sorted(models))}")
                pass
            
            if unrestricted_models:
                # print(f"Found {len(unrestricted_models)} unrestricted models that can run Current firmware:")
                # print(f"  - {', '.join(sorted(unrestricted_models))}")
                pass
            
            return firmware_restrictions, unrestricted_models, last_updated, True
        else:
            # print(f"{YELLOW}Could not parse firmware information from documentation, using fallback{RESET}")
            pass
            return fallback_restrictions, fallback_unrestricted, None, False
            
    except Exception as e:
        # print(f"{RED}Error fetching/parsing documentation: {e}{RESET}")
        # traceback.print_exc()
        
        # Use fallback values but no fallback date
        # print(f"{YELLOW}Using fallback firmware information{RESET}")
        return MS_FIRMWARE_RESTRICTIONS, ["MS130", "MS225", "MS450", "CW9"], None, False

def has_rgb_color(shape, target_rgb):
    """Check if shape has a line with the target RGB color, safely handling None cases."""
    if not hasattr(shape, 'line'):
        return False
    
    # Check if shape has a line
    line = shape.line
    if line is None:
        return False
    
    # Check if line has a color
    if not hasattr(line, 'color') or line.color is None:
        return False
        
    # Check if color has rgb attribute
    if not hasattr(line.color, 'rgb') or line.color.rgb is None:
        return False
    
    # Finally check the RGB value
    return line.color.rgb == target_rgb

# Helper function to extract base model
def get_base_model(model):
    """Extract the base model (e.g., MS225 from MS225-24P)."""
    base_match = re.match(r'(MS\d+|C9300[X-]*)', model)
    return base_match.group(1) if base_match else None

def get_model_firmware_version(model, firmware_restrictions, unrestricted_models):
    """
    Determine if a model has a firmware restriction, and if so, which version.
    
    Args:
        model: The full model string (e.g., MS225-24P)
        firmware_restrictions: Dict of firmware versions and their restricted models
        unrestricted_models: List of models that can run Current firmware
        
    Returns:
        str or None: The firmware version restriction or None if unrestricted
    """
    # Extract the base model
    base_model = get_base_model(model)
    
    if not base_model:
        return None  # Not a recognizable model
    
    # For Catalyst models, treat them as unrestricted
    if base_model.startswith('C9300'):
        return None
    
    # Check if model is explicitly listed as unrestricted
    for um in unrestricted_models:
        if base_model == um or base_model.startswith(um):
            return None
    
    # Check if model has a firmware restriction
    for version, models in firmware_restrictions.items():
        for rm in models:
            if base_model == rm or base_model.startswith(rm):
                return version
    
    # If not found in either list, treat as unrestricted
    return None

async def generate(api_client, template_path, output_path, inventory_devices=None):
    """Generate the MS Firmware Restrictions slide."""
    print(f"\n{GREEN}Generating MS Firmware Restrictions slide (Slide 4)...{RESET}")
    
    start_time = time.time()
    
    # If inventory_devices is provided, use it
    # Otherwise, would need to fetch it (not implemented here)
    if not inventory_devices:
        print(f"{RED}No inventory data provided{RESET}")
        return
    
    #print(f"{BLUE}Using inventory data provided from slide 1{RESET}")
    
    # Get firmware restrictions from documentation (or use hardcoded fallback)
    firmware_restrictions, unrestricted_models, last_updated_date, is_from_doc = get_firmware_restrictions_from_doc()
    
    # Log the source of firmware restrictions
    if is_from_doc:
        if last_updated_date:
            #print(f"{GREEN}Using MS firmware information from documentation (last updated: {last_updated_date}){RESET}")
            pass
        else:
            print(f"{GREEN}Using MS firmware information from documentation (no date found){RESET}")
    else:
        print(f"{YELLOW}Using fallback MS firmware information - documentation unavailable{RESET}")
    
    # Process MS device data
    process_start_time = time.time()
    #print(f"{PURPLE}[{time.strftime('%H:%M:%S')}] Processing MS device data...{RESET}")
    
    # Filter only MS devices and Catalyst 9300 devices
    ms_devices = [device for device in inventory_devices 
                 if device.get('model', '').startswith('MS') or 
                   device.get('model', '').startswith('C9300')]
    
    # Display the firmware restrictions data for verification
    #print(f"{BLUE}Firmware restrictions data:{RESET}")
    for version, models in firmware_restrictions.items():
        #print(f"  - MS {version}: {', '.join(models)}")
        pass
    
    #if unrestricted_models:
        #print(f"{BLUE}Unrestricted models:{RESET}")
        #print(f"  - {', '.join(unrestricted_models)}")
    
    # Count devices by firmware version and model
    restricted_devices = {}
    unrestricted_devices = {}
    total_ms_devices = len(ms_devices)
    
    # Group devices by their firmware restriction and model
    for device in ms_devices:
        model = device.get('model', 'unknown')
        
        # Get the restricted firmware version for this model
        restricted_version = get_model_firmware_version(model, firmware_restrictions, unrestricted_models)
        
        if restricted_version:
            # This model has a firmware restriction
            if restricted_version not in restricted_devices:
                restricted_devices[restricted_version] = {}
            
            if model not in restricted_devices[restricted_version]:
                restricted_devices[restricted_version][model] = 0
            
            restricted_devices[restricted_version][model] += 1
        else:
            # This model doesn't have a specific restriction (is "Current")
            if model not in unrestricted_devices:
                unrestricted_devices[model] = 0
            
            unrestricted_devices[model] += 1
    
    #print(f"{BLUE}MS Device Statistics:{RESET}")
    #print(f"Total MS devices found: {total_ms_devices}")
    
    for version in restricted_devices:
        device_count = sum(restricted_devices[version].values())
        #print(f"MS {version}: {device_count} devices")
        for model, count in sorted(restricted_devices[version].items()):
            #print(f"  - {model}: {count}")
            pass
    
    unrestricted_count = sum(unrestricted_devices.values())
    #print(f"Not Firmware Restricted: {unrestricted_count} devices")
    # Print unrestricted device models
    for model, count in sorted(unrestricted_devices.items()):
        #print(f"  - {model}: {count}")
        pass
    
    process_time = time.time() - process_start_time
    print(f"{BLUE}MS data processing completed in {process_time:.2f} seconds{RESET}")
    
    # Update PowerPoint presentation
    ppt_start_time = time.time()
    print(f"{BLUE}Updating PowerPoint with MS data...{RESET}")
    
    # Load the presentation
    try:
        prs = Presentation(output_path)
        
        # If the slide doesn't exist, add it
        if len(prs.slides) < 4:
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides[3]
        
        # Clear existing shapes except for title
        title_shape = None
        teal_line = None
        black_line = None
        
        # Look for existing title and lines
        for shape in slide.shapes:
            # Find title
            if hasattr(shape, "text_frame") and "MS Firmware Restrictions" in shape.text_frame.text:
                title_shape = shape
                continue
                
            # Find teal horizontal line
            if has_rgb_color(shape, RGBColor(80, 200, 192)):
                teal_line = shape
                continue
                
            # Find black horizontal line
            if has_rgb_color(shape, RGBColor(0, 0, 0)):
                black_line = shape
                continue
        
        # Create title if it doesn't exist
        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_p = title_shape.text_frame.add_paragraph()
            title_p.text = "MS Firmware Restrictions"
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            # print(f"{YELLOW}Added new title: 'MS Firmware Restrictions'{RESET}")
        else:
            # print(f"{BLUE}Found existing textbox title: 'MS Firmware Restrictions'{RESET}")
            pass
        
        # Remove all shapes except title and lines
        shapes_to_preserve = [title_shape, teal_line, black_line]
        shapes_to_remove = []
        
        for shape in slide.shapes:
            if shape not in shapes_to_preserve and shape is not None:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                if hasattr(shape, '_sp'):
                    sp = shape._sp
                    sp.getparent().remove(sp)
            except Exception as e:
                # print(f"{RED}Error removing shape: {e}{RESET}")
                pass
        
        # print(f"{BLUE}Removing {len(shapes_to_remove)} shapes while preserving title and line{RESET}")
        
        # Check if we need to add horizontal lines
        if teal_line is None:
            # print(f"{YELLOW}No teal horizontal line found, this will be added by the template{RESET}")
            pass
        
        if black_line is None:
            # print(f"{YELLOW}No black horizontal line found, this will be added by the template{RESET}")
            pass
        
        # Add last updated date only if we have a date
        if last_updated_date:
            update_text = f"Firmware restriction last updated {last_updated_date}"
            update_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(5), Inches(0.3))
            update_tf = update_box.text_frame
            update_p = update_tf.add_paragraph()
            update_p.text = update_text
            update_p.font.size = Pt(10)
            update_p.font.italic = True
        
        # Add an explanatory note to clarify what "firmware restrictions" means
        explanation_text = "Note: These values represent the maximum firmware versions these devices can run."
        explanation_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(8), Inches(0.3))
        explanation_tf = explanation_box.text_frame
        explanation_p = explanation_tf.add_paragraph()
        explanation_p.text = explanation_text
        explanation_p.font.size = Pt(10)
        explanation_p.font.italic = True
        
        # Define the two-column layout
        left_col_x = Inches(0.65)
        right_col_x = Inches(5.0)
        start_y = Inches(1.9)
        
        # Style settings
        header_size = Pt(16)
        item_size = Pt(12)
        
        # LEFT COLUMN: "Not Firmware Restricted" section
        if unrestricted_devices:
            left_y = start_y
            
            # Add header for unrestricted models
            header = slide.shapes.add_textbox(left_col_x - Inches(0.15), left_y, Inches(4), Inches(0.3))
            tf = header.text_frame
            p = tf.add_paragraph()
            p.text = "Not Firmware Restricted"
            p.font.size = header_size
            p.font.bold = True
            
            left_y += Inches(0.4)
            
            # Handle the Catalyst models separately if they exist
            catalyst_models = {model: count for model, count in unrestricted_devices.items() 
                              if model.startswith('C9300')}
            
            if catalyst_models:
                # Add Catalyst models header
                item = slide.shapes.add_textbox(left_col_x, left_y, Inches(4), Inches(0.25))
                tf = item.text_frame
                p = tf.add_paragraph()
                p.text = "Catalyst Models:"
                p.font.size = item_size
                p.font.bold = True
                
                left_y += Inches(0.3)
                
                # Organize Catalyst models into groups
                catalyst_lines = []
                current_line = ""
                
                for model, count in sorted(catalyst_models.items()):
                    model_text = f"{model} ({count})"
                    if len(current_line) + len(model_text) + 2 <= 40:
                        if current_line:
                            current_line += ", " + model_text
                        else:
                            current_line = model_text
                    else:
                        catalyst_lines.append(current_line)
                        current_line = model_text
                
                if current_line:
                    catalyst_lines.append(current_line)
                
                # Add each Catalyst line
                for line in catalyst_lines:
                    item = slide.shapes.add_textbox(left_col_x, left_y, Inches(4.0), Inches(0.25))
                    tf = item.text_frame
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = item_size
                    
                    left_y += Inches(0.25)
                
                left_y += Inches(0.15)  # Add some spacing after catalyst models
            
            # Handle all MS models
            ms_models = {model: count for model, count in unrestricted_devices.items() 
                        if model.startswith('MS')}
            
            if ms_models:
                # Add Meraki switches header
                item = slide.shapes.add_textbox(left_col_x, left_y, Inches(4), Inches(0.25))
                tf = item.text_frame
                p = tf.add_paragraph()
                p.text = "Meraki Switches:"
                p.font.size = item_size
                p.font.bold = True
                
                left_y += Inches(0.3)
                
                # Group MS models by base model
                ms_groups = {}
                for model, count in ms_models.items():
                    base = re.match(r'(MS\d+)', model)
                    base_model = base.group(1) if base else model
                    
                    if base_model not in ms_groups:
                        ms_groups[base_model] = []
                    
                    ms_groups[base_model].append((model, count))
                
                # Process each MS group
                for base_model, models in sorted(ms_groups.items()):
                    # Format into lines
                    text_lines = []
                    current_line = ""
                    
                    for model, count in sorted(models):
                        model_text = f"{model} ({count})"
                        if len(current_line) + len(model_text) + 2 <= 40:
                            if current_line:
                                current_line += ", " + model_text
                            else:
                                current_line = model_text
                        else:
                            text_lines.append(current_line)
                            current_line = model_text
                    
                    if current_line:
                        text_lines.append(current_line)
                    
                    # Add each line
                    for line in text_lines:
                        item = slide.shapes.add_textbox(left_col_x, left_y, Inches(4.0), Inches(0.25))
                        tf = item.text_frame
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = item_size
                        
                        left_y += Inches(0.25)
        
        sorted_versions = sorted(restricted_devices.keys(), key=lambda x: float(x), reverse=True)
        
        if sorted_versions:
            right_y = start_y
            
            # Process one version at a time in the right column
            for version_index, version in enumerate(sorted_versions):
                # Add firmware version header
                header = slide.shapes.add_textbox(right_col_x - Inches(0.15), right_y, Inches(4), Inches(0.3))
                tf = header.text_frame
                p = tf.add_paragraph()
                p.text = f"MS {version}"
                p.font.size = header_size
                p.font.bold = True
                
                right_y += Inches(0.4)
                
                # Add Meraki switches subtitle
                subtitle = slide.shapes.add_textbox(right_col_x, right_y, Inches(4), Inches(0.25))
                tf = subtitle.text_frame
                p = tf.add_paragraph()
                p.text = "Meraki Switches:"
                p.font.size = item_size
                p.font.bold = True
                
                right_y += Inches(0.3)
                
                # Add models for this version
                models_list = restricted_devices[version]
                
                # Group models for better organization
                model_groups = {}
                for model, count in sorted(models_list.items()):
                    # Group by base model type
                    base_match = re.match(r'(MS\d+|C9300[X-]*)', model)
                    base_model = base_match.group(1) if base_match else model
                    
                    if base_model not in model_groups:
                        model_groups[base_model] = []
                    model_groups[base_model].append((model, count))
                
                # Add each model group on a separate line
                for base_model, models in sorted(model_groups.items()):
                    line_text = ""
                    for model, count in sorted(models):
                        if line_text:
                            line_text += ", "
                        line_text += f"{model} ({count})"
                    
                    # Add the line
                    item = slide.shapes.add_textbox(right_col_x, right_y, Inches(4.25), Inches(0.25))
                    tf = item.text_frame
                    p = tf.add_paragraph()
                    p.text = line_text
                    p.font.size = item_size
                    
                    right_y += Inches(0.25)
                
                # Add spacing between versions if there are more to come
                if version_index < len(sorted_versions) - 1:
                    right_y += Inches(0.3)
        
        # Add total count at the bottom right
        total_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(3), Inches(0.4))
        tf = total_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Total MS Devices: {total_ms_devices}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT
        
        # Add URL to slide notes (visible only to the presenter)
        documentation_urls = [
            "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions#MS",
            "https://documentation.meraki.com/General_Administration/Firmware_Upgrades/Product_Firmware_Version_Restrictions#Cisco_Catalyst"
        ]
        
        if hasattr(slide, 'notes_slide'):
            notes = slide.notes_slide
        else:
            notes = slide.notes_slide = prs.notes_master.clone_master_slide()
        
        # Clear any existing notes
        for shape in notes.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()
        
        # Add the URLs to the slide notes
        notes_text_frame = notes.notes_text_frame
        note_p = notes_text_frame.add_paragraph()
        note_p.text = f"Sources:\n{documentation_urls[0]}\n{documentation_urls[1]}"
        note_p.font.size = Pt(12)
        
        # Save the presentation
        prs.save(output_path)
        #print(f"{GREEN}Updated MS slide (Slide 4) with proper firmware categorization{RESET}")
        
    except Exception as e:
        print(f"{RED}Error updating PowerPoint: {e}{RESET}")
        traceback.print_exc()
    
    ppt_time = time.time() - ppt_start_time
    print(f"{PURPLE}MS Firmware Restrictions slide generation completed in {ppt_time:.2f} seconds{RESET}")
    
    # Calculate total execution time
    total_time = time.time() - start_time
    return total_time

async def main_async(org_ids, template_path=None, output_path=None):
    """
    Standalone async entry point for testing
    """
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
    
    # Create dummy API client
    class DummyApiClient:
        def __init__(self, org_ids):
            self.org_ids = org_ids
    
    api_client = DummyApiClient(org_ids)
    
    # Would need to fetch inventory devices in a real scenario
    # For testing, create some dummy data
    inventory_devices = [
        {"model": "MS120-24", "firmware": "14.32.1", "networkId": "N1"},
        {"model": "MS120-48", "firmware": "14.32.1", "networkId": "N1"},
        {"model": "MS210-24", "firmware": "14.28.2", "networkId": "N2"},
        {"model": "MS420-24", "firmware": "15.23.1", "networkId": "N3"},
        {"model": "MS220-48", "firmware": "14.24.3", "networkId": "N3"},
        {"model": "MS250-24", "firmware": "14.22.1", "networkId": "N4"},
        {"model": "MS350-48", "firmware": "14.32.5", "networkId": "N5"},
        {"model": "C9300-24P", "firmware": "14.32.2", "networkId": "N6"}
    ]
    
    await generate(api_client, template_path, output_path, inventory_devices=inventory_devices)

if __name__ == "__main__":
    # Process command line arguments when run directly
    if len(sys.argv) < 2:
        print("Usage: python ms_firmware_restrictions.py <output_path> [<template_path>]")
        sys.exit(1)
    
    output_path = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else output_path
    
    asyncio.run(main_async(["dummy_org"], template_path, output_path))
