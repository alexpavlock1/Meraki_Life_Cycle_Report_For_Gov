import os
import sys
import asyncio
import datetime
import random
import time
import shutil
import tempfile
import meraki.aio
import logging
from meraki.aio import AsyncRestSession

# PowerPoint libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx import Presentation


# ANSI color codes for terminal output
PURPLE = '\033[95m'   # Timer information
RED = '\033[91m'      # Rate limiting decreases
GREEN = '\033[92m'    # Rate limiting increases
BLUE = '\033[94m'     # General information highlights
YELLOW = '\033[93m'   # Warnings
RESET = '\033[0m'     # Reset to default color

async def get_network_clients_with_logging(aiomeraki, network_id, t0, t1, rate_limiter):
    """Version of get_network_clients with additional logging."""
    logger = logging.getLogger('meraki_api_detailed')
    logger.debug(f"Getting clients for network {network_id}, t0={t0}, t1={t1}")
    
    try:
        start_time = time.time()
        clients = await get_network_clients(aiomeraki, network_id, t0, t1, rate_limiter)
        elapsed_time = time.time() - start_time
        #logger.debug(f"Got {len(clients)} clients in {elapsed_time:.2f}s for network {network_id}")
        return clients
    except Exception as e:
        logger.error(f"Error getting clients for network {network_id}: {str(e)}")
        raise

def get_api_key():
    """Retrieve Meraki API key from environment variable."""
    api_key = os.environ.get("MERAKI_API_KEY")
    if not api_key:
        raise ValueError("MERAKI_API_KEY environment variable is not set")
    return api_key
# Direct manual API implementation as a fallback
async def get_network_clients_manual(network_id, t0, t1, api_key):
    """Manually implement getNetworkClients using aiohttp to match Postman behavior."""
    import aiohttp
    
    # Set up logging
    logger = logging.getLogger('meraki_api_detailed')
    
    url = f"https://api.gov-meraki.com/api/v1/networks/{network_id}/clients"
    
    params = {
        "t0": t0,
        "t1": t1,
        "perPage": 5000
    }
    
    headers = {
        "X-Cisco-Meraki-API-Key": api_key,
        "Content-Type": "application/json",
        "Accept": "application/json",
        "User-Agent": "Python Meraki Client"
    }
    
    logger.debug(f"MANUAL API REQUEST: {url} with t0={t0}, t1={t1}")
    
    async with aiohttp.ClientSession() as session:
        try:
            start_time = time.time()
            async with session.get(url, params=params, headers=headers, timeout=60) as response:
                response_time = time.time() - start_time
                logger.debug(f"MANUAL API RESPONSE: Status {response.status} in {response_time:.2f}s")
                
                if response.status == 200:
                    data = await response.json()
                    logger.debug(f"MANUAL API SUCCESS: Got {len(data)} clients")
                    return data
                else:
                    error_text = await response.text()
                    logger.error(f"MANUAL API ERROR: Status {response.status}, {error_text[:200]}...")
                    raise Exception(f"API error: {response.status}, {error_text[:200]}")
        except asyncio.TimeoutError as e:
            logger.error(f"MANUAL API TIMEOUT: {str(e)}")
            raise e
        except Exception as e:
            logger.error(f"MANUAL API EXCEPTION: {type(e).__name__} - {str(e)}")
            raise e

# Test function to compare with Postman
async def test_postman_equivalent(aiomeraki, network_id, api_key, rate_limiter):
    """Run a test that mimics a Postman request exactly."""
    # Set up logging
    logger = logging.getLogger('meraki_api_detailed')
    
    logger.debug(f"\n===== POSTMAN COMPARISON TEST =====")
    logger.debug(f"Testing network ID: {network_id}")
    
    # Use the exact same parameters as a Postman request
    t0 = "2025-04-01T00:01:00Z"
    t1 = "2025-04-01T23:59:00Z"
    
    # Try with SDK first
    try:
        logger.debug(f"POSTMAN TEST: Using SDK with t0={t0}, t1={t1}")
        
        start_time = time.time()
        clients = await get_network_clients_with_logging(aiomeraki, network_id, t0, t1, rate_limiter)
        end_time = time.time()
        
        logger.debug(f"POSTMAN TEST SDK SUCCESS: Got {len(clients)} clients in {end_time - start_time:.2f}s")
        return True, "sdk", len(clients), end_time - start_time
    except Exception as e:
        logger.error(f"POSTMAN TEST SDK ERROR: {type(e).__name__} - {str(e)}")
        
        # Try with manual implementation
        try:
            logger.debug(f"POSTMAN TEST: Using manual implementation with t0={t0}, t1={t1}")
            
            start_time = time.time()
            clients = await get_network_clients_manual(network_id, t0, t1, api_key)
            end_time = time.time()
            
            logger.debug(f"POSTMAN TEST MANUAL SUCCESS: Got {len(clients)} clients in {end_time - start_time:.2f}s")
            return True, "manual", len(clients), end_time - start_time
        except Exception as e2:
            logger.error(f"POSTMAN TEST MANUAL ERROR: {type(e2).__name__} - {str(e2)}")
            return False, None, 0, 0
def setup_enhanced_logging():
    """Set up enhanced logging for API calls"""
    # Create a logs directory if it doesn't exist
    if not os.path.exists('logs'):
        os.makedirs('logs')
    
    # Configure logging
    logger = logging.getLogger('meraki_api_detailed')
    logger.setLevel(logging.DEBUG)
    
    # Clear any existing handlers
    if logger.handlers:
        logger.handlers = []
    
    # Create handlers
    file_handler = logging.FileHandler('logs/meraki_api_detailed.log')
    file_handler.setLevel(logging.DEBUG)
    
    # Create formatters
    file_format = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    
    # Add formatters to handlers
    file_handler.setFormatter(file_format)
    
    # Add handlers to logger
    logger.addHandler(file_handler)
    
    return logger


class AdaptiveRateLimiter:
    """Simplified rate limiter with minimal complexity to avoid deadlocks."""
    
    def __init__(self, initial_limit=5, min_limit=3, max_limit=10):
        self.current_limit = initial_limit
        self.min_limit = min_limit
        self.max_limit = max_limit
        self.semaphore = asyncio.Semaphore(self.current_limit)
        self.error_count = 0
        self.success_count = 0
    
    async def wait(self):
        """Acquire the semaphore before making an API request."""
        await self.semaphore.acquire()
        try:
            return True
        finally:
            self.semaphore.release()
    
    # Keep your existing methods
    def decrease_limit(self):
        """Decrease the concurrency limit after a 429 error."""
        if self.current_limit > self.min_limit:
            # Simple decrease by 1
            self.current_limit = max(self.min_limit, self.current_limit - 1)
            self.semaphore = asyncio.Semaphore(self.current_limit)
            print(f"{RED}Rate limit hit! Decreasing concurrency limit to {self.current_limit}{RESET}")
    
    def increase_limit(self):
        """Cautiously increase the concurrency limit after many successes."""
        if self.current_limit < self.max_limit:
            # Simple increase by 1
            self.current_limit = min(self.max_limit, self.current_limit + 1)
            self.semaphore = asyncio.Semaphore(self.current_limit)
            print(f"{GREEN}Increased concurrency limit to {self.current_limit}{RESET}")
    
    def check_and_adjust(self):
        """Periodically check and adjust the limit based on success/error ratio."""
        # If we've had a significant number of successes and no errors
        if self.success_count > 100 and self.error_count == 0:
            self.increase_limit()
            self.success_count = 0
        
        # If we've had errors, decrease the limit
        elif self.error_count > 0:
            self.decrease_limit()
            self.error_count = 0
            self.success_count = 0


async def rate_limited_api_call(api_func, rate_limiter, max_retries=4, *args, **kwargs):
    """
    Execute an API call with simpler error handling to avoid stalls.
    
    Args:
        api_func: The API function to call
        rate_limiter: AdaptiveRateLimiter instance
        max_retries: Maximum number of retries on failure
        *args, **kwargs: Arguments to pass to the API function
        
    Returns:
        API call result or raises an exception after retries
    """
    retries = 0
    
    while retries <= max_retries:
        try:
            # Add a small random delay to avoid synchronized requests
            jitter = random.uniform(0.1, 0.3)  # 100-300ms jitter
            await asyncio.sleep(jitter)
            
            async with rate_limiter.semaphore:
                # Set a timeout on the API call to prevent hanging indefinitely
                result = await asyncio.wait_for(api_func(*args, **kwargs), timeout=15.0)
                # Simple success counter, no complex logic
                rate_limiter.success_count += 1
                return result
                
        except asyncio.TimeoutError:
            # Handle timeout errors
            print(f"{RED}Timeout in API call {api_func.__name__}{RESET}")
            retries += 1
            if retries <= max_retries:
                await asyncio.sleep(1.0)  # Simple backoff
            else:
                raise Exception(f"API timeout after {max_retries} retries")
                
        except Exception as e:
            error_str = str(e)
            is_rate_limit = "429" in error_str
            
            if is_rate_limit:
                # Rate limit error - increment counter and retry
                rate_limiter.error_count += 1
                retries += 1
                if retries <= max_retries:
                    # Simple backoff strategy
                    wait_time = 1.0 * retries
                    print(f"{YELLOW}Rate limit hit, retrying in {wait_time:.1f}s (attempt {retries}/{max_retries}){RESET}")
                    await asyncio.sleep(wait_time)
                else:
                    raise Exception(f"Rate limit exceeded after {max_retries} retries")
            else:
                # Other error - pass through
                raise

async def get_organization_names(aiomeraki, org_ids, rate_limiter):
    """Get organization names for the given organization IDs."""
    org_names = {}
    
    # Get all organizations
    try:
        orgs = await rate_limited_api_call(
            aiomeraki.organizations.getOrganizations,
            rate_limiter
        )
        
        # Create a mapping of organization ID to name
        for org in orgs:
            if org['id'] in org_ids:
                org_names[org['id']] = org['name']
        
        # If we couldn't find all org IDs, try to fetch them individually
        missing_ids = [org_id for org_id in org_ids if org_id not in org_names]
        for org_id in missing_ids:
            try:
                org = await rate_limited_api_call(
                    aiomeraki.organizations.getOrganization,
                    rate_limiter,
                    organizationId=org_id
                )
                org_names[org_id] = org['name']
            except Exception as e:
                print(f"{YELLOW}Could not retrieve name for organization {org_id}: {e}{RESET}")
                # Use org ID as name if we can't get the actual name
                org_names[org_id] = f"Organization {org_id}"
    
    except Exception as e:
        print(f"{RED}Error retrieving organization names: {e}{RESET}")
        # Use org IDs as names if we can't get the actual names
        for org_id in org_ids:
            org_names[org_id] = f"Organization {org_id}"
    
    return org_names

async def collect_data(org_ids, days=14):
    """Only collect data without updating PowerPoint."""
    # Run the same data collection code as in main_async, but don't update PowerPoint
    api_key = get_api_key()
    rate_limiter = AdaptiveRateLimiter(initial_limit=5, min_limit=3, max_limit=15)
    
    async with meraki.aio.AsyncDashboardAPI(
        api_key=api_key,
        suppress_logging=False,
        log_file_prefix='meraki_api_log',
        log_path='.',
        maximum_retries=3,
        base_url="https://api.gov-meraki.com/api/v1"
    ) as aiomeraki:
        # Get organization names
        org_names = await get_organization_names(aiomeraki, org_ids, rate_limiter)
        
        # Get all networks
        all_networks = []
        for org_id in org_ids:
            networks = await get_networks(aiomeraki, org_id, rate_limiter)
            all_networks.extend(networks)
        
        # Extract network IDs
        network_ids = [network['id'] for network in all_networks]
        
        # Filter incompatible networks
        valid_network_ids = await filter_incompatible_networks(network_ids, all_networks)
        
        # Get dashboard statistics
        dashboard_stats = await get_dashboard_stats(aiomeraki, org_ids, rate_limiter)
        
        # Get client statistics
        client_stats = await get_client_stats(aiomeraki, valid_network_ids, rate_limiter, days)
        
        # Combine all stats
        combined_stats = {**dashboard_stats, **client_stats}
        
        # Also collect inventory devices for slide2
        all_inventory_devices = []
        for org_id in org_ids:
            devices = await get_inventory_devices(aiomeraki, org_id, rate_limiter)
            all_inventory_devices.extend(devices)
        
        return combined_stats, all_inventory_devices, org_names

async def get_organizations(aiomeraki, rate_limiter):
    """Get list of organizations the API key has access to."""
    return await rate_limited_api_call(
        aiomeraki.organizations.getOrganizations,
        rate_limiter
    )


async def get_networks(aiomeraki, org_id, rate_limiter):
    """Get list of networks for a specific organization with pagination and rate limiting."""
    all_networks = []
    
    # Start with the first page
    current_page = 1
    per_page = 1000  # Maximum allowed by the API
    last_network_id = None
    
    while True:
        # Get a page of networks with rate limiting
        try:
            # Only include startingAfter for pages after the first
            if current_page == 1:
                networks_page = await rate_limited_api_call(
                    aiomeraki.organizations.getOrganizationNetworks,
                    rate_limiter,
                    organizationId=org_id,
                    perPage=per_page
                )
            else:
                networks_page = await rate_limited_api_call(
                    aiomeraki.organizations.getOrganizationNetworks,
                    rate_limiter,
                    organizationId=org_id,
                    perPage=per_page,
                    startingAfter=last_network_id
                )
        except Exception as e:
            print(f"{RED}Error getting networks for org {org_id}, page {current_page}: {e}{RESET}")
            break
            
        # If we got no networks or fewer than requested, we've reached the end
        if not networks_page:
            break
            
        # Add networks to our collection
        all_networks.extend(networks_page)
        
        # If we got fewer networks than requested, we've reached the end
        if len(networks_page) < per_page:
            break
            
        # Store the last network's ID for pagination
        last_network_id = networks_page[-1]['id']
            
        # Move to the next page
        current_page += 1
        #print(f"{BLUE}Retrieved {len(networks_page)} networks from page {current_page-1}, total so far: {len(all_networks)}{RESET}")
    
    return all_networks


async def get_inventory_devices(aiomeraki, org_id, rate_limiter):
    """Get all inventory devices across an organization with pagination and rate limiting."""
    all_devices = []
    
    # Start with the first page
    current_page = 1
    per_page = 1000
    last_device_serial = None
    
    while True:
        # Get a page of devices with rate limiting
        try:
            # Only include startingAfter for pages after the first
            if current_page == 1:
                devices_page = await rate_limited_api_call(
                    aiomeraki.organizations.getOrganizationInventoryDevices,
                    rate_limiter,
                    organizationId=org_id, 
                    perPage=per_page
                )
            else:
                devices_page = await rate_limited_api_call(
                    aiomeraki.organizations.getOrganizationInventoryDevices,
                    rate_limiter,
                    organizationId=org_id, 
                    perPage=per_page,
                    startingAfter=last_device_serial
                )
        except Exception as e:
            print(f"Error getting inventory devices for org {org_id}, page {current_page}: {e}")
            break
            
        # If we got no devices or fewer than requested, we've reached the end
        if not devices_page:
            break
            
        # Add devices to our collection
        all_devices.extend(devices_page)
        
        # If we got fewer devices than requested, we've reached the end
        if len(devices_page) < per_page:
            break
            
        # Store the last device's serial for pagination
        last_device_serial = devices_page[-1]['serial']
            
        # Move to the next page
        current_page += 1
    
    return all_devices


def filter_active_devices(devices):
    """Filter to only include active devices (those with a non-blank networkId)."""
    return [device for device in devices if device.get('networkId')]


async def get_network_clients(aiomeraki, network_id, t0, t1, rate_limiter):
    """Get all clients for a network within the specified time range."""
    all_clients = []
    
    # Start with the first page
    current_page = 1
    per_page = 5000  # Maximum allowed
    last_client_id = None
    
    while True:
        try:
            # Use rate_limited_api_call with appropriate pagination parameters
            if current_page == 1:
                # First page doesn't need startingAfter
                clients_page = await rate_limited_api_call(
                    aiomeraki.networks.getNetworkClients,
                    rate_limiter,
                    networkId=network_id,
                    t0=t0,
                    t1=t1,
                    perPage=per_page
                )
            else:
                # Subsequent pages need startingAfter parameter
                clients_page = await rate_limited_api_call(
                    aiomeraki.networks.getNetworkClients,
                    rate_limiter,
                    networkId=network_id,
                    t0=t0,
                    t1=t1,
                    perPage=per_page,
                    startingAfter=last_client_id
                )
            
            # If we got no clients, we've reached the end
            if not clients_page:
                break
                
            # Add clients to our collection
            all_clients.extend(clients_page)
            
            # If we got fewer clients than requested, we've reached the end
            if len(clients_page) < per_page:
                break

            if clients_page and 'id' in clients_page[-1]:
                last_client_id = clients_page[-1]['id']
            else:
                # If id isn't available, we can't paginate further
                break
                
            # Move to the next page
            current_page += 1
            
        except Exception as e:
            # Let exceptions propagate to the caller for proper handling
            raise e
    
    return all_clients


async def get_dashboard_stats(aiomeraki, org_ids, rate_limiter):
    """Get Meraki dashboard statistics for the given organizations."""
    total_networks = 0
    total_inventory = 0
    total_active_nodes = 0
    
    # Create tasks for concurrent API calls
    network_tasks = [get_networks(aiomeraki, org_id, rate_limiter) for org_id in org_ids]
    inventory_tasks = [get_inventory_devices(aiomeraki, org_id, rate_limiter) for org_id in org_ids]
    
    # Wait for all network tasks to complete
    networks_results = await asyncio.gather(*network_tasks, return_exceptions=True)
    
    # Process network results
    for result in networks_results:
        if isinstance(result, Exception):
            print(f"Error getting networks: {result}")
            continue
        total_networks += len(result)
    
    # Wait for all inventory tasks to complete
    inventory_results = await asyncio.gather(*inventory_tasks, return_exceptions=True)
    
    # Process inventory results
    for result in inventory_results:
        if isinstance(result, Exception):
            print(f"Error getting inventory devices: {result}")
            continue
        
        # Count all devices in inventory
        total_inventory += len(result)
        
        # Count devices with a networkId (active nodes)
        active_devices = filter_active_devices(result)
        total_active_nodes += len(active_devices)
    
    return {
        "total_networks": total_networks,
        "total_inventory": total_inventory,
        "total_active_nodes": total_active_nodes
    }


async def filter_incompatible_networks(network_ids, all_networks):
    """Filter out networks that don't support client API (Systems Manager, Camera, etc.)."""
    # Create a mapping of network ID to network data
    network_map = {net['id']: net for net in all_networks}
    
    valid_network_ids = []
    incompatible_networks = {
        "systemsManager": 0,
        "camera": 0,
        "other": 0
    }
    
    for network_id in network_ids:
        if network_id in network_map:
            network = network_map[network_id]
            product_types = network.get('productTypes', [])
            
            # Check incompatible product types
            if "systemsManager" in product_types:
                incompatible_networks["systemsManager"] += 1
                continue  # Skip this network
            elif "camera" in product_types and len(product_types) == 1:  # Only skip if it's just camera
                incompatible_networks["camera"] += 1
                continue  # Skip this network
        valid_network_ids.append(network_id)
    
    # Print summary of filtered networks
    total_filtered = sum(incompatible_networks.values())
    if total_filtered > 0:
        #print(f"{BLUE}Filtered out {total_filtered} incompatible networks:{RESET}")
        pass
        for network_type, count in incompatible_networks.items():
            if count > 0:
                #print(f"  - {network_type}: {count} networks")
                pass
    
    #print(f"{BLUE}Processing {len(valid_network_ids)} client-supporting networks{RESET}")
    
    return valid_network_ids


async def get_client_stats(aiomeraki, network_ids, rate_limiter, days_back=14):
    """Get client statistics with dynamic day parallelism and smart time chunking."""
    # Set up logging
    logger = setup_enhanced_logging()
    logger.debug("Starting get_client_stats function")
    
    # Get API key for direct API calls if needed
    api_key = os.environ.get("MERAKI_API_KEY")
    
    # Invalid network blacklist - share between all days
    invalid_network_blacklist = set()
    # Track networks that experience timeouts
    timeout_prone_networks = set()
    # Track networks with high client volumes
    high_volume_networks = set()
    # Track client counts by network
    network_client_counts = {}
    # Track timeouts by network
    network_timeout_counts = {}
    
    # Pre-mark problematic networks (like 2006) This is for testing a known network that is failing
    for network_id in network_ids:
        if network_id.endswith('2006'):
            timeout_prone_networks.add(network_id)
            logger.debug(f"Pre-marking network {network_id} as timeout-prone")
    
    # Get current time in UTC
    now = datetime.datetime.now(datetime.timezone.utc)
    
    # Track all results
    all_unique_macs = set()
    total_non_unique_clients = 0
    unique_clients_per_day = []
    non_unique_clients_per_day = []
    day_results = []
    
    # Dynamic day parallelism control
    day_parallelism = 2  # Start with 2 days in parallel
    max_parallelism = 2  # Maximum of 2 days in parallel
    
    # Class to track rate limit status
    class RateLimitTracker:
        def __init__(self):
            self._hit = False
            self._count = 0
            
        def record_hit(self):
            self._hit = True
            self._count += 1
            print(f"{RED}Rate limit hit! ({self._count} total){RESET}")
            
        def check_and_reset(self):
            was_hit = self._hit
            self._hit = False
            return was_hit
            
        @property
        def total_hits(self):
            return self._count
    
    rate_tracker = RateLimitTracker()
    
    # Smart chunking function that adapts to network behavior
    async def get_network_clients_with_smart_chunking(network_id, t0, t1):
        """Use adaptive time chunking based on network behavior."""
        # Log the start of this function call
        logger.debug(f"get_network_clients_with_smart_chunking: network={network_id}, t0={t0}, t1={t1}")
        
        # Check for problematic networks
        if network_id in timeout_prone_networks:
            # Initial chunking size (smaller for timeout-prone networks)
            chunk_hours = 1  # Start with 1-hour chunks for problematic networks
            max_retries = 3
            
            print(f"{BLUE}Using smaller time chunking for timeout-prone network {network_id[-4:]}{RESET}")
            logger.debug(f"Using smaller time chunking for timeout-prone network {network_id}")
            
            # For network 2006, first try the test function that mimics Postman
            if network_id.endswith('2006'):
                logger.debug(f"Running Postman equivalent test for network {network_id}")
                success, method, num_clients, time_taken = await test_postman_equivalent(
                    aiomeraki, network_id, api_key, rate_limiter
                )
                
                if success:
                    print(f"{GREEN}Postman test successful! Got {num_clients} clients in {time_taken:.2f}s using {method}{RESET}")
                    logger.debug(f"Postman test successful with {method} method")
                    
                    # If the manual method worked, use it for all requests to this network
                    if method == "manual":
                        try:
                            logger.debug(f"Using manual implementation for main request t0={t0}, t1={t1}")
                            clients = await get_network_clients_manual(network_id, t0, t1, api_key)
                            logger.debug(f"Manual implementation succeeded with {len(clients)} clients")
                            return clients
                        except Exception as e:
                            logger.error(f"Manual implementation failed: {str(e)}")
                            # Continue with chunking approach below
            
            # Parse start and end times
            start_time = datetime.datetime.strptime(t0, '%Y-%m-%dT%H:%M:%SZ')
            end_time = datetime.datetime.strptime(t1, '%Y-%m-%dT%H:%M:%SZ')
            
            # Create chunks
            current_time = start_time
            all_clients = []
            
            while current_time < end_time:
                next_time = min(current_time + datetime.timedelta(hours=chunk_hours), end_time)
                
                # Format times for API
                chunk_t0 = current_time.strftime('%Y-%m-%dT%H:%M:%SZ')
                chunk_t1 = next_time.strftime('%Y-%m-%dT%H:%M:%SZ')
                
                retry_count = 0
                success = False
                
                # Calculate timeout based on history with this network
                timeout_count = network_timeout_counts.get(network_id, 0)
                initial_timeout = max(15.0, 60.0 / (1 + timeout_count * 0.2))
                
                # Use different approaches based on the network
                if network_id.endswith('2006'):
                    # For network 2006, try manual implementation first
                    try:
                        logger.debug(f"Using manual implementation for chunk t0={chunk_t0}, t1={chunk_t1}")
                        chunk_clients = await get_network_clients_manual(network_id, chunk_t0, chunk_t1, api_key)
                        all_clients.extend(chunk_clients)
                        print(f"{GREEN}Manual API succeeded for {chunk_t0} to {chunk_t1} with {len(chunk_clients)} clients{RESET}")
                        logger.debug(f"Manual API succeeded with {len(chunk_clients)} clients")
                        success = True
                    except Exception as e:
                        logger.error(f"Manual API failed: {str(e)}")
                        # Fall back to the standard approach with retries
                else:
                    # Add progressive backoff for retries
                    while retry_count < max_retries and not success:
                        try:
                            # Add jitter to avoid synchronized requests
                            await asyncio.sleep(random.uniform(0.1, 0.5) * (retry_count + 1))
                            
                            # Adjust timeout based on retry count
                            adjusted_timeout = initial_timeout - (retry_count * 5.0)
                            adjusted_timeout = max(10.0, adjusted_timeout)  # Don't go below 10 seconds
                            
                            logger.debug(f"Attempting chunk with timeout={adjusted_timeout}s (retry {retry_count})")
                            
                            # Use the logging version for problematic networks
                            if network_id in timeout_prone_networks:
                                chunk_clients = await asyncio.wait_for(
                                    get_network_clients_with_logging(aiomeraki, network_id, chunk_t0, chunk_t1, rate_limiter),
                                    timeout=adjusted_timeout
                                )
                            else:
                                # Get clients for this time chunk with a shorter timeout
                                chunk_clients = await asyncio.wait_for(
                                    get_network_clients(aiomeraki, network_id, chunk_t0, chunk_t1, rate_limiter),
                                    timeout=adjusted_timeout
                                )
                            
                            # Success - add clients and track counts
                            all_clients.extend(chunk_clients)
                            
                            # Update high-volume tracking if we got a lot of clients
                            if len(chunk_clients) > 1000:
                                high_volume_networks.add(network_id)
                                network_client_counts[network_id] = network_client_counts.get(network_id, 0) + len(chunk_clients)
                                logger.debug(f"Network {network_id} has many clients ({len(chunk_clients)})")
                            
                            print(f"{GREEN}Successfully retrieved {len(chunk_clients)} clients for chunk {chunk_t0} to {chunk_t1} for network {network_id[-4:]}{RESET}")
                            logger.debug(f"Successfully retrieved {len(chunk_clients)} clients for chunk")
                            success = True
                            
                        except asyncio.TimeoutError:
                            # Update timeout counter for this network
                            network_timeout_counts[network_id] = network_timeout_counts.get(network_id, 0) + 1
                            logger.debug(f"Timeout on chunk (retry {retry_count})")
                            
                            retry_count += 1
                            if retry_count < max_retries:
                                print(f"{YELLOW}Timeout on chunk {chunk_t0} to {chunk_t1} for network {network_id[-4:]} (retry {retry_count}/{max_retries}){RESET}")
                                
                                # If we're on the last retry and using larger chunks, try with a smaller chunk size
                                if retry_count == max_retries - 1 and chunk_hours > 1:
                                    # Exit this retry loop to allow the outer loop to use smaller chunks
                                    break
                            else:
                                print(f"{RED}Giving up on chunk {chunk_t0} to {chunk_t1} for network {network_id[-4:]} after {max_retries} retries{RESET}")
                        except Exception as e:
                            # Handle other exceptions as before
                            error_str = str(e)
                            logger.error(f"Error on chunk: {error_str}")
                            
                            if "Invalid device type" in error_str:
                                invalid_network_blacklist.add(network_id)
                                print(f"{YELLOW}Invalid device type for network {network_id[-4:]}, blacklisting{RESET}")
                                logger.debug(f"Blacklisting network {network_id} due to invalid device type")
                                return all_clients
                            elif "429" in error_str:
                                rate_tracker.record_hit()
                                retry_count += 1
                                await asyncio.sleep(1.0 * retry_count)  # Back off on rate limits
                                print(f"{YELLOW}Rate limit for network {network_id[-4:]} in time chunk (retry {retry_count}/{max_retries}){RESET}")
                                logger.debug(f"Rate limit hit, retry {retry_count}")
                            else:
                                print(f"{YELLOW}Error for network {network_id[-4:]} in time chunk: {str(e)[:100]}{RESET}")
                                retry_count += 1
                
                # If we failed with larger chunks, retry with smaller chunks
                if not success and chunk_hours > 1:
                    # Reduce chunk size and retry this time period with smaller chunks
                    chunk_hours = max(1, chunk_hours // 2)
                    print(f"{BLUE}Reducing chunk size to {chunk_hours} hour(s) for network {network_id[-4:]}{RESET}")
                    logger.debug(f"Reducing chunk size to {chunk_hours} hour(s)")
                    continue  # Retry the same time period with smaller chunks
                
                # If network 2006 still fails, try with very small chunks (15 minutes)
                if not success and network_id.endswith('2006'):
                    logger.debug(f"Network 2006 failed, trying with very small chunks")
                    # Try with 15 minute chunks
                    small_chunk_minutes = 15
                    chunk_start = current_time
                    
                    while chunk_start < next_time and not success:
                        chunk_end = min(chunk_start + datetime.timedelta(minutes=small_chunk_minutes), next_time)
                        small_t0 = chunk_start.strftime('%Y-%m-%dT%H:%M:%SZ')
                        small_t1 = chunk_end.strftime('%Y-%m-%dT%H:%M:%SZ')
                        
                        try:
                            logger.debug(f"Trying very small chunk {small_t0} to {small_t1}")
                            # Try manual implementation for these very small chunks
                            small_clients = await get_network_clients_manual(network_id, small_t0, small_t1, api_key)
                            all_clients.extend(small_clients)
                            print(f"{GREEN}Very small chunk succeeded for {small_t0} to {small_t1} with {len(small_clients)} clients{RESET}")
                            logger.debug(f"Very small chunk succeeded with {len(small_clients)} clients")
                            # Mark as partial success
                            success = True
                        except Exception as e:
                            logger.error(f"Very small chunk failed: {str(e)}")
                            # Just move to the next small chunk
                        
                        # Move to next small chunk
                        chunk_start = chunk_end
                
                # Move to next chunk
                current_time = next_time
                
                # Small break between chunks to avoid rate limiting
                await asyncio.sleep(random.uniform(0.3, 0.7))
            
            logger.debug(f"get_network_clients_with_smart_chunking completed with {len(all_clients)} clients")
            return all_clients
        else:
            # Use standard approach for networks that haven't timed out before
            try:
                # Set a reasonable timeout for the network call
                timeout_value = 60.0
                logger.debug(f"Using standard approach with timeout={timeout_value}s")
                
                return await asyncio.wait_for(
                    get_network_clients(aiomeraki, network_id, t0, t1, rate_limiter),
                    timeout=timeout_value
                )
            except asyncio.TimeoutError:
                # Add this network to the timeout-prone set for future queries
                timeout_prone_networks.add(network_id)
                # Update timeout counter
                network_timeout_counts[network_id] = network_timeout_counts.get(network_id, 0) + 1
                
                print(f"{RED}Timeout getting clients for network {network_id[-4:]}, will use time chunking for future queries{RESET}")
                logger.debug(f"Timeout in standard approach, adding to timeout_prone_networks")
                
                # Fall back to time chunking approach for this call
                return await get_network_clients_with_smart_chunking(network_id, t0, t1)
            except Exception as e:
                error_str = str(e)
                logger.error(f"Error in standard approach: {error_str}")
                
                if "Invalid device type" in error_str:
                    invalid_network_blacklist.add(network_id)
                    logger.debug(f"Blacklisting network {network_id} due to invalid device type")
                elif "429" in error_str:
                    # Directly record the rate limit hit
                    rate_tracker.record_hit()
                    logger.debug(f"Rate limit hit for network {network_id}")
                else:
                    logger.error(f"Error for network {network_id}: {str(e)}")
                return []
    
    # Process day statistics with controlled parallelism
    async def process_day(day):
        try:
            day_start_time = time.time()
            logger.debug(f"process_day: Processing day {day+1} of {days_back}")
            
            # Calculate the date for the day being processed
            target_date = (now - datetime.timedelta(days=day)).date()
            
            # Set time to 12:01 AM for start and 11:59 PM for end
            day_start = datetime.datetime.combine(
                target_date, 
                datetime.time(0, 1, 0),  # 12:01 AM
                tzinfo=datetime.timezone.utc
            )
            # For end time: use current time if it's today, otherwise use 11:59 PM
            if day == 0:  # Today
                day_end = now  # Use current time as end time
            else:
                day_end = datetime.datetime.combine(
                    target_date,
                    datetime.time(23, 59, 0),  # 11:59 PM
                    tzinfo=datetime.timezone.utc
                )
            
            # Convert to ISO 8601 format
            t0 = day_start.strftime('%Y-%m-%dT%H:%M:%SZ')
            t1 = day_end.strftime('%Y-%m-%dT%H:%M:%SZ')
            
            print(f"{PURPLE}Starting day {day+1} of {days_back}... ({t0} to {t1}){RESET}")
            
            # Track unique MACs for this day
            day_macs = set()
            daily_non_unique_count = 0
            
            # Filter out blacklisted networks
            valid_networks = [network_id for network_id in network_ids 
                             if network_id not in invalid_network_blacklist]
            
            # Process networks in smaller chunks
            chunk_size = 6  # Networks per chunk
            
            for i in range(0, len(valid_networks), chunk_size):
                chunk_start = time.time()
                chunk = valid_networks[i:i+chunk_size]
                
                # Print progress
                print(f"{BLUE}Processing networks {i+1}-{min(i+chunk_size, len(valid_networks))} of {len(valid_networks)}{RESET}")
                logger.debug(f"Processing networks {i+1}-{min(i+chunk_size, len(valid_networks))} of {len(valid_networks)}")
                
                # Process networks concurrently with a limit
                clients_by_network = {}
                for j in range(0, len(chunk), 3):  # Process up to 3 networks at once
                    sub_chunk = chunk[j:j+3]
                    # Create tasks for this sub-chunk using smart chunking
                    tasks = [get_network_clients_with_smart_chunking(net_id, t0, t1) for net_id in sub_chunk]
                    
                    # Wait for all tasks with combined timeout
                    try:
                        results = await asyncio.wait_for(asyncio.gather(*tasks), timeout=360.0)  # Increased timeout
                        # Store results by network
                        for k, net_id in enumerate(sub_chunk):
                            clients_by_network[net_id] = results[k]
                    except asyncio.TimeoutError:
                        print(f"{RED}Timeout in sub-chunk processing{RESET}")
                        logger.error(f"Timeout in sub-chunk processing")
                        # Try each network individually to avoid entire batch failing
                        for net_id in sub_chunk:
                            try:
                                logger.debug(f"Trying network {net_id} individually after batch timeout")
                                result = await asyncio.wait_for(
                                    get_network_clients_with_smart_chunking(net_id, t0, t1),
                                    timeout=180.0
                                )
                                clients_by_network[net_id] = result
                            except Exception as e:
                                logger.error(f"Error processing network {net_id} individually: {str(e)}")
                                print(f"{RED}Error processing network {net_id[-4:]} individually: {str(e)[:100]}{RESET}")
                
                # Process results for this chunk
                successful_networks = 0
                for net_id, clients in clients_by_network.items():
                    if clients:
                        successful_networks += 1
                        # Count all clients as non-unique
                        daily_non_unique_count += len(clients)
                        # Add MAC addresses to the unique set
                        for client in clients:
                            if 'mac' in client and client['mac'] is not None:
                                mac = client['mac'].lower()
                                day_macs.add(mac)
                
                chunk_time = time.time() - chunk_start
                print(f"{BLUE}Processed chunk {i//chunk_size + 1}/{(len(valid_networks)+chunk_size-1)//chunk_size} with {successful_networks} successful networks ({chunk_time:.1f}s){RESET}")
                logger.debug(f"Processed chunk with {successful_networks} successful networks in {chunk_time:.1f}s")
                
                # Small pause between chunks
                await asyncio.sleep(0.2)
            
            daily_unique_count = len(day_macs)
            
            # Report day timing
            day_time = time.time() - day_start_time
            print(f"{PURPLE}Day {day+1} completed: {daily_unique_count} unique, {daily_non_unique_count} non-unique clients ({day_time:.2f}s){RESET}")
            logger.debug(f"Day {day+1} completed with {daily_unique_count} unique clients in {day_time:.2f}s")
            
            return {
                'day': day,
                'unique_count': daily_unique_count,
                'non_unique_count': daily_non_unique_count,
                'unique_macs': day_macs
            }
        except Exception as e:
            logger.error(f"Error processing day {day+1}: {str(e)}")
            print(f"{RED}Error processing day {day+1}: {e}{RESET}")
            # Return empty results for this day
            return {
                'day': day,
                'unique_count': 0,
                'non_unique_count': 0,
                'unique_macs': set()
            }
    
    # Process all days with dynamic batch size
    batch_idx = 0
    while batch_idx < days_back:
        # Get current batch based on parallelism
        current_batch = list(range(batch_idx, min(batch_idx + day_parallelism, days_back)))
        batch_size = len(current_batch)
        
        print(f"\n{BLUE}Processing batch: days {current_batch[0]+1}-{current_batch[-1]+1} ({batch_size} days in parallel){RESET}")
        logger.debug(f"Processing batch: days {current_batch[0]+1}-{current_batch[-1]+1} ({batch_size} days in parallel)")
        
        # Process days in this batch concurrently
        batch_tasks = [process_day(day) for day in current_batch]
        batch_results = await asyncio.gather(*batch_tasks)
        
        # Add results to overall collection
        day_results.extend(batch_results)
        
        # Check if rate limits were hit and adjust parallelism
        if rate_tracker.check_and_reset():
            # We hit rate limits, so decrease parallelism to 1 day at a time
            old_parallelism = day_parallelism
            day_parallelism = 1
            if old_parallelism != day_parallelism:
                print(f"{YELLOW}Rate limits hit! Reducing parallelism from {old_parallelism} to {day_parallelism} day{RESET}")
                logger.debug(f"Rate limits hit! Reducing parallelism to {day_parallelism} day")
        elif day_parallelism < max_parallelism:
            # No rate limits, so increase parallelism (up to max)
            old_parallelism = day_parallelism
            day_parallelism = min(max_parallelism, day_parallelism + 1)
            if old_parallelism != day_parallelism:
                print(f"{GREEN}No rate limits! Increasing parallelism from {old_parallelism} to {day_parallelism} days{RESET}")
                logger.debug(f"Increasing parallelism to {day_parallelism} days")
        
        # Move to next batch
        batch_idx += batch_size
    
    # Process the results from all days
    # Sort by day to keep order in output
    day_results.sort(key=lambda x: x['day'])
    
    for result in day_results:
        unique_clients_per_day.append(result['unique_count'])
        non_unique_clients_per_day.append(result['non_unique_count'])
        total_non_unique_clients += result['non_unique_count']
        all_unique_macs.update(result['unique_macs'])
    
    # Calculate total unique clients as the number of unique MAC addresses across all days
    total_unique_clients = len(all_unique_macs)
    
    # Calculate averages
    avg_unique_clients_per_day = sum(unique_clients_per_day) / days_back if days_back > 0 else 0
    avg_non_unique_clients_per_day = sum(non_unique_clients_per_day) / days_back if days_back > 0 else 0
    
    # Print statistics about blacklisted networks
    if len(invalid_network_blacklist) > 0:
        print(f"\n{YELLOW}Skipped {len(invalid_network_blacklist)} networks due to API errors{RESET}")
        logger.info(f"Skipped {len(invalid_network_blacklist)} networks due to API errors")
    
    # Print statistics about timeout-prone networks
    if len(timeout_prone_networks) > 0:
        print(f"\n{YELLOW}Used time chunking for {len(timeout_prone_networks)} timeout-prone networks{RESET}")
        logger.info(f"Used time chunking for {len(timeout_prone_networks)} timeout-prone networks")
        for net_id in timeout_prone_networks:
            timeout_count = network_timeout_counts.get(net_id, 0)
            logger.info(f"  - Network {net_id}: {timeout_count} timeouts")
    
    # Print statistics about high-volume networks
    if len(high_volume_networks) > 0:
        print(f"\n{YELLOW}Detected {len(high_volume_networks)} high-volume networks{RESET}")
        logger.info(f"Detected {len(high_volume_networks)} high-volume networks")
        for net_id in high_volume_networks:
            client_count = network_client_counts.get(net_id, 0)
            print(f"  - Network {net_id[-4:]}: ~{client_count} clients")
            logger.info(f"  - Network {net_id}: ~{client_count} clients")
    
    # Print rate limit statistics
    print(f"\n{YELLOW}Total rate limits encountered: {rate_tracker.total_hits}{RESET}")
    logger.info(f"Total rate limits encountered: {rate_tracker.total_hits}")
    
    # Final log entry
    logger.debug("get_client_stats completed successfully")
    
    return {
        "total_unique_clients": total_unique_clients,
        "total_non_unique_clients": total_non_unique_clients,
        "avg_unique_clients_per_day": round(avg_unique_clients_per_day),
        "avg_non_unique_clients_per_day": round(avg_non_unique_clients_per_day)
    }




def create_or_update_presentation(stats, template_path, output_path, days=14, org_names=None):
    """Create or update PowerPoint presentation with dashboard statistics."""
    try:      
        #print(f"{BLUE}Opening template: {template_path}{RESET}")
        #print(f"{BLUE}Will save to: {output_path}{RESET}")
        #print(f"{BLUE}Stats to update: {stats}{RESET}")
        
        # Create temporary working directory
        temp_dir = tempfile.mkdtemp()
        temp_file = os.path.join(temp_dir, "temp_output.pptx")
        
        # Copy the template to the temp location
        shutil.copy2(template_path, temp_file)
        
        # Open the template
        prs = Presentation(temp_file)
        #print(f"{GREEN}Successfully opened template with {len(prs.slides)} slides{RESET}")
        
        # Update slide 1 (title slide) if org_names are provided
        if org_names and len(prs.slides) > 0:
            title_slide = prs.slides[0]
            
            # Find the title placeholder
            title_shape = None
            for shape in title_slide.shapes:
                if hasattr(shape, 'text') and 'Company - Meraki Bi-Weekly Report' in shape.text:
                    title_shape = shape
                    break
            
            # If we found the title shape, update it
            if title_shape:
                # Create organization names string
                if len(org_names) == 1:
                    org_names_str = next(iter(org_names.values()))
                else:
                    org_names_str = " & ".join(org_names.values())
                
                # Limit length to avoid overflow
                if len(org_names_str) > 60:
                    org_names_str = org_names_str[:57] + "..."
                
                # Get current date for the report
                current_date = datetime.datetime.now().strftime("%B %d, %Y")
                
                # Replace the text
                title_shape.text = title_shape.text.replace(
                    "Company - Meraki Bi-Weekly Report", 
                    f"{org_names_str} - Meraki Bi-Weekly Report"
                )
                
                # Also update the date if it's in the title
                if "March 22, 2025" in title_shape.text:
                    title_shape.text = title_shape.text.replace(
                        "March 22, 2025",
                        current_date
                    )
                
                #print(f"{GREEN}Updated title slide with organization name(s): {org_names_str}{RESET}")
            else:
                print(f"{YELLOW}Could not find title text to update on slide 1{RESET}")
        
        # Create a brand new slide 2 with updated data
        if len(prs.slides) >= 2:
            # Create a copy of slide 1 (usually the title slide)
            title_slide = prs.slides[0]
            
            # Remove slide 2 (index 1)
            slide_to_remove = prs.slides[1]
            r_id = prs.slides._sldIdLst.index(slide_to_remove._element)
            prs.slides._sldIdLst.remove(slide_to_remove._element)
            
            # Add new slide 2
            slide_layout = prs.slide_layouts[0]  # Usually the title layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Move the new slide to position 2 (index 1)
            new_rid = prs.slides._sldIdLst.index(new_slide._element)
            prs.slides._sldIdLst.insert(1, prs.slides._sldIdLst[new_rid])
            prs.slides._sldIdLst.pop(new_rid + 1)  # +1 because we inserted before
            
            # Add title
            title = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_tf = title.text_frame
            title_tf.text = "Overview Stats"
            title_p = title_tf.paragraphs[0]
            title_p.font.bold = True
            title_p.font.size = Pt(32)
            
            # Add Cisco Confidential at bottom
            conf = new_slide.shapes.add_textbox(Inches(8), Inches(6.5), Inches(2), Inches(0.5))
            conf_tf = conf.text_frame
            conf_tf.text = "Cisco Confidential"
            
            # Add date
            date = new_slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5), Inches(0.5))
            date_tf = date.text_frame
            current_date = datetime.datetime.now().strftime("%B %d, %Y")
            date_tf.text = f"Data as of {current_date}"
            
            # NETWORKS - top left
            networks = new_slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(3), Inches(1.5))
            networks_tf = networks.text_frame
            networks_tf.word_wrap = True
            
            networks_p = networks_tf.add_paragraph()
            # Format number with commas
            networks_p.text = f"{stats['total_networks']:,}"
            networks_p.font.size = Pt(48)
            networks_p.font.bold = True
            networks_p.font.color.rgb = RGBColor(0, 150, 0)  # Green
            
            networks_label = networks_tf.add_paragraph()
            networks_label.text = "Networks"
            networks_label.font.size = Pt(18)
            
            # INVENTORY - top center
            inventory = new_slide.shapes.add_textbox(Inches(4.0), Inches(2.0), Inches(3), Inches(1.5))
            inventory_tf = inventory.text_frame
            inventory_tf.word_wrap = True
            
            inventory_p = inventory_tf.add_paragraph()
            # Format number with commas
            inventory_p.text = f"{stats['total_inventory']:,}"
            inventory_p.font.size = Pt(48)
            inventory_p.font.bold = True
            inventory_p.font.color.rgb = RGBColor(0, 0, 200)  # Bright blue
            
            inventory_label = inventory_tf.add_paragraph()
            inventory_label.text = "Total Inventory"
            inventory_label.font.size = Pt(18)
            
            # ACTIVE NODES - top right
            nodes = new_slide.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(3), Inches(1.5))
            nodes_tf = nodes.text_frame
            nodes_tf.word_wrap = True
            
            nodes_p = nodes_tf.add_paragraph()
            # Format number with commas
            nodes_p.text = f"{stats['total_active_nodes']:,}"
            nodes_p.font.size = Pt(48)
            nodes_p.font.bold = True
            nodes_p.font.color.rgb = RGBColor(200, 0, 0)  # Bright red
            
            nodes_label = nodes_tf.add_paragraph()
            nodes_label.text = "Total Active Nodes"
            nodes_label.font.size = Pt(18)
            
            # CLIENT STATS - bottom half
            if 'total_unique_clients' in stats:
                # Add a label for clients section - updated to use dynamic days value
                clients_label = new_slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.4))
                clients_label_tf = clients_label.text_frame
                clients_label_tf.text = f"Clients (for last {days} days)"
                clients_label_tf.paragraphs[0].font.size = Pt(14)
                clients_label_tf.paragraphs[0].font.bold = True
                
                # Unique clients total - bottom left
                unique_total = new_slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(2.5), Inches(1.5))
                unique_total_tf = unique_total.text_frame
                unique_total_tf.word_wrap = True
                
                unique_total_p = unique_total_tf.add_paragraph()
                # Format number with commas
                unique_total_p.text = f"{stats['total_unique_clients']:,}"
                unique_total_p.font.size = Pt(36)
                unique_total_p.font.bold = True
                unique_total_p.font.color.rgb = RGBColor(0, 150, 150)  # Teal
                
                unique_total_label = unique_total_tf.add_paragraph()
                unique_total_label.text = "Unique clients total"
                unique_total_label.font.size = Pt(14)
                
                # Unique clients daily - bottom center-left
                unique_daily = new_slide.shapes.add_textbox(Inches(3.0), Inches(4.3), Inches(2.5), Inches(1.5))
                unique_daily_tf = unique_daily.text_frame
                unique_daily_tf.word_wrap = True
                
                unique_daily_p = unique_daily_tf.add_paragraph()
                # Format number with commas
                unique_daily_p.text = f"{stats['avg_unique_clients_per_day']:,}"
                unique_daily_p.font.size = Pt(36)
                unique_daily_p.font.bold = True
                unique_daily_p.font.color.rgb = RGBColor(0, 150, 150)
                
                unique_daily_label = unique_daily_tf.add_paragraph()
                unique_daily_label.text = "Avg unique clients per day"
                unique_daily_label.font.size = Pt(14)
                
                # Non-unique clients total - bottom center-right
                nonunique_total = new_slide.shapes.add_textbox(Inches(5.5), Inches(4.3), Inches(2.5), Inches(1.5))
                nonunique_total_tf = nonunique_total.text_frame
                nonunique_total_tf.word_wrap = True
                
                nonunique_total_p = nonunique_total_tf.add_paragraph()
                # Format number with commas
                nonunique_total_p.text = f"{stats['total_non_unique_clients']:,}"
                nonunique_total_p.font.size = Pt(36)
                nonunique_total_p.font.bold = True
                nonunique_total_p.font.color.rgb = RGBColor(150, 75, 0)  # Orange-brown
                
                nonunique_total_label = nonunique_total_tf.add_paragraph()
                nonunique_total_label.text = "Non-unique clients total"
                nonunique_total_label.font.size = Pt(14)
                
                # Non-unique clients daily - bottom right
                nonunique_daily = new_slide.shapes.add_textbox(Inches(8.0), Inches(4.3), Inches(2.5), Inches(1.5))
                nonunique_daily_tf = nonunique_daily.text_frame
                nonunique_daily_tf.word_wrap = True
                
                nonunique_daily_p = nonunique_daily_tf.add_paragraph()
                # Format number with commas
                nonunique_daily_p.text = f"{stats['avg_non_unique_clients_per_day']:,}"
                nonunique_daily_p.font.size = Pt(36)
                nonunique_daily_p.font.bold = True
                nonunique_daily_p.font.color.rgb = RGBColor(150, 75, 0)
                
                nonunique_daily_label = nonunique_daily_tf.add_paragraph()
                nonunique_daily_label.text = "Non-unique clients per day"
                nonunique_daily_label.font.size = Pt(14)
            
            #print(f"{GREEN}Successfully built new dashboard slide{RESET}")
            
            # Save the presentation
            try:
                prs.save(temp_file)
                #print(f"{GREEN}Saved to temporary file{RESET}")
                
                # Copy back to the output location
                shutil.copy2(temp_file, output_path)
                #print(f"{GREEN}Copied to final location: {output_path}{RESET}")
            except Exception as e:
                print(f"{RED}Error saving file: {e}{RESET}")
                import traceback
                traceback.print_exc()
        
        # Clean up the temp directory
        try:
            shutil.rmtree(temp_dir)
        except:
            pass
        
        return output_path
        
    except Exception as e:
        print(f"{RED}Error in create_or_update_presentation: {e}{RESET}")
        import traceback
        traceback.print_exc()
        return output_path

async def main_async(org_ids, days=14, template_path=None, output_path=None):
    """Async main function with conservative settings to avoid rate limits."""
    # Use default paths if not provided
    if template_path is None:
        template_path = "template.pptx"
    if output_path is None:
        output_path = "meraki_report.pptx"
        
    # Start timer
    start_time = time.time()
    
    # Get API key
    api_key = get_api_key()
    
    # Create a conservative rate limiter that starts with low limits
    rate_limiter = AdaptiveRateLimiter(initial_limit=5, min_limit=3, max_limit=15)
    
    #print(f"{BLUE}Starting with conservative concurrency limit of {rate_limiter.current_limit}{RESET}")
    
    # Clean up old log files before starting
    try:
        log_files = [f for f in os.listdir('.') if f.startswith('meraki_api_log') and f.endswith('.log')]
        log_files.sort(key=lambda x: os.path.getmtime(x), reverse=True)
        for old_log in log_files[1:]:  # Keep only the most recent log
            os.remove(old_log)
            print(f"Removed old log file: {old_log}")
    except Exception as e:
        print(f"Error cleaning up log files: {e}")
    
    # Initialize async Meraki client 
    async with meraki.aio.AsyncDashboardAPI(
        api_key=api_key,
        suppress_logging=False,  # Keep this False if you want logs
        log_file_prefix='meraki_api_log',  # Fixed name prefix
        log_path='.',  # Current directory
        maximum_retries=3,
        base_url="https://api.gov-meraki.com/api/v1"
    ) as aiomeraki:
        try:
            # Get organization names first
            #print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Getting organization names...{RESET}")
            org_names = await get_organization_names(aiomeraki, org_ids, rate_limiter)
            #print(f"Organizations: {org_names}")
            
            # Get all networks for the specified organizations
            #print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Getting networks...{RESET}")
            network_start_time = time.time()
            
            # Process organizations sequentially to avoid rate limits
            all_networks = []
            for org_id in org_ids:
                try:
                    networks = await get_networks(aiomeraki, org_id, rate_limiter)
                    all_networks.extend(networks)
                    #print(f"Found {len(networks)} networks in organization {org_id}")
                except Exception as e:
                    print(f"{RED}Error retrieving networks for organization {org_id}: {e}{RESET}")
            
            network_time = time.time() - network_start_time
            #print(f"{PURPLE}Network retrieval completed in {network_time:.2f} seconds{RESET}")
            
            # Print network product types for diagnostic purposes
            product_types = {}
            for network in all_networks:
                for product in network.get('productTypes', []):
                    product_types[product] = product_types.get(product, 0) + 1
            
            #print("\nNetwork product types:")
            for product, count in sorted(product_types.items(), key=lambda x: x[1], reverse=True):
                #print(f"  - {product}: {count} networks")
                pass
            
            # Extract network IDs
            network_ids = [network['id'] for network in all_networks]
            #print(f"{BLUE}Total networks to process: {len(network_ids)}{RESET}")
            
            # Filter out incompatible networks (Systems Manager, Camera, etc.)
            filter_start_time = time.time()
            #print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Filtering networks...{RESET}")
            valid_network_ids = await filter_incompatible_networks(network_ids, all_networks)
            filter_time = time.time() - filter_start_time
            #print(f"{PURPLE}Network filtering completed in {filter_time:.2f} seconds{RESET}")
            
            # Get dashboard statistics
            dashboard_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Getting dashboard statistics...{RESET}")
            dashboard_stats = await get_dashboard_stats(aiomeraki, org_ids, rate_limiter)
            dashboard_time = time.time() - dashboard_start_time
            print(f"{PURPLE}Dashboard statistics completed in {dashboard_time:.2f} seconds{RESET}")
            
            # Get client statistics - pass filtered network IDs
            client_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Getting client statistics...{RESET}")
            client_stats = await get_client_stats(aiomeraki, valid_network_ids, rate_limiter, days)
            client_time = time.time() - client_start_time
            print(f"{PURPLE}Client statistics completed in {client_time:.2f} seconds{RESET}")
            
            # End timer for data collection
            data_collection_time = time.time() - start_time
            print(f"\n{PURPLE}Total data collection completed in {data_collection_time:.2f} seconds{RESET}")
            
            # Combine all stats
            combined_stats = {**dashboard_stats, **client_stats}
            
            # Start timer for PowerPoint generation
            ppt_start_time = time.time()
            print(f"\n{PURPLE}[{time.strftime('%H:%M:%S')}] Creating PowerPoint presentation...{RESET}")
            
            # Create/update PowerPoint - pass the days parameter and org_names
            output_file = create_or_update_presentation(
                combined_stats, 
                template_path, 
                output_path, 
                days,
                org_names
            )
            
            ppt_time = time.time() - ppt_start_time
            print(f"{PURPLE}PowerPoint generation completed in {ppt_time:.2f} seconds{RESET}")
            
            # Calculate total script execution time
            total_time = time.time() - start_time
            print(f"\n{PURPLE}Total script execution time: {total_time:.2f} seconds{RESET}")
            
            #print(f"\n{BLUE}Dashboard Report created successfully at {output_file}{RESET}")
            #print(f"Total Networks: {dashboard_stats['total_networks']}")
            #print(f"Total Inventory Devices: {dashboard_stats['total_inventory']}")
            #print(f"Total Active Nodes: {dashboard_stats['total_active_nodes']}")
            #print(f"Total Unique Clients: {client_stats['total_unique_clients']}")
            #print(f"Total Non-Unique Clients: {client_stats['total_non_unique_clients']}")
            #print(f"Avg Unique Clients Per Day: {client_stats['avg_unique_clients_per_day']}")
            #print(f"Avg Non-Unique Clients Per Day: {client_stats['avg_non_unique_clients_per_day']}")
            
            # Add timing breakdown
            #print(f"\n{PURPLE}Timing Breakdown:{RESET}")
            #print(f"{PURPLE}  - Network Retrieval: {network_time:.2f}s ({network_time/total_time*100:.1f}%){RESET}")
            #print(f"{PURPLE}  - Network Filtering: {filter_time:.2f}s ({filter_time/total_time*100:.1f}%){RESET}")
            #print(f"{PURPLE}  - Dashboard Statistics: {dashboard_time:.2f}s ({dashboard_time/total_time*100:.1f}%){RESET}")
            #print(f"{PURPLE}  - Client Statistics: {client_time:.2f}s ({client_time/total_time*100:.1f}%){RESET}")
            #print(f"{PURPLE}  - PowerPoint Generation: {ppt_time:.2f}s ({ppt_time/total_time*100:.1f}%){RESET}")
            
            return combined_stats
            
        except Exception as e:
            print(f"{RED}Error in main processing: {e}{RESET}")
            # Return empty stats as a fallback
            return {
                "total_networks": 0,
                "total_inventory": 0, 
                "total_active_nodes": 0,
                "total_unique_clients": 0,
                "total_non_unique_clients": 0,
                "avg_unique_clients_per_day": 0,
                "avg_non_unique_clients_per_day": 0
            }
